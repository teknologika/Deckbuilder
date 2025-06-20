import os
import shutil
import json
import yaml
import re
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from placeholder_types import (
    is_title_placeholder, 
    is_subtitle_placeholder, 
    is_content_placeholder,
    is_media_placeholder
)
try:
    from table_styles import TABLE_HEADER_STYLES, TABLE_ROW_STYLES, TABLE_BORDER_STYLES
except ImportError:
    # Fallback values if modules don't exist
    TABLE_HEADER_STYLES = {
        "dark_blue_white_text": {"bg": RGBColor(46, 89, 132), "text": RGBColor(255, 255, 255)}
    }
    TABLE_ROW_STYLES = {
        "alternating_light_gray": {"primary": RGBColor(255, 255, 255), "alt": RGBColor(240, 240, 240)}
    }
    TABLE_BORDER_STYLES = {
        "thin_gray": {"width": Pt(1), "color": RGBColor(128, 128, 128), "style": "all"}
    }

def singleton(cls):
    instances = {}
    def get_instance(*args, **kwargs):
        if cls not in instances:
            instances[cls] = cls(*args, **kwargs)
        return instances[cls]
    return get_instance

@singleton
class Deckbuilder:

    def __init__(self):
        self.template_path = os.getenv('DECK_TEMPLATE_FOLDER')
        self.template_name = os.getenv('DECK_TEMPLATE_NAME')
        self.output_folder = os.getenv('DECK_OUTPUT_FOLDER')
        self.prs = Presentation()
        self.layout_mapping = None
        
        # Ensure default template exists in templates folder
        self._check_template_exists(self.template_name or 'default')
        
    def _check_template_exists(self, templateName: str):
        """Check if template exists in the templates folder and copy if needed."""

        # Use self.template_name if available, otherwise use default
        if not templateName or templateName == 'default':
            templateName = self.template_name or 'default'

        # Ensure templateName ends with .pptx
        if not templateName.endswith('.pptx'):
            templateName += '.pptx'
        
        if self.template_path:
            try:
                # Create templates folder if it doesn't exist
                os.makedirs(self.template_path, exist_ok=True)
                
                # Check if template exists in templates folder
                default_template = os.path.join(self.template_path, templateName)
                if not os.path.exists(default_template):
                    # Copy from src/default.pptx
                    src_template = os.path.join(os.path.dirname(__file__), 'default.pptx')
                    if os.path.exists(src_template):
                        shutil.copy2(src_template, default_template)
                
                # Also copy the corresponding JSON mapping file
                base_name = templateName.replace('.pptx', '')
                json_template = os.path.join(self.template_path, base_name + '.json')
                if not os.path.exists(json_template):
                    # Copy from src/default.json
                    src_json = os.path.join(os.path.dirname(__file__), base_name + '.json')
                    if os.path.exists(src_json):
                        shutil.copy2(src_json, json_template)
            except (OSError, IOError):
                # Handle file operation errors silently
                pass
    
    def _load_layout_mapping(self, templateName: str):
        """Load layout mapping from JSON file."""
        if not templateName.endswith('.json'):
            templateName += '.json'
        
        # Try to load from template folder first
        if self.template_path:
            mapping_path = os.path.join(self.template_path, templateName)
            if os.path.exists(mapping_path):
                try:
                    with open(mapping_path, 'r', encoding='utf-8') as f:
                        self.layout_mapping = json.load(f)
                        return
                except:
                    pass
        
        # Fallback to src folder
        src_mapping_path = os.path.join(os.path.dirname(__file__), templateName)
        if os.path.exists(src_mapping_path):
            try:
                with open(src_mapping_path, 'r', encoding='utf-8') as f:
                    self.layout_mapping = json.load(f)
                    return
            except:
                pass
        
        # Use fallback mapping if JSON not found
        self.layout_mapping = {
            "layouts": {"Title and Content": {"index": 1}},
            "aliases": {"content": "Title and Content", "title": "Title Slide"}
        }
    
    def create_presentation(self, templateName: str = "default", fileName: str = "Sample_Presentation") -> str:
        # Check template exists
        self._check_template_exists(templateName)
        
        # Load layout mapping
        base_name = templateName.replace('.pptx', '') if templateName.endswith('.pptx') else templateName
        self._load_layout_mapping(base_name)
        
        # Create deck with template
        if not templateName.endswith('.pptx'):
            templateName += '.pptx'
        if self.template_path:
            template_path = os.path.join(self.template_path, templateName)
            self.prs = Presentation(template_path) if os.path.exists(template_path) else Presentation()
        else:
            self.prs = Presentation()
        
        self._clear_slides()

        return f"Creating presentation: {fileName}"

    def write_presentation(self, fileName: str = "Sample_Presentation") -> str:
        """Writes the generated presentation to disk with ISO timestamp."""
        from datetime import datetime
        
        # Get output folder from environment or use default
        output_folder = self.output_folder or '.'
        
        # Ensure output folder exists
        os.makedirs(output_folder, exist_ok=True)
        
        # Create filename with ISO timestamp and .g.pptx extension for generated files
        timestamp = datetime.now().strftime("%Y-%m-%d_%H%M")
        generated_filename = f"{fileName}.{timestamp}.g.pptx"
        output_file = os.path.join(output_folder, generated_filename)
        
        # Save the presentation (overwrites if same timestamp exists)
        self.prs.save(output_file)
        
        return f"Successfully created presentation: {os.path.basename(output_file)}"

    def add_slide_from_json(self, json_data) -> str:
        """
        Add a slide to the presentation using JSON data.
        
        Args:
            json_data: JSON string or dictionary containing slide data
            
        Returns:
            Success message
        """
        try:
            # Handle both string and dictionary inputs
            if isinstance(json_data, str):
                # Parse JSON data - handle potential double encoding
                data = json.loads(json_data)
                
                # If the result is still a string, parse it again
                if isinstance(data, str):
                    data = json.loads(data)
            else:
                # Already a dictionary
                data = json_data
            
            # Handle different JSON formats
            if "slides" in data:
                # Multiple slides format
                for slide_data in data["slides"]:
                    self._add_slide(slide_data)
            elif "presentation" in data and "slides" in data["presentation"]:
                # Presentation wrapper format
                for slide_data in data["presentation"]["slides"]:
                    self._add_slide(slide_data)
            else:
                # Single slide format
                self._add_slide(data)
                
            return "Successfully added slide(s) from JSON data"
            
        except json.JSONDecodeError as e:
            return f"Error parsing JSON: {str(e)}"
        except Exception as e:
            return f"Error adding slide: {str(e)}"

    def _clear_slides(self):
        """Clear all slides from the presentation."""
        slide_count = len(self.prs.slides)
        for i in range(slide_count - 1, -1, -1):
            rId = self.prs.slides._sldIdLst[i].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[i]

    def _add_slide(self, slide_data: dict):
        """
        Add a single slide to the presentation based on slide data.
        
        Args:
            slide_data: Dictionary containing slide information
        """
        # Auto-parse JSON formatting for inline formatting support
        slide_data = self._auto_parse_json_formatting(slide_data)
        
        # Get slide type and determine layout using JSON mapping
        slide_type = slide_data.get("type", "content")
        
        # Use layout mapping if available
        if self.layout_mapping:
            aliases = self.layout_mapping.get("aliases", {})
            layouts = self.layout_mapping.get("layouts", {})
            
            # Get layout name from aliases
            layout_name = aliases.get(slide_type, slide_type)
            
            # Get layout index
            layout_info = layouts.get(layout_name, {})
            layout_index = layout_info.get("index", 1)
        else:
            # Fallback
            layout_index = 1
        
        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add content to placeholders using template mapping + semantic detection
        self._apply_content_to_mapped_placeholders(slide, slide_data, layout_name)
        
        # Handle rich content
        if "rich_content" in slide_data:
            self._add_rich_content_to_slide(slide, slide_data["rich_content"])
        elif "content" in slide_data:
            # Fallback to simple content (backwards compatibility)
            self._add_simple_content_to_slide(slide, slide_data["content"])
        
        # Add table if provided
        if "table" in slide_data:
            self._add_table_to_slide(slide, slide_data["table"])

    def _apply_content_to_mapped_placeholders(self, slide, slide_data, layout_name):
        """
        Apply content to placeholders using template JSON mappings + semantic detection.
        
        This unified method works with both JSON input and markdown frontmatter input:
        1. Looks up layout in template JSON mappings
        2. For each field in slide_data, finds corresponding placeholder index
        3. Gets actual placeholder and determines its semantic type
        4. Applies content using appropriate semantic handler
        
        Args:
            slide: PowerPoint slide object
            slide_data: Dictionary containing slide content (from JSON or markdown)
            layout_name: Name of the PowerPoint layout
        """
        if not self.layout_mapping:
            # Fallback to basic semantic detection if no mapping available
            self._add_content_to_placeholders_fallback(slide, slide_data)
            return
        
        # Get layout info from template mapping
        layouts = self.layout_mapping.get("layouts", {})
        layout_info = layouts.get(layout_name, {})
        placeholder_mappings = layout_info.get("placeholders", {})
        
        # Create reverse mapping: field_name -> placeholder_index
        field_to_index = {}
        for placeholder_idx, field_name in placeholder_mappings.items():
            field_to_index[field_name] = int(placeholder_idx)
        
        # Process each field in slide_data using semantic detection
        for field_name, field_value in slide_data.items():
            # Skip non-content fields
            if field_name in ['type', 'rich_content', 'table']:
                continue
                
            # Skip formatted variants (handled automatically)
            if field_name.endswith('_formatted'):
                continue
            
            # Find placeholder using semantic detection
            target_placeholder = None
            
            # Handle title placeholders
            if field_name == "title":
                for placeholder in slide.placeholders:
                    if is_title_placeholder(placeholder.placeholder_format.type):
                        target_placeholder = placeholder
                        break
            
            # Handle subtitle placeholders
            elif field_name == "subtitle":
                for placeholder in slide.placeholders:
                    if is_subtitle_placeholder(placeholder.placeholder_format.type):
                        target_placeholder = placeholder
                        break
            
            # Handle content placeholders
            elif field_name == "content":
                for placeholder in slide.placeholders:
                    if is_content_placeholder(placeholder.placeholder_format.type):
                        target_placeholder = placeholder
                        break
            
            # Handle other fields by checking if they match placeholder names in JSON mapping
            else:
                # Try to find by exact field name match in JSON mapping
                if field_name in field_to_index:
                    placeholder_idx = field_to_index[field_name]
                    for placeholder in slide.placeholders:
                        if placeholder.placeholder_format.idx == placeholder_idx:
                            target_placeholder = placeholder
                            break
            
            if target_placeholder:
                # Apply content based on placeholder's semantic type
                self._apply_content_by_semantic_type(target_placeholder, field_name, field_value, slide_data)

    def _add_content_to_placeholders_fallback(self, slide, slide_data):
        """
        Fallback method for basic semantic placeholder detection when no JSON mapping available.
        """
        for shape in slide.placeholders:
            placeholder_type = shape.placeholder_format.type
            
            # Handle title placeholders
            if "title" in slide_data and is_title_placeholder(placeholder_type):
                if "title_formatted" in slide_data:
                    self._apply_formatted_segments_to_shape(shape, slide_data["title_formatted"])
                else:
                    shape.text = slide_data["title"]
            
            # Handle subtitle placeholders
            elif "subtitle" in slide_data and is_subtitle_placeholder(placeholder_type):
                if "subtitle_formatted" in slide_data:
                    self._apply_formatted_segments_to_shape(shape, slide_data["subtitle_formatted"])
                else:
                    shape.text = slide_data["subtitle"]
            
            # Handle main content placeholders (for simple content)
            elif "content" in slide_data and is_content_placeholder(placeholder_type):
                # Only use simple content if rich_content is not available
                if "rich_content" not in slide_data:
                    self._add_simple_content_to_placeholder(shape, slide_data["content"])

    def _apply_content_by_semantic_type(self, placeholder, field_name, field_value, slide_data):
        """
        Apply content to a placeholder based on its semantic type and the content type.
        """
        placeholder_type = placeholder.placeholder_format.type
        
        # Check for formatted version of the field
        formatted_field = field_name + '_formatted'
        has_formatted = formatted_field in slide_data
        
        # Apply content based on placeholder semantic type
        if is_title_placeholder(placeholder_type) or is_subtitle_placeholder(placeholder_type):
            # Title/subtitle placeholders - simple text with formatting
            if has_formatted:
                self._apply_formatted_segments_to_shape(placeholder, slide_data[formatted_field])
            else:
                placeholder.text = str(field_value)
                
        elif is_content_placeholder(placeholder_type):
            # Content placeholders - can handle text, lists, etc.
            if has_formatted:
                self._apply_formatted_segments_to_shape(placeholder, slide_data[formatted_field])
            elif isinstance(field_value, (list, tuple)):
                self._add_simple_content_to_placeholder(placeholder, field_value)
            else:
                self._add_simple_content_to_placeholder(placeholder, str(field_value))
                
        else:
            # Other placeholder types - simple text for now
            if has_formatted:
                self._apply_formatted_segments_to_shape(placeholder, slide_data[formatted_field])
            else:
                placeholder.text = str(field_value)

    def _add_simple_content_to_placeholder(self, placeholder, content):
        """Add simple content to a content placeholder with inline formatting support."""
        if not hasattr(placeholder, 'text_frame'):
            return
            
        text_frame = placeholder.text_frame
        text_frame.clear()
        
        if isinstance(content, str):
            p = text_frame.paragraphs[0]
            self._apply_inline_formatting(content, p)
        elif isinstance(content, list):
            for i, line in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]  # Use existing first paragraph
                else:
                    p = text_frame.add_paragraph()
                self._apply_inline_formatting(line, p)

    def _parse_inline_formatting(self, text):
        """Parse inline formatting and return structured formatting data"""
        import re
        
        if not text:
            return [{"text": "", "format": {}}]
        
        # Patterns in order of precedence (longest patterns first to avoid conflicts)
        patterns = [
            (r'\*\*\*___(.*?)___\*\*\*', {'bold': True, 'italic': True, 'underline': True}),  # ***___text___***
            (r'___\*\*\*(.*?)\*\*\*___', {'bold': True, 'italic': True, 'underline': True}),  # ___***text***___
            (r'\*\*\*(.*?)\*\*\*', {'bold': True, 'italic': True}),                          # ***text***
            (r'___(.*?)___', {'underline': True}),                                            # ___text___
            (r'\*\*(.*?)\*\*', {'bold': True}),                                              # **text**
            (r'\*(.*?)\*', {'italic': True})                                                 # *text*
        ]
        
        # Find all matches and their positions
        all_matches = []
        for pattern, format_dict in patterns:
            for match in re.finditer(pattern, text):
                all_matches.append((match.start(), match.end(), match.group(1), format_dict))
        
        # Sort matches by position
        all_matches.sort(key=lambda x: x[0])
        
        # Remove overlapping matches (keep the first one found)
        filtered_matches = []
        last_end = 0
        for start, end, content, format_dict in all_matches:
            if start >= last_end:
                filtered_matches.append((start, end, content, format_dict))
                last_end = end
        
        # Build the formatted text segments
        segments = []
        last_pos = 0
        
        for start, end, content, format_dict in filtered_matches:
            # Add plain text before the formatted text
            if start > last_pos:
                plain_text = text[last_pos:start]
                if plain_text:
                    segments.append({"text": plain_text, "format": {}})
            
            # Add formatted text
            segments.append({"text": content, "format": format_dict})
            last_pos = end
        
        # Add any remaining plain text
        if last_pos < len(text):
            remaining_text = text[last_pos:]
            if remaining_text:
                segments.append({"text": remaining_text, "format": {}})
        
        # If no formatting found, return the original text
        if not segments:
            segments = [{"text": text, "format": {}}]
            
        return segments

    def _apply_inline_formatting(self, text, paragraph):
        """Apply inline formatting to paragraph using parsed formatting data."""
        # Clear any existing text
        paragraph.text = ""
        
        # Parse the formatting
        segments = self._parse_inline_formatting(text)
        
        # Apply each segment to the paragraph
        for segment in segments:
            run = paragraph.add_run()
            run.text = segment["text"]
            
            # Apply formatting
            format_dict = segment["format"]
            if format_dict.get("bold"):
                run.font.bold = True
            if format_dict.get("italic"):
                run.font.italic = True
            if format_dict.get("underline"):
                run.font.underline = True

    def _apply_formatted_segments_to_shape(self, shape, segments):
        """Apply formatted text segments to a shape's text frame."""
        if not hasattr(shape, 'text_frame'):
            # For shapes that don't have text_frame, fall back to simple text
            shape.text = ''.join(segment["text"] for segment in segments)
            return
            
        text_frame = shape.text_frame
        text_frame.clear()
        
        # Use the first paragraph or create one
        if text_frame.paragraphs:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
        
        paragraph.text = ""
        
        # Apply each segment
        for segment in segments:
            run = paragraph.add_run()
            run.text = segment["text"]
            
            # Apply formatting
            format_dict = segment["format"]
            if format_dict.get("bold"):
                run.font.bold = True
            if format_dict.get("italic"):
                run.font.italic = True
            if format_dict.get("underline"):
                run.font.underline = True

    def _apply_formatted_segments_to_cell(self, cell, segments):
        """Apply formatted text segments to a table cell."""
        text_frame = cell.text_frame
        text_frame.clear()
        
        # Create first paragraph
        paragraph = text_frame.paragraphs[0]
        paragraph.text = ""
        
        # Apply each segment
        for segment in segments:
            run = paragraph.add_run()
            run.text = segment["text"]
            
            # Apply formatting
            format_dict = segment["format"]
            if format_dict.get("bold"):
                run.font.bold = True
            if format_dict.get("italic"):
                run.font.italic = True
            if format_dict.get("underline"):
                run.font.underline = True

    def _add_rich_content_to_slide(self, slide, rich_content: list):
        """Add rich content blocks to a slide with improved formatting"""
        # Find the content placeholder
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Content placeholder
                content_placeholder = shape
                break
        
        if not content_placeholder:
            return
        
        # Clear existing content
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        
        # Set margins for better spacing
        text_frame.margin_left = Cm(0.25)
        text_frame.margin_right = Cm(0.25)
        text_frame.margin_top = Cm(0.25)
        text_frame.margin_bottom = Cm(0.25)
        
        # Add each content block with proper hierarchy
        first_content = True
        for block in rich_content:
            if "heading" in block:
                if first_content:
                    p = text_frame.paragraphs[0]  # Use existing first paragraph
                else:
                    p = text_frame.add_paragraph()
                self._apply_inline_formatting(block["heading"], p)
                # Apply bold to all runs in the heading paragraph
                for run in p.runs:
                    run.font.bold = True
                p.space_after = Pt(6)
                p.space_before = Pt(12) if not first_content else Pt(0)
                
            elif "paragraph" in block:
                if first_content:
                    p = text_frame.paragraphs[0]  # Use existing first paragraph
                else:
                    p = text_frame.add_paragraph()
                self._apply_inline_formatting(block["paragraph"], p)
                p.space_after = Pt(6)
                p.space_before = Pt(3)
                
            elif "bullets" in block:
                # Get bullet levels if available, otherwise default to level 1
                bullet_levels = block.get("bullet_levels", [1] * len(block["bullets"]))
                
                for bullet_idx, bullet in enumerate(block["bullets"]):
                    if first_content and bullet_idx == 0:
                        p = text_frame.paragraphs[0]  # Use existing first paragraph for first bullet
                    else:
                        p = text_frame.add_paragraph()
                    self._apply_inline_formatting(bullet, p)
                    
                    # Use the parsed bullet level
                    bullet_level = bullet_levels[bullet_idx] if bullet_idx < len(bullet_levels) else 1
                    p.level = bullet_level
                    
                    # Set spacing based on level
                    if bullet_level == 1:
                        p.space_after = Pt(3)
                    else:  # Level 2+ (sub-bullets)
                        p.space_after = Pt(2)
            
            first_content = False

    def _add_simple_content_to_slide(self, slide, content):
        """Add simple content to slide with inline formatting support (backwards compatibility)"""
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Content placeholder
                text_frame = shape.text_frame
                text_frame.clear()
                
                if isinstance(content, str):
                    p = text_frame.paragraphs[0]
                    self._apply_inline_formatting(content, p)
                elif isinstance(content, list):
                    for i, line in enumerate(content):
                        if i == 0:
                            p = text_frame.paragraphs[0]  # Use existing first paragraph
                        else:
                            p = text_frame.add_paragraph()
                        self._apply_inline_formatting(line, p)
                break

    def _add_table_to_slide(self, slide, table_data):
        """
        Add a styled table to a slide.
        
        Args:
            slide: The slide to add the table to
            table_data: Dictionary containing table data and styling options
        """
        # Get table data - support both 'data' and 'rows' keys for backwards compatibility
        data = table_data.get("data", table_data.get("rows", []))
        if not data:
            return
        
        # Get styling options
        header_style = table_data.get("header_style", "dark_blue_white_text")
        row_style = table_data.get("row_style", "alternating_light_gray")
        border_style = table_data.get("border_style", "thin_gray")
        custom_colors = table_data.get("custom_colors", {})
        
        # Find content placeholder or create table in available space
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Content placeholder
                content_placeholder = shape
                break
        
        if content_placeholder:
            # Remove placeholder and create table in its place
            left = content_placeholder.left
            top = content_placeholder.top
            width = content_placeholder.width
            height = content_placeholder.height
            
            # Remove the placeholder
            sp = content_placeholder._element
            sp.getparent().remove(sp)
        else:
            # Default positioning if no placeholder found
            left = Cm(2.5)
            top = Cm(5)
            width = Cm(20)
            height = Cm(12)
        
        # Create the table
        rows = len(data)
        if data:
            # Handle both old (list of strings) and new (list of dicts) formats
            first_row = data[0]
            if isinstance(first_row, list):
                cols = len(first_row)
            else:
                cols = 1  # Fallback
        else:
            cols = 1
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Apply table data with formatting support
        for row_idx, row_data in enumerate(data):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                
                # Handle both old (string) and new (formatted) cell data
                if isinstance(cell_data, dict) and "formatted" in cell_data:
                    # New formatted cell data
                    self._apply_formatted_segments_to_cell(cell, cell_data["formatted"])
                else:
                    # Old string cell data
                    cell.text = str(cell_data)
        
        # Apply styling
        self._apply_table_styling(table, header_style, row_style, border_style, custom_colors)

    def _apply_table_styling(self, table, header_style, row_style, border_style, custom_colors):
        """
        Apply styling to a table.
        
        Args:
            table: The table object to style
            header_style: Header style name
            row_style: Row style name
            border_style: Border style name
            custom_colors: Dictionary of custom color overrides
        """
        # Apply header styling
        if header_style in TABLE_HEADER_STYLES:
            header_colors = TABLE_HEADER_STYLES[header_style]
            
            # Override with custom colors if provided
            bg_color = self._parse_custom_color(custom_colors.get("header_bg")) or header_colors["bg"]
            text_color = self._parse_custom_color(custom_colors.get("header_text")) or header_colors["text"]
            
            # Style header row (first row)
            for col_idx in range(len(table.columns)):
                cell = table.cell(0, col_idx)
                # Set background color
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_color
                
                # Set text color and formatting
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = text_color
                        run.font.bold = True
        
        # Apply row styling
        if row_style in TABLE_ROW_STYLES and len(table.rows) > 1:
            row_colors = TABLE_ROW_STYLES[row_style]
            
            # Override with custom colors if provided
            primary_color = self._parse_custom_color(custom_colors.get("primary_row")) or row_colors["primary"]
            alt_color = self._parse_custom_color(custom_colors.get("alt_row")) or row_colors["alt"]
            
            # Style data rows (skip header row)
            for row_idx in range(1, len(table.rows)):
                is_alt_row = (row_idx - 1) % 2 == 1
                bg_color = alt_color if is_alt_row else primary_color
                
                if bg_color is not None:
                    for col_idx in range(len(table.columns)):
                        cell = table.cell(row_idx, col_idx)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = bg_color
        
        # Apply border styling
        if border_style in TABLE_BORDER_STYLES:
            self._apply_table_borders(table, TABLE_BORDER_STYLES[border_style], custom_colors)

    def _apply_table_borders(self, table, border_config, custom_colors):
        """
        Apply border styling to a table.
        
        Args:
            table: The table object
            border_config: Border configuration dictionary
            custom_colors: Custom color overrides
        """
        border_width = border_config["width"]
        border_color = self._parse_custom_color(custom_colors.get("border_color")) or border_config["color"]
        border_style = border_config["style"]
        
        if border_style == "none" or border_width.cm == 0:
            return
        
        # Apply borders based on style
        for row_idx in range(len(table.rows)):
            for col_idx in range(len(table.columns)):
                cell = table.cell(row_idx, col_idx)
                
                if border_style == "all":
                    # All borders
                    self._set_cell_borders(cell, border_width, border_color, all_sides=True)
                elif border_style == "header" and row_idx == 0:
                    # Only header bottom border
                    self._set_cell_borders(cell, border_width, border_color, bottom=True)
                elif border_style == "outer":
                    # Only outer borders
                    is_top = row_idx == 0
                    is_bottom = row_idx == len(table.rows) - 1
                    is_left = col_idx == 0
                    is_right = col_idx == len(table.columns) - 1
                    
                    self._set_cell_borders(cell, border_width, border_color, 
                                         top=is_top, bottom=is_bottom, 
                                         left=is_left, right=is_right)

    def _set_cell_borders(self, cell, width, color, all_sides=False, 
                         top=False, bottom=False, left=False, right=False):
        """
        Set borders for a table cell.
        
        Args:
            cell: The table cell
            width: Border width
            color: Border color
            all_sides: Apply to all sides
            top, bottom, left, right: Apply to specific sides
        """
        if color is None:
            return
            
        if all_sides:
            top = bottom = left = right = True
        
        # Note: python-pptx has limited border support
        # This is a simplified implementation
        try:
            if hasattr(cell, 'border'):
                if top and hasattr(cell.border, 'top'):
                    cell.border.top.color.rgb = color
                    cell.border.top.width = width
                if bottom and hasattr(cell.border, 'bottom'):
                    cell.border.bottom.color.rgb = color
                    cell.border.bottom.width = width
                if left and hasattr(cell.border, 'left'):
                    cell.border.left.color.rgb = color
                    cell.border.left.width = width
                if right and hasattr(cell.border, 'right'):
                    cell.border.right.color.rgb = color
                    cell.border.right.width = width
        except (AttributeError, Exception):
            # Borders not fully supported in python-pptx, skip silently
            pass

    def _parse_custom_color(self, color_value):
        """
        Parse a custom color value (hex string) to RGBColor.
        
        Args:
            color_value: Hex color string (e.g., "#FF0000")
            
        Returns:
            RGBColor object or None if invalid
        """
        if not color_value or not isinstance(color_value, str):
            return None
        
        try:
            # Remove # if present
            color_value = color_value.lstrip('#')
            
            # Convert hex to RGB
            if len(color_value) == 6:
                r = int(color_value[0:2], 16)
                g = int(color_value[2:4], 16)
                b = int(color_value[4:6], 16)
                return RGBColor(r, g, b)
        except (ValueError, TypeError):
            pass
        
        return None

    def parse_markdown_with_frontmatter(self, markdown_content: str) -> list:
        """
        Parse markdown content with frontmatter into slide data.
        
        Args:
            markdown_content: Markdown string with frontmatter slide definitions
            
        Returns:
            List of slide dictionaries ready for _add_slide()
        """
        slides = []
        
        # Split content by frontmatter boundaries
        slide_blocks = re.split(r'^---\s*$', markdown_content, flags=re.MULTILINE)
        
        i = 0
        while i < len(slide_blocks):
            # Skip empty blocks
            if not slide_blocks[i].strip():
                i += 1
                continue
                
            # Look for frontmatter + content pairs
            if i + 1 < len(slide_blocks):
                try:
                    frontmatter_raw = slide_blocks[i].strip()
                    content_raw = slide_blocks[i + 1].strip() if i + 1 < len(slide_blocks) else ""
                    
                    # Parse frontmatter with error handling for special characters
                    try:
                        slide_config = yaml.safe_load(frontmatter_raw) or {}
                    except yaml.YAMLError:
                        # If YAML parsing fails due to special characters, try pre-processing
                        slide_config = self._parse_frontmatter_safe(frontmatter_raw)
                    
                    # Parse markdown content into slide data
                    slide_data = self._parse_slide_content(content_raw, slide_config)
                    slides.append(slide_data)
                    
                    i += 2  # Skip both frontmatter and content blocks
                except yaml.YAMLError:
                    # If YAML parsing fails, treat as regular content
                    content_raw = slide_blocks[i].strip()
                    slide_data = self._parse_slide_content(content_raw, {})
                    slides.append(slide_data)
                    i += 1
            else:
                # Single block without frontmatter
                content_raw = slide_blocks[i].strip()
                slide_data = self._parse_slide_content(content_raw, {})
                slides.append(slide_data)
                i += 1
        
        return slides

    def _parse_frontmatter_safe(self, frontmatter_raw: str) -> dict:
        """
        Parse frontmatter safely by handling special characters that break YAML.
        
        This method processes frontmatter line by line to handle values with
        markdown formatting characters (*, _, etc.) that would break YAML parsing.
        """
        config = {}
        for line in frontmatter_raw.split('\n'):
            line = line.strip()
            if not line or line.startswith('#'):
                continue
                
            if ':' in line:
                key, value = line.split(':', 1)
                key = key.strip()
                value = value.strip()
                
                # Remove quotes if present
                if value.startswith('"') and value.endswith('"'):
                    value = value[1:-1]
                elif value.startswith("'") and value.endswith("'"):
                    value = value[1:-1]
                
                config[key] = value
        
        return config

    def _parse_slide_content(self, content: str, config: dict) -> dict:
        """Convert markdown content + config into slide data dict with mixed content support"""
        slide_data = {
            "type": config.get("layout", "content"),
            **config  # Include all frontmatter as slide properties
        }
        
        if not content.strip():
            return slide_data
            
        lines = content.split('\n')
        
        # Extract title (first # header) 
        title_found = False
        content_lines = []
        
        for line in lines:
            if line.startswith('# ') and not title_found:
                title_text = line[2:].strip()
                slide_data["title"] = title_text
                slide_data["title_formatted"] = self._parse_inline_formatting(title_text)
                title_found = True
            elif line.startswith('## ') and slide_data["type"] == "title":
                subtitle_text = line[3:].strip()
                slide_data["subtitle"] = subtitle_text
                slide_data["subtitle_formatted"] = self._parse_inline_formatting(subtitle_text)
            else:
                content_lines.append(line)
        
        # Parse mixed content based on slide type
        if slide_data["type"] == "table":
            slide_data["table"] = self._parse_markdown_table('\n'.join(content_lines), config)
        elif slide_data["type"] != "title":  # Content slides get rich content
            rich_content = self._parse_rich_content('\n'.join(content_lines))
            if rich_content:
                slide_data["rich_content"] = rich_content
        
        return slide_data

    def _parse_rich_content(self, content: str) -> list:
        """Parse mixed markdown content into structured content blocks with better hierarchy"""
        blocks = []
        lines = content.split('\n')
        current_block = None
        
        for line in lines:
            original_line = line
            line = line.strip()
            if not line:
                continue
                
            # Handle nested bullet points by preserving indentation
            if line.startswith('- ') or line.startswith('* '):
                # Determine indentation level
                indent_level = len(original_line) - len(original_line.lstrip())
                bullet_text = line[2:].strip()
                
                if not current_block or "bullets" not in current_block:
                    if current_block:
                        blocks.append(current_block)
                    current_block = {"bullets": [], "bullet_levels": []}
                
                current_block["bullets"].append(bullet_text)
                # Map indentation to bullet levels (0 indent = level 1, 2+ spaces = level 2, etc.)
                level = 1 if indent_level < 2 else 2
                current_block["bullet_levels"].append(level)
                
            elif line.startswith('## '):  # Subheading
                if current_block:
                    blocks.append(current_block)
                current_block = {
                    "heading": line[3:].strip(),
                    "level": 2
                }
                
            elif line.startswith('### '):  # Sub-subheading
                if current_block:
                    blocks.append(current_block)
                current_block = {
                    "heading": line[4:].strip(),
                    "level": 3
                }
                
            else:  # Regular paragraph
                if not current_block or "paragraph" not in current_block:
                    if current_block:
                        blocks.append(current_block)
                    current_block = {"paragraph": line}
                else:
                    current_block["paragraph"] += " " + line
        
        if current_block:
            blocks.append(current_block)
            
        return blocks

    def _parse_markdown_table(self, content: str, config: dict) -> dict:
        """Extract table from markdown and apply styling config"""
        table_data = {
            "data": [],
            "header_style": config.get("style", "dark_blue_white_text"),
            "row_style": config.get("row_style", "alternating_light_gray"),
            "border_style": config.get("border_style", "thin_gray"),
            "custom_colors": config.get("custom_colors", {})
        }
        
        lines = [line.strip() for line in content.split('\n') if line.strip()]
        
        for line in lines:
            if line.startswith('|') and line.endswith('|'):
                # Parse table row with inline formatting
                cells = [cell.strip() for cell in line[1:-1].split('|')]
                formatted_cells = []
                for cell in cells:
                    formatted_cells.append({
                        "text": cell,
                        "formatted": self._parse_inline_formatting(cell)
                    })
                table_data["data"].append(formatted_cells)
            elif '|' in line and not line.startswith('|'):
                # Handle tables without outer pipes with inline formatting
                cells = [cell.strip() for cell in line.split('|')]
                formatted_cells = []
                for cell in cells:
                    formatted_cells.append({
                        "text": cell,
                        "formatted": self._parse_inline_formatting(cell)
                    })
                table_data["data"].append(formatted_cells)
            elif line.startswith('---') or line.startswith('==='):
                # Skip separator lines
                continue
        
        return table_data

    def _auto_parse_json_formatting(self, slide_data):
        """Auto-parse inline formatting in JSON slide data."""
        # Create a copy to avoid modifying original
        processed_data = slide_data.copy()
        
        # Parse title if present
        if "title" in processed_data and processed_data["title"]:
            title_text = processed_data["title"]
            processed_data["title_formatted"] = self._parse_inline_formatting(title_text)
        
        # Parse subtitle if present
        if "subtitle" in processed_data and processed_data["subtitle"]:
            subtitle_text = processed_data["subtitle"]
            processed_data["subtitle_formatted"] = self._parse_inline_formatting(subtitle_text)
        
        # Parse content list if present
        if "content" in processed_data and isinstance(processed_data["content"], list):
            # Convert simple content to rich content with formatting
            rich_content = []
            for item in processed_data["content"]:
                if isinstance(item, str):
                    # Treat as paragraph text
                    rich_content.append({
                        "paragraph": item
                    })
            processed_data["rich_content"] = rich_content
            # Remove old content key to avoid conflicts
            del processed_data["content"]
        
        # Parse table data if present
        if "table" in processed_data and "data" in processed_data["table"]:
            table_data = processed_data["table"]
            if isinstance(table_data["data"], list):
                formatted_data = []
                for row in table_data["data"]:
                    if isinstance(row, list):
                        formatted_row = []
                        for cell in row:
                            if isinstance(cell, str):
                                formatted_row.append({
                                    "text": cell,
                                    "formatted": self._parse_inline_formatting(cell)
                                })
                            else:
                                # Keep non-string cells as-is
                                formatted_row.append(cell)
                        formatted_data.append(formatted_row)
                    else:
                        # Keep non-list rows as-is
                        formatted_data.append(row)
                processed_data["table"]["data"] = formatted_data
        
        return processed_data

    def create_presentation_from_markdown(self, markdown_content: str, fileName: str = "Sample_Presentation", templateName: str = "default") -> str:
        """Create presentation from formatted markdown with frontmatter"""
        try:
            slides = self.parse_markdown_with_frontmatter(markdown_content)
            
            # Create presentation
            self.create_presentation(templateName, fileName)
            
            # Add all slides to the presentation
            for slide_data in slides:
                self._add_slide(slide_data)
            
            # Automatically save the presentation to disk after creation
            write_result = self.write_presentation(fileName)
                
            return f"Successfully created presentation with {len(slides)} slides from markdown. {write_result}"
        except Exception as e:
            return f"Error creating presentation from markdown: {str(e)}"

def get_deckbuilder_client():
    # Return singleton instance of Deckbuilder
    return Deckbuilder()