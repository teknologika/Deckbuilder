import os
import json
from typing import Dict, Optional
from pptx import Presentation


class TemplateAnalyzer:
    """Analyzes PowerPoint templates to extract raw layout and placeholder information."""
    
    def __init__(self):
        self.template_path = os.getenv('DECK_TEMPLATE_FOLDER')
        self.output_folder = os.getenv('DECK_OUTPUT_FOLDER')
    
    def analyze_pptx_template(self, template_name: str) -> Dict:
        """
        Analyze a PowerPoint template and extract raw layout information.
        
        Args:
            template_name: Name of the template file (with or without .pptx extension)
            
        Returns:
            Dictionary containing template structure with placeholder indices
        """
        # Ensure template name has .pptx extension
        if not template_name.endswith('.pptx'):
            template_name += '.pptx'
        
        # Build template path
        if not self.template_path:
            raise RuntimeError("DECK_TEMPLATE_FOLDER environment variable not set")
        
        template_path = os.path.join(self.template_path, template_name)
        
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template file not found: {template_path}")
        
        try:
            prs = Presentation(template_path)
            
            # Extract basic template info
            base_name = os.path.splitext(template_name)[0]
            template_info = {
                "name": base_name.replace("_", " ").title(),
                "version": "1.0"
            }
            
            # Extract raw layout data
            layouts = self._extract_layouts(prs)
            
            # Generate basic aliases structure (empty for user to fill)
            aliases = self._generate_aliases_template()
            
            result = {
                "template_info": template_info,
                "layouts": layouts,
                "aliases": aliases
            }
            
            # Save to output folder as .g.json
            self._save_json_mapping(base_name, result)
            
            return result
            
        except Exception as e:
            raise RuntimeError(f"Error analyzing template: {str(e)}")
    
    def _extract_layouts(self, presentation: Presentation) -> Dict:
        """Extract raw layout data from all slide layouts."""
        layouts = {}
        
        for idx, layout in enumerate(presentation.slide_layouts):
            layout_info = self._extract_single_layout(layout, idx)
            if layout_info:
                # Use generic layout names based on index
                layout_name = f"layout_{idx}"
                layouts[layout_name] = layout_info
        
        return layouts
    
    def _extract_single_layout(self, layout, index: int) -> Optional[Dict]:
        """
        Extract raw placeholder data from a single slide layout.
        
        Args:
            layout: PowerPoint slide layout object
            index: Layout index in the template
            
        Returns:
            Dictionary with raw layout information
        """
        try:
            placeholders = {}
            
            # Extract each placeholder's index
            for shape in layout.placeholders:
                placeholder_idx = shape.placeholder_format.idx
                # Leave placeholder value empty for user to fill in
                placeholders[str(placeholder_idx)] = f"placeholder_{placeholder_idx}"
            
            return {
                "index": index,
                "placeholders": placeholders
            }
            
        except Exception as e:
            print(f"Warning: Could not analyze layout {index}: {str(e)}")
            return None
    
    def _generate_aliases_template(self) -> Dict:
        """Generate empty aliases template for user configuration."""
        return {
            "table": "layout_0",
            "bullets": "layout_0"
        }
    
    def _save_json_mapping(self, template_name: str, data: Dict) -> None:
        """
        Save the JSON mapping to output folder as .g.json file.
        
        Args:
            template_name: Base name of the template (without extension)
            data: Dictionary to save as JSON
        """
        if not self.output_folder:
            print("Warning: DECK_OUTPUT_FOLDER not set, saving to current directory")
            output_folder = '.'
        else:
            output_folder = self.output_folder
        
        # Ensure output folder exists
        os.makedirs(output_folder, exist_ok=True)
        
        # Create output filename
        output_filename = f"{template_name}.g.json"
        output_path = os.path.join(output_folder, output_filename)
        
        # Save JSON (overwrite if exists)
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=2, ensure_ascii=False)
        
        print(f"Template mapping saved to: {output_path}")
        print(f"Rename to {template_name}.json when ready to use with deckbuilder")


def analyze_pptx_template(template_name: str) -> Dict:
    """
    Convenience function to analyze a PowerPoint template.
    
    Args:
        template_name: Name of the template file (with or without .pptx extension)
        
    Returns:
        Dictionary containing raw template structure for user mapping
    """
    analyzer = TemplateAnalyzer()
    return analyzer.analyze_pptx_template(template_name)


def test_with_default_template():
    """Test the analyzer with the default template."""
    try:
        result = analyze_pptx_template("default")
        print("Raw Template Structure:")
        print(json.dumps(result, indent=2))
        return result
    except Exception as e:
        print(f"Error analyzing template: {str(e)}")
        return None


if __name__ == "__main__":
    test_with_default_template()