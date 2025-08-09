import json
import os
import sys
from pathlib import Path
from typing import Dict, Optional

from pptx import Presentation

# Add the parent directory to Python path for imports
sys.path.insert(0, str(Path(__file__).parent.parent))
from deckbuilder.utils.path import path_manager  # noqa: E402


class TemplateAnalyzer:
    """Analyzes PowerPoint templates to extract raw layout and placeholder information."""

    def __init__(self):
        self.template_path = str(path_manager.get_template_folder())
        self.output_folder = str(path_manager.get_output_folder())

    def analyze_pptx_template(self, template_name: str) -> Dict:
        """
        Analyze a PowerPoint template and extract raw layout information.

        Args:
            template_name: Name of the template file (with or without .pptx extension)

        Returns:
            Dictionary containing template structure with placeholder indices
        """
        # Ensure template name has .pptx extension
        if not template_name.endswith(".pptx"):
            template_name += ".pptx"

        # Build template path using PathManager
        if not path_manager.validate_template_folder_exists():
            raise RuntimeError(f"Template folder not found: {self.template_path}")

        template_path = str(path_manager.get_template_file_path(template_name))

        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template file not found: {template_path}")

        try:
            prs = Presentation(template_path)

            # Extract basic template info
            base_name = os.path.splitext(template_name)[0]
            template_info = {"name": base_name.replace("_", " ").title(), "version": "1.0"}

            # Extract raw layout data
            layouts = self._extract_layouts(prs)

            # Validate template and generate warnings
            validation_results = self._validate_template(layouts)

            # Generate basic aliases structure (empty for user to fill)
            aliases = self._generate_aliases_template()

            result = {"template_info": template_info, "layouts": layouts, "aliases": aliases}

            # Add validation results
            if validation_results["warnings"] or validation_results["errors"]:
                result["validation"] = validation_results

            # Save to output folder as .g.json
            self._save_json_mapping(base_name, result)

            # Print validation results
            self._print_validation_results(validation_results)

            return result

        except Exception as e:
            raise RuntimeError(f"Error analyzing template: {str(e)}")

    def _extract_layouts(self, presentation: Presentation) -> Dict:
        """Extract raw layout data from all slide layouts."""
        layouts = {}

        for idx, layout in enumerate(presentation.slide_layouts):
            layout_info = self._extract_single_layout(layout, idx)
            if layout_info:
                # Use actual PowerPoint layout name
                layout_name = f"layout_{idx}"  # fallback
                try:
                    if hasattr(layout, "name") and layout.name:
                        layout_name = layout.name
                except Exception:  # nosec B110
                    # Layout name unavailable, keep fallback name
                    pass

                layouts[layout_name] = layout_info

        return layouts

    def _extract_single_layout(self, layout, index: int) -> Optional[Dict]:
        """
        Extract raw placeholder data from a single slide layout.

        Args:
            layout: PowerPoint slide layout object
            index: Layout index in the template

        Returns:
            Dictionary with raw layout information
        """
        try:
            placeholders = {}

            # Extract each placeholder's index and available properties
            for shape in layout.placeholders:
                placeholder_idx = shape.placeholder_format.idx

                # Try to get more information about the placeholder
                placeholder_info = f"placeholder_{placeholder_idx}"

                # Check if placeholder has a name or type information
                try:
                    if hasattr(shape, "name") and shape.name:
                        placeholder_info = shape.name
                    elif hasattr(shape.placeholder_format, "type"):
                        placeholder_type = shape.placeholder_format.type
                        placeholder_info = f"type_{placeholder_type}"
                except Exception:  # nosec B110
                    # Shape properties unavailable, keep default name
                    pass

                placeholders[str(placeholder_idx)] = placeholder_info

            return {"index": index, "placeholders": placeholders}

        except Exception as e:
            print(f"Warning: Could not analyze layout {index}: {str(e)}")
            return None

    def _generate_aliases_template(self) -> Dict:
        """Generate basic aliases template for user configuration."""
        return {
            "table": "Title and Content",
            "bullets": "Title and Content",
            "content": "Title and Content",
            "title": "Title Slide",
        }

    def _validate_template(self, layouts: Dict) -> Dict:
        """
        Validate template layouts and detect common issues.

        Args:
            layouts: Dictionary of layout definitions

        Returns:
            Dictionary containing validation results with warnings and errors
        """
        validation_results = {"warnings": [], "errors": [], "layout_analysis": {}}

        all_placeholder_names = []

        for layout_name, layout_info in layouts.items():
            placeholders = layout_info.get("placeholders", {})
            layout_warnings = []
            layout_errors = []

            # Check for duplicate placeholder names within layout
            placeholder_names = list(placeholders.values())
            unique_names = set(placeholder_names)

            if len(placeholder_names) != len(unique_names):
                for name in unique_names:
                    count = placeholder_names.count(name)
                    if count > 1:
                        layout_errors.append(f"Duplicate placeholder name '{name}' appears {count} times")

            # Check for column layouts with inconsistent naming
            if "column" in layout_name.lower():
                self._validate_column_layout(layout_name, placeholders, layout_warnings, layout_errors)

            # Check for comparison layouts
            if "comparison" in layout_name.lower():
                self._validate_comparison_layout(layout_name, placeholders, layout_warnings, layout_errors)

            # Track all placeholder names across layouts
            all_placeholder_names.extend(placeholder_names)

            # Store layout-specific validation results
            if layout_warnings or layout_errors:
                validation_results["layout_analysis"][layout_name] = {
                    "warnings": layout_warnings,
                    "errors": layout_errors,
                }
                validation_results["warnings"].extend([f"{layout_name}: {w}" for w in layout_warnings])
                validation_results["errors"].extend([f"{layout_name}: {e}" for e in layout_errors])

        # Global validation checks
        self._validate_global_consistency(all_placeholder_names, validation_results, layouts)

        return validation_results

    def _validate_column_layout(self, layout_name: str, placeholders: Dict, warnings: list, errors: list) -> None:
        """Validate column-based layouts for consistent naming patterns."""
        placeholder_names = list(placeholders.values())

        # Check for expected column patterns
        col_titles = [name for name in placeholder_names if "col" in name.lower() and "title" in name.lower()]
        col_contents = [name for name in placeholder_names if "col" in name.lower() and ("text" in name.lower() or "content" in name.lower())]

        # Extract column numbers from names and track specific placeholders that need fixing
        title_cols = []
        content_cols = []
        fix_suggestions = []

        for name in col_titles:
            try:
                # Extract number from names like "Col 1 Title" or "Col 2 Title Placeholder"
                parts = name.lower().split()
                for i, part in enumerate(parts):
                    if part == "col" and i + 1 < len(parts):
                        col_num = int(parts[i + 1])
                        title_cols.append((col_num, name))
                        break
            except (ValueError, IndexError):
                warnings.append(f"Could not parse column number from title placeholder: '{name}'")

        for name in col_contents:
            try:
                parts = name.lower().split()
                for i, part in enumerate(parts):
                    if part == "col" and i + 1 < len(parts):
                        col_num = int(parts[i + 1])
                        content_cols.append((col_num, name))
                        break
            except (ValueError, IndexError):
                warnings.append(f"Could not parse column number from content placeholder: '{name}'")

        # Check for consistent column numbering
        title_nums = sorted([col[0] for col in title_cols])
        content_nums = sorted([col[0] for col in content_cols])

        # For layouts with titles, check title/content pairs match
        if "title" in layout_name.lower() and title_cols:
            if title_nums != content_nums:
                # Find content placeholders that need fixing
                expected_content_nums = title_nums
                for expected_num in expected_content_nums:
                    if expected_num not in content_nums:
                        # Find what column number this content actually has
                        for actual_num, content_name in content_cols:
                            if actual_num != expected_num and expected_num not in [c[0] for c in content_cols]:
                                # This content placeholder has wrong number
                                correct_name = content_name.replace(f"Col {actual_num}", f"Col {expected_num}")
                                fix_suggestions.append(f"In PowerPoint: Rename '{content_name}' → '{correct_name}'")
                                break

                # Also check for duplicate column numbers in content
                content_num_counts = {}
                for num, name in content_cols:
                    if num not in content_num_counts:
                        content_num_counts[num] = []
                    content_num_counts[num].append(name)

                for num, names in content_num_counts.items():
                    if len(names) > 1:
                        # Multiple content placeholders have same column number
                        for i, name in enumerate(names[1:], start=2):  # Skip first one, fix the rest
                            # Find the next available column number
                            next_num = num + i - 1
                            while next_num in content_num_counts and next_num != num:
                                next_num += 1
                            if next_num <= len(title_nums):
                                correct_name = name.replace(f"Col {num}", f"Col {next_num}")
                                fix_suggestions.append(f"In PowerPoint: Rename '{name}' → '{correct_name}'")

                error_msg = f"Column title numbers {title_nums} don't match content numbers {content_nums}"
                if fix_suggestions:
                    error_msg += f". Required fixes in '{layout_name}' layout: {'; '.join(fix_suggestions)}"
                errors.append(error_msg)

            # Check for proper sequential numbering
            if title_nums and title_nums != list(range(1, len(title_nums) + 1)):
                expected_nums = list(range(1, len(title_nums) + 1))
                warnings.append(f"Column numbers not sequential: {title_nums} (expected: {expected_nums})")

        # Check for proper content numbering in content-only layouts
        elif content_cols:
            if content_nums != list(range(1, len(content_nums) + 1)):
                expected_nums = list(range(1, len(content_nums) + 1))
                warnings.append(f"Column numbers not sequential: {content_nums} (expected: {expected_nums})")

    def _validate_comparison_layout(self, layout_name: str, placeholders: Dict, warnings: list, errors: list) -> None:
        """Validate comparison layouts for proper left/right structure."""
        placeholder_names = list(placeholders.values())

        text_placeholders = [name for name in placeholder_names if "text" in name.lower() and "placeholder" in name.lower()]
        content_placeholders = [name for name in placeholder_names if "content" in name.lower() and "placeholder" in name.lower()]

        if len(text_placeholders) < 2:
            errors.append("Comparison layout should have at least 2 text placeholders for left/right titles")

        if len(content_placeholders) < 2:
            msg = "Comparison layout should have at least 2 content placeholders " "for left/right content"
            errors.append(msg)

    def _validate_global_consistency(self, all_placeholder_names: list, validation_results: Dict, layouts: Dict = None) -> None:
        """Validate global consistency across all layouts."""
        # Track patterns by layout for more specific reporting
        layout_patterns = {}
        unique_patterns = set()

        if layouts:
            for layout_name, layout_info in layouts.items():
                layout_patterns[layout_name] = set()
                placeholders = layout_info.get("placeholders", {})

                for name in placeholders.values():
                    if "placeholder" in str(name).lower():
                        # Extract pattern like "Text Placeholder", "Content Placeholder"
                        parts = str(name).split()
                        if len(parts) >= 2:
                            pattern = f"{parts[0]} {parts[1]}"
                            layout_patterns[layout_name].add(pattern)
                            unique_patterns.add(pattern)
        else:
            # Fallback to old method if layouts not provided
            for name in all_placeholder_names:
                if "placeholder" in name.lower():
                    parts = name.split()
                    if len(parts) >= 2:
                        pattern = f"{parts[0]} {parts[1]}"
                        unique_patterns.add(pattern)

        # Check for mixed naming conventions
        if len(unique_patterns) > 1:
            warning_msg = f"Multiple placeholder naming patterns detected: {sorted(unique_patterns)}"

            # Add specific layout information if available
            if layouts and layout_patterns:
                layouts_with_patterns = []
                for layout_name, patterns in layout_patterns.items():
                    if patterns:  # Only include layouts that have patterns
                        pattern_list = sorted(patterns)
                        layouts_with_patterns.append(f"{layout_name} ({', '.join(pattern_list)})")

                if layouts_with_patterns:
                    warning_msg += f". Affected layouts: {'; '.join(layouts_with_patterns)}"

            validation_results["warnings"].append(warning_msg)

    def _print_validation_results(self, validation_results: Dict) -> None:
        """Print validation results to console with formatting."""
        if not validation_results["warnings"] and not validation_results["errors"]:
            print("\n✓ Template validation passed - no issues detected")
            return

        print("\n" + "=" * 60)
        print("TEMPLATE VALIDATION RESULTS")
        print("=" * 60)

        if validation_results["errors"]:
            print(f"\n❌ ERRORS ({len(validation_results['errors'])}):")
            for i, error in enumerate(validation_results["errors"], 1):
                print(f"  {i}. {error}")

        if validation_results["warnings"]:
            print(f"\n⚠️  WARNINGS ({len(validation_results['warnings'])}):")
            for i, warning in enumerate(validation_results["warnings"], 1):
                print(f"  {i}. {warning}")

        print("\n" + "=" * 60)
        print("RECOMMENDED ACTIONS:")
        print("=" * 60)

        if validation_results["errors"]:
            print("\n🔧 Fix these errors in your PowerPoint template:")
            print("   • Open your PowerPoint template file")
            print("   • Open View > Slide Master to edit the template layouts")
            print("   • On Mac: Open Arrange > Selection Pane to see all placeholder objects")
            print("   • Select the placeholder objects and rename them in the Selection Pane")
            print("   • Rename placeholders as specified in the error messages above")
            print("   • Ensure column layouts have consistent numbering (Col 1, Col 2, etc.)")
            print("   • Verify comparison layouts have proper left/right structure")
            print("   • Close Slide Master view when finished")

        if validation_results["warnings"]:
            print("\n💡 Consider these improvements:")
            print("   • Use consistent placeholder naming patterns")
            print("   • Ensure column numbers are sequential (1, 2, 3, 4)")
            print("   • Follow naming conventions like 'Col 1 Title Placeholder 2'")

        print("\n📝 After fixing placeholder names in PowerPoint, regenerate the template mapping:")
        print("   python src/deckbuilder/cli_tools.py analyze default --verbose")
        print("\n💡 The analyzer will show ✅ validation passed when all issues are resolved")
        print("=" * 60)

    def _save_json_mapping(self, template_name: str, data: Dict) -> None:
        """
        Save the JSON mapping to output folder as .g.json file.

        Args:
            template_name: Base name of the template (without extension)
            data: Dictionary to save as JSON
        """
        if not self.output_folder:
            print("Warning: DECK_OUTPUT_FOLDER not set, saving to current directory")
            output_folder = "."
        else:
            output_folder = self.output_folder

        # Ensure output folder exists
        os.makedirs(output_folder, exist_ok=True)

        # Create output filename
        output_filename = f"{template_name}.g.json"
        output_path = os.path.join(output_folder, output_filename)

        # Save JSON (overwrite if exists)
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(data, f, indent=2, ensure_ascii=False)

        print(f"\nTemplate mapping saved to: {output_path}")
        print(f"Rename to {template_name}.json when ready to use with deckbuilder")


def analyze_pptx_template(template_name: str) -> Dict:
    """
    Convenience function to analyze a PowerPoint template.

    Args:
        template_name: Name of the template file (with or without .pptx extension)

    Returns:
        Dictionary containing raw template structure for user mapping
    """
    analyzer = TemplateAnalyzer()
    return analyzer.analyze_pptx_template(template_name)


def test_with_default_template():
    """Test the analyzer with the default template."""
    try:
        result = analyze_pptx_template("default")
        print("Raw Template Structure:")
        print(json.dumps(result, indent=2))
        return result
    except Exception as e:
        print(f"Error analyzing template: {str(e)}")
        return None


if __name__ == "__main__":
    test_with_default_template()
