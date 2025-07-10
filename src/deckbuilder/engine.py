# import json
import sys
from datetime import datetime
from pathlib import Path
from typing import Dict, Any

from pptx import Presentation

from .path_manager import path_manager, PathManager
from .presentation_builder import PresentationBuilder
from .content_processor import ContentProcessor
from .template_manager import TemplateManager
from .image_handler import ImageHandler

# Import PlaceKitten from parent directory
sys.path.insert(0, str(Path(__file__).parent.parent))
from placekitten import PlaceKitten  # noqa: E402


def singleton(cls):
    instances = {}

    def get_instance(*args, **kwargs):
        if cls not in instances:
            instances[cls] = cls(*args, **kwargs)
        return instances[cls]

    def reset():
        """Reset the singleton instance for testing purposes"""
        instances.clear()

    # Allow external access to clear instances for testing
    get_instance._instances = instances
    get_instance.reset = reset
    cls._instances = instances
    cls.reset = reset

    return get_instance


@singleton
class Deckbuilder:
    def __init__(self, path_manager_instance: PathManager = None):
        # Use provided path manager or default global instance
        self._path_manager = path_manager_instance or path_manager

        self.output_folder = str(self._path_manager.get_output_folder())
        self.prs = Presentation()

        # Initialize components
        self.template_manager = TemplateManager(self._path_manager)
        self.content_processor = ContentProcessor()
        self.presentation_builder = PresentationBuilder(self._path_manager)

        # Initialize image-related components
        self.image_handler = ImageHandler()
        self.placekitten = PlaceKitten()

        # Ensure default template exists in templates folder
        template_name = self._path_manager.get_template_name() or "default"
        self.template_manager.check_template_exists(template_name)

        # Store template info for tests
        self.template_name = template_name
        self.template_path, _ = self.template_manager.prepare_template(template_name)

    def _initialize_presentation(self, templateName: str = "default") -> None:
        # Prepare template and get path and layout mapping
        template_path, layout_mapping = self.template_manager.prepare_template(templateName)

        # Store template path for tests
        self.template_path = template_path

        # Update components with layout mapping
        self.content_processor.layout_mapping = layout_mapping
        self.presentation_builder.layout_mapping = layout_mapping

        # Load template or create empty presentation
        if template_path:
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()

        self.presentation_builder.clear_slides(self.prs)

    def create_presentation(
        self,
        presentation_data: Dict[str, Any],
        fileName: str = "Sample_Presentation",
        templateName: str = "default",
    ) -> str:
        """
        Creates a presentation from the canonical JSON data model.
        Only accepts canonical format: {"slides": [{"layout": "...", "placeholders": {...}, "content": [...]}]}

        Includes built-in end-to-end validation to prevent layout regressions.
        """
        # Import validation here to avoid circular imports
        from .validation import PresentationValidator

        self._initialize_presentation(templateName)

        # Strict validation for canonical JSON format only
        if not isinstance(presentation_data, dict):
            raise ValueError("Input must be a dictionary containing canonical JSON data.")

        if "slides" not in presentation_data:
            raise ValueError("Canonical JSON data must contain a 'slides' array at root level.")

        if not isinstance(presentation_data["slides"], list):
            raise ValueError("'slides' must be an array of slide objects.")

        if len(presentation_data["slides"]) == 0:
            raise ValueError("At least one slide is required.")

        # Validate each slide has required canonical structure
        for i, slide_data in enumerate(presentation_data["slides"]):
            if not isinstance(slide_data, dict):
                raise ValueError(f"Slide {i + 1} must be a dictionary.")

            if "layout" not in slide_data:
                raise ValueError(f"Slide {i + 1} must have a 'layout' field.")

            # Ensure canonical structure exists (placeholders and content are optional but must be correct types)
            if "placeholders" in slide_data and not isinstance(slide_data["placeholders"], dict):
                raise ValueError(f"Slide {i + 1} 'placeholders' must be a dictionary.")

            if "content" in slide_data and not isinstance(slide_data["content"], list):
                raise ValueError(f"Slide {i + 1} 'content' must be an array.")

        # STEP 1: Pre-generation validation (JSON ↔ Template alignment)
        template_folder = str(self._path_manager.get_template_folder())
        validator = PresentationValidator(presentation_data, templateName, template_folder)
        validator.validate_pre_generation()

        # STEP 2: Process slides using canonical format
        for slide_data in presentation_data["slides"]:
            self.presentation_builder.add_slide(self.prs, slide_data)

        # STEP 3: Save the presentation to disk
        write_result = self.write_presentation(fileName)

        # Extract the file path from write_result for post-generation validation
        # write_result format: "Successfully created presentation: filename.pptx"
        if "Successfully created presentation:" in write_result:
            file_path = write_result.split("Successfully created presentation: ")[1].strip()
            full_path = str(self._path_manager.get_output_folder() / file_path)

            # STEP 4: Post-generation validation (PPTX ↔ JSON verification)
            validator.validate_post_generation(full_path)

        # Show completion summary
        from .logging_config import success_print

        slide_count = len(presentation_data["slides"])
        file_name = write_result.split("Successfully created presentation: ")[1].strip() if "Successfully created presentation:" in write_result else "presentation.pptx"
        success_print(f"✅ Presentation complete: {file_name} ({slide_count} slides)")

        return f"Successfully created presentation with {slide_count} slides. {write_result}"

    def write_presentation(self, fileName: str = "Sample_Presentation") -> str:
        """Writes the generated presentation to disk with ISO timestamp."""
        import os

        # Get output folder from environment or use default
        output_folder = self.output_folder or "."

        # Ensure output folder exists
        os.makedirs(output_folder, exist_ok=True)

        # Create filename with ISO timestamp and .g.pptx extension for generated files
        timestamp = datetime.now().strftime("%Y-%m-%d_%H%M")
        generated_filename = f"{fileName}.{timestamp}.g.pptx"
        output_file = os.path.join(output_folder, generated_filename)

        # Save the presentation (overwrites if same timestamp exists)
        self.prs.save(output_file)

        return f"Successfully created presentation: {os.path.basename(output_file)}"


def get_deckbuilder_client():
    # Return Deckbuilder instance with MCP context
    from .path_manager import create_mcp_path_manager

    return Deckbuilder(path_manager_instance=create_mcp_path_manager())
