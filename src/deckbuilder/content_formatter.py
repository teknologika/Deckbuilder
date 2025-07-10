import re

try:
    from .formatting_support import FormattingSupport, get_default_language, get_default_font
except ImportError:
    # Fallback if formatting support is not available
    class FormattingSupport:
        def apply_language_to_run(self, run, language_code):
            return False

        def apply_font_to_run(self, run, font_name):
            return False

    def get_default_language():
        return None

    def get_default_font():
        return None


class ContentFormatter:
    """Handles text formatting, rich content processing, and inline formatting."""

    def __init__(self):
        """Initialize the content formatter."""
        # Initialize formatting support
        self.formatting_support = FormattingSupport()
        self.default_language = get_default_language()
        self.default_font = get_default_font()

    def add_simple_content_to_placeholder(self, placeholder, content):
        """Add content to a placeholder with support for rich content blocks and formatted lists."""
        if not hasattr(placeholder, "text_frame"):
            return

        text_frame = placeholder.text_frame
        text_frame.clear()

        # Debug logging to track content processing pipeline decisions
        content_type = type(content).__name__
        self._debug_log(f"Processing content type: {content_type}")

        # Priority 1: Check for rich content blocks (from content_formatting.py)
        if isinstance(content, dict):
            # Check for formatted content with rich content blocks first
            if "rich_content_blocks" in content:
                self._debug_log("Processing rich content blocks from content_formatting")
                self._add_rich_content_blocks_to_placeholder(text_frame, content["rich_content_blocks"])
                return
            elif "formatted_list" in content:
                self._debug_log("Processing formatted_list from content_formatting")
                # This case means content is a dict like {'formatted_list': [...]}
                # We need to pass the list itself to the handler
                self._add_rich_content_list_to_placeholder(text_frame, content["formatted_list"])
                return
            elif any(key in content for key in ["heading", "paragraph", "bullets"]):
                self._debug_log("Processing direct rich content structure (single block)")
                # Wrap single rich content block in a list for consistent handling
                self._add_rich_content_list_to_placeholder(text_frame, [content])
                return
            elif "text" in content and "formatted" in content:
                self._debug_log("Processing formatted content segments (single text field)")
                formatted_segments = content["formatted"]
                p = text_frame.paragraphs[0]
                self.apply_formatted_segments_to_paragraph(formatted_segments, p)
                return
            elif content.get("type") == "table":
                self._debug_log("Processing table content")
                # Handle table content using TableBuilder
                slide = getattr(self, "_current_slide", None)
                if slide:
                    from .table_builder import TableBuilder

                    table_builder = TableBuilder(self)
                    table_builder.add_table_to_slide(slide, content)
                    # Remove the placeholder since table has been added
                    self._remove_placeholder_from_slide(slide, placeholder)
                else:
                    self._debug_log("No slide context for table, using text fallback")
                    p = text_frame.paragraphs[0]
                    p.text = "Table content (could not render)"
                return
            else:
                # Fallback for other dict types - avoid string conversion if possible
                self._debug_log(f"Unknown dict content structure, attempting text extraction: {list(content.keys())}")
                p = text_frame.paragraphs[0]
                if "text" in content:
                    self.apply_inline_formatting(content["text"], p)
                else:
                    p.text = str(content)  # Final fallback to string conversion
                return

        # Priority 2: Check for list of content (can be rich content blocks or simple strings)
        elif isinstance(content, list):
            if content and isinstance(content[0], dict):
                # Check if this is a list of formatted segments (e.g., from parse_inline_formatting)
                if "text" in content[0] and "format" in content[0]:
                    self._debug_log("Processing list of formatted segments")
                    p = text_frame.paragraphs[0]
                    self.apply_formatted_segments_to_paragraph(content, p)
                    return
                # Check if this is a list of rich content blocks
                elif any(key in content[0] for key in ["heading", "paragraph", "bullets"]):
                    self._debug_log("Processing list of rich content blocks")
                    self._add_rich_content_list_to_placeholder(text_frame, content)
                    return
                else:
                    # Other dict list - treat as rich content list for now
                    self._debug_log("Processing dict list as rich content (general)")
                    self._add_rich_content_list_to_placeholder(text_frame, content)
                    return
            else:
                # Handle list of simple strings
                self._debug_log("Processing simple string list")
                self._add_rich_content_list_to_placeholder(text_frame, content)
                return

        # Priority 3: Handle plain text strings
        elif isinstance(content, str):
            self._debug_log("Processing plain text string")
            p = text_frame.paragraphs[0]
            self.apply_inline_formatting(content, p)
            return

        # Fallback for unexpected content types
        else:
            self._debug_log(f"Fallback: Converting {content_type} to string")
            p = text_frame.paragraphs[0]
            p.text = str(content)

    def _debug_log(self, message):
        """Debug logging for content processing pipeline"""
        from .logging_config import content_processor_print

        content_processor_print(f"[ContentFormatter] {message}")

    def _add_rich_content_list_to_placeholder(self, text_frame, content_list):
        """Add list content with proper formatting and bullet support."""
        paragraph_added = False

        for item in content_list:
            if isinstance(item, str):
                # Simple string item - apply inline formatting
                if not paragraph_added:
                    p = text_frame.paragraphs[0]
                    paragraph_added = True
                else:
                    p = text_frame.add_paragraph()
                self.apply_inline_formatting(item, p)

            elif isinstance(item, dict):
                # Check if this is a rich content block (heading, paragraph, bullets)
                if any(key in item for key in ["heading", "paragraph", "bullets"]):
                    self._debug_log(f"Processing rich content block in list: {list(item.keys())}")
                    # Process as rich content block using the specialized handler
                    self._add_single_rich_content_block_to_placeholder(text_frame, item, paragraph_added)
                    paragraph_added = True
                elif "text" in item:
                    # Simple text item
                    if not paragraph_added:
                        p = text_frame.paragraphs[0]
                        paragraph_added = True
                    else:
                        p = text_frame.add_paragraph()
                    self.apply_inline_formatting(item["text"], p)
                elif "formatted" in item:
                    # Apply formatted content segments
                    if not paragraph_added:
                        p = text_frame.paragraphs[0]
                        paragraph_added = True
                    else:
                        p = text_frame.add_paragraph()
                    self.apply_formatted_segments_to_paragraph(item["formatted"], p)
                else:
                    # Unknown dict structure - extract text if possible
                    self._debug_log(f"Unknown dict in list, keys: {list(item.keys())}")
                    if not paragraph_added:
                        p = text_frame.paragraphs[0]
                        paragraph_added = True
                    else:
                        p = text_frame.add_paragraph()
                    # Try to find any text content to avoid string conversion
                    text_content = item.get("text", str(item))
                    self.apply_inline_formatting(text_content, p)

    def _add_single_rich_content_block_to_placeholder(self, text_frame, content_block, paragraph_added):
        """Add a single rich content block (heading, paragraph, or bullets) to placeholder."""
        # Handle heading
        if "heading" in content_block:
            p = text_frame.paragraphs[0] if not paragraph_added else text_frame.add_paragraph()
            heading_text = content_block["heading"]
            self.apply_inline_formatting(heading_text, p)
            # Make heading bold by default
            for run in p.runs:
                run.font.bold = True
            self._debug_log(f"Added heading: '{heading_text}'")

        # Handle paragraph
        if "paragraph" in content_block:
            p = text_frame.paragraphs[0] if not paragraph_added else text_frame.add_paragraph()
            paragraph_text = content_block["paragraph"]
            self.apply_inline_formatting(paragraph_text, p)
            self._debug_log(f"Added paragraph: '{paragraph_text[:50]}...'")

        # Handle bullets with proper level support
        if "bullets" in content_block and isinstance(content_block["bullets"], list):
            bullet_levels = content_block.get("bullet_levels", [])
            for i, bullet_text in enumerate(content_block["bullets"]):
                p = text_frame.paragraphs[0] if not paragraph_added else text_frame.add_paragraph()
                self.apply_inline_formatting(bullet_text, p)

                # Set bullet level from bullet_levels array or default to level 1
                if i < len(bullet_levels):
                    # Convert level (1-based) to PowerPoint level (0-based)
                    p.level = max(0, bullet_levels[i] - 1)
                else:
                    p.level = 0  # Default to top level bullets

                self._debug_log(f"Added bullet: '{bullet_text}' at level {p.level}")
                paragraph_added = True

    def _add_rich_content_blocks_to_placeholder(self, text_frame, content_dict):
        """Add rich content blocks (headings, paragraphs, bullets) to placeholder."""
        # Check if this is formatted content from content_formatting.py
        if "formatted_list" in content_dict:
            # Handle formatted list content
            formatted_list = content_dict["formatted_list"]
            for i, item in enumerate(formatted_list):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                if "text" in item and "formatted" in item:
                    # Apply formatted segments
                    self.apply_formatted_segments_to_paragraph(item["formatted"], p)
                else:
                    # Fallback to text
                    text = item.get("text", str(item))
                    self.apply_inline_formatting(text, p)
        elif "text" in content_dict and "formatted" in content_dict:
            # Handle single formatted field content like {'text': '...', 'formatted': [...]}
            p = text_frame.paragraphs[0]
            formatted_segments = content_dict["formatted"]
            self.apply_formatted_segments_to_paragraph(formatted_segments, p)
        else:
            # Handle direct rich content blocks
            paragraph_added = False

            # Handle heading
            if "heading" in content_dict:
                p = text_frame.paragraphs[0] if not paragraph_added else text_frame.add_paragraph()
                heading_text = content_dict["heading"]
                self.apply_inline_formatting(heading_text, p)
                # Make heading bold by default
                for run in p.runs:
                    run.font.bold = True
                paragraph_added = True

            # Handle paragraph
            if "paragraph" in content_dict:
                p = text_frame.paragraphs[0] if not paragraph_added else text_frame.add_paragraph()
                paragraph_text = content_dict["paragraph"]
                self.apply_inline_formatting(paragraph_text, p)
                paragraph_added = True

            # Handle bullets with proper level support
            if "bullets" in content_dict and isinstance(content_dict["bullets"], list):
                bullet_levels = content_dict.get("bullet_levels", [])
                for i, bullet_text in enumerate(content_dict["bullets"]):
                    p = text_frame.paragraphs[0] if not paragraph_added else text_frame.add_paragraph()
                    self.apply_inline_formatting(bullet_text, p)

                    # Set bullet level from bullet_levels array or default to level 1
                    if i < len(bullet_levels):
                        # Convert level (1-based) to PowerPoint level (0-based)
                        p.level = max(0, bullet_levels[i] - 1)
                    else:
                        p.level = 0  # Default to top level bullets

                    self._debug_log(f"Added bullet: '{bullet_text}' at level {p.level}")
                    paragraph_added = True

    def apply_formatted_segments_to_paragraph(self, formatted_segments, paragraph):
        """Apply formatted text segments to a paragraph."""
        if not formatted_segments:
            return

        # Clear existing runs
        paragraph.clear()

        for segment in formatted_segments:
            if isinstance(segment, dict) and "text" in segment:
                text = segment["text"]
                format_info = segment.get("format", {})

                run = paragraph.add_run()
                run.text = text

                # Apply formatting
                if format_info.get("bold"):
                    run.font.bold = True
                if format_info.get("italic"):
                    run.font.italic = True
                if format_info.get("underline"):
                    run.font.underline = True

    def parse_inline_formatting(self, text):
        """Parse inline formatting and return structured formatting data"""
        if not text:
            return [{"text": "", "format": {}}]

        # Patterns in order of precedence (longest patterns first to avoid conflicts)
        patterns = [
            (
                r"\*\*\*___(.*?)___\*\*\*",
                {"bold": True, "italic": True, "underline": True},
            ),  # ***___text___***
            (
                r"___\*\*\*(.*?)\*\*\*___",
                {"bold": True, "italic": True, "underline": True},
            ),  # ___***text***___
            (r"\*\*\*(.*?)\*\*\*", {"bold": True, "italic": True}),  # ***text***
            (r"___(.*?)___", {"underline": True}),  # ___text___
            (r"\*\*(.*?)\*\*", {"bold": True}),  # **text**
            (r"\*(.*?)\*", {"italic": True}),  # *text*
        ]

        # Find all matches and their positions
        all_matches = []
        for pattern, format_dict in patterns:
            for match in re.finditer(pattern, text):
                all_matches.append((match.start(), match.end(), match.group(1), format_dict))

        # Sort matches by position
        all_matches.sort(key=lambda x: x[0])

        # Remove overlapping matches (keep the first one found)
        filtered_matches = []
        last_end = 0
        for start, end, content, format_dict in all_matches:
            if start >= last_end:
                filtered_matches.append((start, end, content, format_dict))
                last_end = end

        # Build the formatted text segments
        segments = []
        last_pos = 0

        for start, end, content, format_dict in filtered_matches:
            # Add plain text before the formatted text
            if start > last_pos:
                plain_text = text[last_pos:start]
                if plain_text:
                    segments.append({"text": plain_text, "format": {}})

            # Add formatted text
            segments.append({"text": content, "format": format_dict})
            last_pos = end

        # Add any remaining plain text
        if last_pos < len(text):
            remaining_text = text[last_pos:]
            if remaining_text:
                segments.append({"text": remaining_text, "format": {}})

        # If no formatting found, return the original text
        if not segments:
            segments = [{"text": text, "format": {}}]

        return segments

    def apply_inline_formatting(self, text, paragraph):
        """Apply inline formatting to paragraph using parsed formatting data."""
        # Clear any existing text
        paragraph.text = ""

        # Parse the formatting
        segments = self.parse_inline_formatting(text)

        # Apply each segment to the paragraph
        for segment in segments:
            run = paragraph.add_run()
            run.text = segment["text"]

            # Apply formatting
            format_dict = segment["format"]
            if format_dict.get("bold"):
                run.font.bold = True
            if format_dict.get("italic"):
                run.font.italic = True
            if format_dict.get("underline"):
                run.font.underline = True

            # Apply default language and font settings
            if self.default_language:
                self.formatting_support.apply_language_to_run(run, self.default_language)
            if self.default_font:
                self.formatting_support.apply_font_to_run(run, self.default_font)

    def apply_formatted_segments_to_shape(self, shape, segments):
        """Apply formatted text segments to a shape's text frame."""
        if not hasattr(shape, "text_frame"):
            # For shapes that don't have text_frame, fall back to simple text
            shape.text = "".join(segment["text"] for segment in segments)
            return

        text_frame = shape.text_frame
        text_frame.clear()

        # Use the first paragraph or create one
        if text_frame.paragraphs:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()

        paragraph.text = ""

        # Apply each segment
        for segment in segments:
            run = paragraph.add_run()
            run.text = segment["text"]

            # Apply formatting
            format_dict = segment["format"]
            if format_dict.get("bold"):
                run.font.bold = True
            if format_dict.get("italic"):
                run.font.italic = True
            if format_dict.get("underline"):
                run.font.underline = True

    def apply_formatted_segments_to_cell(self, cell, segments):
        """Apply formatted text segments to a table cell."""
        text_frame = cell.text_frame
        text_frame.clear()

        # Create first paragraph
        paragraph = text_frame.paragraphs[0]
        paragraph.text = ""

        # Apply each segment
        for segment in segments:
            run = paragraph.add_run()
            run.text = segment["text"]

            # Apply formatting
            format_dict = segment["format"]
            if format_dict.get("bold"):
                run.font.bold = True
            if format_dict.get("italic"):
                run.font.italic = True
            if format_dict.get("underline"):
                run.font.underline = True

    def add_rich_content_to_slide(self, slide, rich_content: list):
        """DEPRECATED: Use add_content_to_slide instead"""
        print("Warning: add_rich_content_to_slide is deprecated, using unified add_content_to_slide")

        # Convert legacy rich content format to canonical format
        canonical_content = []
        for block in rich_content:
            if isinstance(block, dict):
                if "heading" in block:
                    canonical_content.append(
                        {
                            "type": "heading",
                            "text": block["heading"],
                            "level": block.get("level", 2),
                        }
                    )
                elif "paragraph" in block:
                    canonical_content.append({"type": "paragraph", "text": block["paragraph"]})
                elif "bullets" in block:
                    # Convert bullets format
                    items = []
                    bullets = block["bullets"]
                    bullet_levels = block.get("bullet_levels", [1] * len(bullets))

                    for bullet, level in zip(bullets, bullet_levels):
                        items.append({"text": bullet, "level": level})

                    canonical_content.append({"type": "bullets", "items": items})

        self.add_content_to_slide(slide, canonical_content)

    def add_content_to_slide(self, slide, content):
        """
        Unified content processing for canonical JSON content blocks only.

        Uses semantic placeholder detection instead of hardcoded indices.
        Expects all content to be in canonical JSON format: [{"type": "...", "text": "..."}, ...]
        """
        # Store slide context for table creation
        self._current_slide = slide

        from .placeholder_types import is_content_placeholder, is_title_placeholder

        self._debug_log(f"Processing {type(content).__name__} content with {len(content) if isinstance(content, list) else 'N/A'} blocks")

        # Enhanced content analysis debugging
        if isinstance(content, list):
            for i, block in enumerate(content):
                if isinstance(block, dict):
                    block_type = block.get("type", "unknown")
                    block_keys = list(block.keys())
                    self._debug_log(f"  Block {i + 1}: type='{block_type}', keys={block_keys}")
                else:
                    self._debug_log(f"  Block {i + 1}: {type(block).__name__} = {str(block)[:50]}...")
        elif isinstance(content, dict):
            content_keys = list(content.keys())
            self._debug_log(f"  Content dict keys: {content_keys}")
        else:
            self._debug_log(f"  Content: {str(content)[:100]}...")

        # Find title placeholder to check if it needs content extracted
        title_placeholder = None
        for shape in slide.placeholders:
            if is_title_placeholder(shape.placeholder_format.type):
                if hasattr(shape, "text_frame") and shape.text_frame is not None:
                    title_placeholder = shape
                    break

        # Check if title placeholder is empty and needs first heading extracted
        title_needs_content = False
        if title_placeholder:
            has_title_content = False
            if title_placeholder.text_frame.paragraphs:
                for paragraph in title_placeholder.text_frame.paragraphs:
                    if paragraph.text.strip():
                        has_title_content = True
                        break
            title_needs_content = not has_title_content
            self._debug_log(f"Title placeholder analysis: empty={title_needs_content}, needs_content={title_needs_content}")

        # Extract first heading to title if needed
        processed_content = content
        if title_needs_content and isinstance(content, list) and content:
            # Look for first heading in content blocks
            first_heading_idx = None
            for i, item in enumerate(content):
                if isinstance(item, dict) and item.get("type") == "heading":
                    first_heading_idx = i
                    break

            if first_heading_idx is not None:
                heading_block = content[first_heading_idx]
                heading_text = heading_block.get("text", "")
                self._debug_log(f"Extracting first heading to title: '{heading_text}'")

                # Add heading text to title placeholder
                p = title_placeholder.text_frame.paragraphs[0]
                self.apply_inline_formatting(heading_text, p)

                # Remove the heading from content blocks
                processed_content = content[:first_heading_idx] + content[first_heading_idx + 1 :]
                self._debug_log(f"Content blocks remaining after title extraction: {len(processed_content)}")

        # Find content placeholders using semantic detection
        content_placeholders = []
        self._debug_log("Analyzing slide placeholders for content placement:")

        for shape in slide.placeholders:
            placeholder_type = shape.placeholder_format.type
            try:
                placeholder_name = getattr(shape.element.nvSpPr.cNvPr, "name", "unnamed")
            except AttributeError:
                placeholder_name = "unnamed"
            type_name = placeholder_type.name if hasattr(placeholder_type, "name") else str(placeholder_type)

            self._debug_log(f"  Placeholder {shape.placeholder_format.idx}: {type_name} ('{placeholder_name}')")

            if is_content_placeholder(placeholder_type):
                # Skip if converted to image placeholder
                if hasattr(shape, "text_frame") and shape.text_frame is not None:
                    content_placeholders.append(shape)
                    self._debug_log("    -> SELECTED for content placement")
                else:
                    self._debug_log("    -> SKIPPED (no text_frame)")
            else:
                self._debug_log("    -> SKIPPED (not content type)")

        self._debug_log(f"Content placeholder summary: {len(content_placeholders)} selected")

        if not content_placeholders:
            self._debug_log("No content placeholders found - trying fallback options")
            # For slides without content placeholders (e.g., Title Slide), try to use other placeholder types
            from .placeholder_types import is_subtitle_placeholder

            subtitle_placeholders = []
            for shape in slide.placeholders:
                if is_subtitle_placeholder(shape.placeholder_format.type):
                    if hasattr(shape, "text_frame") and shape.text_frame is not None:
                        # Check if subtitle placeholder already has content (from explicit subtitle field)
                        has_existing_content = False
                        if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
                            for paragraph in shape.text_frame.paragraphs:
                                if paragraph.text.strip():
                                    has_existing_content = True
                                    break

                        if not has_existing_content:
                            subtitle_placeholders.append(shape)
                            self._debug_log(f"  Fallback option: subtitle placeholder {shape.placeholder_format.idx}")
                        else:
                            self._debug_log(f"  Subtitle {shape.placeholder_format.idx} has existing content, skipping")

            if subtitle_placeholders:
                self._debug_log("Using subtitle placeholders as content fallback")
                content_placeholders = subtitle_placeholders
            else:
                self._debug_log("CRITICAL: No suitable placeholders found for content - content will be lost!")
                return

        # Expect only canonical JSON content blocks
        if isinstance(processed_content, list) and processed_content:
            # Validate all items are canonical JSON blocks
            for i, item in enumerate(processed_content):
                if not isinstance(item, dict) or "type" not in item:
                    raise ValueError(f"Content item {i} must be canonical JSON block with 'type' field. Got: {item}")

            self._debug_log("Starting canonical content block processing")
            self._process_canonical_content_blocks(content_placeholders, processed_content)
        elif processed_content:
            self._debug_log(f"WARNING: Expected list of canonical JSON content blocks, got: {type(processed_content)}")
            self._debug_log(f"Content value: {str(processed_content)[:200]}...")
        else:
            self._debug_log("No content blocks to process")

    def _process_canonical_content_blocks(self, content_placeholders, content_blocks):
        """Process canonical JSON content blocks like {"type": "paragraph", "text": "..."}"""
        self._debug_log(f"Processing {len(content_blocks)} canonical content blocks")

        # Show content block types for debugging
        block_types = [block.get("type", "unknown") for block in content_blocks]
        self._debug_log(f"Block types: {block_types}")

        # Use first content placeholder for now
        # TODO: Support multi-placeholder layouts (Two Content, Four Columns)
        placeholder = content_placeholders[0]
        text_frame = placeholder.text_frame
        text_frame.clear()

        try:
            placeholder_name = getattr(placeholder.element.nvSpPr.cNvPr, "name", "unnamed")
        except AttributeError:
            placeholder_name = "unnamed"
        self._debug_log(f"Using placeholder {placeholder.placeholder_format.idx} ('{placeholder_name}') for content")

        first_block = True
        for i, block in enumerate(content_blocks):
            block_type = block.get("type", "")
            block_text = block.get("text", "")
            self._debug_log(f"  Block {i + 1}/{len(content_blocks)}: type='{block_type}', text_length={len(block_text)}")

            if block_type == "paragraph":
                if first_block:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                text_content = block.get("text", "")
                # print(f"    Paragraph: {text_content[:50]}...")
                self.apply_inline_formatting(text_content, p)

            elif block_type == "heading":
                if first_block:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                text_content = block.get("text", "")
                self._debug_log(f"    Adding heading: '{text_content[:50]}{'...' if len(text_content) > 50 else ''}'")
                self.apply_inline_formatting(text_content, p)
                # Make headings bold
                for run in p.runs:
                    run.font.bold = True

            elif block_type == "bullets":
                items = block.get("items", [])
                self._debug_log(f"    Adding bullet list: {len(items)} items")
                for item in items:
                    if first_block:
                        p = text_frame.paragraphs[0]
                        first_block = False
                    else:
                        p = text_frame.add_paragraph()

                    # Handle both simple and complex bullet items
                    if isinstance(item, dict):
                        text = item.get("text", "")
                        level = item.get("level", 1) - 1  # Convert to 0-based
                    else:
                        text = str(item)
                        level = 0

                    p.level = level
                    # print(f"      - {text[:30]}...")
                    self.apply_inline_formatting(text, p)

            elif block_type == "columns":
                columns = block.get("columns", [])
                self._debug_log(f"    Processing multi-column layout: {len(columns)} columns")

                # Distribute columns across available content placeholders
                for col_idx, column in enumerate(columns):
                    if col_idx < len(content_placeholders):
                        placeholder = content_placeholders[col_idx]
                        text_frame = placeholder.text_frame
                        text_frame.clear()

                        self._debug_log(f"      Column {col_idx + 1} -> placeholder {placeholder.placeholder_format.idx}")

                        # Process column content
                        column_content = column.get("content", [])
                        first_item = True
                        for content_item in column_content:
                            if isinstance(content_item, dict) and "type" in content_item:
                                item_type = content_item.get("type", "")
                                item_text = content_item.get("text", "")

                                if first_item:
                                    p = text_frame.paragraphs[0]
                                    first_item = False
                                else:
                                    p = text_frame.add_paragraph()

                                # self._debug_log(f"        {item_type}: {item_text[:30]}...")
                                self.apply_inline_formatting(item_text, p)

                                # Apply type-specific formatting
                                if item_type == "heading":
                                    for run in p.runs:
                                        run.font.bold = True
                    else:
                        self._debug_log(f"      Column {col_idx + 1}: no available placeholder (only {len(content_placeholders)} total)")

                # Mark that we've processed multi-column content, don't process individual blocks
                return

            elif block_type == "table":
                self._debug_log("    Processing table content block")
                # Get slide from content placeholders context - we know the slide from the calling context
                slide = getattr(self, "_current_slide", None)
                if slide:
                    self._create_table_in_slide(slide, placeholder, block)
                else:
                    self._debug_log("      No slide context available, using text fallback")
                    self._fallback_table_as_text(placeholder, block)

            else:
                # Unknown block type - treat as paragraph
                self._debug_log(f"    Unknown block type '{block_type}', treating as paragraph")
                if first_block:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                text = str(block.get("text", block))
                # print(f"      Text: {text[:50]}...")
                self.apply_inline_formatting(text, p)

            first_block = False

    def _process_simple_content_list(self, placeholder, content_list):
        """Process simple content list (legacy format)"""
        text_frame = placeholder.text_frame
        text_frame.clear()

        for i, line in enumerate(content_list):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            self.apply_inline_formatting(str(line), p)

    def _process_simple_content_string(self, placeholder, content_string):
        """Process simple string content"""
        text_frame = placeholder.text_frame
        text_frame.clear()

        p = text_frame.paragraphs[0]
        self.apply_inline_formatting(content_string, p)

    def add_simple_content_to_slide(self, slide, content):
        """DEPRECATED: Use add_content_to_slide instead"""
        print("Warning: add_simple_content_to_slide is deprecated, using unified add_content_to_slide")
        self.add_content_to_slide(slide, content)

    def auto_parse_json_formatting(self, slide_data):
        """Auto-parse inline formatting in JSON slide data."""
        # Type validation: ensure slide_data is a dictionary
        if not isinstance(slide_data, dict):
            raise TypeError(f"slide_data must be a dictionary, got {type(slide_data).__name__}: {slide_data}")

        # Create a copy to avoid modifying original
        processed_data = slide_data.copy()

        # Parse title if present
        if "title" in processed_data and processed_data["title"]:
            title_text = processed_data["title"]
            processed_data["title_formatted"] = self.parse_inline_formatting(title_text)

        # Parse subtitle if present
        if "subtitle" in processed_data and processed_data["subtitle"]:
            subtitle_text = processed_data["subtitle"]
            processed_data["subtitle_formatted"] = self.parse_inline_formatting(subtitle_text)

        # Parse content list - convert legacy strings to canonical JSON format
        if "content" in processed_data and isinstance(processed_data["content"], list):
            content_list = processed_data["content"]
            canonical_content = []

            for item in content_list:
                if isinstance(item, str):
                    # Convert legacy string to canonical paragraph block
                    canonical_content.append({"type": "paragraph", "text": item})
                    print(
                        # "Converted string to paragraph"
                    )
                elif isinstance(item, dict) and "type" in item:
                    # Already canonical JSON format - keep as is
                    canonical_content.append(item)
                    print(
                        # f"Preserved {item.get('type')} block"
                    )
                else:
                    # Unknown format - convert to paragraph
                    canonical_content.append({"type": "paragraph", "text": str(item)})
                    print(
                        # f"Converted unknown item to paragraph: {item}"
                    )

            # Replace with canonical format
            processed_data["content"] = canonical_content

        # Parse table data if present
        if "table" in processed_data and "data" in processed_data["table"]:
            table_data = processed_data["table"]
            if isinstance(table_data["data"], list):
                formatted_data = []
                for row in table_data["data"]:
                    if isinstance(row, list):
                        formatted_row = []
                        for cell in row:
                            if isinstance(cell, str):
                                formatted_row.append(
                                    {
                                        "text": cell,
                                        "formatted": self.parse_inline_formatting(cell),
                                    }
                                )
                            else:
                                # Keep non-string cells as-is
                                formatted_row.append(cell)
                        formatted_data.append(formatted_row)
                    else:
                        # Keep non-list rows as-is
                        formatted_data.append(row)
                processed_data["table"]["data"] = formatted_data

        # Note: Removed complex formatting preprocessing - formatting now handled at render time
        return processed_data

    def _create_table_in_slide(self, slide, placeholder, table_block):
        """
        Create a PowerPoint table in the given slide, replacing the placeholder.

        Args:
            slide: PowerPoint slide object
            placeholder: PowerPoint placeholder shape to replace
            table_block: Table block with header and rows data
        """
        from pptx.util import Inches  # noqa: F401

        try:
            # Extract table data from block
            header = table_block.get("header", [])
            rows = table_block.get("rows", [])

            if not header and not rows:
                print("      No table data found")
                return

            # Calculate table dimensions
            if header:
                table_data = [header] + rows
            else:
                table_data = rows

            if not table_data:
                print("      Empty table data")
                return

            rows_count = len(table_data)
            cols_count = max(len(row) for row in table_data) if table_data else 0

            if rows_count == 0 or cols_count == 0:
                print(f"      Invalid table dimensions: {rows_count}x{cols_count}")
                return

            print(f"      Creating {rows_count}x{cols_count} table")

            # Get placeholder dimensions and position
            left = placeholder.left
            top = placeholder.top
            width = placeholder.width
            height = placeholder.height

            # Remove the placeholder shape and create table in its place
            sp = placeholder._element
            parent = sp.getparent()
            parent.remove(sp)

            # Add table to slide
            table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
            table = table_shape.table

            # Populate table with data
            for row_idx, row_data in enumerate(table_data):
                for col_idx, cell_data in enumerate(row_data):
                    if col_idx < len(table.rows[row_idx].cells):
                        cell = table.rows[row_idx].cells[col_idx]
                        cell_text = str(cell_data) if cell_data is not None else ""

                        # Apply inline formatting to cell text
                        cell.text_frame.clear()
                        p = cell.text_frame.paragraphs[0]
                        self.apply_inline_formatting(cell_text, p)

                        # Apply header formatting
                        if row_idx == 0 and header:
                            for run in p.runs:
                                run.font.bold = True

            print(f"      Table created successfully with {rows_count} rows and {cols_count} columns")

        except Exception as e:
            print(f"      Error creating table: {e}")
            # Fallback: add table data as text
            self._fallback_table_as_text(placeholder, table_block)

    def _fallback_table_as_text(self, placeholder, table_block):
        """Fallback method to display table as text when table creation fails."""
        try:
            header = table_block.get("header", [])
            rows = table_block.get("rows", [])

            if hasattr(placeholder, "text_frame") and placeholder.text_frame:
                text_frame = placeholder.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]

                # Convert table to text representation
                table_text = "Table data:\n"
                if header:
                    table_text += " | ".join(str(cell) for cell in header) + "\n"
                    table_text += "-" * 50 + "\n"
                for row in rows:
                    table_text += " | ".join(str(cell) for cell in row) + "\n"

                self.apply_inline_formatting(table_text, p)
                print("      Table displayed as text fallback")
            else:
                print("      Could not display table - no text frame available")
        except Exception as e:
            print(f"      Error in table text fallback: {e}")

    def _remove_placeholder_from_slide(self, slide, placeholder):
        """Remove a placeholder from the slide after it's been replaced by a table."""
        try:
            # Find the placeholder shape in the slide and remove it
            for shape in slide.shapes:
                if shape == placeholder:
                    # Remove the shape element from the slide
                    sp = shape._element
                    parent = sp.getparent()
                    if parent is not None:
                        parent.remove(sp)
                    break
            print("      Placeholder removed from slide")
        except Exception as e:
            print(f"      Error removing placeholder: {e}")
