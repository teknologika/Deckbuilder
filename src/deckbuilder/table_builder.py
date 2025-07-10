from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor

try:
    from .table_styles import TABLE_BORDER_STYLES, TABLE_HEADER_STYLES, TABLE_ROW_STYLES
except ImportError:
    # Fallback values if modules don't exist
    TABLE_HEADER_STYLES = {"dark_blue_white_text": {"bg": RGBColor(46, 89, 132), "text": RGBColor(255, 255, 255)}}
    TABLE_ROW_STYLES = {
        "alternating_light_gray": {
            "primary": RGBColor(255, 255, 255),
            "alt": RGBColor(240, 240, 240),
        }
    }
    TABLE_BORDER_STYLES = {"thin_gray": {"width": Pt(1), "color": RGBColor(128, 128, 128), "style": "all"}}


class TableBuilder:
    """Handles table creation, styling, and formatting for PowerPoint presentations."""

    def __init__(self, content_formatter=None):
        """
        Initialize the table builder.

        Args:
            content_formatter: ContentFormatter instance for text formatting
        """
        self.content_formatter = content_formatter

    def add_table_to_slide(self, slide, table_data):
        """
        Add a styled table to a slide.

        Args:
            slide: The slide to add the table to
            table_data: Dictionary containing table data and styling options
        """
        # Get table data - support both 'data' and 'rows' keys for backwards compatibility
        data = table_data.get("data", table_data.get("rows", []))
        if not data:
            return

        # Get styling options
        header_style = table_data.get("header_style", "dark_blue_white_text")
        row_style = table_data.get("row_style", "alternating_light_gray")
        border_style = table_data.get("border_style", "thin_gray")
        custom_colors = table_data.get("custom_colors", {})

        # Find content placeholder or create table in available space
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Content placeholder
                content_placeholder = shape
                break

        if content_placeholder:
            # Remove placeholder and create table in its place
            left = content_placeholder.left
            top = content_placeholder.top
            width = content_placeholder.width
            height = content_placeholder.height

            # Remove the placeholder
            sp = content_placeholder._element
            sp.getparent().remove(sp)
        else:
            # Default positioning if no placeholder found
            left = Cm(2.5)
            top = Cm(5)
            width = Cm(20)
            height = Cm(12)

        # Create the table
        rows = len(data)
        if data:
            # Handle both old (list of strings) and new (list of dicts) formats
            first_row = data[0]
            if isinstance(first_row, list):
                cols = len(first_row)
            else:
                cols = 1  # Fallback
        else:
            cols = 1

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Apply table data with formatting support
        for row_idx, row_data in enumerate(data):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)

                # Handle both old (string) and new (formatted) cell data
                if isinstance(cell_data, dict) and "formatted" in cell_data:
                    # New formatted cell data
                    if self.content_formatter:
                        self.content_formatter.apply_formatted_segments_to_cell(cell, cell_data["formatted"])
                    else:
                        cell.text = cell_data.get("text", str(cell_data))
                else:
                    # Old string cell data
                    cell.text = str(cell_data)

        # Apply styling
        self._apply_table_styling(table, header_style, row_style, border_style, custom_colors)

    def _apply_table_styling(self, table, header_style, row_style, border_style, custom_colors):
        """
        Apply styling to a table.

        Args:
            table: The table object to style
            header_style: Header style name
            row_style: Row style name
            border_style: Border style name
            custom_colors: Dictionary of custom color overrides
        """
        # Apply header styling
        if header_style in TABLE_HEADER_STYLES:
            header_colors = TABLE_HEADER_STYLES[header_style]

            # Override with custom colors if provided
            bg_color = self._parse_custom_color(custom_colors.get("header_bg")) or header_colors["bg"]
            text_color = self._parse_custom_color(custom_colors.get("header_text")) or header_colors["text"]

            # Style header row (first row)
            for col_idx in range(len(table.columns)):
                cell = table.cell(0, col_idx)
                # Set background color
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_color

                # Set text color and formatting
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = text_color
                        run.font.bold = True

        # Apply row styling
        if row_style in TABLE_ROW_STYLES and len(table.rows) > 1:
            row_colors = TABLE_ROW_STYLES[row_style]

            # Override with custom colors if provided
            primary_color = self._parse_custom_color(custom_colors.get("primary_row")) or row_colors["primary"]
            alt_color = self._parse_custom_color(custom_colors.get("alt_row")) or row_colors["alt"]

            # Style data rows (skip header row)
            for row_idx in range(1, len(table.rows)):
                is_alt_row = (row_idx - 1) % 2 == 1
                bg_color = alt_color if is_alt_row else primary_color

                if bg_color is not None:
                    for col_idx in range(len(table.columns)):
                        cell = table.cell(row_idx, col_idx)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = bg_color

        # Apply border styling
        if border_style in TABLE_BORDER_STYLES:
            self._apply_table_borders(table, TABLE_BORDER_STYLES[border_style], custom_colors)

    def _apply_table_borders(self, table, border_config, custom_colors):
        """
        Apply border styling to a table.

        Args:
            table: The table object
            border_config: Border configuration dictionary
            custom_colors: Custom color overrides
        """
        border_width = border_config["width"]
        border_color = self._parse_custom_color(custom_colors.get("border_color")) or border_config["color"]
        border_style = border_config["style"]

        if border_style == "none" or border_width.cm == 0:
            return

        # Apply borders based on style
        for row_idx in range(len(table.rows)):
            for col_idx in range(len(table.columns)):
                cell = table.cell(row_idx, col_idx)

                if border_style == "all":
                    # All borders
                    self._set_cell_borders(cell, border_width, border_color, all_sides=True)
                elif border_style == "header" and row_idx == 0:
                    # Only header bottom border
                    self._set_cell_borders(cell, border_width, border_color, bottom=True)
                elif border_style == "outer":
                    # Only outer borders
                    is_top = row_idx == 0
                    is_bottom = row_idx == len(table.rows) - 1
                    is_left = col_idx == 0
                    is_right = col_idx == len(table.columns) - 1

                    self._set_cell_borders(
                        cell,
                        border_width,
                        border_color,
                        top=is_top,
                        bottom=is_bottom,
                        left=is_left,
                        right=is_right,
                    )

    def _set_cell_borders(self, cell, width, color, all_sides=False, top=False, bottom=False, left=False, right=False):
        """
        Set borders for a table cell.

        Args:
            cell: The table cell
            width: Border width
            color: Border color
            all_sides: Apply to all sides
            top, bottom, left, right: Apply to specific sides
        """
        if color is None:
            return

        if all_sides:
            top = bottom = left = right = True

        # Note: python-pptx has limited border support
        # This is a simplified implementation
        try:
            if hasattr(cell, "border"):
                if top and hasattr(cell.border, "top"):
                    cell.border.top.color.rgb = color
                    cell.border.top.width = width
                if bottom and hasattr(cell.border, "bottom"):
                    cell.border.bottom.color.rgb = color
                    cell.border.bottom.width = width
                if left and hasattr(cell.border, "left"):
                    cell.border.left.color.rgb = color
                    cell.border.left.width = width
                if right and hasattr(cell.border, "right"):
                    cell.border.right.color.rgb = color
                    cell.border.right.width = width
        except Exception:
            # Borders not fully supported in python-pptx, skip silently
            return  # nosec - Skip border styling if not supported

    def _parse_custom_color(self, color_value):
        """
        Parse a custom color value (hex string) to RGBColor.

        Args:
            color_value: Hex color string (e.g., "#FF0000")

        Returns:
            RGBColor object or None if invalid
        """
        if not color_value or not isinstance(color_value, str):
            return None

        try:
            # Remove # if present
            color_value = color_value.lstrip("#")

            # Convert hex to RGB
            if len(color_value) == 6:
                r = int(color_value[0:2], 16)
                g = int(color_value[2:4], 16)
                b = int(color_value[4:6], 16)
                return RGBColor(r, g, b)
        except (ValueError, TypeError):
            pass

        return None
