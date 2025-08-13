from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from pptx.util import Cm

from ..content.placeholder_types import (
    is_content_placeholder,
    is_media_placeholder,
    is_subtitle_placeholder,
    is_title_placeholder,
)
from ..utils.logging import slide_builder_print, debug_print, error_print


class SlideBuilder:
    """Handles core slide creation, layout mapping, and placeholder management."""

    def __init__(self, layout_mapping=None):
        """
        Initialize the slide builder.

        Args:
            layout_mapping: Optional layout mapping dictionary
        """
        self.layout_mapping = layout_mapping
        self._current_slide_index = 0

    def clear_slides(self, prs):
        """Clear all slides from the presentation."""
        slide_count = len(prs.slides)
        for i in range(slide_count - 1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]

        # Reset slide index for consistent image selection
        self._current_slide_index = 0

    def add_slide(self, prs, slide_data: dict, content_formatter, image_placeholder_handler):
        """
        Add a single slide to the presentation based on slide data.

        Args:
            prs: PowerPoint presentation object
            slide_data: Dictionary containing slide information
            content_formatter: ContentFormatter instance for handling content
            image_placeholder_handler: ImagePlaceholderHandler instance for handling images
        """
        # Type validation: ensure slide_data is a dictionary
        if not isinstance(slide_data, dict):
            raise TypeError(f"slide_data must be a dictionary, got {type(slide_data).__name__}: {slide_data}")

        # Track slide index for consistent image selection
        self._current_slide_index = getattr(self, "_current_slide_index", 0) + 1

        # Auto-parse JSON formatting for inline formatting support
        slide_data = content_formatter.format_slide_data(slide_data)

        # Get slide type and determine layout using JSON mapping
        # Prefer explicit "layout" field over "type" field
        layout_or_type = slide_data.get("layout", slide_data.get("type", "content"))

        # Use layout mapping if available
        if self.layout_mapping:
            aliases = self.layout_mapping.get("aliases", {})
            layouts = self.layout_mapping.get("layouts", {})

            # Get layout name from aliases (or use direct layout name if it exists in layouts)
            if layout_or_type in layouts:
                layout_name = layout_or_type
            else:
                layout_name = aliases.get(layout_or_type, layout_or_type)

            # Get layout index
            layout_info = layouts.get(layout_name, {})
            layout_index = layout_info.get("index", 1)
        else:
            # Fallback
            layout_name = layout_or_type  # Use the original layout name as fallback
            layout_index = 1

        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)

        # Copy descriptive placeholder names from template mapping
        self._copy_placeholder_names_from_mapping(slide, layout_name)

        # Add content to placeholders using template mapping + semantic detection
        # For dynamic shapes, content placeholder already contains first text segment only
        self._apply_content_to_mapped_placeholders(slide, slide_data, layout_name, content_formatter, image_placeholder_handler)

        # Add speaker notes if they exist in slide_data or placeholders
        speaker_notes = None
        if "speaker_notes" in slide_data:
            speaker_notes = slide_data["speaker_notes"]
        elif "placeholders" in slide_data and "speaker_notes" in slide_data["placeholders"]:
            speaker_notes = slide_data["placeholders"]["speaker_notes"]

        if speaker_notes:
            self.add_speaker_notes(slide, speaker_notes, content_formatter)

        # ENHANCED: Handle dynamic multi-shape creation for mixed content
        if slide_data.get("_requires_dynamic_shapes") and slide_data.get("_content_segments"):
            # CRITICAL: Clean up content placeholder before dynamic shape creation
            content_placeholder = self._find_content_placeholder(slide)
            if content_placeholder and hasattr(content_placeholder, "text_frame"):
                # Get first text segment for content placeholder
                first_text_segment = next((seg for seg in slide_data["_content_segments"] if seg["type"] == "text"), None)
                if first_text_segment:
                    # Clear any existing mixed content and set only first text segment
                    content_placeholder.text_frame.clear()
                    content_placeholder.text_frame.text = first_text_segment["content"].strip()

                    # CRITICAL: Apply consistent font sizing to match title placeholder
                    title_placeholder = self._find_title_placeholder(slide)
                    if title_placeholder and hasattr(title_placeholder, "text_frame"):
                        # Get font size from title to ensure consistency
                        title_font_size = self._get_placeholder_font_size(title_placeholder)

                        # Apply same font size to content placeholder paragraphs and runs
                        for paragraph in content_placeholder.text_frame.paragraphs:
                            # Set paragraph level font
                            if paragraph.font:
                                paragraph.font.size = title_font_size

                            # Set run level fonts
                            for run in paragraph.runs:
                                if run.font:
                                    run.font.size = title_font_size

                        # Also set default character formatting for the text frame
                        if hasattr(content_placeholder.text_frame, "auto_size"):
                            from pptx.enum.text import MSO_AUTO_SIZE

                            content_placeholder.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                        debug_print(f"    Applied title font size ({title_font_size.pt if hasattr(title_font_size, 'pt') else title_font_size}) to content placeholder")

                    debug_print(f"    Cleaned content placeholder: set to first text segment only ({len(first_text_segment['content'])} chars)")
                else:
                    # No text segments, clear the placeholder
                    content_placeholder.text_frame.clear()
                    debug_print("    Cleaned content placeholder: cleared (no text segments)")

            self._create_dynamic_content_shapes(slide, slide_data, content_formatter)

        # All content should be processed through placeholders only - no legacy content blocks
        debug_print("  Slide completed using structured frontmatter placeholders only")

        return slide

    def add_speaker_notes(self, slide, notes_content, content_formatter):
        """
        Adds speaker notes to the slide.

        Args:
            slide: The slide to add notes to.
            notes_content: The content of the speaker notes.
            content_formatter: ContentFormatter instance for handling content
        """
        if notes_content:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
            content_formatter.apply_inline_formatting(notes_content, p)

    def _find_placeholder_by_name(self, slide, field_name):
        """
        Find a placeholder by its name (after names have been updated to match field names).

        Args:
            slide: The slide to search
            field_name: The field name to look for

        Returns:
            The placeholder object if found, None otherwise
        """
        for placeholder in slide.placeholders:
            try:
                placeholder_name = placeholder.element.nvSpPr.cNvPr.name
                if placeholder_name == field_name:
                    return placeholder
            except AttributeError:
                # Some placeholders don't have accessible names
                continue
        return None

    def add_slide_with_direct_mapping(self, prs, slide_data: dict, content_formatter, image_placeholder_handler):
        """
        Add slide using direct field mapping (no markdown conversion).

        Args:
            prs: PowerPoint presentation object
            slide_data: Dictionary containing slide information
            content_formatter: ContentFormatter instance
            image_placeholder_handler: ImagePlaceholderHandler instance
        """
        # Import the universal formatting module
        from ..content.formatting_support import content_formatter as cf

        # Apply universal formatting to slide data
        formatted_slide = cf.format_slide_data(slide_data)

        # Add slide using the standard method
        return self.add_slide(prs, formatted_slide, content_formatter, image_placeholder_handler)

    def _copy_placeholder_names_from_mapping(self, slide, layout_name):
        """
        Copy descriptive placeholder names from template mapping to slide placeholders.

        This enhances the PowerPoint editing experience by providing meaningful placeholder
        names like "Col 1 Title Placeholder 2" instead of generic "Text Placeholder 2".

        Args:
            slide: PowerPoint slide object
            layout_name: Name of the PowerPoint layout
        """
        if not self.layout_mapping:
            return

        # Get layout info from template mapping
        layouts = self.layout_mapping.get("layouts", {})
        layout_info = layouts.get(layout_name, {})
        placeholder_mappings = layout_info.get("placeholders", {})

        # Update placeholder names to match template mapping
        for placeholder in slide.placeholders:
            placeholder_idx = str(placeholder.placeholder_format.idx)
            if placeholder_idx in placeholder_mappings:
                descriptive_name = placeholder_mappings[placeholder_idx]
                try:
                    # Update the placeholder name
                    try:
                        placeholder.element.nvSpPr.cNvPr.name = descriptive_name
                    except AttributeError:
                        pass  # Some placeholder types don't support name changes
                except Exception:
                    # Fallback: some placeholder types might not allow name changes
                    pass  # nosec - Continue processing other placeholders

    def _apply_content_to_mapped_placeholders(self, slide, slide_data, layout_name, content_formatter, image_placeholder_handler):
        """
        Apply content to placeholders using template JSON mappings + semantic detection.

        This unified method works with both JSON input and markdown frontmatter input:
        1. Looks up layout in template JSON mappings
        2. For each field in slide_data, finds corresponding placeholder index
        3. Gets actual placeholder and determines its semantic type
        4. Applies content using appropriate semantic handler

        Args:
            slide: PowerPoint slide object
            slide_data: Dictionary containing slide content (from JSON or markdown)
            layout_name: Name of the PowerPoint layout
            content_formatter: ContentFormatter instance
            image_placeholder_handler: ImagePlaceholderHandler instance
        """
        if not self.layout_mapping:
            # Fallback to basic semantic detection if no mapping available
            self._apply_content_to_placeholder_unified(slide, slide_data, content_formatter, image_placeholder_handler, fallback_mode=True)
            return

        # Get layout info from template mapping
        layouts = self.layout_mapping.get("layouts", {})
        layout_info = layouts.get(layout_name, {})
        placeholder_mappings = layout_info.get("placeholders", {})

        # Create reverse mapping: field_name -> placeholder_index
        field_to_index = {}
        for placeholder_idx, field_name in placeholder_mappings.items():
            field_to_index[field_name] = int(placeholder_idx)

        # Enhanced debugging: Show all available placeholders and template mapping
        slide_builder_print(f"Layout '{layout_name}' - Template Mapping Analysis:")
        slide_builder_print(f"  Available placeholders in template: {list(field_to_index.keys())}")

        # Show actual PowerPoint placeholders
        actual_placeholders = []
        for ph in slide.placeholders:
            ph_type = ph.placeholder_format.type.name if hasattr(ph.placeholder_format.type, "name") else str(ph.placeholder_format.type)
            try:
                ph_name = getattr(ph.element.nvSpPr.cNvPr, "name", "unnamed")
            except AttributeError:
                ph_name = "unnamed"
            actual_placeholders.append(f"{ph.placeholder_format.idx}:{ph_type}({ph_name})")
        slide_builder_print(f"  PowerPoint placeholders found: {actual_placeholders}")

        # Process each field in slide_data using semantic detection
        # For canonical JSON format, process the placeholders object if it exists
        content_data = slide_data.get("placeholders", {}) if "placeholders" in slide_data else slide_data
        slide_builder_print(f"  Content fields to map: {list(content_data.keys())}")
        successful_mappings = []
        failed_mappings = []

        for field_name, field_value in content_data.items():
            # Skip non-content fields
            if field_name in ["type", "table", "layout"]:
                continue

            slide_builder_print(f"    Mapping field '{field_name}' (value: {str(field_value)[:50]}...)")

            # Find placeholder using semantic detection
            target_placeholder = None
            mapping_method = None

            # Handle title placeholders
            if field_name == "title":
                slide_builder_print("      Method: Title field resolution")

                # First try to find title field in template mapping
                title_field = field_to_index.get("title")
                if title_field is not None:
                    slide_builder_print(f"        Trying template mapping 'title' -> idx {title_field}")
                    for placeholder in slide.placeholders:
                        if placeholder.placeholder_format.idx == title_field:
                            target_placeholder = placeholder
                            mapping_method = f"Template mapping title -> idx {title_field}"
                            break
                else:
                    slide_builder_print("        No 'title' field in template mapping, trying title_top")
                    # Try title_top which is common in templates
                    title_top_idx = field_to_index.get("title_top")
                    if title_top_idx is not None:
                        slide_builder_print(f"        Trying template mapping 'title_top' -> idx {title_top_idx}")
                        for placeholder in slide.placeholders:
                            if placeholder.placeholder_format.idx == title_top_idx:
                                target_placeholder = placeholder
                                mapping_method = f"Template mapping title_top -> idx {title_top_idx}"
                                break

                # Final fallback to semantic title detection
                if not target_placeholder:
                    slide_builder_print("        Template mapping failed, trying semantic title detection")
                    for placeholder in slide.placeholders:
                        if is_title_placeholder(placeholder.placeholder_format.type):
                            target_placeholder = placeholder
                            mapping_method = f"Semantic title (idx {placeholder.placeholder_format.idx})"
                            break

            # Handle subtitle placeholders
            elif field_name == "subtitle":
                slide_builder_print("      Method: Semantic subtitle detection")
                for placeholder in slide.placeholders:
                    if is_subtitle_placeholder(placeholder.placeholder_format.type):
                        target_placeholder = placeholder
                        mapping_method = f"Semantic subtitle (idx {placeholder.placeholder_format.idx})"
                        break

            # Handle content placeholders
            elif field_name == "content":
                slide_builder_print("      Method: Content field resolution")

                # First try to find content field in template mapping
                content_field = field_to_index.get("content")
                if content_field is not None:
                    slide_builder_print(f"        Trying template mapping 'content' -> idx {content_field}")
                    for placeholder in slide.placeholders:
                        if placeholder.placeholder_format.idx == content_field:
                            target_placeholder = placeholder
                            mapping_method = f"Template mapping content -> idx {content_field}"
                            break
                else:
                    slide_builder_print("        No 'content' field in template mapping, trying content_1")
                    # Fallback: try to find content_1 specifically (for layouts like vertical)
                    content_1_idx = field_to_index.get("content_1")
                    if content_1_idx is not None:
                        slide_builder_print(f"        Trying template mapping 'content_1' -> idx {content_1_idx}")
                        for placeholder in slide.placeholders:
                            if placeholder.placeholder_format.idx == content_1_idx:
                                target_placeholder = placeholder
                                mapping_method = f"Template mapping content_1 -> idx {content_1_idx}"
                                break

                # Final fallback to semantic content placeholder detection
                if not target_placeholder:
                    slide_builder_print("        Template mapping failed, trying semantic content detection")
                    for placeholder in slide.placeholders:
                        if is_content_placeholder(placeholder.placeholder_format.type):
                            target_placeholder = placeholder
                            mapping_method = f"Semantic content (idx {placeholder.placeholder_format.idx})"
                            break

            # Handle image_path fields and image placeholder fields - find PICTURE placeholders
            elif field_name == "image_path" or field_name.endswith(".image_path") or "image" in field_name.lower():
                slide_builder_print("      Method: Image field detection")
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.PICTURE:
                        target_placeholder = placeholder
                        mapping_method = f"Image placeholder (idx {placeholder.placeholder_format.idx})"
                        slide_builder_print(f"        Found PICTURE placeholder at idx {placeholder.placeholder_format.idx}")
                        break
                if not target_placeholder:
                    slide_builder_print(f"        No PICTURE placeholder found for image field '{field_name}'")

            # Handle other fields by checking if they match placeholder names in JSON mapping
            else:
                slide_builder_print("      Method: Template mapping lookup")
                # Try to find by exact field name match in JSON mapping
                target_field = field_name

                # Handle common field name variations - be flexible for users
                resolved_field = self._resolve_field_name_variations(field_name, field_to_index)

                if resolved_field != field_name:
                    slide_builder_print(f"        Field name resolved: '{field_name}' -> '{resolved_field}'")
                    target_field = resolved_field
                else:
                    slide_builder_print(f"        Using field name as-is: '{field_name}'")

                # Try name-based lookup first (after placeholder names have been updated)
                slide_builder_print(f"        Trying name-based lookup for field: '{target_field}'")
                target_placeholder = self._find_placeholder_by_name(slide, target_field)
                if target_placeholder:
                    mapping_method = f"Name-based mapping '{target_field}'"
                    placeholder_type = target_placeholder.placeholder_format.type
                    slide_builder_print(f"        SUCCESS: Found placeholder by name: '{target_field}', type: {placeholder_type}")
                elif target_field in field_to_index:
                    # Fallback to index-based lookup
                    placeholder_idx = field_to_index[target_field]
                    slide_builder_print(f"        Template mapping found: '{target_field}' -> idx {placeholder_idx}")
                    for placeholder in slide.placeholders:
                        if placeholder.placeholder_format.idx == placeholder_idx:
                            target_placeholder = placeholder
                            mapping_method = f"Template mapping '{target_field}' -> idx {placeholder_idx}"
                            slide_builder_print(f"        Matched placeholder at idx {placeholder_idx}")
                            break
                    if not target_placeholder:
                        slide_builder_print(f"        Template mapping failed: idx {placeholder_idx} not found in slide")
                else:
                    slide_builder_print(f"        No template mapping found for field '{target_field}'")

            if target_placeholder:
                slide_builder_print(f"    SUCCESS: '{field_name}' mapped using {mapping_method}")
                successful_mappings.append(f"{field_name} -> {mapping_method}")

                # Apply content based on placeholder's semantic type
                self._apply_content_to_single_placeholder(
                    slide,
                    target_placeholder,
                    field_name,
                    field_value,
                    slide_data,
                    content_formatter,
                    image_placeholder_handler,
                )
            else:
                error_print(f"    FAILED: '{field_name}' could not be mapped to any placeholder")
                failed_mappings.append(field_name)

        # Summary of mapping results
        slide_builder_print("  Mapping Summary:")
        slide_builder_print(f"    Successful: {len(successful_mappings)} fields")
        for mapping in successful_mappings:
            slide_builder_print(f"      {mapping}")
        if failed_mappings:
            error_print(f"    Failed: {len(failed_mappings)} fields")
            for field in failed_mappings:
                error_print(f"      {field} (no suitable placeholder found)")

        # Process nested structures like media.image_path
        self._process_nested_image_fields(slide, slide_data, image_placeholder_handler)

        # Process table data if present - ONLY if not using dynamic shapes
        if not slide_data.get("_requires_dynamic_shapes"):
            self._process_table_data(slide, slide_data, content_formatter)
        else:
            debug_print("  Skipping static table processing - using dynamic multi-shape creation")

    def _resolve_field_name_variations(self, field_name: str, field_to_index: dict) -> str:
        """
        Resolve field name variations to handle user flexibility.

        Users shouldn't have to know exact template field names - the engine should
        be smart enough to map common variations automatically.
        """
        # Return original if exact match exists
        if field_name in field_to_index:
            return field_name

        # Common variations mapping
        variations = {
            # Caption variations
            "text_caption": ["text_caption_1", "caption", "caption_1"],
            "caption": ["text_caption_1", "text_caption", "caption_1"],
            # Title variations - CRITICAL: map "title" to "title_top" for template compatibility
            "title": ["title_top", "title_top_1", "main_title"],
            "title_top": ["title", "title_top_1", "main_title"],
            "title_left": ["title_left_1", "left_title", "title_col1"],
            "title_right": ["title_right_1", "right_title", "title_col2"],
            # Content variations
            "content_left": ["content_left_1", "left_content", "content_col1"],
            "content_right": ["content_right_1", "right_content", "content_col2"],
            "content": ["content_1", "main_content", "body"],
            # Image variations
            "image": ["image_1", "image_path", "picture"],
            "image_1": ["image", "image_path", "picture"],
            "image_path": ["image", "image_1", "picture"],
            # Column variations
            "content_col1": ["content_left", "content_left_1", "col1_content"],
            "content_col2": ["content_right", "content_right_1", "col2_content"],
            "content_col3": ["content_col3_1", "col3_content"],
            "content_col4": ["content_col4_1", "col4_content"],
            "title_col1": ["title_left", "title_left_1", "col1_title"],
            "title_col2": ["title_right", "title_right_1", "col2_title"],
            "title_col3": ["title_col3_1", "col3_title"],
            "title_col4": ["title_col4_1", "col4_title"],
            # Item variations (for agenda, lists)
            "content_item1": ["content_item1_1", "item1_content", "item_1"],
            "content_item2": ["content_item2_1", "item2_content", "item_2"],
            "content_item3": ["content_item3_1", "item3_content", "item_3"],
            "content_item4": ["content_item4_1", "item4_content", "item_4"],
            "content_item5": ["content_item5_1", "item5_content", "item_5"],
            "content_item6": ["content_item6_1", "item6_content", "item_6"],
            # Number variations (for agenda)
            "number_item1": ["number_item1_1", "item1_number", "num_1"],
            "number_item2": ["number_item2_1", "item2_number", "num_2"],
            "number_item3": ["number_item3_1", "item3_number", "num_3"],
            "number_item4": ["number_item4_1", "item4_number", "num_4"],
            "number_item5": ["number_item5_1", "item5_number", "num_5"],
            "number_item6": ["number_item6_1", "item6_number", "num_6"],
            # SWOT Analysis variations
            "content_top_left": ["content_16", "strengths", "strength"],
            "content_top_right": ["content_17", "weaknesses", "weakness"],
            "content_bottom_left": ["content_18", "opportunities", "opportunity"],
            "content_bottom_right": ["content_19", "threats", "threat"],
            "content_16": ["content_top_left", "strengths", "strength"],
            "content_17": ["content_top_right", "weaknesses", "weakness"],
            "content_18": ["content_bottom_left", "opportunities", "opportunity"],
            "content_19": ["content_bottom_right", "threats", "threat"],
        }

        # Check if field_name has variations to try
        if field_name in variations:
            for variant in variations[field_name]:
                if variant in field_to_index:
                    return variant

        # Reverse lookup - check if template has a field that maps to this user field
        for template_field in field_to_index.keys():
            if template_field in variations:
                if field_name in variations[template_field]:
                    return template_field

        # Smart suffix handling - try adding/removing _1 suffix
        if field_name.endswith("_1"):
            base_name = field_name[:-2]
            if base_name in field_to_index:
                return base_name
        else:
            suffixed_name = field_name + "_1"
            if suffixed_name in field_to_index:
                return suffixed_name

        # Partial matching for complex fields
        for template_field in field_to_index.keys():
            # Check if field names are similar (contain each other)
            if field_name.lower() in template_field.lower() or template_field.lower() in field_name.lower():
                return template_field

        # Return original if no variations found
        return field_name

    def _apply_content_to_placeholder_unified(self, slide, slide_data, content_formatter, image_placeholder_handler, fallback_mode=False):
        """
        Unified method that handles both standard and fallback content application.

        Args:
            slide: PowerPoint slide object
            slide_data: Slide content data
            content_formatter: ContentFormatter instance
            image_placeholder_handler: ImagePlaceholderHandler instance
            fallback_mode: If True, use semantic detection fallback; if False, use template mapping
        """
        if fallback_mode:
            # Fallback mode: Apply content using semantic detection
            for placeholder in slide.placeholders:
                placeholder_type = placeholder.placeholder_format.type

                # Apply content based on semantic type detection
                if is_title_placeholder(placeholder_type):
                    if "title" in slide_data:
                        self._apply_content_to_single_placeholder(
                            slide,
                            placeholder,
                            "title",
                            slide_data["title"],
                            slide_data,
                            content_formatter,
                            image_placeholder_handler,
                        )
                elif is_subtitle_placeholder(placeholder_type):
                    if "subtitle" in slide_data:
                        self._apply_content_to_single_placeholder(
                            slide,
                            placeholder,
                            "subtitle",
                            slide_data["subtitle"],
                            slide_data,
                            content_formatter,
                            image_placeholder_handler,
                        )
                elif is_content_placeholder(placeholder_type):
                    if "content" in slide_data and not slide_data.get("_requires_dynamic_shapes"):
                        # Only process content normally if NOT using dynamic shape creation
                        self._apply_content_to_single_placeholder(
                            slide,
                            placeholder,
                            "content",
                            slide_data["content"],
                            slide_data,
                            content_formatter,
                            image_placeholder_handler,
                        )
        else:
            # Standard mode: Use template mapping (handled by caller)
            # This method is called when template mapping is available
            pass

    def _apply_content_to_single_placeholder(
        self,
        slide,
        placeholder,
        field_name,
        field_value,
        slide_data,
        content_formatter,
        image_placeholder_handler,
    ):
        """
        Apply content to a single placeholder based on its semantic type and content type.
        This consolidates the logic from both _apply_content_by_semantic_type and fallback methods.
        """
        # Set slide context for content formatter (needed for table processing)
        content_formatter._current_slide = slide_data.get("slide_object", slide)

        placeholder_type = placeholder.placeholder_format.type

        # Apply content based on placeholder semantic type
        if is_title_placeholder(placeholder_type) or is_subtitle_placeholder(placeholder_type):
            # Title/subtitle placeholders - apply inline formatting directly
            if hasattr(placeholder, "text_frame") and placeholder.text_frame:
                # Use text frame for formatting support
                text_frame = placeholder.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                # Handle formatted content properly for titles
                if isinstance(field_value, list) and field_value and isinstance(field_value[0], dict) and "text" in field_value[0]:
                    content_formatter.apply_formatted_segments_to_paragraph(field_value, p)
                else:
                    content_formatter.apply_inline_formatting(str(field_value), p)
            else:
                # Fallback to simple text
                placeholder.text = str(field_value)

        elif is_content_placeholder(placeholder_type):
            # Check if this is table data that should be handled specially
            if isinstance(field_value, dict) and field_value.get("type") == "table":
                # This is table data - route to table creation instead of text placeholder
                slide_builder_print(f"    SPECIAL HANDLING: Table data detected in content field '{field_name}'")

                # Import TableBuilder for table creation
                from .table_builder import TableBuilder

                table_builder = TableBuilder(content_formatter)

                # Create table on slide
                table_builder.add_table_to_slide(slide, field_value)

                # Clear the placeholder to avoid showing placeholder text
                if hasattr(placeholder, "text_frame") and placeholder.text_frame:
                    placeholder.text_frame.clear()
                    # Optionally add a note that table was created
                    p = placeholder.text_frame.paragraphs[0] if placeholder.text_frame.paragraphs else placeholder.text_frame.add_paragraph()
                    p.text = ""  # Leave empty since table replaces this content
            else:
                # Content placeholders - handle text, lists, etc. with inline formatting
                # ENHANCED: Only process if not table content
                content_formatter.add_content_to_placeholder(placeholder, field_value)

        elif is_media_placeholder(placeholder_type):
            # Media placeholders - handle images, charts, tables, etc.
            if placeholder_type == PP_PLACEHOLDER_TYPE.PICTURE:
                image_placeholder_handler.handle_image_placeholder(placeholder, field_name, field_value, slide_data)
            elif placeholder_type == PP_PLACEHOLDER_TYPE.TABLE:
                # TABLE placeholders - handle table data
                self._handle_table_placeholder(placeholder, field_name, field_value, slide_data, slide)
            elif placeholder_type == PP_PLACEHOLDER_TYPE.OBJECT and hasattr(placeholder, "text_frame"):
                # OBJECT placeholders with text_frame should be treated as content placeholders
                content_formatter.add_content_to_placeholder(placeholder, field_value)
            else:
                # Generic object placeholder - try text fallback
                if hasattr(placeholder, "text"):
                    placeholder.text = str(field_value)

    def _process_nested_image_fields(self, slide, slide_data, image_placeholder_handler):
        """
        Process nested image fields like media.image_path from structured frontmatter.

        Note: This method handles raw frontmatter with nested media structures.
        Structured frontmatter conversion flattens media.image_path to image_1,
        which is already handled by the main field processing loop.

        Args:
            slide: PowerPoint slide object
            slide_data: Dictionary containing slide content
            image_placeholder_handler: ImagePlaceholderHandler instance
        """
        # Skip if this appears to be structured frontmatter that was already converted
        # (indicated by presence of flattened image fields like image_1, image_path)
        has_converted_image_fields = any(field_name == "image_path" or field_name.endswith("_1") and "image" in field_name.lower() for field_name in slide_data.keys())

        if has_converted_image_fields:
            # Already processed by structured frontmatter conversion
            return

        # Check for media structure with image_path (raw frontmatter)
        if "media" in slide_data and isinstance(slide_data["media"], dict):
            media_data = slide_data["media"]
            image_path = media_data.get("image_path")

            if image_path:
                # Find the first PICTURE placeholder
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.PICTURE:
                        image_placeholder_handler.handle_image_placeholder(placeholder, "media.image_path", image_path, slide_data)
                        break

    def _handle_table_placeholder(self, placeholder, field_name, field_value, slide_data, slide):
        """
        Handle TABLE placeholder by creating a table directly in the placeholder space.

        Args:
            placeholder: PowerPoint TABLE placeholder
            field_name: Name of the field (e.g., 'table_data')
            field_value: Table data (dict with table structure)
            slide_data: Full slide data for styling
            slide: PowerPoint slide object
        """
        debug_print(f"  Handling TABLE placeholder for field: {field_name}")
        debug_print(f"    Field value type: {type(field_value)}")
        debug_print(f"    Field value: {field_value}")

        if not isinstance(field_value, dict) or field_value.get("type") != "table":
            debug_print(f"    Skipping - not table data: {type(field_value)}")
            return

        # Use the table builder to create the table
        from ..core.table_builder import TableBuilder

        table_builder = TableBuilder()

        # Create table configuration combining field data with slide styling
        table_config = field_value.copy()

        # Add slide-level styling if not specified in field
        for style_field in ["header_style", "row_style", "border_style", "row_height", "table_width", "column_widths"]:
            if style_field not in table_config and style_field in slide_data:
                table_config[style_field] = slide_data[style_field]

        # Get placeholder bounds for positioning
        left = placeholder.left
        top = placeholder.top
        width = placeholder.width
        height = placeholder.height

        # Create table shape manually at placeholder position
        if "data" in table_config and table_config["data"]:
            rows = len(table_config["data"])
            cols = len(table_config["data"][0]) if table_config["data"] else 1

            # Create table shape
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table

            # Apply data and styling using table builder
            table_builder._apply_table_data_and_styling(table, table_config)

            debug_print(f"    Created table in TABLE placeholder: {rows}x{cols}")
        else:
            debug_print(f"    No table data found in: {table_config}")

    def _process_table_data(self, slide, slide_data, content_formatter):
        """
        ENHANCED: Complete table processing - handles formatting AND configuration.

        This is now the SINGLE method responsible for ALL table creation.
        It finds table content anywhere in slide_data, preserves formatting from markdown,
        applies configuration from frontmatter, and creates one complete table.

        Args:
            slide: PowerPoint slide object
            slide_data: Dictionary containing slide content (may include table data)
            content_formatter: ContentFormatter instance for table creation
        """
        # ENHANCED: Look for table content in multiple places (content, table, table_data fields)
        table_content = self._find_table_content_in_slide_data(slide_data)
        if not table_content:
            return

        debug_print(f"  Found table content in: {table_content['source_field']}")

        # Enhanced table detection to prevent duplicate table creation
        existing_tables = self._detect_existing_tables(slide)
        if existing_tables:
            debug_print(f"  Table already exists ({len(existing_tables)} found), skipping duplicate table creation")
            return

        # ENHANCED: Parse table markdown with formatting preservation
        if table_content["markdown"]:
            debug_print("  Parsing table markdown with formatting preservation")
            formatted_table_data = content_formatter.parse_table_markdown_with_formatting(table_content["markdown"])
        else:
            # Fallback for non-markdown table data
            formatted_table_data = table_content.get("table_data", {"data": [], "type": "table"})

        # ENHANCED: Read configuration from frontmatter (dimensions, styling)
        table_config = {
            "header_style": slide_data.get("style", "dark_blue_white_text"),
            "row_style": slide_data.get("row_style", "alternating_light_gray"),
            "border_style": slide_data.get("border_style", "thin_gray"),
            "column_widths": slide_data.get("column_widths", []),
            "row_height": slide_data.get("row_height", None),
            "table_width": slide_data.get("table_width", None),
            "header_font_size": slide_data.get("header_font_size", 12),
            "data_font_size": slide_data.get("data_font_size", 10),
            "custom_colors": slide_data.get("custom_colors", {}),
        }

        # ENHANCED: Merge formatted data with configuration for complete table
        complete_table_data = {**formatted_table_data, **table_config}
        debug_print(f"  Creating complete table with {len(complete_table_data.get('data', []))} rows")

        # ENHANCED: Clear ALL table-related placeholders before creating table
        self._clear_all_table_placeholders(slide, slide_data, table_content)

        # Set slide context for content formatter
        content_formatter._current_slide = slide

        # Create the complete table with formatting AND configuration
        from .table_builder import TableBuilder

        table_builder = TableBuilder(content_formatter)
        table_builder.add_table_to_slide(slide, complete_table_data)

        # ENHANCED: Mark content as processed to prevent content processing from also creating tables
        self._mark_table_content_as_processed(slide_data, table_content["source_field"])
        debug_print("  Complete table creation finished")

    def _clear_table_content_from_placeholders(self, slide, slide_data):
        """
        Clear table text content from placeholders when a table object will be created.

        This prevents duplication when both content placeholders contain table markdown
        AND a separate table object exists with table data.

        Args:
            slide: PowerPoint slide object
            slide_data: Dictionary containing slide content including placeholders
        """
        from ..utils.logging import debug_print

        # Check both placeholders data and direct slide_data fields
        placeholders_data = slide_data.get("placeholders", {})

        # Also check direct fields that might contain table content (like 'content' field)
        all_fields = {**placeholders_data}
        for key in ["content", "body", "text"]:
            if key in slide_data:
                all_fields[key] = slide_data[key]

        debug_print(f"  Checking {len(all_fields)} fields for table content: {list(all_fields.keys())}")

        for field_name, field_value in all_fields.items():
            if isinstance(field_value, str) and self._is_table_markdown(field_value):
                debug_print(f"  Found table markdown in field '{field_name}'")
                # Find and clear ALL placeholders that contain this table content
                for shape in slide.placeholders:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        shape_text = shape.text_frame.text if shape.text_frame.text else ""
                        if shape_text and self._is_table_markdown(shape_text):
                            debug_print(f"  Clearing table content from placeholder idx {shape.placeholder_format.idx}")
                            shape.text_frame.clear()
                            # Leave completely empty since table object will replace this content
                            break

    def _is_table_markdown(self, text_content):
        """
        Enhanced detection for markdown table syntax.

        Args:
            text_content: String content to check

        Returns:
            bool: True if content appears to contain table markdown
        """
        if not isinstance(text_content, str) or not text_content.strip():
            return False

        lines = [line.strip() for line in text_content.split("\n") if line.strip()]
        if len(lines) < 2:
            return False

        table_rows = 0

        for line in lines:
            # Table separator line (like |---|---|---| or | --- | --- | --- |)
            if "|" in line and all(c in "|-:= \t" for c in line.replace("|", "").strip()):
                continue  # Skip separator lines
            # Table data row (contains | but has actual content)
            elif "|" in line and not all(c in "|-:= \t" for c in line.strip()):
                table_rows += 1

        # Valid table: at least 2 data rows (header + content) and optionally a separator
        return table_rows >= 2

    def _detect_existing_tables(self, slide):
        """
        Enhanced method to detect existing tables on a slide using multiple approaches.

        Args:
            slide: PowerPoint slide object

        Returns:
            List of existing table shapes
        """
        from ..utils.logging import debug_print

        existing_tables = []
        debug_print(f"  Checking {len(slide.shapes)} shapes for existing tables")

        for i, shape in enumerate(slide.shapes):
            try:
                shape_info = f"Shape {i}: "

                # Method 1: Check shape_type for TABLE (MSO_SHAPE_TYPE.TABLE = 19)
                if hasattr(shape, "shape_type"):
                    shape_type_str = str(shape.shape_type)
                    shape_info += f"type={shape_type_str} "

                    # PowerPoint table shape type is 19
                    if "19" in shape_type_str or "TABLE" in shape_type_str.upper():
                        existing_tables.append(shape)
                        debug_print(f"    {shape_info}-> FOUND TABLE (shape_type)")
                        continue

                # Method 2: Check for table attribute directly
                if hasattr(shape, "table"):
                    try:
                        # Try to access table properties to confirm it's a real table
                        rows = len(shape.table.rows)
                        cols = len(shape.table.columns)
                        existing_tables.append(shape)
                        debug_print(f"    {shape_info}-> FOUND TABLE (table attr, {rows}x{cols})")
                        continue
                    except Exception:  # nosec B110
                        pass

                # Method 3: Check element tag name (backup method)
                if hasattr(shape, "element") and hasattr(shape.element, "tag"):
                    tag_name = shape.element.tag if hasattr(shape.element, "tag") else ""
                    if "tbl" in tag_name.lower():
                        existing_tables.append(shape)
                        debug_print(f"    {shape_info}-> FOUND TABLE (element tag)")
                        continue

                debug_print(f"    {shape_info}-> not a table")

            except Exception as e:
                debug_print(f"    Shape {i}: Error checking - {str(e)}")
                pass  # Skip shapes that can't be inspected

        debug_print(f"  Total tables found: {len(existing_tables)}")
        return existing_tables

    def _find_table_content_in_slide_data(self, slide_data):
        """
        ENHANCED: Find table content anywhere in slide_data (content, table, table_data fields).

        Args:
            slide_data: Dictionary containing slide content

        Returns:
            Dictionary with table content info or None if no table found
        """
        from ..utils.logging import debug_print

        # Check multiple possible locations for table content
        content_fields_to_check = [
            ("content", slide_data.get("content", "")),
            ("table_data", slide_data.get("table_data", "")),
            ("table", slide_data.get("table", "")),
            ("body", slide_data.get("body", "")),
            ("text", slide_data.get("text", "")),
        ]

        # Also check placeholders object
        placeholders = slide_data.get("placeholders", {})
        for field_name, field_value in placeholders.items():
            if isinstance(field_value, str):
                content_fields_to_check.append((f"placeholders.{field_name}", field_value))

        debug_print(f"  Checking {len(content_fields_to_check)} fields for table content")

        for field_name, field_value in content_fields_to_check:
            if isinstance(field_value, str) and self._is_table_markdown(field_value):
                debug_print(f"  Table markdown found in field: {field_name}")
                return {"source_field": field_name, "markdown": field_value, "table_data": None}
            elif isinstance(field_value, dict) and field_value.get("type") == "table":
                debug_print(f"  Table data object found in field: {field_name}")
                return {"source_field": field_name, "markdown": None, "table_data": field_value}

        debug_print("  No table content found in slide_data")
        return None

    def _clear_all_table_placeholders(self, slide, slide_data, table_content):
        """
        ENHANCED: Clear ALL placeholders that might contain table content.

        Args:
            slide: PowerPoint slide object
            slide_data: Dictionary containing slide content
            table_content: Information about where table content was found
        """
        from ..utils.logging import debug_print

        debug_print(f"  Clearing table placeholders (source: {table_content['source_field']})")

        # Clear all content placeholders that might have table content
        for shape in slide.placeholders:
            if hasattr(shape, "text_frame") and shape.text_frame:
                try:
                    # Check if this placeholder contains table content
                    shape_text = shape.text_frame.text if shape.text_frame.text else ""
                    if shape_text and self._is_table_markdown(shape_text):
                        debug_print(f"    Clearing table content from placeholder idx {shape.placeholder_format.idx}")
                        shape.text_frame.clear()
                        # Ensure completely empty
                        if hasattr(shape, "text"):
                            shape.text = ""
                except Exception as e:
                    debug_print(f"    Warning: Could not clear placeholder {getattr(shape.placeholder_format, 'idx', 'unknown')}: {e}")
                    continue

    def _mark_table_content_as_processed(self, slide_data, source_field):
        """
        ENHANCED: Mark field as processed to prevent content processing from also creating tables.

        Args:
            slide_data: Dictionary containing slide content
            source_field: Field name that contained the table content
        """
        # Add a marker to prevent content processing from processing this field
        if "_processed_table_fields" not in slide_data:
            slide_data["_processed_table_fields"] = set()
        slide_data["_processed_table_fields"].add(source_field)

    def _contains_table_content(self, text_content):
        """
        Legacy method name - redirects to _is_table_markdown for backwards compatibility.
        """
        return self._is_table_markdown(text_content)

    def _create_dynamic_content_shapes(self, slide, slide_data, content_formatter):
        """
        Create dynamic text shapes and tables for mixed content segments.

        This method processes content segments created by converter's intelligent splitting,
        creating additional text shapes and properly positioned tables as needed.

        Args:
            slide: PowerPoint slide object
            slide_data: Dictionary containing _content_segments
            content_formatter: ContentFormatter instance for text formatting
        """
        from .table_builder import TableBuilder

        content_segments = slide_data.get("_content_segments", [])
        if not content_segments:
            debug_print("    No content segments found for dynamic shape creation")
            return

        debug_print(f"    Creating dynamic shapes for {len(content_segments)} content segments")

        # Find the content placeholder to use as reference for positioning
        content_placeholder = None
        for shape in slide.placeholders:
            if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.BODY:
                content_placeholder = shape
                break

        if not content_placeholder:
            # Fallback: look for any content-type placeholder
            for shape in slide.placeholders:
                if hasattr(shape, "text_frame"):
                    content_placeholder = shape
                    break

        if not content_placeholder:
            error_print("    Warning: No content placeholder found for dynamic shape positioning")
            return

        # Calculate positioning parameters using placeholder properties
        start_left = content_placeholder.left
        start_top = content_placeholder.top
        available_width = content_placeholder.width

        # Get font-based sizing from placeholder
        base_font_size = self._get_placeholder_font_size(content_placeholder)
        line_height = base_font_size * 1.2  # Standard 1.2 line spacing

        # Dynamic spacing based on font size
        text_shape_height = line_height * 2  # Height for 2 lines of text
        spacing_between_shapes = Cm(0.5)  # 0.5cm gap as requested

        # Calculate initial positioning - start after the content placeholder
        # The content placeholder now contains the first text segment, so position after it
        placeholder_bottom = start_top + content_placeholder.height
        current_top = placeholder_bottom + spacing_between_shapes

        # Process each content segment
        table_builder = TableBuilder(content_formatter)
        segment_index = 0

        for segment in content_segments:
            if segment["type"] == "text":
                segment_index += 1

                # First text segment is handled by the main content placeholder
                # but we need to account for its space when positioning subsequent elements
                if segment_index == 1:
                    # Calculate height needed for this first text segment
                    text_lines = segment["content"].count("\n") + 1
                    first_segment_height = line_height * text_lines
                    current_top += first_segment_height + spacing_between_shapes
                    debug_print(f"    First text segment in main placeholder, height: {first_segment_height}, next position: {current_top}")
                    continue

                # Create additional text shape for subsequent text segments (ensure integer EMU values)
                text_shape = slide.shapes.add_textbox(int(start_left), int(current_top), int(available_width), int(text_shape_height))

                # Apply content formatting
                content_formatter.apply_inline_formatting(segment["content"], text_shape.text_frame.paragraphs[0])

                debug_print(f"    Created text shape {segment_index} at position {current_top}")
                current_top += text_shape_height + spacing_between_shapes

            elif segment["type"] == "table":
                # Create table shape positioned at current location
                table_data = segment["table_data"]

                # Create positioned table creation function

                def positioned_table_creation(slide, table_data, table_top):
                    """Override table positioning for dynamic layout"""
                    # Get table data
                    data = table_data.get("data", table_data.get("rows", []))
                    if not data:
                        return

                    # Calculate table dimensions based on content and row height
                    rows = len(data)
                    cols = len(data[0]) if data else 1

                    # Get row height from table data or use font-based default
                    configured_row_height = table_data.get("row_height")

                    if configured_row_height:
                        try:
                            row_height_cm = float(configured_row_height)
                            table_height = Cm(row_height_cm * rows + 0.5)  # Add padding
                        except (ValueError, TypeError):
                            table_height = line_height * rows * 1.5  # Font-based fallback
                    else:
                        table_height = line_height * rows * 1.5  # Font-based default

                    # Create table at calculated position (ensure integer EMU values)
                    table = slide.shapes.add_table(rows, cols, int(start_left), int(table_top), int(available_width), int(table_height)).table

                    # Apply table data and styling with font-aware sizing
                    # Use placeholder font size as base for table fonts
                    font_size_pt = int(base_font_size.pt) if hasattr(base_font_size, "pt") else 12

                    # Add font sizing to table data if not specified
                    if not table_data.get("header_font_size"):
                        table_data["header_font_size"] = font_size_pt
                    if not table_data.get("data_font_size"):
                        table_data["data_font_size"] = max(8, font_size_pt - 2)  # Slightly smaller for data

                    # Apply row height if configured
                    if configured_row_height:
                        try:
                            row_height_emu = Cm(float(configured_row_height))
                            table_builder._apply_row_heights(table, row_height_emu)
                        except (ValueError, TypeError) as e:
                            debug_print(f"    Warning: Could not apply row height: {e}")

                    table_builder._apply_table_data_and_styling(table, table_data)

                    return table_height

                # Create the table
                try:
                    table_height = positioned_table_creation(slide, table_data, current_top)
                    debug_print(f"    Created table at position {current_top}, height: {table_height}")
                    current_top += table_height + spacing_between_shapes
                except Exception as e:
                    error_print(f"    Error creating dynamic table: {e}")
                    # Fallback to standard table creation
                    table_builder.add_table_to_slide(slide, table_data)

        debug_print(f"    Dynamic shape creation completed, final position: {current_top}")

    def _apply_content_to_mapped_placeholders_selective(self, slide, slide_data, layout_name, content_formatter, image_placeholder_handler, skip_content=False):
        """
        Selective version of _apply_content_to_mapped_placeholders that can skip content placeholders.
        Used when dynamic shape processing handles content but we still need to process title, subtitle, etc.
        """
        if not self.layout_mapping:
            return

        # Get layout info from template mapping
        layouts = self.layout_mapping.get("layouts", {})
        layout_info = layouts.get(layout_name, {})
        placeholder_mappings = layout_info.get("placeholders", {})

        # Create reverse mapping: field_name -> placeholder_index
        field_to_placeholder = {}
        for field_name, placeholder_index in placeholder_mappings.items():
            field_to_placeholder[field_name] = placeholder_index

        # Get available placeholders on the slide
        placeholders_by_index = {}
        for shape in slide.placeholders:
            placeholders_by_index[shape.placeholder_format.idx] = shape

        # Process each field in slide_data, mapping to placeholders
        processed_fields = set()

        # First priority: Process placeholders section
        if "placeholders" in slide_data:
            placeholders = slide_data["placeholders"]
            for field_name, content in placeholders.items():
                # For dynamic shapes, completely skip content placeholders during template mapping
                # The cleanup process will handle content properly
                if skip_content and field_name == "content":
                    debug_print(f"    Skipping content field '{field_name}' - will be handled by dynamic shape cleanup")
                    continue

                placeholder_index = field_to_placeholder.get(field_name)
                if placeholder_index is not None and placeholder_index in placeholders_by_index:
                    placeholder = placeholders_by_index[placeholder_index]
                    try:
                        self._apply_content_to_single_placeholder(slide, placeholder, field_name, content, slide_data, content_formatter, image_placeholder_handler)
                        processed_fields.add(field_name)
                    except Exception as e:
                        error_print(f"    Error applying content to placeholder {field_name}: {e}")

        # Second priority: Process direct fields in slide_data
        for field_name, content in slide_data.items():
            if field_name in ["placeholders", "layout", "style", "_content_segments", "_requires_dynamic_shapes", "_table_count"] or field_name.startswith("_"):
                continue
            if field_name in processed_fields:
                continue
            if skip_content and field_name == "content":
                debug_print(f"    Skipping content field '{field_name}' - will be handled by dynamic shape cleanup")
                continue

            placeholder_index = field_to_placeholder.get(field_name)
            if placeholder_index is not None and placeholder_index in placeholders_by_index:
                placeholder = placeholders_by_index[placeholder_index]
                try:
                    self._apply_content_to_single_placeholder(slide, placeholder, field_name, content, slide_data, content_formatter, image_placeholder_handler)
                except Exception as e:
                    error_print(f"    Error applying content to placeholder {field_name}: {e}")

    def _find_content_placeholder(self, slide):
        """
        Find the content placeholder in the slide.

        Args:
            slide: PowerPoint slide object

        Returns:
            Content placeholder shape or None if not found
        """
        for shape in slide.placeholders:
            try:
                if is_content_placeholder(shape.placeholder_format.type):
                    return shape
            except Exception:  # nosec B112
                continue
        return None

    def _find_title_placeholder(self, slide):
        """
        Find the title placeholder in the slide.

        Args:
            slide: PowerPoint slide object

        Returns:
            Title placeholder shape or None if not found
        """
        for shape in slide.placeholders:
            try:
                if is_title_placeholder(shape.placeholder_format.type):
                    return shape
            except Exception:  # nosec B112
                continue
        return None

    def _get_placeholder_font_size(self, placeholder):
        """
        Extract font size from placeholder for consistent sizing.

        Args:
            placeholder: PowerPoint placeholder shape

        Returns:
            Font size in EMU units
        """
        from pptx.util import Pt

        try:
            if hasattr(placeholder, "text_frame") and placeholder.text_frame:
                # Try to get font size from runs first
                for paragraph in placeholder.text_frame.paragraphs:
                    if paragraph.runs:
                        for run in paragraph.runs:
                            if run.font.size:
                                debug_print(f"    Found font size from run: {run.font.size.pt}pt")
                                return run.font.size

                # Try to get from paragraph font
                for paragraph in placeholder.text_frame.paragraphs:
                    if paragraph.font and paragraph.font.size:
                        debug_print(f"    Found font size from paragraph: {paragraph.font.size.pt}pt")
                        return paragraph.font.size

                # Try template default - check slide master
                try:
                    # This might give us template defaults, but it's complex
                    _ = placeholder._element.part.slide_master
                    debug_print("    No explicit font size found, using intelligent default")
                except Exception:  # nosec B110
                    pass

        except Exception as e:
            debug_print(f"    Error getting placeholder font size: {e}")

        # Intelligent fallback - use a reasonable size for presentations
        debug_print("    Using fallback font size: 18pt")
        return Pt(18)  # Slightly larger default for presentations

    def _estimate_placeholder_content_height(self, placeholder):
        """
        Estimate the height occupied by content in a placeholder.

        Args:
            placeholder: PowerPoint placeholder shape

        Returns:
            Estimated height in EMU units
        """
        if not hasattr(placeholder, "text_frame") or not placeholder.text_frame.text:
            return Cm(1)  # Minimal height for empty placeholder

        # Rough estimation based on text length and line breaks
        text_content = placeholder.text_frame.text
        line_count = len(text_content.split("\n"))

        # Estimate: ~0.5cm per line of text
        estimated_height = Cm(0.5 * max(line_count, 1) + 0.5)  # +0.5cm padding

        return estimated_height
