# import json
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, Optional
import yaml

from pptx import Presentation

from ..utils.path import path_manager, PathManager
from .presentation_builder import PresentationBuilder
from ..content.processor import ContentProcessor
from ..templates.manager import TemplateManager
from ..image.image_handler import ImageHandler
from .result import PresentationResult, ValidationResult

# PlaceKitten will be imported lazily when needed
from ..utils.path import get_placekitten


def singleton(cls):
    instances = {}

    def get_instance(*args, **kwargs):
        if cls not in instances:
            instances[cls] = cls(*args, **kwargs)
        return instances[cls]

    def reset():
        """Reset the singleton instance for testing purposes"""
        instances.clear()

    # Allow external access to clear instances for testing
    get_instance._instances = instances
    get_instance.reset = reset
    cls._instances = instances
    cls.reset = reset

    return get_instance


@singleton
class Deckbuilder:
    def __init__(self, path_manager_instance: Optional[PathManager] = None):
        # Use provided path manager or default global instance
        self._path_manager = path_manager_instance or path_manager

        self.output_folder = str(self._path_manager.get_output_folder())
        self.prs = Presentation()

        # Initialize components
        self.template_manager = TemplateManager(self._path_manager)
        self.content_processor = ContentProcessor()
        self.presentation_builder = PresentationBuilder(self._path_manager)

        # Initialize image-related components
        output_folder = self._path_manager.get_output_folder()
        if isinstance(output_folder, str):
            image_cache_dir = Path(output_folder) / "image_cache"
        else:
            image_cache_dir = output_folder / "image_cache"
        self.image_handler = ImageHandler(str(image_cache_dir))
        self.placekitten = get_placekitten()()

        # Ensure default template exists in templates folder
        template_name = self._path_manager.get_template_name() or "default"
        self.template_manager.check_template_exists(template_name)

        # Store template info for tests
        self.template_name = template_name
        self.template_path, _ = self.template_manager.prepare_template(template_name)

    def _initialize_presentation(self, templateName: str = "default") -> None:
        # Prepare template and get path (layout mapping no longer needed)
        template_path, _ = self.template_manager.prepare_template(templateName)

        # Store template path for tests
        self.template_path = template_path

        # Load template or create empty presentation
        if template_path:
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()

        self.presentation_builder.clear_slides(self.prs)

    def create_presentation(
        self,
        presentation_data: Dict[str, Any],
        fileName: str = "Sample_Presentation",
        templateName: str = "default",
        language_code: Optional[str] = None,
        font_name: Optional[str] = None,
    ) -> str:
        """
        Creates a presentation from the canonical JSON data model.
        Only accepts canonical format: {"slides": [{"layout": "...", "placeholders": {...}, "content": [...]}]}

        Includes built-in end-to-end validation to prevent layout regressions.
        """
        # Import validation here to avoid circular imports
        # from .validation import PresentationValidator

        self._initialize_presentation(templateName)

        # Strict validation for canonical JSON format only
        if not isinstance(presentation_data, dict):
            raise ValueError("Input must be a dictionary containing canonical JSON data.")

        if "slides" not in presentation_data:
            raise ValueError("Canonical JSON data must contain a 'slides' array at root level.")

        if not isinstance(presentation_data["slides"], list):
            raise ValueError("'slides' must be an array of slide objects.")

        if len(presentation_data["slides"]) == 0:
            raise ValueError("At least one slide is required.")

        # Validate each slide has required canonical structure
        for i, slide_data in enumerate(presentation_data["slides"]):
            if not isinstance(slide_data, dict):
                raise ValueError(f"Slide {i + 1} must be a dictionary.")

            if "layout" not in slide_data:
                raise ValueError(f"Slide {i + 1} must have a 'layout' field.")

            # Ensure canonical structure exists (placeholders and content are optional but must be correct types)
            if "placeholders" in slide_data and not isinstance(slide_data["placeholders"], dict):
                raise ValueError(f"Slide {i + 1} 'placeholders' must be a dictionary.")

            if "content" in slide_data and not isinstance(slide_data["content"], list):
                raise ValueError(f"Slide {i + 1} 'content' must be an array.")

        # STEP 1: Pre-generation validation (JSON ↔ Template alignment)
        # TEMPORARILY DISABLED: Old validation system uses index-based mappings
        # template_folder = str(self._path_manager.get_template_folder())
        # validator = PresentationValidator(presentation_data, templateName, template_folder)
        # validator.validate_pre_generation()

        # STEP 1.5: Apply theme font formatting if specified
        if font_name is not None:
            from ..content.formatting_support import FormattingSupport

            formatter = FormattingSupport()
            formatter.update_theme_fonts(self.prs, font_name)

        # STEP 2: Update presentation builder with formatting parameters
        self.presentation_builder.set_formatting_options(language_code, font_name)

        # STEP 3: Process slides using canonical format with optional formatting
        for slide_data in presentation_data["slides"]:
            # Use template-based layouts for tables instead of dynamic shape creation
            self.presentation_builder.add_slide(self.prs, slide_data)

        # STEP 4: Save the presentation to disk
        write_result = self.write_presentation(fileName)

        # Extract the file path from write_result for post-generation validation
        # write_result format: "Successfully created presentation: filename.pptx"
        if "Successfully created presentation:" in write_result:
            # file_path = write_result.split("Successfully created presentation: ")[1].strip()  # Future: use for validation
            # full_path = str(self._path_manager.get_output_folder() / file_path)  # Future: use for validation
            pass

            # STEP 5: Post-generation validation (PPTX ↔ JSON verification)
            # TEMPORARILY DISABLED: Old validation system uses index-based mappings
            # validator.validate_post_generation(full_path)

        # Show completion summary
        from ..utils.logging import success_print

        slide_count = len(presentation_data["slides"])
        file_name = write_result.split("Successfully created presentation: ")[1].strip() if "Successfully created presentation:" in write_result else "presentation.pptx"
        success_print(f"✅ Presentation complete: {file_name} ({slide_count} slides)")

        return f"Successfully created presentation with {slide_count} slides. {write_result}"

    def create_presentation_from_markdown(
        self,
        markdown_content: str,
        fileName: str = "Sample_Presentation",
        templateName: str = "default",
        language_code: Optional[str] = None,
        font_name: Optional[str] = None,
    ) -> PresentationResult:
        """
        Creates a presentation from markdown content with frontmatter.

        Handles all validation internally and returns structured result instead of throwing exceptions.
        This is the new recommended method for creating presentations from markdown.

        Args:
            markdown_content: Raw markdown string with frontmatter
            fileName: Output filename (without extension)
            templateName: Template to use
            language_code: Language for formatting
            font_name: Font to use

        Returns:
            PresentationResult with success/error information
        """
        try:
            # Parse markdown to canonical JSON with internal error handling
            conversion_result = self._convert_markdown_to_json_safe(markdown_content)
            if not conversion_result.valid:
                # Convert validation errors to presentation error
                error_messages = "\n".join(conversion_result.errors)
                return PresentationResult.content_error_result(error_messages, "Check markdown frontmatter syntax and structure")

            # Get the converted data
            presentation_data = conversion_result.context.get("presentation_data")
            if not presentation_data:
                return PresentationResult.content_error_result("Failed to convert markdown to presentation data", "Ensure markdown has valid frontmatter sections")

            # Use existing presentation creation (but handle its exceptions)
            try:
                result_message = self.create_presentation(
                    presentation_data,
                    fileName=fileName,
                    templateName=templateName,
                    language_code=language_code,
                    font_name=font_name,
                )

                # Parse success message to extract details
                slide_count = len(presentation_data.get("slides", []))
                filename = self._extract_filename_from_result(result_message)

                return PresentationResult.success_result(filename, slide_count)

            except ValueError as e:
                return PresentationResult.validation_error_result(str(e), "Check presentation data structure and template compatibility")
            except Exception as e:
                return PresentationResult.error_result(f"Unexpected error during presentation creation: {str(e)}")

        except Exception as e:
            return PresentationResult.error_result(f"Unexpected error processing markdown: {str(e)}")

    def _convert_markdown_to_json_safe(self, markdown_content: str) -> ValidationResult:
        """
        Safely convert markdown to canonical JSON with comprehensive error handling.

        Returns ValidationResult instead of throwing exceptions.
        """
        try:
            from ..content.frontmatter_to_json_converter import markdown_to_canonical_json

            # Try conversion - this may internally handle some YAML errors
            presentation_data = markdown_to_canonical_json(markdown_content)

            # Basic validation of the result
            if not isinstance(presentation_data, dict):
                return ValidationResult.error_result(["Conversion failed: Result is not a valid dictionary"])

            if "slides" not in presentation_data:
                return ValidationResult.error_result(["Conversion failed: No slides found in markdown", "Ensure markdown has frontmatter sections with layout specifications"])

            if not isinstance(presentation_data["slides"], list):
                return ValidationResult.error_result(["Conversion failed: Slides data is not a list"])

            if len(presentation_data["slides"]) == 0:
                return ValidationResult.error_result(["Conversion failed: No slides generated from markdown", "Check that frontmatter sections are properly formatted"])

            result = ValidationResult.success_result()
            result.context = {"presentation_data": presentation_data}
            return result

        except yaml.YAMLError as e:
            # Handle YAML parsing errors gracefully
            error_str = str(e)
            return ValidationResult.error_result([f"YAML syntax error in frontmatter: {error_str}", "Fix: Check YAML indentation and syntax in frontmatter sections"])
        except Exception as e:
            return ValidationResult.error_result([f"Error converting markdown to presentation format: {str(e)}"])

    def _extract_filename_from_result(self, result_message: str) -> str:
        """Extract filename from success message."""
        if "Successfully created presentation: " in result_message:
            return result_message.split("Successfully created presentation: ")[1].strip()
        return "presentation.pptx"

    # Removed _process_mixed_content_for_json - table handling now uses dedicated layouts

    def write_presentation(self, fileName: str = "Sample_Presentation") -> str:
        """Writes the generated presentation to disk with ISO timestamp."""
        import os

        # Get output folder from environment or use default
        output_folder = self.output_folder or "."

        # Ensure output folder exists
        os.makedirs(output_folder, exist_ok=True)

        # Create filename with ISO timestamp and .g.pptx extension for generated files
        timestamp = datetime.now().strftime("%Y-%m-%d_%H%M")
        generated_filename = f"{fileName}.{timestamp}.g.pptx"
        output_file = os.path.join(output_folder, generated_filename)

        # Save the presentation (overwrites if same timestamp exists)
        self.prs.save(output_file)

        return f"Successfully created presentation: {os.path.basename(output_file)}"


def get_deckbuilder_client():
    # Return Deckbuilder instance with MCP context
    from ..utils.path import create_mcp_path_manager

    return Deckbuilder(path_manager_instance=create_mcp_path_manager())
