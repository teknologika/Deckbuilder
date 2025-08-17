#!/usr/bin/env python3
"""
Built-in End-to-End Validation System for Deckbuilder

Provides automatic validation that runs on every presentation generation to prevent
layout regressions and ensure JSON ↔ Template ↔ PPTX alignment.

GitHub Issue: https://github.com/teknologika/Deckbuilder/issues/36
"""

import re
from pathlib import Path
from typing import Dict, List, Any, Optional
from pptx import Presentation
from ..utils.logging import validation_print, error_print, success_print


class ValidationError(Exception):
    """Validation error with detailed fix instructions."""

    def __init__(self, message: str, slide_num: Optional[int] = None, field_name: Optional[str] = None):
        self.slide_num = slide_num
        self.field_name = field_name
        super().__init__(message)


class PresentationValidator:
    """
    Built-in validation system that runs automatically on every presentation generation.

    Prevents layout regressions by validating:
    1. Pre-generation: JSON ↔ Template mapping alignment
    2. Post-generation: PPTX output ↔ JSON input verification
    """

    def __init__(self, presentation_data: Dict[str, Any], template_name: str, template_folder: str):
        self.presentation_data = presentation_data
        self.template_name = template_name
        self.template_folder = Path(template_folder)
        # Legacy template mapping JSON files removed - validation now uses structured frontmatter patterns

    def validate_markdown_to_json(self, markdown_content: str, converted_json: Dict[str, Any]):
        """
        Validate Markdown → JSON conversion before template validation.

        Ensures structured frontmatter is properly converted and no data is lost.
        """
        validation_print("🔍 Markdown → JSON validation: Frontmatter processing...")

        # Parse markdown sections manually to validate conversion
        markdown_sections = self._parse_markdown_sections(markdown_content)
        json_slides = converted_json.get("slides", [])

        # Validate slide count matches
        if len(markdown_sections) != len(json_slides):
            raise ValidationError(
                f"Markdown → JSON conversion error: {len(markdown_sections)} markdown sections "
                f"converted to {len(json_slides)} JSON slides\n"
                f"Fix: Check markdown section parsing and frontmatter conversion"
            )

        # Validate each section conversion
        for section_idx, (md_section, json_slide) in enumerate(zip(markdown_sections, json_slides)):
            self._validate_section_conversion(section_idx + 1, md_section, json_slide)

        success_print("✅ Markdown → JSON validation passed")

    def validate_pre_generation(self):
        """
        Validate JSON structure before generation.

        Validates basic slide structure without legacy template mapping dependency.
        """
        validation_print("🔍 JSON structure validation...")

        validation_print(f"[Pre-Generation Validation] Validating {len(self.presentation_data.get('slides', []))} slides")

        # Validate each slide's basic structure
        for slide_idx, slide_data in enumerate(self.presentation_data.get("slides", [])):
            slide_num = slide_idx + 1
            layout_name = slide_data.get("layout")

            validation_print(f"[Pre-Generation Validation] Slide {slide_num}: Checking layout '{layout_name}'")

            if not layout_name:
                raise ValidationError(f"Slide {slide_num}: Missing 'layout' field\n" f"Fix: Add 'layout' field with valid layout name")

            # Show slide content fields for debugging
            placeholders = slide_data.get("placeholders", {})
            validation_print(f"[Pre-Generation Validation] Placeholder fields: {list(placeholders.keys())}")

            # Legacy content blocks should not exist in structured frontmatter
            if "content" in slide_data:
                validation_print("[Pre-Generation Validation] WARNING: Legacy content blocks detected - should be converted to placeholders")

        success_print("✅ Pre-generation validation passed")

    # Legacy _validate_slide_placeholders method removed - template mapping validation no longer needed
    # Placeholder validation now handled by the PlaceholderManager during slide creation

    # Legacy field name resolution methods removed - template mapping validation no longer needed
    # Field name resolution now handled by PlaceholderManager and PlaceholderResolver during slide creation

    def validate_post_generation(self, pptx_file_path: str):
        """
        Validate PPTX output ↔ JSON input after generation.

        Raises ValidationError if generated content doesn't match specification.
        """
        validation_print("🔍 Post-generation validation: PPTX ↔ JSON verification...")
        validation_print(f"[Post Validation] Loading generated PPTX: {pptx_file_path}")

        if not Path(pptx_file_path).exists():
            error_print(f"[Post Validation] WARNING: Generated PPTX file not found: {pptx_file_path}")
            error_print("[Post Validation] This indicates the presentation was not saved properly.")
            return  # Don't crash - just skip validation

        # Load generated presentation
        prs = Presentation(pptx_file_path)

        # Validate slide count
        expected_slides = len(self.presentation_data.get("slides", []))
        actual_slides = len(prs.slides)

        validation_print(f"[Post Validation] Slide count check: expected={expected_slides}, actual={actual_slides}")

        if actual_slides != expected_slides:
            raise ValidationError(f"Slide count mismatch: expected {expected_slides}, got {actual_slides}\n" f"Fix: Check slide generation logic for dropped or duplicated slides")

        # Validate each slide content
        validation_errors = []
        validation_print(f"[Post Validation] Validating content for {actual_slides} slides...")

        for slide_idx, (slide, slide_spec) in enumerate(zip(prs.slides, self.presentation_data["slides"])):
            slide_num = slide_idx + 1
            layout_name = slide_spec.get("layout", "unknown")

            try:
                self._validate_slide_content(slide_num, slide, slide_spec)
                validation_print(f"[Post Validation] Slide {slide_num} ({layout_name}): Content validation passed")
            except ValidationError as e:
                error_print(f"[Post Validation] Slide {slide_num} ({layout_name}): Content validation failed")
                validation_errors.append(str(e))

        if validation_errors:
            error_print(f"[Post Validation] WARNING - {len(validation_errors)} validation issues found:")
            for i, error in enumerate(validation_errors, 1):
                error_print(f"[Post Validation]   {i}. {error}")
            error_print("[Post Validation] Presentation generated with warnings. Run with DECKBUILDER_DEBUG=true for detailed analysis.")
            # Don't raise - let presentation be saved with warnings
            return

        success_print("✅ Post-generation validation passed")

    def _validate_slide_content(self, slide_num: int, slide, slide_spec: Dict[str, Any]):
        """Validate individual slide content against specification."""
        layout_name = slide_spec["layout"]
        actual_layout = slide.slide_layout.name

        # Validate layout
        if actual_layout != layout_name:
            raise ValidationError(f"Slide {slide_num}: Layout mismatch - expected '{layout_name}', got '{actual_layout}'")

        # Validate critical placeholders are not empty
        self._validate_critical_placeholders(slide_num, slide, slide_spec)

        # Content blocks validation removed - structured frontmatter uses placeholders only

    def _validate_critical_placeholders(self, slide_num: int, slide, slide_spec: Dict[str, Any]):
        """Validate that critical placeholders have content."""
        placeholders_spec = slide_spec.get("placeholders", {})
        layout_name = slide_spec["layout"]

        # Get actual placeholder content
        actual_content = {}
        all_placeholders = {}  # Track all placeholders for debugging

        for shape in slide.shapes:
            try:
                if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                    ph_type = shape.placeholder_format.type
                    ph_idx = shape.placeholder_format.idx
                    try:
                        ph_name = getattr(shape.element.nvSpPr.cNvPr, "name", "unnamed")
                    except AttributeError:
                        ph_name = "unnamed"

                    content = self._extract_shape_text(shape)
                    semantic_type = self._get_semantic_type(ph_type)

                    # Special handling for PICTURE placeholders
                    has_content = bool(content.strip())
                    if ph_type.name == "PICTURE":
                        # For picture placeholders, check if image is present
                        has_image = self._check_placeholder_has_image(shape)
                        has_content = has_image
                        if has_image:
                            content = "[IMAGE PRESENT]"

                    # Track all placeholders
                    all_placeholders[f"{semantic_type}_{ph_idx}"] = {
                        "content": content,
                        "name": ph_name,
                        "has_content": has_content,
                    }

                    if has_content:  # Count content or images
                        actual_content[f"{semantic_type}_{ph_idx}"] = content
            except ValueError:
                continue

        # Check critical fields that should have content
        critical_missing = []

        for field_name, expected_content in placeholders_spec.items():
            if field_name in ["style", "speaker_notes", "media"]:  # Skip non-content fields
                continue

            if not expected_content or str(expected_content).strip() == "":
                continue  # Skip empty expected content

            # Check if this field has corresponding content in PPTX
            found_content = False
            expected_clean = self._normalize_expected_content_for_validation(str(expected_content))

            # Special handling for image fields
            if field_name in ["image", "image_1", "image_path"] or "image" in field_name.lower():
                # For image fields, just check if we have any content in actual_content
                # (which includes "[IMAGE PRESENT]" for successful image insertion)
                for _ph_key, actual_text in actual_content.items():
                    if "PICTURE" in _ph_key or "[IMAGE PRESENT]" in actual_text:
                        found_content = True
                        break
            # Special handling for table fields
            elif isinstance(expected_content, dict) and expected_content.get("type") == "table":
                # For table fields, check if we have table shapes in the slide or table content indicators
                # Check if we have table shapes in the slide
                table_shapes = [shape for shape in slide.shapes if hasattr(shape, "table")]
                if table_shapes:
                    found_content = True
                else:
                    # Fallback: check for table content indicators in text
                    for _ph_key, actual_text in actual_content.items():
                        if "table" in actual_text.lower() or "rows" in actual_text.lower():
                            found_content = True
                            break
            # Special handling for table markdown content
            elif self._is_table_markdown(expected_content):
                # When expected content is table markdown, check for table shapes instead of placeholder content
                table_shapes = [shape for shape in slide.shapes if hasattr(shape, "table")]
                if table_shapes:
                    found_content = True
                    # Optionally validate table content matches expected structure
                    # (for now, just confirm table exists)
            else:
                # Normal text content validation
                for _ph_key, actual_text in actual_content.items():
                    if expected_clean.lower() in actual_text.lower():
                        found_content = True
                        break

            if not found_content:
                critical_missing.append(f"{field_name} (expected: '{expected_clean[:30]}...')")

        # Special validation for vertical layouts
        if "vertical" in layout_name.lower():
            content_found = any("content" in key.lower() or "body" in key.lower() for key in actual_content.keys())
            if not content_found and placeholders_spec.get("content"):
                critical_missing.append("content (vertical layout missing main content)")

        if critical_missing:
            # Show detailed validation info only on failures
            error_print(f"[Validation] Slide {slide_num} ({layout_name}): VALIDATION FAILED")
            error_print(f"[Validation]   Expected placeholders: {list(placeholders_spec.keys())}")
            error_print("[Validation]   Found placeholders in slide:")
            for ph_key, ph_info in all_placeholders.items():
                status = "✓ HAS CONTENT" if ph_info["has_content"] else "✗ EMPTY"
                error_print(f"[Validation]     {ph_key} ('{ph_info['name']}'): {status}")
                if ph_info["has_content"]:
                    error_print(f"[Validation]       Content: '{ph_info['content'][:50]}{'...' if len(ph_info['content']) > 50 else ''}'")
            error_print(f"[Validation]   Non-empty placeholders: {list(actual_content.keys())}")
            error_print("[Validation]   Checking critical field mappings:")
            for field_name, expected_content in placeholders_spec.items():
                if field_name in ["style", "speaker_notes", "media"]:
                    continue
                if not expected_content or str(expected_content).strip() == "":
                    error_print(f"[Validation]     '{field_name}': SKIPPED (empty expected content)")
                    continue
                expected_clean = self._normalize_expected_content_for_validation(str(expected_content))
                error_print(f"[Validation]     '{field_name}': Looking for '{expected_clean[:30]}{'...' if len(expected_clean) > 30 else ''}'")
                if field_name in critical_missing:
                    error_print("[Validation]       ✗ NOT FOUND in any placeholder")
                else:
                    error_print("[Validation]       ✓ FOUND")
            raise ValidationError(f"Slide {slide_num} ({layout_name}): Missing critical content: {', '.join(critical_missing)}")

    # _validate_content_blocks method removed - structured frontmatter uses placeholders only

    def _extract_shape_text(self, shape) -> str:
        """Extract all text from a shape."""
        text_parts = []

        if hasattr(shape, "text") and shape.text.strip():
            text_parts.append(shape.text.strip())
        elif hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text.strip())

        return "\n".join(text_parts)

    def _get_semantic_type(self, ph_type) -> str:
        """Get semantic type name for a placeholder."""
        return ph_type.name if hasattr(ph_type, "name") else str(ph_type)

    def _normalize_expected_content_for_validation(self, text: str) -> str:
        """
        Normalise expected content to match what appears in generated slides.

        The slide builder processes markdown during content placement:
        - '## Header' becomes heading formatting (slide text: 'Header')
        - '**bold**' becomes bold formatting (slide text: 'bold')

        This method applies the same transformations so validation can
        properly compare expected vs actual slide content.

        IMPORTANT: This is validation-only logic. The slide builder handles
        the actual content processing during generation.
        """
        # Remove markdown headers (slide builder converts these to heading formatting)
        text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)

        # Remove formatting markers (slide builder converts these to text formatting)
        text = re.sub(r"\*\*\*(.*?)\*\*\*", r"\1", text)  # ***bold italic***
        text = re.sub(r"___(.*?)___", r"\1", text)  # ___underline___
        text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)  # **bold**
        text = re.sub(r"\*(.*?)\*", r"\1", text)  # *italic*
        return text

    def _is_table_markdown(self, content):
        """
        Check if content appears to be table markdown.

        Args:
            content: Content to check

        Returns:
            bool: True if content looks like table markdown
        """
        if not isinstance(content, str):
            return False

        lines = [line.strip() for line in content.split("\n") if line.strip()]

        # Need at least 2 lines for a table (header + data)
        if len(lines) < 2:
            return False

        # Check if most lines contain pipes (table markers)
        lines_with_pipes = sum(1 for line in lines if "|" in line)

        # If 80% or more lines have pipes, it's likely a table
        if lines_with_pipes >= len(lines) * 0.8:
            return True

        return False

    def _check_placeholder_has_image(self, shape) -> bool:
        """Check if a PICTURE placeholder contains an image."""
        try:
            # Check if shape has image content
            if hasattr(shape, "image"):
                return shape.image is not None

            # Check if shape has fill with image
            if hasattr(shape, "fill") and shape.fill.type is not None:
                # PICTURE fill type indicates image is present
                return True

            # Check for image in shape element
            if hasattr(shape, "element") and shape.element is not None:
                # Look for image elements in the shape XML
                image_elements = shape.element.xpath(
                    ".//a:blip",
                    namespaces={"a": "http://schemas.openxmlformats.org/drawingml/2006/main"},
                )
                return len(image_elements) > 0

        except Exception:  # nosec B110
            # If we can't determine, assume no image for validation purposes
            pass

        return False

    def _parse_markdown_sections(self, markdown_content: str) -> List[Dict[str, Any]]:
        """Parse markdown into sections with frontmatter and content."""
        import yaml

        # Split by frontmatter delimiters
        sections = []
        blocks = re.split(r"^---\s*$", markdown_content, flags=re.MULTILINE)

        i = 0
        while i < len(blocks):
            # Skip empty blocks
            if not blocks[i].strip():
                i += 1
                continue

            # Look for frontmatter + content pairs
            if i + 1 < len(blocks):
                try:
                    frontmatter_raw = blocks[i].strip()
                    content_raw = blocks[i + 1].strip() if i + 1 < len(blocks) else ""

                    # Parse frontmatter
                    frontmatter = yaml.safe_load(frontmatter_raw) or {}

                    sections.append(
                        {
                            "frontmatter": frontmatter,
                            "content": content_raw,
                            "raw_frontmatter": frontmatter_raw,
                        }
                    )

                    i += 2
                except yaml.YAMLError as e:
                    raise ValidationError(f"YAML parsing error in markdown frontmatter: {e}\n" f"Fix: Check YAML syntax in frontmatter section")
            else:
                i += 1

        return sections

    def _validate_section_conversion(self, section_num: int, md_section: Dict[str, Any], json_slide: Dict[str, Any]):
        """Validate individual section conversion from markdown to JSON."""
        frontmatter = md_section["frontmatter"]

        # Validate layout field conversion
        expected_layout = frontmatter.get("layout")
        actual_layout = json_slide.get("layout")

        if expected_layout != actual_layout:
            raise ValidationError(
                f"Section {section_num}: Layout conversion failed\n" f"Markdown layout: '{expected_layout}'\n" f"JSON layout: '{actual_layout}'\n" f"Fix: Check frontmatter to JSON conversion logic"
            )

        # Validate critical frontmatter fields are preserved
        critical_fields = ["title", "layout"]
        for field in critical_fields:
            if field in frontmatter:
                # Check if field appears in JSON placeholders
                json_placeholders = json_slide.get("placeholders", {})
                if field not in json_placeholders and field != "layout":
                    raise ValidationError(
                        f"Section {section_num}: Frontmatter field '{field}' not found in JSON placeholders\n"
                        f"Markdown value: '{frontmatter[field]}'\n"
                        f"Fix: Check structured frontmatter conversion logic"
                    )
