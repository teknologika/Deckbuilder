"""
TableHandler

Handles table detection, parsing, creation, and font configuration.
Tables only support plain text. Markdown formatting is not supported within tables.
Enhanced with font logic moved from table_integration.py and table_builder.py.
"""

from typing import List, Dict, Any, Optional, Tuple

# pptx.util.Pt removed - not used in this module
from ..utils.logging import debug_print, error_print


class TableHandler:
    """
    Handles table detection, parsing, and creation for PowerPoint slides.

    Focuses on plain text processing to eliminate complexity of markdown within
    table cells.
    """

    def __init__(self):
        """Initialize the TableHandler."""
        pass

    def detect_table_content(self, text_content: str) -> bool:
        """
        Detection for markdown table syntax within slide content

        Args:
            text_content: String content to check

        Returns:
            bool: True if content appears to contain table markdown
        """
        if not isinstance(text_content, str) or not text_content.strip():
            return False

        lines = [line.strip() for line in text_content.split("\n") if line.strip()]
        if len(lines) < 2:
            return False

        table_rows = 0

        for line in lines:
            # Table separator line (like |---|---|---| or | --- | --- | --- |)
            if "|" in line and all(c in "|-:= \t" for c in line.replace("|", "").strip()):
                continue  # Skip separator lines
            # Table data row (contains | but has actual content)
            elif "|" in line and not all(c in "|-:= \t" for c in line.strip()):
                table_rows += 1

        # Valid table: at least 2 data rows (header + content) and optionally a separator
        return table_rows >= 2

    def parse_table_structure(self, markdown_content: str) -> List[List[str]]:
        """
        Parse table markdown into plain text structure.

        IMPORTANT: This method processes table cells as PLAIN TEXT ONLY.
        No markdown formatting is parsed within cells - all content is treated
        as literal text to improve performance and reduce complexity.

        Args:
            markdown_content: Markdown table content

        Returns:
            List of lists representing table rows and columns (plain text only)
        """
        if not self.detect_table_content(markdown_content):
            return []

        lines = [line.strip() for line in markdown_content.split("\n") if line.strip()]
        table_data = []

        for line in lines:
            # Skip separator lines (like |---|---|---|)
            if "|" in line and all(c in "|-:= \t" for c in line.replace("|", "").strip()):
                continue

            # Process data rows
            if "|" in line:
                # Split by | and clean up cells - PLAIN TEXT ONLY
                cells = [cell.strip() for cell in line.split("|")]
                # Remove empty cells at start/end (from leading/trailing |)
                if cells and not cells[0]:
                    cells = cells[1:]
                if cells and not cells[-1]:
                    cells = cells[:-1]

                # Store as plain text - no markdown processing
                if cells:
                    table_data.append(cells)

        return table_data

    def create_table_from_data(self, slide, table_data: List[List[str]], position: Tuple[float, float], size: Tuple[float, float]) -> Any:
        """
        Create a PowerPoint table from plain text data.

        Args:
            slide: PowerPoint slide object
            table_data: List of lists containing plain text cell data
            position: (left, top) position tuple in PowerPoint units
            size: (width, height) size tuple in PowerPoint units

        Returns:
            Created table shape object
        """
        if not table_data or len(table_data) < 1:
            error_print("Cannot create table: no table data provided")
            return None

        try:
            # Calculate table dimensions
            rows = len(table_data)
            cols = max(len(row) for row in table_data) if table_data else 0

            if cols == 0:
                error_print("Cannot create table: no columns detected")
                return None

            # Create table shape
            left, top = position
            width, height = size

            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table

            # Populate table with plain text data
            for row_idx, row_data in enumerate(table_data):
                for col_idx, cell_text in enumerate(row_data):
                    if col_idx < len(table.rows[row_idx].cells):
                        # Set cell text as plain text - no formatting applied
                        cell = table.rows[row_idx].cells[col_idx]
                        cell.text = str(cell_text).strip()

            debug_print(f"Created table with {rows} rows and {cols} columns")
            return table_shape

        except Exception as e:
            error_print(f"Error creating table: {e}")
            return None

    def position_table_on_slide(self, slide, content_height: float = None) -> Tuple[float, float]:
        """
        Calculate intelligent table positioning on slide.

        Positions table below existing content with appropriate spacing,
        avoiding overlap with text placeholders.

        Args:
            slide: PowerPoint slide object
            content_height: Height of content above table (optional)

        Returns:
            (left, top) position tuple for table placement
        """
        from pptx.util import Cm

        # Default positioning
        default_left = Cm(1.0)  # 1cm from left edge
        default_top = Cm(8.0)  # 8cm from top (middle of slide)

        if content_height is None:
            return (default_left, default_top)

        # Calculate position based on content height
        # Add spacing below content (minimum 0.5cm)
        spacing = max(Cm(0.5), content_height * 0.1)  # 10% of content height or 0.5cm
        calculated_top = content_height + spacing

        # Ensure table fits on slide (assume slide height ~19cm)
        max_top = Cm(17.0)  # Leave room for table
        final_top = min(calculated_top, max_top)

        debug_print(f"Positioning table at left: {default_left}, top: {final_top}")
        return (default_left, final_top)

    def detect_existing_tables(self, slide) -> List[Any]:
        """
        Ddetect existing tables on a slide.

        Args:
            slide: PowerPoint slide object

        Returns:
            List of existing table shapes
        """
        existing_tables = []
        debug_print(f"Checking {len(slide.shapes)} shapes for existing tables")

        for i, shape in enumerate(slide.shapes):
            try:
                # Check shape_type for TABLE (MSO_SHAPE_TYPE.TABLE = 19)
                if hasattr(shape, "shape_type") and shape.shape_type == 19:
                    debug_print(f"  Found table shape at index {i}")
                    existing_tables.append(shape)

                # Additional check: has table attribute
                elif hasattr(shape, "table"):
                    debug_print(f"  Found shape with table attribute at index {i}")
                    existing_tables.append(shape)

            except Exception as e:
                debug_print(f"  Error checking shape {i}: {e}")
                continue

        debug_print(f"Found {len(existing_tables)} existing tables on slide")
        return existing_tables

    # ===== FONT HANDLING METHODS (moved from table_integration.py and table_builder.py) =====

    def get_default_fonts(self, slide_data: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        Get table font sizes from slide data when explicitly specified.

        Returns only explicitly configured font sizes - does not provide fallbacks.
        Template fonts are preserved when no explicit sizes are specified.

        Args:
            slide_data: Optional slide data containing font specifications

        Returns:
            Dictionary with explicitly specified font sizes (may be empty)
        """
        fonts = {}

        if slide_data:
            # Only include explicitly specified font sizes
            if "header_font_size" in slide_data and slide_data["header_font_size"] is not None:
                fonts["header_font_size"] = slide_data["header_font_size"]
            if "data_font_size" in slide_data and slide_data["data_font_size"] is not None:
                fonts["data_font_size"] = slide_data["data_font_size"]

        if fonts:
            debug_print(f"Explicit table fonts: {fonts}")
        else:
            debug_print("No explicit table fonts - using template defaults")

        return fonts

    def validate_font_sizes(self, font_config: Dict[str, Any]) -> Dict[str, Any]:
        """
        Validate and sanitize font size configuration.

        Moved from table_integration.py for centralized font handling.

        Args:
            font_config: Dictionary containing font size specifications

        Returns:
            Validated font configuration dictionary
        """
        validated_config = font_config.copy()

        # Validate header and data font sizes
        for font_field in ["header_font_size", "data_font_size"]:
            if font_field in validated_config and validated_config[font_field] is not None:
                try:
                    font_size = int(validated_config[font_field])
                    # Reasonable bounds: 8pt to 24pt for headers, 8pt to 20pt for data
                    if font_field == "header_font_size":
                        font_size = max(8, min(24, font_size))
                    else:  # data_font_size
                        font_size = max(8, min(20, font_size))

                    validated_config[font_field] = font_size
                    debug_print(f"Validated {font_field}: {font_size}pt")

                except (ValueError, TypeError):
                    # Invalid font size - remove it to use defaults
                    debug_print(f"Invalid {font_field}: {validated_config[font_field]} - removing")
                    del validated_config[font_field]

        return validated_config

    def calculate_table_fonts(self, slide_data: Dict[str, Any], base_font_size: Optional[Any] = None) -> Dict[str, Any]:
        """
        Calculate appropriate font sizes for table headers and data.

        Moved logic from slide_builder_legacy.py for centralized font handling.

        Args:
            slide_data: Slide data containing font specifications
            base_font_size: Optional base font size from placeholder (for relative sizing)

        Returns:
            Dictionary with calculated font sizes (may be empty to preserve template fonts)
        """
        # Start with explicitly specified fonts
        font_config = self.get_default_fonts(slide_data)

        # Validate any explicit font configurations
        if font_config:
            font_config = self.validate_font_sizes(font_config)

        # If base font size provided and no explicit sizes, calculate relative sizes
        if base_font_size and not font_config.get("header_font_size") and not font_config.get("data_font_size"):
            try:
                base_size_pt = int(base_font_size.pt) if hasattr(base_font_size, "pt") else int(base_font_size)

                # Use base font size for headers, slightly smaller for data
                font_config["header_font_size"] = base_size_pt
                font_config["data_font_size"] = max(8, base_size_pt - 2)  # At least 8pt

                debug_print(f"Calculated relative table fonts from base {base_size_pt}pt: header={font_config['header_font_size']}pt, data={font_config['data_font_size']}pt")

            except (ValueError, AttributeError):
                # No valid base font size - preserve template fonts by returning empty config
                debug_print(f"Invalid base font size: {base_font_size} - preserving template fonts")
                font_config = {}

        return font_config

    def get_font_size_for_row(self, row_idx: int, header_font_size: int, data_font_size: int) -> int:
        """
        Get appropriate font size for table cell based on row type.

        Moved from table_builder.py _get_font_size method.

        Args:
            row_idx: Row index (0 is header)
            header_font_size: Font size for header rows
            data_font_size: Font size for data rows

        Returns:
            Font size in points
        """
        # Validate font sizes with reasonable bounds
        header_size = max(8, min(24, int(header_font_size))) if header_font_size else 12
        data_size = max(8, min(20, int(data_font_size))) if data_font_size else 10

        return header_size if row_idx == 0 else data_size

    def find_table_content_in_slide_data(self, slide_data: Dict[str, Any]) -> Optional[Dict[str, Any]]:
        """
        Find table content in slide data, checking both markdown and structured data.

        Args:
            slide_data: Dictionary containing slide data

        Returns:
            Dictionary with table information or None if no table found
        """
        debug_print("Searching for table content in slide data")

        # Get all content fields to check
        content_fields_to_check = []

        # Check direct fields
        for field_name, field_value in slide_data.items():
            if field_name not in ["layout", "type"]:
                content_fields_to_check.append((field_name, field_value))

        # Check placeholders if they exist
        if "placeholders" in slide_data:
            for field_name, field_value in slide_data["placeholders"].items():
                content_fields_to_check.append((field_name, field_value))

        # Look for table content
        for field_name, field_value in content_fields_to_check:
            if isinstance(field_value, str) and self.detect_table_content(field_value):
                debug_print(f"Table markdown found in field: {field_name}")
                return {"source_field": field_name, "markdown": field_value, "table_data": None}
            elif isinstance(field_value, dict) and field_value.get("type") == "table":
                debug_print(f"Table data object found in field: {field_name}")
                return {"source_field": field_name, "markdown": None, "table_data": field_value}

        debug_print("No table content found in slide data")
        return None

    # TODO: With correct table parsing logic, this should be unnecessry.
    def clear_table_content_from_placeholders(self, slide, slide_data: Dict[str, Any]) -> None:
        """
        Clear table markdown content from text placeholders.

        This prevents duplication when table objects are created to replace
        the markdown content in placeholders.

        Args:
            slide: PowerPoint slide object
            slide_data: Dictionary containing slide data
        """
        debug_print("Clearing table content from placeholders to prevent duplication")

        # Get all fields that might contain table content
        all_fields = dict(slide_data)
        if "placeholders" in slide_data:
            all_fields.update(slide_data["placeholders"])

        # Find fields with table markdown
        for field_name, field_value in all_fields.items():
            if isinstance(field_value, str) and self.detect_table_content(field_value):
                debug_print(f"Found table markdown in field '{field_name}'")
                # Find and clear ALL placeholders that contain this table content
                for shape in slide.placeholders:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        shape_text = shape.text_frame.text if shape.text_frame.text else ""
                        if shape_text and self.detect_table_content(shape_text):
                            debug_print(f"Clearing table content from placeholder idx {shape.placeholder_format.idx}")
                            shape.text_frame.clear()
                            # Leave completely empty since table object will replace this content
                            break
