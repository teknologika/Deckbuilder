"""Placeholder resolution utilities for name-based lookups."""

from typing import List, Dict, Optional
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE


class PlaceholderResolver:
    """Provides name-based placeholder resolution functionality."""

    @staticmethod
    def get_placeholder_by_name(slide, placeholder_name: str) -> Optional[SlidePlaceholder]:
        """
        Find placeholder by exact name match.

        Args:
            slide: PowerPoint slide object
            placeholder_name: Exact placeholder name to find

        Returns:
            SlidePlaceholder object if found, None otherwise
        """
        for placeholder in slide.placeholders:
            try:
                # Get placeholder name from element
                actual_name = placeholder.element.nvSpPr.cNvPr.name
                if actual_name == placeholder_name:
                    return placeholder
            except AttributeError:
                # Some placeholder types might not have accessible names
                continue
        return None

    @staticmethod
    def find_placeholder_by_name_with_suggestions(slide, placeholder_name: str) -> SlidePlaceholder:
        """
        Find placeholder by exact name match with helpful error messages.

        Args:
            slide: PowerPoint slide object
            placeholder_name: Exact placeholder name to find

        Returns:
            SlidePlaceholder object if found

        Raises:
            ValueError: If placeholder not found with helpful suggestions
        """
        placeholder = PlaceholderResolver.get_placeholder_by_name(slide, placeholder_name)
        if placeholder:
            return placeholder

        # Generate helpful error with suggestions
        available = PlaceholderResolver.list_placeholder_names(slide)
        similar = [name for name in available if placeholder_name.lower() in name.lower()]

        error_msg = f"Placeholder '{placeholder_name}' not found.\\n"
        if similar:
            error_msg += f"Similar placeholders: {similar}\\n"
        error_msg += f"Available placeholders: {available}"

        raise ValueError(error_msg)

    @staticmethod
    def list_placeholder_names(slide) -> List[str]:
        """
        List all placeholder names in a slide.

        Args:
            slide: PowerPoint slide object

        Returns:
            List of placeholder names
        """
        names = []
        for placeholder in slide.placeholders:
            try:
                name = placeholder.element.nvSpPr.cNvPr.name
                names.append(name)
            except AttributeError:
                # Fallback to index if name not accessible
                names.append(f"Placeholder_{placeholder.placeholder_format.idx}")
        return names

    @staticmethod
    def get_placeholder_summary(slide) -> List[Dict[str, str]]:
        """
        Get detailed summary of all placeholders in a slide.

        Args:
            slide: PowerPoint slide object

        Returns:
            List of dictionaries with placeholder details
        """
        summary = []
        for placeholder in slide.placeholders:
            try:
                name = placeholder.element.nvSpPr.cNvPr.name
            except AttributeError:
                name = f"Placeholder_{placeholder.placeholder_format.idx}"

            # Get placeholder type
            ph_type = placeholder.placeholder_format.type
            type_name = PlaceholderResolver._get_placeholder_type_name(ph_type)

            summary.append({"name": name, "index": placeholder.placeholder_format.idx, "type": type_name, "type_code": str(ph_type)})

        return summary

    @staticmethod
    def find_placeholders_by_type(slide, placeholder_type: PP_PLACEHOLDER_TYPE) -> List[SlidePlaceholder]:
        """
        Find all placeholders of a specific type.

        Args:
            slide: PowerPoint slide object
            placeholder_type: PP_PLACEHOLDER_TYPE enum value

        Returns:
            List of SlidePlaceholder objects matching the type
        """
        matching_placeholders = []
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == placeholder_type:
                matching_placeholders.append(placeholder)
        return matching_placeholders

    @staticmethod
    def get_placeholder_type_summary(slide) -> Dict[str, List[str]]:
        """
        Group placeholders by type for debugging.

        Args:
            slide: PowerPoint slide object

        Returns:
            Dictionary mapping type names to lists of placeholder names
        """
        type_summary = {}
        for placeholder in slide.placeholders:
            ph_type = placeholder.placeholder_format.type
            type_name = PlaceholderResolver._get_placeholder_type_name(ph_type)

            try:
                name = placeholder.element.nvSpPr.cNvPr.name
            except AttributeError:
                name = f"Placeholder_{placeholder.placeholder_format.idx}"

            if type_name not in type_summary:
                type_summary[type_name] = []
            type_summary[type_name].append(name)

        return type_summary

    @staticmethod
    def validate_placeholder_exists(slide, placeholder_name: str) -> bool:
        """
        Check if placeholder exists without raising exception.

        Args:
            slide: PowerPoint slide object
            placeholder_name: Placeholder name to check

        Returns:
            True if placeholder exists, False otherwise
        """
        return PlaceholderResolver.get_placeholder_by_name(slide, placeholder_name) is not None

    @staticmethod
    def _get_placeholder_type_name(ph_type: PP_PLACEHOLDER_TYPE) -> str:
        """
        Convert placeholder type enum to readable name.

        Args:
            ph_type: PP_PLACEHOLDER_TYPE enum value

        Returns:
            Human-readable type name
        """
        type_mapping = {
            PP_PLACEHOLDER_TYPE.TITLE: "TITLE",
            PP_PLACEHOLDER_TYPE.BODY: "BODY",
            PP_PLACEHOLDER_TYPE.SUBTITLE: "SUBTITLE",
            PP_PLACEHOLDER_TYPE.TABLE: "TABLE",
            PP_PLACEHOLDER_TYPE.PICTURE: "PICTURE",
            PP_PLACEHOLDER_TYPE.MEDIA_CLIP: "MEDIA_CLIP",
            PP_PLACEHOLDER_TYPE.OBJECT: "OBJECT",
            PP_PLACEHOLDER_TYPE.CHART: "CHART",
            PP_PLACEHOLDER_TYPE.DATE: "DATE",
            PP_PLACEHOLDER_TYPE.FOOTER: "FOOTER",
            PP_PLACEHOLDER_TYPE.HEADER: "HEADER",
            PP_PLACEHOLDER_TYPE.SLIDE_NUMBER: "SLIDE_NUMBER",
            PP_PLACEHOLDER_TYPE.CENTER_TITLE: "CENTER_TITLE",
            PP_PLACEHOLDER_TYPE.BITMAP: "BITMAP",
            PP_PLACEHOLDER_TYPE.ORG_CHART: "ORG_CHART",
            PP_PLACEHOLDER_TYPE.SLIDE_IMAGE: "SLIDE_IMAGE",
            PP_PLACEHOLDER_TYPE.VERTICAL_BODY: "VERTICAL_BODY",
            PP_PLACEHOLDER_TYPE.VERTICAL_TITLE: "VERTICAL_TITLE",
            PP_PLACEHOLDER_TYPE.VERTICAL_OBJECT: "VERTICAL_OBJECT",
            PP_PLACEHOLDER_TYPE.MIXED: "MIXED",
        }
        return type_mapping.get(ph_type, f"UNKNOWN_{ph_type}")
