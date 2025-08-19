from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor

try:
    from .table_styles import TABLE_BORDER_STYLES, TABLE_HEADER_STYLES, TABLE_ROW_STYLES
except ImportError:
    # Fallback values if modules don't exist
    TABLE_HEADER_STYLES = {"dark_blue_white_text": {"bg": RGBColor(46, 89, 132), "text": RGBColor(255, 255, 255)}}
    TABLE_ROW_STYLES = {
        "alternating_light_gray": {
            "primary": RGBColor(255, 255, 255),
            "alt": RGBColor(240, 240, 240),
        }
    }
    TABLE_BORDER_STYLES = {"thin_gray": {"width": Pt(1), "color": RGBColor(128, 128, 128), "style": "all"}}


class TableBuilder:
    """Handles table creation, styling, and formatting for PowerPoint presentations."""

    def __init__(self, content_formatter=None):
        """
        Initialize the table builder.

        Args:
            content_formatter: ContentFormatter instance for text formatting
        """
        self.content_formatter = content_formatter

    def add_table_to_slide(self, slide, table_data):
        """
        Add a styled table to a slide.

        Args:
            slide: The slide to add the table to
            table_data: Dictionary containing table data and styling options
        """
        # Get table data - support both 'data' and 'rows' keys for backwards compatibility
        data = table_data.get("data", table_data.get("rows", []))
        if not data:
            return

        # Get styling options with enhanced configuration support
        header_style = table_data.get("header_style", "dark_blue_white_text")
        row_style = table_data.get("row_style", "alternating_light_gray")
        border_style = table_data.get("border_style", "thin_gray")
        custom_colors = table_data.get("custom_colors", {})

        # Font sizing configuration
        header_font_size = table_data.get("header_font_size", 12)  # 12pt default for headers
        data_font_size = table_data.get("data_font_size", 10)  # 10pt default for data

        # Parse dimension options
        dimensions = self._parse_dimensions(table_data, len(data[0]) if data else 0, len(data) if data else 0)

        # Find content placeholder or create table in available space
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Content placeholder
                content_placeholder = shape
                break

        if content_placeholder and hasattr(content_placeholder, "text_frame"):
            # Insert table within the content placeholder bounds

            # Calculate smart positioning based on content line count
            content_offset = self._calculate_content_offset(content_placeholder)
            left = content_placeholder.left
            top = content_placeholder.top + content_offset
            width = dimensions["table_width"] or (content_placeholder.width - Cm(1))
            height = dimensions["table_height"] or self._calculate_smart_table_height(data, header_font_size, data_font_size)

        else:
            # Default positioning if no content placeholder found
            left = Cm(2.5)
            top = Cm(5)
            width = dimensions["table_width"] or Cm(20)
            height = dimensions["table_height"] or self._calculate_smart_table_height(data, header_font_size, data_font_size)

        # Create the table
        rows = len(data)
        if data:
            # Handle both old (list of strings) and new (list of dicts) formats
            first_row = data[0]
            if isinstance(first_row, list):
                cols = len(first_row)
            else:
                cols = 1  # Fallback
        else:
            cols = 1

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Apply table data with formatting support
        for row_idx, row_data in enumerate(data):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)

                # Handle both old (string) and new (formatted) cell data
                if isinstance(cell_data, dict) and "formatted" in cell_data:
                    # New formatted cell data
                    if self.content_formatter:
                        self.content_formatter.apply_formatted_segments_to_cell(cell, cell_data["formatted"])
                    else:
                        cell.text = cell_data.get("text", str(cell_data))
                else:
                    # Old string cell data
                    cell.text = str(cell_data)

        # Apply dimension controls
        if dimensions["column_widths"]:
            self._apply_column_widths(table, dimensions["column_widths"])
        if dimensions["row_height"]:
            self._apply_row_heights(table, dimensions["row_height"])

        # Apply styling with font size control
        self._apply_table_styling(table, header_style, row_style, border_style, custom_colors, table_data, header_font_size, data_font_size)

        # Apply per-cell color overrides based on configuration
        cell_color_mode = table_data.get("cell_color_mode", "auto")
        if cell_color_mode in ["auto", "enabled"]:
            self._apply_per_cell_colors(table, data, cell_color_mode)

    def _apply_table_styling(self, table, header_style, row_style, border_style, custom_colors, table_data=None, header_font_size=12, data_font_size=10):
        """
        Apply styling to a table with font size control.

        Args:
            table: The table object to style
            header_style: Header style name
            row_style: Row style name
            border_style: Border style name
            custom_colors: Dictionary of custom color overrides
        """
        # Apply header styling
        if header_style in TABLE_HEADER_STYLES:
            header_colors = TABLE_HEADER_STYLES[header_style]

            # Override with custom colors if provided
            bg_color = self._parse_custom_color(custom_colors.get("header_bg")) or header_colors["bg"]
            text_color = self._parse_custom_color(custom_colors.get("header_text")) or header_colors["text"]

            # Style header row (first row)
            for col_idx in range(len(table.columns)):
                cell = table.cell(0, col_idx)
                # Set background color
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_color

                # Set text color, formatting, and font size
                font_size = self._get_font_size(0, header_font_size, data_font_size)  # Row 0 is header
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = text_color
                        run.font.bold = True
                        run.font.size = Pt(font_size)

        # Apply row styling
        if row_style in TABLE_ROW_STYLES and len(table.rows) > 1:
            row_colors = TABLE_ROW_STYLES[row_style]

            # Override with custom colors if provided
            primary_color = self._parse_custom_color(custom_colors.get("primary_row")) or row_colors["primary"]
            alt_color = self._parse_custom_color(custom_colors.get("alt_row")) or row_colors["alt"]

            # Style data rows (skip header row)
            for row_idx in range(1, len(table.rows)):
                is_alt_row = (row_idx - 1) % 2 == 1
                bg_color = alt_color if is_alt_row else primary_color

                if bg_color is not None:
                    font_size = self._get_font_size(row_idx, header_font_size, data_font_size)  # Data row
                    for col_idx in range(len(table.columns)):
                        cell = table.cell(row_idx, col_idx)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = bg_color

                        # Apply font size to data cells
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(font_size)

        # Apply border styling
        if border_style in TABLE_BORDER_STYLES:
            self._apply_table_borders(table, TABLE_BORDER_STYLES[border_style], custom_colors)

    def _apply_table_borders(self, table, border_config, custom_colors):
        """
        Apply border styling to a table.

        Args:
            table: The table object
            border_config: Border configuration dictionary
            custom_colors: Custom color overrides
        """
        border_width = border_config["width"]
        border_color = self._parse_custom_color(custom_colors.get("border_color")) or border_config["color"]
        border_style = border_config["style"]

        if border_style == "none" or border_width.cm == 0:
            return

        # Apply borders based on style
        for row_idx in range(len(table.rows)):
            for col_idx in range(len(table.columns)):
                cell = table.cell(row_idx, col_idx)

                if border_style == "all":
                    # All borders
                    self._set_cell_borders(cell, border_width, border_color, all_sides=True)
                elif border_style == "header" and row_idx == 0:
                    # Only header bottom border
                    self._set_cell_borders(cell, border_width, border_color, bottom=True)
                elif border_style == "outer":
                    # Only outer borders
                    is_top = row_idx == 0
                    is_bottom = row_idx == len(table.rows) - 1
                    is_left = col_idx == 0
                    is_right = col_idx == len(table.columns) - 1

                    self._set_cell_borders(
                        cell,
                        border_width,
                        border_color,
                        top=is_top,
                        bottom=is_bottom,
                        left=is_left,
                        right=is_right,
                    )

    def _set_cell_borders(self, cell, width, color, all_sides=False, top=False, bottom=False, left=False, right=False):
        """
        Set borders for a table cell.

        Args:
            cell: The table cell
            width: Border width
            color: Border color
            all_sides: Apply to all sides
            top, bottom, left, right: Apply to specific sides
        """
        if color is None:
            return

        if all_sides:
            top = bottom = left = right = True

        # Note: python-pptx has limited border support
        # This is a simplified implementation
        try:
            if hasattr(cell, "border"):
                if top and hasattr(cell.border, "top"):
                    cell.border.top.color.rgb = color
                    cell.border.top.width = width
                if bottom and hasattr(cell.border, "bottom"):
                    cell.border.bottom.color.rgb = color
                    cell.border.bottom.width = width
                if left and hasattr(cell.border, "left"):
                    cell.border.left.color.rgb = color
                    cell.border.left.width = width
                if right and hasattr(cell.border, "right"):
                    cell.border.right.color.rgb = color
                    cell.border.right.width = width
        except Exception:
            # Borders not fully supported in python-pptx, skip silently
            return  # nosec - Skip border styling if not supported

    def _parse_custom_color(self, color_value):
        """
        Parse a custom color value using HTML color names or 'transparent'.

        Args:
            color_value: HTML color name (e.g., "red", "navy", "lightgray") or "transparent"

        Returns:
            RGBColor object or None for transparent/invalid
        """
        if not color_value or not isinstance(color_value, str):
            return None

        try:
            import webcolors

            # Normalize the color name
            color_name = color_value.strip().lower()

            # Handle transparent specially
            if color_name == "transparent":
                return None

            # Convert HTML color name to RGB using webcolors
            rgb_tuple = webcolors.name_to_rgb(color_name)
            return RGBColor(rgb_tuple.red, rgb_tuple.green, rgb_tuple.blue)

        except (ValueError, TypeError, AttributeError):
            # webcolors raises ValueError for unknown color names
            return None

    def _apply_per_cell_colors(self, table, data, mode="auto"):
        """
        Apply per-cell color overrides based on cell content detection.

        This creates status tables where cell content like "GREEN", "RED", "BLUE"
        automatically gets colored backgrounds with matching text for invisible text effect.

        Args:
            table: PowerPoint table object
            data: Table data array
            mode: Color detection mode - "auto" (default), "enabled", or "disabled"
        """
        for row_idx, row_data in enumerate(data):
            for col_idx, cell_data in enumerate(row_data):
                if row_idx >= len(table.rows) or col_idx >= len(table.columns):
                    continue

                cell = table.cell(row_idx, col_idx)

                # Get cell text content
                if isinstance(cell_data, dict) and "text" in cell_data:
                    cell_text = str(cell_data["text"]).strip()
                else:
                    cell_text = str(cell_data).strip()

                # Check if cell text is a valid color name
                detected_color = self._parse_custom_color(cell_text)
                if detected_color is not None:
                    # Apply background color (same as detected color)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = detected_color

                    # Apply text color (same as background for invisible text effect)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = detected_color

                elif cell_text.upper() == "TRANSPARENT":
                    # Make cell completely transparent
                    cell.fill.background()

                    # Make text white or black based on overall table styling for visibility
                    text_color = self._get_transparent_cell_text_color()
                    if text_color:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = text_color

    def _get_transparent_cell_text_color(self):
        """
        Get appropriate text color for transparent cells.

        Returns white for dark themes, black for light themes.
        """
        from pptx.dml.color import RGBColor

        # Default to dark text for transparent cells (works on most backgrounds)
        return RGBColor(64, 64, 64)  # Dark gray for good visibility

    def _calculate_smart_table_height(self, data, header_font_size=12, data_font_size=10):
        """
        Calculate smart table height based on content analysis.

        Estimates appropriate height based on:
        - Text content length and wrapping
        - Font sizes
        - Number of rows
        - Content complexity

        Args:
            data: Table data array
            header_font_size: Font size for header row
            data_font_size: Font size for data rows

        Returns:
            Cm object for estimated table height
        """
        if not data:
            return Cm(3)  # Minimal height for empty table

        total_estimated_height = 0

        for row_idx, row_data in enumerate(data):
            # Determine if this is header row (first row)
            is_header = row_idx == 0
            font_size = header_font_size if is_header else data_font_size

            # Estimate row height based on content
            row_height = self._estimate_row_height(row_data, font_size, is_header)
            total_estimated_height += row_height

        return Cm(total_estimated_height)

    def _estimate_row_height(self, row_data, font_size, is_header=False):
        """
        Estimate height needed for a single row based on content.

        Args:
            row_data: List of cell data for the row
            font_size: Font size in points
            is_header: Whether this is a header row

        Returns:
            Height in cm (float)
        """
        # Base height for font size (rough estimate: font_size in pt * 0.035 = height in cm)
        base_height = font_size * 0.04  # Slightly generous for readability

        # Header rows get extra padding
        if is_header:
            base_height += 0.3

        # Find the longest content in this row
        max_content_length = 0
        longest_text = ""

        for cell_data in row_data:
            if isinstance(cell_data, dict) and "text" in cell_data:
                text = str(cell_data["text"]).strip()
            else:
                text = str(cell_data).strip()

            if len(text) > max_content_length:
                max_content_length = len(text)
                longest_text = text

        # Estimate text wrapping (approximate: chars per line varies by table width)
        # Conservative estimate for table columns (narrower columns = fewer chars per line)
        estimated_chars_per_line = 60  # More conservative for typical table columns
        estimated_lines = max(1, (max_content_length + estimated_chars_per_line - 1) // estimated_chars_per_line)  # Ceiling division

        # Check for explicit line breaks
        line_breaks = longest_text.count("\n") + longest_text.count("\\n")
        estimated_lines = max(estimated_lines, line_breaks + 1)

        # Calculate height based on lines needed
        if estimated_lines > 1:
            # Multi-line content needs more space
            line_height = font_size * 0.05  # Line spacing (slightly more generous)
            total_height = estimated_lines * line_height + 0.3  # Extra padding for multi-line
        else:
            # Single line content
            total_height = base_height

        # Apply minimum and maximum constraints
        min_height = 0.6  # Minimum readable height
        max_height = 5.0  # Maximum before it gets unwieldy (increased for long content)

        return max(min_height, min(max_height, total_height))

    def _parse_dimensions(self, table_data, column_count, row_count=None):
        """
        Parse and validate table dimension parameters.

        Args:
            table_data: Dictionary containing table configuration
            column_count: Number of columns in the table data
            row_count: Number of rows in the table data (for row_heights validation)

        Returns:
            Dictionary with parsed dimensions or None values
        """
        dimensions = {
            "table_width": None,
            "table_height": None,
            "column_widths": None,
            "row_height": None,
            "row_heights": None,  # New: per-row height array
        }

        # Parse table width
        if "table_width" in table_data:
            try:
                width_value = float(table_data["table_width"])
                if width_value > 0:
                    dimensions["table_width"] = Cm(width_value)
            except (ValueError, TypeError):
                print(f"Warning: Invalid table_width '{table_data['table_width']}', using default")

        # Parse table height
        if "table_height" in table_data:
            try:
                height_value = float(table_data["table_height"])
                if height_value > 0:
                    dimensions["table_height"] = Cm(height_value)
            except (ValueError, TypeError):
                print(f"Warning: Invalid table_height '{table_data['table_height']}', using default")

        # Parse column widths
        if "column_widths" in table_data:
            column_widths_raw = table_data["column_widths"]
            if isinstance(column_widths_raw, list):
                try:
                    column_widths = [float(w) for w in column_widths_raw if float(w) > 0]

                    # Validate column count
                    if len(column_widths) < column_count:
                        # Extend with default width (average of specified widths or 5cm)
                        default_width = sum(column_widths) / len(column_widths) if column_widths else 5.0
                        missing_count = column_count - len(column_widths)
                        column_widths.extend([default_width] * missing_count)
                        print(f"Warning: Only {len(column_widths_raw)} column widths specified for {column_count} columns. Extended with {default_width}cm default.")

                    elif len(column_widths) > column_count:
                        # Truncate to match actual columns
                        column_widths = column_widths[:column_count]
                        print(f"Warning: {len(column_widths_raw)} column widths specified for {column_count} columns. Truncated to match.")

                    # Convert to Cm objects
                    dimensions["column_widths"] = [Cm(w) for w in column_widths]

                    # Calculate total table width from column widths
                    dimensions["table_width"] = Cm(sum(column_widths))

                except (ValueError, TypeError) as e:
                    print(f"Warning: Invalid column_widths '{column_widths_raw}', using default: {e}")
            else:
                print(f"Warning: column_widths must be a list, got {type(column_widths_raw)}")

        # Parse row height
        if "row_height" in table_data:
            try:
                row_height_value = float(table_data["row_height"])
                if row_height_value > 0:
                    dimensions["row_height"] = Cm(row_height_value)
            except (ValueError, TypeError):
                print(f"Warning: Invalid row_height '{table_data['row_height']}', using default")

        # Parse per-row heights array
        if "row_heights" in table_data and row_count is not None:
            row_heights_raw = table_data["row_heights"]
            if isinstance(row_heights_raw, list):
                try:
                    row_heights = [float(h) for h in row_heights_raw if float(h) > 0]

                    # Validate row count
                    if len(row_heights) < row_count:
                        # Extend with default height (average of specified heights or 0.8cm)
                        default_height = sum(row_heights) / len(row_heights) if row_heights else 0.8
                        missing_count = row_count - len(row_heights)
                        row_heights.extend([default_height] * missing_count)
                        print(f"Warning: Only {len(row_heights_raw)} row heights specified for {row_count} rows. Extended with {default_height}cm default.")

                    elif len(row_heights) > row_count:
                        # Truncate to match actual rows
                        row_heights = row_heights[:row_count]
                        print(f"Warning: {len(row_heights_raw)} row heights specified for {row_count} rows. Truncated to match.")

                    # Convert to Cm objects
                    dimensions["row_heights"] = [Cm(h) for h in row_heights]

                except (ValueError, TypeError) as e:
                    print(f"Warning: Invalid row_heights '{row_heights_raw}', using default: {e}")
            else:
                print(f"Warning: row_heights must be a list, got {type(row_heights_raw)}")

        return dimensions

    def _calculate_content_offset(self, content_placeholder):
        """
        Calculate intelligent positioning offset based on content analysis.

        Args:
            content_placeholder: The content placeholder shape

        Returns:
            Cm object representing the offset from placeholder top
        """
        try:
            if not hasattr(content_placeholder, "text_frame") or not content_placeholder.text_frame:
                return Cm(0.8)  # Default fallback for placeholder without text frame

            # Analyze the actual text content
            total_text_length = 0
            line_count = 0
            has_bullets = False
            has_long_lines = False

            for paragraph in content_placeholder.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:  # Only count non-empty paragraphs
                    line_count += 1
                    total_text_length += len(text)

                    # Check for bullet points or list items
                    if text.startswith(("•", "-", "*", "1.", "2.", "3.")) or "\\n•" in text or "\\n-" in text:
                        has_bullets = True

                    # Check for long lines that might wrap
                    if len(text) > 80:  # Approximate wrapping threshold
                        has_long_lines = True
                        # Add extra line for wrapping
                        line_count += len(text) // 80

            # Enhanced offset calculation based on content characteristics:
            base_offset = 0.5  # Minimum spacing

            if line_count == 0:
                # Empty placeholder - minimal offset
                return Cm(base_offset)

            # Calculate offset based on estimated content height
            line_height = 0.6  # Approximate cm per line
            content_height = line_count * line_height

            # Add extra spacing for different content types
            if has_bullets:
                content_height += 0.3  # Extra space for bullet formatting

            if has_long_lines:
                content_height += 0.4  # Extra space for text wrapping

            # Add base spacing plus content-based spacing
            total_offset = base_offset + content_height + 0.5  # 0.5cm buffer between content and table

            # Cap the maximum offset to prevent tables from going off-slide
            return Cm(min(total_offset, 6.0))  # Max 6cm offset

        except Exception:
            # Fallback to safe default if analysis fails
            return Cm(1.2)

    def _get_font_size(self, row_idx, header_font_size, data_font_size):
        """
        Get appropriate font size for table cell based on row type.

        Args:
            row_idx: Row index (0 is header)
            header_font_size: Font size for header rows
            data_font_size: Font size for data rows

        Returns:
            Font size in points
        """
        # Validate font sizes
        header_size = max(8, min(24, int(header_font_size))) if header_font_size else 12
        data_size = max(8, min(20, int(data_font_size))) if data_font_size else 10

        return header_size if row_idx == 0 else data_size

    def _apply_font_size_to_cell(self, cell, font_size_pt):
        """
        Apply font size to all text in a table cell.

        Args:
            cell: Table cell object
            font_size_pt: Font size in points
        """
        try:
            from pptx.util import Pt

            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size_pt)
        except Exception:  # nosec B110
            # Ignore font sizing errors to prevent breaking table creation
            # This is intentional - font errors shouldn't break table functionality
            pass

    def _apply_column_widths(self, table, column_widths):
        """
        Apply individual column widths to a table.

        Args:
            table: The table object
            column_widths: List of Cm objects for column widths
        """
        try:
            for i, width in enumerate(column_widths):
                if i < len(table.columns):
                    table.columns[i].width = width
        except Exception as e:
            print(f"Warning: Failed to apply column widths: {e}")

    def _apply_row_heights(self, table, row_height):
        """
        Apply uniform row height to all table rows.

        Args:
            table: The table object
            row_height: Cm object for row height
        """
        try:
            for i in range(len(table.rows)):
                table.rows[i].height = row_height
        except Exception as e:
            print(f"Warning: Failed to apply row height: {e}")

    def _apply_table_data_and_styling(self, table, table_data):
        """
        Apply data and styling to an existing table shape.

        This method is used by dynamic shape creation to apply data and styling
        to a table that has already been created with specific positioning.

        Args:
            table: PowerPoint table object
            table_data: Dictionary containing table data and styling options
        """
        # Get table data
        data = table_data.get("data", table_data.get("rows", []))
        if not data:
            return

        # Get styling options
        header_style = table_data.get("header_style", "dark_blue_white_text")
        row_style = table_data.get("row_style", "alternating_light_gray")
        border_style = table_data.get("border_style", "thin_gray")
        custom_colors = table_data.get("custom_colors", {})
        header_font_size = table_data.get("header_font_size", 12)
        data_font_size = table_data.get("data_font_size", 10)

        # Apply data to table cells
        for row_idx, row_data in enumerate(data):
            if row_idx >= len(table.rows):
                break  # Safety check
            for col_idx, cell_data in enumerate(row_data):
                if col_idx >= len(table.columns):
                    break  # Safety check
                cell = table.cell(row_idx, col_idx)

                # Handle both old (string) and new (formatted) cell data
                if isinstance(cell_data, dict) and "formatted" in cell_data:
                    # New formatted cell data
                    if self.content_formatter:
                        self.content_formatter.apply_formatted_segments_to_cell(cell, cell_data["formatted"])
                    else:
                        cell.text = cell_data.get("text", str(cell_data))
                else:
                    # Old string cell data
                    cell.text = str(cell_data)

        # Apply styling using existing method
        self._apply_table_styling(table, header_style, row_style, border_style, custom_colors, table_data, header_font_size, data_font_size)

        # Apply per-cell color overrides based on configuration
        cell_color_mode = table_data.get("cell_color_mode", "auto")
        if cell_color_mode in ["auto", "enabled"]:
            self._apply_per_cell_colors(table, data, cell_color_mode)
