"""
ContentProcessor Module

Handles content application to PowerPoint placeholders with minimal font interference.
Focuses on content placement accuracy and proper newline support.

ARCHITECTURE DECISION:
- NO font override logic (preserve template fonts)
- Only handle semantic content placement
- Table font logic moved to TableHandler
- Newline support: convert \n → PowerPoint paragraphs
"""

from typing import Any, Dict
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE

from ..content.placeholder_types import (
    is_content_placeholder,
    is_media_placeholder,
    is_subtitle_placeholder,
    is_title_placeholder,
)
from ..utils.logging import debug_print


class ContentProcessor:
    """
    Handles content application to PowerPoint placeholders.

    DESIGN PRINCIPLES:
    - Preserve template fonts (no font overrides)
    - Focus on accurate content placement
    - Proper newline handling (\n → paragraphs)
    - Clean semantic type detection
    - Table processing delegates to TableHandler
    """

    def __init__(self):
        """Initialize the ContentProcessor."""
        pass

    def apply_content_to_placeholder(
        self,
        slide,
        placeholder,
        field_name: str,
        field_value: Any,
        slide_data: Dict[str, Any],
        content_formatter,
        image_placeholder_handler,
    ) -> None:
        """
        Apply content to a single placeholder based on its semantic type.

        SIMPLIFIED VERSION: Extracted from _apply_content_to_single_placeholder
        with font override logic removed.

        Args:
            slide: PowerPoint slide object
            placeholder: PowerPoint placeholder shape
            field_name: Field name being processed
            field_value: Content value to apply
            slide_data: Complete slide data dictionary
            content_formatter: ContentFormatter instance for formatting
            image_placeholder_handler: ImagePlaceholderHandler for images
        """
        # Set slide context for content formatter (needed for table processing)
        content_formatter._current_slide = slide_data.get("slide_object", slide)

        placeholder_type = placeholder.placeholder_format.type

        # Apply content based on placeholder semantic type
        if is_title_placeholder(placeholder_type) or is_subtitle_placeholder(placeholder_type):
            self._apply_title_content(placeholder, field_value, content_formatter)

        elif is_content_placeholder(placeholder_type):
            self._apply_content_placeholder_content(slide, placeholder, field_name, field_value, content_formatter)

        elif is_media_placeholder(placeholder_type):
            self._apply_media_placeholder_content(slide, placeholder, field_name, field_value, slide_data, image_placeholder_handler)

    def _apply_title_content(self, placeholder, field_value: Any, content_formatter) -> None:
        """
        Apply content to title/subtitle placeholders with proper formatting.

        NO FONT OVERRIDES: Preserve template font sizes.

        Args:
            placeholder: Title or subtitle placeholder
            field_value: Content to apply
            content_formatter: ContentFormatter for inline formatting
        """
        if not (hasattr(placeholder, "text_frame") and placeholder.text_frame):
            # Fallback to simple text for placeholders without text frames
            placeholder.text = str(field_value)
            return

        # Use text frame for formatting support
        text_frame = placeholder.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()

        # Handle formatted content properly for titles
        if isinstance(field_value, list) and field_value and isinstance(field_value[0], dict) and "text" in field_value[0]:
            # Pre-formatted segments
            content_formatter.apply_formatted_segments_to_paragraph(field_value, paragraph)
        else:
            # Raw text with potential inline formatting and newlines
            text_content = str(field_value)
            self._apply_content_with_newlines(text_content, text_frame, content_formatter)

        debug_print(f"    Applied title content: {len(str(field_value))} chars (template fonts preserved)")

    def _apply_content_placeholder_content(self, slide, placeholder, field_name: str, field_value: Any, content_formatter) -> None:
        """
        Apply content to content placeholders, handling tables and text content.

        Args:
            slide: PowerPoint slide object
            placeholder: Content placeholder
            field_name: Field name being processed
            field_value: Content to apply
            content_formatter: ContentFormatter instance
        """
        # Import TableHandler for table detection
        from .table_handler import TableHandler

        # Check if this content contains table data (especially for OBJECT placeholders)
        table_handler = TableHandler()

        # Detect if field_value contains table data (markdown string or structured dict)
        is_table_content = False

        if isinstance(field_value, str) and table_handler.detect_table_content(field_value):
            # Raw markdown table content
            is_table_content = True
            debug_print(f"    DETECTED MARKDOWN TABLE CONTENT in content placeholder: {field_name}")
        elif isinstance(field_value, dict) and field_value.get("type") == "table":
            # Already parsed table data structure
            is_table_content = True
            debug_print(f"    DETECTED STRUCTURED TABLE CONTENT in content placeholder: {field_name}")

        if is_table_content:
            # This content placeholder contains table data - route to table processing
            self._handle_table_in_content_placeholder(slide, placeholder, field_name, field_value, table_handler)
        else:
            # Regular content - text, lists, etc. with inline formatting and newline support
            content_formatter.add_content_to_placeholder(placeholder, field_value)
            debug_print(f"    Applied content to placeholder: {field_name} ({type(field_value).__name__})")

    def _apply_media_placeholder_content(
        self,
        slide,
        placeholder,
        field_name: str,
        field_value: Any,
        slide_data: Dict[str, Any],
        image_placeholder_handler,
    ) -> None:
        """
        Apply content to media placeholders (images, tables, objects).

        Args:
            slide: PowerPoint slide object
            placeholder: Media placeholder
            field_name: Field name being processed
            field_value: Content to apply
            slide_data: Complete slide data
            image_placeholder_handler: Handler for image processing
        """
        placeholder_type = placeholder.placeholder_format.type

        if placeholder_type == PP_PLACEHOLDER_TYPE.PICTURE:
            image_placeholder_handler.handle_image_placeholder(placeholder, field_name, field_value, slide_data)

        elif placeholder_type == PP_PLACEHOLDER_TYPE.TABLE:
            # TABLE placeholders - handle table data
            self._handle_table_placeholder(placeholder, field_name, field_value, slide_data, slide)

        elif placeholder_type == PP_PLACEHOLDER_TYPE.OBJECT and hasattr(placeholder, "text_frame"):
            # OBJECT placeholders with text_frame should be treated as content placeholders
            if hasattr(placeholder, "text_frame") and placeholder.text_frame:
                text_content = str(field_value)
                self._apply_content_with_newlines(text_content, placeholder.text_frame, None)
            else:
                placeholder.text = str(field_value)

        else:
            # Generic object placeholder - try text fallback
            if hasattr(placeholder, "text"):
                placeholder.text = str(field_value)
            debug_print(f"    Applied media content: {field_name} to {placeholder_type}")

    def _handle_table_placeholder(self, placeholder, field_name: str, field_value: Any, slide_data: Dict[str, Any], slide) -> None:
        """
        Handle TABLE placeholder types with proper table creation.

        Args:
            placeholder: Table placeholder
            field_name: Field name
            field_value: Table content (can be markdown or structured data)
            slide_data: Complete slide data
            slide: PowerPoint slide object
        """
        # Import TableHandler for proper table processing
        from .table_handler import TableHandler

        table_handler = TableHandler()

        # Use the field_value directly instead of searching slide_data
        table_content = field_value

        if table_content and table_handler.detect_table_content(table_content):
            table_data = table_handler.parse_table_structure(table_content)
            if table_data:
                # Get number of rows and columns from parsed data
                num_rows = len(table_data)
                num_cols = len(table_data[0]) if table_data else 0

                if num_rows > 0 and num_cols > 0:
                    # Use placeholder.insert_table() to insert table INTO the placeholder
                    try:
                        graphic_frame = placeholder.insert_table(rows=num_rows, cols=num_cols)
                        table = graphic_frame.table

                        # Populate the table with our data
                        for row_idx, row_data in enumerate(table_data):
                            for col_idx, cell_data in enumerate(row_data):
                                if row_idx < len(table.rows) and col_idx < len(table.columns):
                                    table.cell(row_idx, col_idx).text = str(cell_data)

                        debug_print(f"    ✅ Table inserted into placeholder: {num_rows}x{num_cols}")

                    except Exception as e:
                        debug_print(f"    Failed to insert table into placeholder: {e}")
                        # Fallback to text content
                        if hasattr(placeholder, "text"):
                            placeholder.text = str(field_value)
                else:
                    debug_print("    Invalid table dimensions - no data to insert")
                    if hasattr(placeholder, "text"):
                        placeholder.text = str(field_value)
            else:
                debug_print("    Failed to parse table structure from placeholder")
                if hasattr(placeholder, "text"):
                    placeholder.text = str(field_value)
        else:
            # No table content detected - treat as regular text
            if hasattr(placeholder, "text"):
                placeholder.text = str(field_value)
            debug_print("    TABLE placeholder set as text content (no table detected)")

    def _handle_table_in_content_placeholder(self, slide, placeholder, field_name: str, field_value, table_handler) -> None:
        """
        Handle table content found within a content placeholder (typically OBJECT type).

        This addresses the issue where table markdown appears in content placeholders
        but needs to be processed as actual PowerPoint table shapes.

        Args:
            slide: PowerPoint slide object
            placeholder: Content placeholder containing table data
            field_name: Field name containing table content
            field_value: String content with table markdown OR dict with structured table data
            table_handler: TableHandler instance for processing
        """
        debug_print(f"    Processing table content in content placeholder: {field_name}")

        # Check if slide already has tables to avoid duplication
        existing_tables = table_handler.detect_existing_tables(slide)
        if existing_tables:
            debug_print(f"    Skipping table creation - slide already has {len(existing_tables)} table(s)")
            # Clear the placeholder to avoid showing table markdown as text
            if hasattr(placeholder, "text_frame") and placeholder.text_frame:
                placeholder.text_frame.clear()
            elif hasattr(placeholder, "text"):
                placeholder.text = ""
            return

        # Get table data depending on format
        if isinstance(field_value, dict) and field_value.get("type") == "table":
            # Already structured table data
            table_data = field_value.get("data")
            debug_print(f"    Using structured table data with {len(table_data) if table_data else 0} rows")
        else:
            # Parse table structure from the markdown content
            table_data = table_handler.parse_table_structure(field_value)

        if table_data:
            # Calculate position and size for the table
            # Default position in the center-lower area of the slide
            position = (1.0, 3.0)  # inches from left, inches from top
            size = (8.0, 4.0)  # width, height in inches

            # Create the table shape
            table_shape = table_handler.create_table_from_data(slide, table_data, position, size)

            if table_shape:
                # Table positioning is handled in create_table_from_data with position/size
                debug_print(f"    Table created at position {position} with size {size}")

                # Clear table content from the original placeholder to avoid duplication
                table_handler.clear_table_content_from_placeholders(slide, field_value)

                debug_print(f"    ✅ Created table with {len(table_data)} rows from content placeholder")
            else:
                debug_print("    ❌ Failed to create table from content placeholder")
                # Fallback: keep original content in placeholder
                if hasattr(placeholder, "text_frame") and placeholder.text_frame:
                    placeholder.text_frame.text = str(field_value)
                elif hasattr(placeholder, "text"):
                    placeholder.text = str(field_value)
        else:
            debug_print("    ❌ Failed to parse table structure from content placeholder")
            # Fallback: keep original content in placeholder
            if hasattr(placeholder, "text_frame") and placeholder.text_frame:
                placeholder.text_frame.text = str(field_value)
            elif hasattr(placeholder, "text"):
                placeholder.text = str(field_value)

    def convert_newlines_to_paragraphs(self, text_content: str, text_frame, content_formatter=None) -> None:
        """
        Convert newline characters (\n) to actual PowerPoint paragraphs.

        CORE FEATURE: Proper newline support for PowerPoint presentations.

        Args:
            text_content: Text with potential \n characters
            text_frame: PowerPoint text frame to populate
            content_formatter: Optional ContentFormatter for inline formatting
        """
        if not text_content or "\n" not in text_content:
            # No newlines - use single paragraph with formatting
            if content_formatter:
                paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                content_formatter.apply_inline_formatting(text_content, paragraph)
            else:
                text_frame.text = text_content
            return

        # Split by newlines and create separate paragraphs
        lines = text_content.split("\n")
        text_frame.clear()

        for i, line in enumerate(lines):
            # Create paragraph for each line
            if i == 0:
                # Use existing first paragraph
                paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
            else:
                # Add new paragraph for subsequent lines
                paragraph = text_frame.add_paragraph()

            # Apply content with inline formatting if available
            if content_formatter and line.strip():
                content_formatter.apply_inline_formatting(line, paragraph)
            else:
                paragraph.text = line

        debug_print(f"    Converted {len(lines)} lines with newlines to PowerPoint paragraphs")

    def _apply_content_with_newlines(self, text_content: str, text_frame, content_formatter) -> None:
        """
        Apply text content with proper newline handling.

        Args:
            text_content: Text content to apply
            text_frame: PowerPoint text frame
            content_formatter: ContentFormatter for inline formatting (can be None)
        """
        self.convert_newlines_to_paragraphs(text_content, text_frame, content_formatter)
