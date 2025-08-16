"""Single field processing system using name-based placeholder resolution."""

from typing import Dict, Any, List
from ..utils.logging import slide_builder_print, debug_print, error_print
from .placeholder_resolver import PlaceholderResolver
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE


class FieldProcessor:
    """Handles single-path field processing using name-based placeholder resolution."""

    def __init__(self):
        """Initialize the field processor."""
        pass

    def process_slide_fields(self, slide, slide_data: dict, content_formatter, image_placeholder_handler) -> Dict[str, str]:
        """
        Process all fields in slide_data using name-based placeholder resolution.

        This is the single processing path that replaces dual field-driven and semantic processing.

        Args:
            slide: PowerPoint slide object
            slide_data: Dictionary containing slide content (from JSON or markdown)
            content_formatter: ContentFormatter instance
            image_placeholder_handler: ImagePlaceholderHandler instance

        Returns:
            Dictionary of processing results {field_name: status}
        """
        processing_results = {}
        failed_fields = []

        # Extract fields to process - look in placeholders object first, then root level
        fields_to_process = {}

        # Primary source: placeholders object (structured frontmatter format)
        if "placeholders" in slide_data and isinstance(slide_data["placeholders"], dict):
            fields_to_process.update(slide_data["placeholders"])
            debug_print(f"    Found {len(slide_data['placeholders'])} fields in placeholders object")

        # Secondary source: root level fields (direct format)
        for field_name, field_value in slide_data.items():
            if field_name.startswith("_") or field_name in ["layout", "type", "style", "placeholders", "speaker_notes", "content"]:
                # Skip internal fields, layout metadata, and already processed placeholders
                continue
            if field_name not in fields_to_process:  # Don't override placeholders data
                fields_to_process[field_name] = field_value

        # Also check for special table fields at root level (table_data)
        for table_field in ["table_data", "table", "content"]:
            if table_field in slide_data and table_field not in fields_to_process:
                fields_to_process[table_field] = slide_data[table_field]

        slide_builder_print(f"  Processing {len(fields_to_process)} fields using name-based resolution")
        debug_print(f"    Fields to process: {list(fields_to_process.keys())}")

        # Get all available placeholders for reference
        available_placeholders = PlaceholderResolver.get_placeholder_summary(slide)
        placeholder_names = [ph["name"] for ph in available_placeholders]

        debug_print(f"    Available placeholders: {placeholder_names}")

        # Process each field
        for field_name, field_value in fields_to_process.items():
            result = self._process_single_field(slide, field_name, field_value, content_formatter, image_placeholder_handler, available_placeholders, slide_data)

            processing_results[field_name] = result
            if result.startswith("FAILED"):
                failed_fields.append(field_name)

        # Report results
        if failed_fields:
            error_print(f"    FAILED: {len(failed_fields)} fields could not be mapped")
            error_print(f"    Failed fields: {failed_fields}")
            for field in failed_fields:
                error_print(f"      {field} ({processing_results[field]})")
        else:
            slide_builder_print(f"    SUCCESS: All {len(processing_results)} fields processed")

        return processing_results

    def _process_single_field(self, slide, field_name: str, field_value: Any, content_formatter, image_placeholder_handler, available_placeholders: List[Dict], slide_data: dict = None) -> str:
        """
        Process a single field using name-based placeholder resolution.

        Args:
            slide: PowerPoint slide object
            field_name: Name of the field to process
            field_value: Value of the field
            content_formatter: ContentFormatter instance
            image_placeholder_handler: ImagePlaceholderHandler instance
            available_placeholders: List of placeholder summaries

        Returns:
            Status string describing the processing result
        """
        debug_print(f"      Processing field '{field_name}' = {str(field_value)[:50]}...")

        # Try direct name matching first
        placeholder = PlaceholderResolver.get_placeholder_by_name(slide, field_name)
        if placeholder:
            return self._apply_field_to_placeholder(placeholder, field_name, field_value, content_formatter, image_placeholder_handler, slide_data)

        # Try semantic field name mapping
        mapped_placeholder_name = self._resolve_semantic_field_name(field_name, available_placeholders)
        if mapped_placeholder_name:
            placeholder = PlaceholderResolver.get_placeholder_by_name(slide, mapped_placeholder_name)
            if placeholder:
                debug_print(f"        Semantic mapping: '{field_name}' -> '{mapped_placeholder_name}'")
                return self._apply_field_to_placeholder(placeholder, field_name, field_value, content_formatter, image_placeholder_handler, slide_data)

        # No suitable placeholder found
        available_names = [ph["name"] for ph in available_placeholders]
        similar_names = [name for name in available_names if field_name.lower() in name.lower()]

        error_msg = "no suitable placeholder found"
        if similar_names:
            error_msg += f" (similar: {similar_names})"

        return f"FAILED: {error_msg}"

    # TODO: This method should be able to be removed as we map names in structured_frontmatter_patterns
    # Exact name matching should be used with out ANY layout hard-coding
    # this logic causes a lot of bugs.
    def _resolve_semantic_field_name(self, field_name: str, available_placeholders: List[Dict]) -> str:
        """
        Map semantic field names to actual placeholder names.

        Args:
            field_name: Field name from slide_data
            available_placeholders: List of placeholder summaries

        Returns:
            Actual placeholder name or None if no mapping found
        """
        field_lower = field_name.lower()

        # First try exact name matching (case insensitive)
        for placeholder in available_placeholders:
            placeholder_name = placeholder["name"]
            if placeholder_name.lower() == field_lower:
                return placeholder_name

        # Then try pattern matching based on layout structure
        content_placeholders = [ph for ph in available_placeholders if ph["type"] in ["BODY", "OBJECT"]]
        content_placeholders.sort(key=lambda x: x["index"])  # Sort by index for consistent ordering

        # Multi-column layouts: map by position
        if field_name.startswith("content_col"):
            try:
                col_num = int(field_name.split("_col")[1])  # Extract column number
                if 1 <= col_num <= len(content_placeholders):
                    return content_placeholders[col_num - 1]["name"]  # 1-based to 0-based
            except (ValueError, IndexError):
                pass

        # Two-column layouts: left/right mapping
        elif field_name in ["content_left", "content_right"]:
            if len(content_placeholders) >= 2:
                if field_name == "content_left":
                    return content_placeholders[0]["name"]  # First content placeholder
                elif field_name == "content_right":
                    return content_placeholders[1]["name"]  # Second content placeholder

        # Specialized layout fields: content_item1-6, Title1-6, etc.
        elif field_name.startswith("content_item") or field_name.startswith("Title") or field_name.startswith("Content"):
            # Extract number from field name if present
            import re

            number_match = re.search(r"(\d+)", field_name)
            if number_match:
                item_num = int(number_match.group(1))
                if 1 <= item_num <= len(content_placeholders):
                    return content_placeholders[item_num - 1]["name"]  # 1-based to 0-based
            # Fallback to first content placeholder
            elif content_placeholders:
                return content_placeholders[0]["name"]

        # SWOT Analysis quadrant mapping
        elif field_name.startswith("content_") and any(quad in field_name for quad in ["top_left", "top_right", "bottom_left", "bottom_right"]):
            # Map quadrant fields to content placeholders by position
            quadrant_map = {"content_top_left": 0, "content_top_right": 1, "content_bottom_left": 2, "content_bottom_right": 3}
            if field_name in quadrant_map and quadrant_map[field_name] < len(content_placeholders):
                return content_placeholders[quadrant_map[field_name]]["name"]

        # Generic content mapping (first available content placeholder)
        elif field_name == "content" and content_placeholders:
            return content_placeholders[0]["name"]

        # Legacy content field mapping
        elif field_name in ["rich_content", "rich_content_formatted"] and content_placeholders:
            return content_placeholders[0]["name"]

        # Title matching
        elif field_name == "title":
            title_placeholders = [ph for ph in available_placeholders if ph["type"] == "TITLE"]
            if title_placeholders:
                return title_placeholders[0]["name"]

        # Table data matching
        elif field_name in ["table_data", "table", "table_formatted"]:
            table_placeholders = [ph for ph in available_placeholders if ph["type"] == "TABLE"]
            if table_placeholders:
                return table_placeholders[0]["name"]
            # Fallback to content placeholders for table data if no TABLE placeholder exists
            elif content_placeholders:
                return content_placeholders[0]["name"]

        # Image matching
        elif field_name in ["image", "media"]:
            image_placeholders = [ph for ph in available_placeholders if ph["type"] in ["PICTURE", "MEDIA_CLIP"]]
            if image_placeholders:
                return image_placeholders[0]["name"]

        # Fallback: Try partial name matching for any remaining fields
        for placeholder in available_placeholders:
            placeholder_name = placeholder["name"]
            placeholder_lower = placeholder_name.lower()

            # Check if any part of the field name appears in the placeholder name
            field_parts = field_name.lower().split("_")
            if any(part in placeholder_lower for part in field_parts if len(part) > 2):
                return placeholder_name

        return None

    def _apply_field_to_placeholder(self, placeholder, field_name: str, field_value: Any, content_formatter, image_placeholder_handler, slide_data: dict = None) -> str:
        """
        Apply field value to a specific placeholder.

        Args:
            placeholder: PowerPoint placeholder object
            field_name: Name of the field
            field_value: Value to apply
            content_formatter: ContentFormatter instance
            image_placeholder_handler: ImagePlaceholderHandler instance

        Returns:
            Status string describing the result
        """
        try:
            placeholder_type = placeholder.placeholder_format.type
            placeholder_name = placeholder.element.nvSpPr.cNvPr.name

            debug_print(f"        Applying to placeholder '{placeholder_name}' (type: {PlaceholderResolver._get_placeholder_type_name(placeholder_type)})")

            # Handle different placeholder types
            if placeholder_type == PP_PLACEHOLDER_TYPE.TITLE:
                return self._handle_title_placeholder(placeholder, field_value, content_formatter)
            elif placeholder_type == PP_PLACEHOLDER_TYPE.BODY:
                return self._handle_content_placeholder(placeholder, field_value, content_formatter)
            elif placeholder_type == PP_PLACEHOLDER_TYPE.TABLE:
                return self._handle_table_placeholder(placeholder, field_value, content_formatter)
            elif placeholder_type == PP_PLACEHOLDER_TYPE.PICTURE:
                return self._handle_image_placeholder(placeholder, field_name, field_value, image_placeholder_handler, slide_data)
            elif placeholder_type == PP_PLACEHOLDER_TYPE.SUBTITLE:
                return self._handle_title_placeholder(placeholder, field_value, content_formatter)  # Same as title
            else:
                # Default to text content handling
                return self._handle_content_placeholder(placeholder, field_value, content_formatter)

        except Exception as e:
            error_print(f"        Error applying field '{field_name}': {e}")
            return f"FAILED: error during application - {str(e)}"

    def _handle_title_placeholder(self, placeholder, field_value: Any, content_formatter) -> str:
        """Handle title placeholder content."""
        try:
            if hasattr(placeholder, "text_frame"):
                content_formatter.add_content_to_placeholder(placeholder, field_value)
                return "SUCCESS: title content applied"
            else:
                return "FAILED: placeholder has no text_frame"
        except Exception as e:
            return f"FAILED: title application error - {str(e)}"

    def _handle_content_placeholder(self, placeholder, field_value: Any, content_formatter) -> str:
        """Handle content/body placeholder content."""
        try:
            if hasattr(placeholder, "text_frame"):
                content_formatter.add_content_to_placeholder(placeholder, field_value)
                return "SUCCESS: content applied"
            else:
                return "FAILED: placeholder has no text_frame"
        except Exception as e:
            return f"FAILED: content application error - {str(e)}"

    def _handle_table_placeholder(self, placeholder, field_value: Any, content_formatter) -> str:
        """Handle table placeholder content using proper python-pptx table insertion."""
        try:
            # This is the critical table processing logic
            debug_print(f"        Processing TABLE placeholder with data: {str(field_value)[:100]}...")

            if hasattr(placeholder, "insert_table"):
                # Proper TABLE placeholder - parse markdown and create real table
                return self._create_real_table(placeholder, field_value)
            elif hasattr(placeholder, "text_frame"):
                # Fallback to text processing for table data
                content_formatter.add_content_to_placeholder(placeholder, field_value)
                return "SUCCESS: table content applied as text (fallback)"
            else:
                return "FAILED: TABLE placeholder has no insert_table or text_frame"
        except Exception as e:
            error_print(f"        TABLE processing error: {e}")
            return f"FAILED: table application error - {str(e)}"

    def _create_real_table(self, placeholder, table_data: str) -> str:
        """Create a real PowerPoint table from markdown table data."""
        try:
            # Parse markdown table
            lines = [line.strip() for line in table_data.strip().split("\n") if line.strip()]
            if not lines:
                return "FAILED: empty table data"

            # Filter out separator lines (lines with only |, -, and spaces)
            data_lines = []
            for line in lines:
                if not all(c in "|-: " for c in line):
                    data_lines.append(line)

            if not data_lines:
                return "FAILED: no data rows found in table"

            # Parse rows and columns
            rows_data = []
            for line in data_lines:
                # Split by | and clean up
                cells = [cell.strip() for cell in line.split("|")]
                # Remove empty cells from start/end (common in markdown)
                while cells and not cells[0]:
                    cells.pop(0)
                while cells and not cells[-1]:
                    cells.pop()
                if cells:
                    rows_data.append(cells)

            if not rows_data:
                return "FAILED: no valid rows parsed from table data"

            num_rows = len(rows_data)
            num_cols = max(len(row) for row in rows_data)

            debug_print(f"        Creating {num_rows}x{num_cols} table...")

            # Insert table using python-pptx
            graphic_frame = placeholder.insert_table(rows=num_rows, cols=num_cols)
            table = graphic_frame.table

            # Fill table with data
            for row_idx, row_data in enumerate(rows_data):
                for col_idx in range(num_cols):
                    cell_text = row_data[col_idx] if col_idx < len(row_data) else ""
                    table.cell(row_idx, col_idx).text = cell_text

            debug_print("        Successfully created real PowerPoint table!")
            return f"SUCCESS: created {num_rows}x{num_cols} PowerPoint table"

        except Exception as e:
            error_print(f"        Real table creation error: {e}")
            return f"FAILED: real table creation error - {str(e)}"

    def _handle_image_placeholder(self, placeholder, field_name: str, field_value: Any, image_placeholder_handler, slide_data: dict = None) -> str:
        """Handle image placeholder content."""
        try:
            if image_placeholder_handler:
                # Use image handler for proper image processing
                result = image_placeholder_handler.handle_image_placeholder(placeholder, field_name, field_value, slide_data or {})
                return "SUCCESS: image applied" if result else "FAILED: image application failed"
            else:
                return "FAILED: no image placeholder handler available"
        except Exception as e:
            return f"FAILED: image application error - {str(e)}"
