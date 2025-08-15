"""
Layout Resolver Module

Enhanced layout resolution module that USES existing layout resolution systems.
Provides clean interface for layout-to-pattern mapping with proper error handling.

Part of the slide_builder.py refactor - Phase 2: PlaceholderManager Dependencies.
"""

from pptx import Presentation
from pptx.slide import SlideLayout

from ..core.layout_resolver import LayoutResolver as CoreLayoutResolver


class LayoutResolver:
    """
    Enhanced layout resolver that USES existing core layout resolution.

    DESIGN PRINCIPLE: ENHANCE existing code, don't create alternate paths.
    Uses existing core/layout_resolver.py with clean interface.
    """

    def __init__(self):
        """Initialize layout resolver using existing core functionality."""
        # USE existing CoreLayoutResolver - no alternate implementation
        self.core_resolver = CoreLayoutResolver()

    def resolve_layout_by_name(self, prs: Presentation, layout_name: str) -> SlideLayout:
        """
        Resolve PowerPoint layout by name using existing core functionality.

        Args:
            prs: PowerPoint presentation object
            layout_name: Layout name to find

        Returns:
            SlideLayout object

        Raises:
            ValueError: If layout not found with helpful error message and suggestions
        """
        try:
            return self.core_resolver.get_layout_by_name(prs, layout_name)
        except ValueError as e:
            # Enhanced error handling with available layouts and suggestions
            available_layouts = self.get_available_layouts(prs)

            # Find similar layout names for suggestions
            suggestions = self._find_similar_layout_names(layout_name, available_layouts)

            error_msg = f"Layout '{layout_name}' not found in presentation."
            error_msg += f"\n\nAvailable layouts ({len(available_layouts)}):"
            for i, layout in enumerate(available_layouts, 1):
                error_msg += f"\n  {i}. '{layout}'"

            if suggestions:
                error_msg += "\n\nDid you mean one of these?"
                for suggestion in suggestions[:3]:  # Show top 3 suggestions
                    error_msg += f"\n  - '{suggestion}'"

            error_msg += f"\n\nOriginal error: {str(e)}"

            raise ValueError(error_msg)

    def resolve_layout_safely(self, prs: Presentation, layout_name: str) -> dict:
        """
        Resolve layout without crashing - returns success/error information.

        Args:
            prs: PowerPoint presentation object
            layout_name: Layout name to find

        Returns:
            Dictionary with success status, layout (if found), or error info
        """
        try:
            layout = self.core_resolver.get_layout_by_name(prs, layout_name)
            return {"success": True, "layout": layout, "message": "Successfully resolved layout '{}'".format(layout_name)}
        except ValueError as e:
            # Return error information without crashing
            available_layouts = self.get_available_layouts(prs)
            suggestions = self._find_similar_layout_names(layout_name, available_layouts)

            return {
                "success": False,
                "layout": None,
                "error": "Layout not found",
                "requested_layout": layout_name,
                "available_layouts": available_layouts,
                "suggestions": suggestions[:3],  # Top 3 suggestions
                "message": f"Layout '{layout_name}' not found. Found {len(available_layouts)} available layouts.",
                "original_error": str(e),
            }

    def get_available_layouts(self, prs: Presentation) -> list[str]:
        """
        Get list of all available layout names.

        Args:
            prs: PowerPoint presentation object

        Returns:
            List of layout names available in the presentation
        """
        return self.core_resolver.list_available_layouts(prs)

    def validate_layout_compatibility(self, prs: Presentation, layout_name: str) -> bool:
        """
        Validate that a layout exists and is compatible.

        Args:
            prs: PowerPoint presentation object
            layout_name: Layout name to validate

        Returns:
            True if layout exists and is usable, False otherwise
        """
        return self.core_resolver.validate_layout_exists(prs, layout_name)

    def get_layout_info(self, prs: Presentation, layout_name: str) -> dict:
        """
        Get detailed information about a layout for debugging.

        Args:
            prs: PowerPoint presentation object
            layout_name: Layout name to analyze

        Returns:
            Dictionary with layout information
        """
        try:
            layout = self.resolve_layout_by_name(prs, layout_name)
            return {
                "name": layout.name,
                "placeholders": len(layout.placeholders),
                "exists": True,
                "placeholder_details": [
                    {
                        "idx": ph.placeholder_format.idx,
                        "type": ph.placeholder_format.type,
                        "name": getattr(ph, "name", "unnamed"),
                    }
                    for ph in layout.placeholders
                ],
            }
        except ValueError:
            return {
                "name": layout_name,
                "exists": False,
                "available_layouts": self.get_available_layouts(prs),
                "error": f"Layout '{layout_name}' not found",
            }

    def create_layout_to_pattern_mapping(self, prs: Presentation, layout_name: str, pattern_info: dict) -> dict:
        """
        Create a clean mapping between layout and pattern for validation.

        Args:
            prs: PowerPoint presentation object
            layout_name: Layout name to map
            pattern_info: Pattern information from PatternLoader

        Returns:
            Dictionary with layout-pattern mapping details
        """
        layout_info = self.get_layout_info(prs, layout_name)

        if not layout_info.get("exists"):
            return {
                "valid": False,
                "error": "Layout not found",
                "layout_info": layout_info,
                "pattern_info": pattern_info,
            }

        if not pattern_info.get("found"):
            return {
                "valid": False,
                "error": "Pattern not found",
                "layout_info": layout_info,
                "pattern_info": pattern_info,
            }

        # Extract placeholder names from layout
        layout_placeholder_names = [detail.get("name", f"idx_{detail['idx']}") for detail in layout_info["placeholder_details"]]

        # Extract expected fields from pattern
        pattern_fields = pattern_info.get("expected_fields", [])

        # Create mapping analysis
        mapping_analysis = {
            "layout_placeholder_names": layout_placeholder_names,
            "pattern_expected_fields": pattern_fields,
            "potential_matches": [],
            "missing_placeholders": [],
            "extra_placeholders": [],
        }

        # Analyze potential matches
        for field in pattern_fields:
            if field in ["layout", "style"]:  # Skip metadata
                continue
            if field in layout_placeholder_names:
                mapping_analysis["potential_matches"].append(field)
            else:
                mapping_analysis["missing_placeholders"].append(field)

        # Find extra placeholders
        content_fields = [f for f in pattern_fields if f not in ["layout", "style"]]
        for name in layout_placeholder_names:
            if name not in content_fields:
                mapping_analysis["extra_placeholders"].append(name)

        return {
            "valid": len(mapping_analysis["missing_placeholders"]) == 0,
            "layout_name": layout_name,
            "layout_info": layout_info,
            "pattern_info": pattern_info,
            "mapping_analysis": mapping_analysis,
        }

    def _find_similar_layout_names(self, target_name: str, available_names: list[str]) -> list[str]:
        """
        Find layout names similar to the target name for error suggestions.

        Args:
            target_name: The layout name that wasn't found
            available_names: List of available layout names

        Returns:
            List of similar layout names, sorted by similarity
        """
        if not available_names:
            return []

        suggestions = []
        target_lower = target_name.lower()

        # First pass: exact substring matches
        for name in available_names:
            name_lower = name.lower()
            if target_lower in name_lower or name_lower in target_lower:
                suggestions.append(name)

        # Second pass: word matches if no substring matches
        if not suggestions:
            target_words = set(target_lower.split())
            for name in available_names:
                name_words = set(name.lower().split())
                # If any words match
                if target_words.intersection(name_words):
                    suggestions.append(name)

        # Third pass: fuzzy matching if still no matches
        if not suggestions:
            import difflib

            suggestions = difflib.get_close_matches(target_name, available_names, n=3, cutoff=0.3)

        # Remove exact match (shouldn't happen, but just in case)
        if target_name in suggestions:
            suggestions.remove(target_name)

        return suggestions

    def resolve_layout_with_fallback(self, prs: Presentation, primary_layout: str, fallback_layouts: list[str]) -> SlideLayout:
        """
        Resolve layout with fallback options for robustness.

        Args:
            prs: PowerPoint presentation object
            primary_layout: Preferred layout name
            fallback_layouts: List of fallback layout names to try

        Returns:
            SlideLayout object

        Raises:
            ValueError: If none of the layouts can be resolved
        """
        layouts_to_try = [primary_layout] + fallback_layouts

        for layout_name in layouts_to_try:
            try:
                layout = self.core_resolver.get_layout_by_name(prs, layout_name)
                if layout_name != primary_layout:
                    print(f"⚠️  Using fallback layout '{layout_name}' instead of '{primary_layout}'")
                return layout
            except ValueError:
                continue

        # If all layouts failed, provide comprehensive error
        available_layouts = self.get_available_layouts(prs)
        error_msg = f"None of the specified layouts could be found: {layouts_to_try}"
        error_msg += f"\n\nAvailable layouts ({len(available_layouts)}):"
        for i, layout in enumerate(available_layouts, 1):
            error_msg += f"\n  {i}. '{layout}'"

        raise ValueError(error_msg)

    def resolve_layout_with_fallback_safely(self, prs: Presentation, primary_layout: str, fallback_layouts: list[str]) -> dict:
        """
        Resolve layout with fallbacks without crashing.

        Args:
            prs: PowerPoint presentation object
            primary_layout: Preferred layout name
            fallback_layouts: List of fallback layout names to try

        Returns:
            Dictionary with success status, layout (if found), or error info
        """
        layouts_to_try = [primary_layout] + fallback_layouts
        tried_layouts = []

        for layout_name in layouts_to_try:
            try:
                layout = self.core_resolver.get_layout_by_name(prs, layout_name)
                used_fallback = layout_name != primary_layout

                return {
                    "success": True,
                    "layout": layout,
                    "used_layout": layout_name,
                    "requested_layout": primary_layout,
                    "used_fallback": used_fallback,
                    "tried_layouts": tried_layouts + [layout_name],
                    "message": f"{'Used fallback layout' if used_fallback else 'Successfully resolved'} '{layout_name}'",
                }
            except ValueError:
                tried_layouts.append(layout_name)
                continue

        # If all layouts failed, return comprehensive error info
        available_layouts = self.get_available_layouts(prs)
        suggestions = self._find_similar_layout_names(primary_layout, available_layouts)

        return {
            "success": False,
            "layout": None,
            "requested_layout": primary_layout,
            "fallback_layouts": fallback_layouts,
            "tried_layouts": tried_layouts,
            "available_layouts": available_layouts,
            "suggestions": suggestions[:3],
            "error": "No layouts found",
            "message": f"None of the {len(layouts_to_try)} layouts could be found: {layouts_to_try}",
        }
