"""
Placeholder Normalizer Module

Normalizes instantiated slide placeholder names to match template names.
Ensures consistent naming between template layouts and instantiated slides.

Part of the slide_builder.py refactor - Phase 2: PlaceholderManager support.
"""

from typing import Dict
from pptx.slide import Slide
from pptx.slide import SlideLayout


class PlaceholderNormalizer:
    """
    Normalizes placeholder names on instantiated slides to match template names.

    This ensures that patterns using clean names like 'title_top' and 'content'
    can find matching placeholders on instantiated slides.
    """

    def __init__(self):
        """Initialize the placeholder normalizer."""
        # Mapping from template layout names to clean normalized names
        self._template_name_mapping = {}

    def normalize_slide_placeholder_names(self, slide: Slide, layout: SlideLayout) -> Dict[str, str]:
        """
        Normalize placeholder names on an instantiated slide to match template names.

        Args:
            slide: Instantiated slide object
            layout: Source layout object used to create the slide

        Returns:
            Dictionary mapping old names to new names
        """
        name_changes = {}

        try:
            # Update placeholder names on the instantiated slide
            for slide_ph in slide.placeholders:
                slide_idx = slide_ph.placeholder_format.idx

                # Find corresponding placeholder in layout by index
                layout_ph = self._find_layout_placeholder_by_index(layout, slide_idx)
                if layout_ph:
                    try:
                        template_name = layout_ph.element.nvSpPr.cNvPr.name
                        current_name = slide_ph.element.nvSpPr.cNvPr.name

                        if template_name != current_name:
                            # Update the slide placeholder name to match template
                            slide_ph.element.nvSpPr.cNvPr.name = template_name
                            name_changes[current_name] = template_name

                    except AttributeError:
                        # Some placeholders might not have accessible names
                        continue

        except Exception as e:
            # Log the error but don't fail the whole process
            print(f"Warning: Could not normalize placeholder names: {e}")

        return name_changes

    def _get_layout_name_mapping(self, layout: SlideLayout) -> Dict[int, str]:
        """
        Get mapping from placeholder index to template name for a layout.

        Args:
            layout: Layout object to extract names from

        Returns:
            Dictionary mapping placeholder index to template name
        """
        mapping = {}

        try:
            for ph in layout.placeholders:
                idx = ph.placeholder_format.idx
                try:
                    name = ph.element.nvSpPr.cNvPr.name
                    mapping[idx] = name
                except AttributeError:
                    # Placeholder has no name - use default based on type/index
                    mapping[idx] = f"placeholder_{idx}"

        except Exception:
            # If we can't read layout names, return empty mapping
            pass

        return mapping

    def _find_layout_placeholder_by_index(self, layout: SlideLayout, target_idx: int):
        """
        Find placeholder in layout by index.

        Args:
            layout: Layout to search
            target_idx: Placeholder index to find

        Returns:
            Placeholder object if found, None otherwise
        """
        try:
            for ph in layout.placeholders:
                if ph.placeholder_format.idx == target_idx:
                    return ph
        except Exception:
            pass
        return None

    def get_placeholder_name_preview(self, layout: SlideLayout) -> Dict[str, str]:
        """
        Preview what placeholder names would be after normalization.

        Args:
            layout: Layout to preview names for

        Returns:
            Dictionary with placeholder information
        """
        preview = {"layout_name": getattr(layout, "name", "Unknown"), "placeholders": []}

        try:
            for ph in layout.placeholders:
                idx = ph.placeholder_format.idx
                ph_type = ph.placeholder_format.type

                try:
                    template_name = ph.element.nvSpPr.cNvPr.name
                except AttributeError:
                    template_name = f"placeholder_{idx}"

                preview["placeholders"].append({"index": idx, "type": str(ph_type), "template_name": template_name, "would_be_slide_name": f"Placeholder {idx + 1}"})  # PowerPoint default pattern

        except Exception as e:
            preview["error"] = str(e)

        return preview
