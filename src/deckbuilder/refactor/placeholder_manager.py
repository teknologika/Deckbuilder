"""
Placeholder Manager Module

Enhanced placeholder management using name-based resolution with PatternLoader.
Maps structured frontmatter fields to PowerPoint placeholders by name.

Part of the slide_builder.py refactor - Phase 2: PlaceholderManager.
"""

from typing import Dict, Any, Optional, List
from pptx.slide import Slide
from pptx.shapes.placeholder import SlidePlaceholder

from ..templates.pattern_loader import PatternLoader
from ..core.placeholder_resolver import PlaceholderResolver
from .layout_resolver import LayoutResolver


class PlaceholderManager:
    """
    Enhanced placeholder manager using name-based resolution.

    DESIGN PRINCIPLE: ENHANCE existing code, don't create alternate paths.
    - Uses existing PatternLoader from templates/pattern_loader.py
    - Uses existing PlaceholderResolver from core/placeholder_resolver.py
    - Primary approach: Name-based placeholder matching
    - Semantic types: Only for exceptions during content application
    - NO fallbacks - SUCCESS or CLEAR ERRORS only
    """

    def __init__(self, layout_resolver: Optional[LayoutResolver] = None):
        """
        Initialize PlaceholderManager with existing systems.

        Args:
            layout_resolver: LayoutResolver instance (creates default if None)
        """
        # USE existing PatternLoader
        self.pattern_loader = PatternLoader()

        # USE existing PlaceholderResolver
        self.placeholder_resolver = PlaceholderResolver()

        # USE existing LayoutResolver
        self.layout_resolver = layout_resolver or LayoutResolver()

    def map_fields_to_placeholders(self, slide: Slide, slide_data: Dict[str, Any], layout_name: str, layout=None) -> Dict[str, SlidePlaceholder]:
        """
        Map structured frontmatter fields to placeholders by name.

        HYBRID FLOW:
        1. Get layout object (passed in or resolve from presentation)
        2. Normalize slide placeholder names to match template (using PlaceholderNormalizer)
        3. Name-based resolution finds placeholders (now guaranteed to match)

        Args:
            slide: PowerPoint slide object
            slide_data: Dictionary with slide data fields
            layout_name: Layout name for pattern lookup
            layout: Optional layout object (if not provided, will be resolved)

        Returns:
            Dictionary mapping field names to placeholder objects

        Raises:
            ValueError: Clear error message if pattern/placeholder not found
        """
        # Step 1: Get pattern from existing PatternLoader
        pattern = self.pattern_loader.get_pattern_for_layout(layout_name)
        if not pattern:
            available_layouts = self.pattern_loader.get_layout_names()
            raise ValueError(f"No pattern found for layout '{layout_name}' - check structured_frontmatter_patterns/. " f"Available layouts: {available_layouts}")

        # Step 2: Get expected fields from pattern
        expected_fields = pattern.get("yaml_pattern", {})
        if not expected_fields:
            raise ValueError(f"Pattern for layout '{layout_name}' has no yaml_pattern fields defined")

        # Step 3: Get layout object for normalization (if not provided)
        if not layout:
            # Get presentation from slide's slide layout
            prs = slide.slide_layout.part.package.presentation_part.presentation
            layout_result = self.layout_resolver.resolve_layout_safely(prs, layout_name)

            if not layout_result["success"]:
                # Handle layout resolution failure gracefully
                error_msg = f"Cannot normalize placeholders - {layout_result['message']}"
                if layout_result.get("suggestions"):
                    error_msg += f"\nSuggestions: {', '.join(layout_result['suggestions'])}"
                raise ValueError(error_msg)

            layout = layout_result["layout"]

        # Step 4: NORMALIZE SLIDE PLACEHOLDER NAMES - happens every slide build
        from .placeholder_normalizer import PlaceholderNormalizer

        normalizer = PlaceholderNormalizer()
        name_changes = normalizer.normalize_slide_placeholder_names(slide, layout)

        if name_changes:
            # Placeholder names were normalized
            pass

        # Step 5: Map each field to its placeholder by name (now works reliably)
        mapped_placeholders = {}

        # BUGFIX: Extract placeholder data from nested structure
        placeholder_data = slide_data.get("placeholders", slide_data)

        for field_name, _field_value in placeholder_data.items():
            # Skip metadata fields that aren't content placeholders
            if field_name in ["layout", "style", "speaker_notes"]:
                continue

            # Only process fields that are expected by the pattern
            if field_name in expected_fields:
                # Use name-based resolution (now guaranteed to work after normalization)
                placeholder = self.placeholder_resolver.get_placeholder_by_name(slide, field_name)

                if placeholder:
                    mapped_placeholders[field_name] = placeholder
                else:
                    # Clear error - should be rare after normalization
                    available_names = self.placeholder_resolver.list_placeholder_names(slide)
                    placeholder_summary = self.placeholder_resolver.get_placeholder_summary(slide)

                    raise ValueError(
                        f"Cannot find placeholder named '{field_name}' on slide for layout '{layout_name}' "
                        f"(after normalization). Available placeholder names: {available_names}. "
                        f"Placeholder details: {placeholder_summary}"
                    )

        return mapped_placeholders

    def validate_pattern_compatibility(self, slide: Slide, layout_name: str, slide_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Validate that slide data is compatible with pattern and layout.

        Args:
            slide: PowerPoint slide object
            layout_name: Layout name to validate against
            slide_data: Slide data to validate

        Returns:
            Dictionary with validation results
        """
        try:
            # Test the mapping without errors
            mapped = self.map_fields_to_placeholders(slide, slide_data, layout_name)

            return {
                "valid": True,
                "layout_name": layout_name,
                "mapped_fields": list(mapped.keys()),
                "field_count": len(mapped),
                "placeholder_names": self.placeholder_resolver.list_placeholder_names(slide),
                "placeholder_summary": self.placeholder_resolver.get_placeholder_summary(slide),
            }

        except ValueError as e:
            return {
                "valid": False,
                "layout_name": layout_name,
                "error": str(e),
                "slide_fields": list(slide_data.keys()),
                "placeholder_names": self.placeholder_resolver.list_placeholder_names(slide),
                "placeholder_summary": self.placeholder_resolver.get_placeholder_summary(slide),
            }

    def get_pattern_info(self, layout_name: str) -> Dict[str, Any]:
        """
        Get detailed pattern information for debugging.

        Args:
            layout_name: Layout name to get pattern for

        Returns:
            Dictionary with pattern information
        """
        pattern = self.pattern_loader.get_pattern_for_layout(layout_name)

        if not pattern:
            return {
                "found": False,
                "layout_name": layout_name,
                "available_layouts": self.pattern_loader.get_layout_names(),
            }

        return {
            "found": True,
            "layout_name": layout_name,
            "yaml_pattern": pattern.get("yaml_pattern", {}),
            "validation": pattern.get("validation", {}),
            "description": pattern.get("description", "No description"),
            "expected_fields": list(pattern.get("yaml_pattern", {}).keys()),
        }

    def get_placeholder_details(self, slide: Slide) -> Dict[str, Any]:
        """
        Get detailed placeholder information for debugging.

        Args:
            slide: PowerPoint slide object

        Returns:
            Dictionary with placeholder details
        """
        return {
            "placeholder_count": len(slide.placeholders),
            "placeholder_names": self.placeholder_resolver.list_placeholder_names(slide),
            "placeholder_summary": self.placeholder_resolver.get_placeholder_summary(slide),
            "placeholder_by_type": self.placeholder_resolver.get_placeholder_type_summary(slide),
        }

    def clear_pattern_cache(self) -> None:
        """Clear pattern cache to force reloading."""
        self.pattern_loader.clear_cache()

    def get_available_layouts(self) -> List[str]:
        """Get list of all available layout names."""
        return self.pattern_loader.get_layout_names()

    def find_placeholder_by_name_with_suggestions(self, slide: Slide, placeholder_name: str) -> SlidePlaceholder:
        """
        Find placeholder by name with helpful error messages.

        Delegates to existing PlaceholderResolver functionality.

        Args:
            slide: PowerPoint slide object
            placeholder_name: Placeholder name to find

        Returns:
            SlidePlaceholder object

        Raises:
            ValueError: With helpful suggestions if not found
        """
        return self.placeholder_resolver.find_placeholder_by_name_with_suggestions(slide, placeholder_name)

    def validate_placeholder_exists(self, slide: Slide, placeholder_name: str) -> bool:
        """
        Check if placeholder exists by name.

        Args:
            slide: PowerPoint slide object
            placeholder_name: Placeholder name to check

        Returns:
            True if placeholder exists, False otherwise
        """
        return self.placeholder_resolver.validate_placeholder_exists(slide, placeholder_name)
