#!/usr/bin/env python3
"""
Universal Content Formatting Module

Handles ALL formatting types independently of structure processing.
This module can be used by both JSON and Markdown processing without circular dependencies.

Supports:
- Inline formatting: **bold**, *italic*, ___underline___, ***bold italic***
- Rich content blocks: headings, paragraphs, bullets with levels
- Table cell formatting: maintaining structure while applying formatting
- Complex combinations: mixed formatting within same text

Separated from structure handling to enable clean architecture.
"""

import re
from typing import List, Dict, Any, Union


class ContentFormatter:
    """Universal content formatting processor."""

    def __init__(self, language_code=None, font_name=None):
        """Initialize the content formatter with optional language and font settings."""
        self.language_code = language_code
        self.font_name = font_name

        # Only import FormattingSupport if formatting needed
        self.formatter = None
        if language_code is not None or font_name is not None:
            from .formatting_support import FormattingSupport

            self.formatter = FormattingSupport()

    def parse_inline_formatting(self, text: str) -> List[Dict[str, Any]]:
        """
        Parse inline formatting and return structured formatting data.

        Args:
            text: Text with inline formatting markers

        Returns:
            List of segments with text and formatting attributes
        """
        if not text:
            return [{"text": "", "format": {}}]

        # Patterns in order of precedence (longest patterns first to avoid conflicts)
        patterns = [
            (
                r"\*\*\*___(.*?)___\*\*\*",
                {"bold": True, "italic": True, "underline": True},
            ),  # ***___text___***
            (
                r"___\*\*\*(.*?)\*\*\*___",
                {"bold": True, "italic": True, "underline": True},
            ),  # ___***text***___
            (r"\*\*\*(.*?)\*\*\*", {"bold": True, "italic": True}),  # ***text***
            (r"___(.*?)___", {"underline": True}),  # ___text___
            (r"\*\*(.*?)\*\*", {"bold": True}),  # **text**
            (r"\*(.*?)\*", {"italic": True}),  # *text*
        ]

        # Find all matches and their positions
        all_matches = []
        for pattern, format_dict in patterns:
            for match in re.finditer(pattern, text):
                all_matches.append((match.start(), match.end(), match.group(1), format_dict))

        # Sort matches by position
        all_matches.sort(key=lambda x: x[0])

        # Remove overlapping matches (keep the first one found)
        filtered_matches = []
        last_end = 0
        for start, end, content, format_dict in all_matches:
            if start >= last_end:
                filtered_matches.append((start, end, content, format_dict))
                last_end = end

        # Build the formatted text segments
        segments = []
        last_pos = 0

        for start, end, content, format_dict in filtered_matches:
            # Add plain text before the formatted text
            if start > last_pos:
                plain_text = text[last_pos:start]
                if plain_text:
                    segments.append({"text": plain_text, "format": {}})

            # Add formatted text
            segments.append({"text": content, "format": format_dict})
            last_pos = end

        # Add any remaining plain text
        if last_pos < len(text):
            remaining_text = text[last_pos:]
            if remaining_text:
                segments.append({"text": remaining_text, "format": {}})

        # If no formatting found, return the original text
        if not segments:
            segments = [{"text": text, "format": {}}]

        return segments

    def apply_inline_formatting(self, text: str, paragraph) -> None:
        """
        Apply inline formatting to a paragraph, including heading detection and font scaling.

        Args:
            text: The text to format.
            paragraph: The paragraph to apply the formatting to.
        """
        # Check if this text contains heading syntax - process line by line for headings
        lines = text.split("\n")

        if len(lines) > 1:
            # Multi-line text - check for headings
            self._apply_multiline_formatting_with_headings(text, paragraph)
        else:
            # Single line - check if it's a heading
            line = lines[0].strip()
            heading_level = self._detect_heading_level(line)

            if heading_level:
                # Extract heading text without the # markers
                heading_text = line[heading_level + 1 :].strip()  # +1 for the space after #
                self._apply_heading_formatting(paragraph, heading_text, heading_level)
            else:
                # Regular text - apply normal inline formatting
                segments = self.parse_inline_formatting(text)
                for segment in segments:
                    run = paragraph.add_run()
                    run.text = segment["text"]
                    if "bold" in segment["format"]:
                        run.font.bold = True
                    if "italic" in segment["format"]:
                        run.font.italic = True
                    if "underline" in segment["format"]:
                        run.font.underline = True
                    # Apply language and font formatting if specified
                    self._apply_language_font_formatting(run)

    def apply_formatted_segments_to_paragraph(self, formatted_segments, paragraph):
        """Apply formatted text segments to a paragraph."""
        if not formatted_segments:
            return

        # Clear existing runs
        paragraph.clear()

        for segment in formatted_segments:
            if isinstance(segment, dict) and "text" in segment:
                text = segment["text"]
                format_info = segment.get("format", {})

                run = paragraph.add_run()
                run.text = text

                # Apply inline formatting
                if format_info.get("bold"):
                    run.font.bold = True
                if format_info.get("italic"):
                    run.font.italic = True
                if format_info.get("underline"):
                    run.font.underline = True

                # Apply language and font formatting if specified
                self._apply_language_font_formatting(run)

    def _apply_language_font_formatting(self, run):
        """Apply language and font formatting to a text run if specified."""
        if self.formatter is None:
            return  # No formatting specified

        if self.language_code is not None:
            self.formatter.apply_language_to_run(run, self.language_code)

        if self.font_name is not None:
            self.formatter.apply_font_to_run(run, self.font_name)

    def apply_formatted_segments_to_cell(self, cell, segments):
        """Apply formatted text segments to a table cell."""
        text_frame = cell.text_frame
        text_frame.clear()

        # Create first paragraph
        paragraph = text_frame.paragraphs[0]
        paragraph.text = ""

        # Apply each segment
        for segment in segments:
            run = paragraph.add_run()
            run.text = segment["text"]

            # Apply formatting
            format_dict = segment["format"]
            if format_dict.get("bold"):
                run.font.bold = True
            if format_dict.get("italic"):
                run.font.italic = True
            if format_dict.get("underline"):
                run.font.underline = True
            # Apply language and font formatting if specified
            self._apply_language_font_formatting(run)

    def add_content_to_placeholder(self, placeholder, content):
        """
        Add content to a placeholder, handling lists and strings.

        Args:
            placeholder: The placeholder to add content to.
            content: The content to add.
        """
        if isinstance(content, list):
            for item in content:
                p = placeholder.text_frame.add_paragraph()
                self.apply_inline_formatting(item, p)
        else:
            self.apply_inline_formatting(str(content), placeholder.text_frame.paragraphs[0])

    def format_simple_content_list(self, content_list: List[str]) -> List[Dict[str, Any]]:
        """
        Convert simple content list to rich content with formatting.

        Args:
            content_list: List of strings with possible inline formatting

        Returns:
            Rich content list with formatting preserved
        """
        if not content_list:
            return []

        rich_content = []
        for item in content_list:
            if isinstance(item, str):
                # Treat as paragraph text with inline formatting
                rich_content.append({"paragraph": item, "formatted": self.parse_inline_formatting(item)})
            else:
                # Keep non-string items as-is
                rich_content.append(item)

        return rich_content

    def format_field_content(self, content: Union[str, List[str], Dict[str, Any]]) -> Dict[str, Any]:
        """
        Format any field content type with inline formatting preservation.

        Args:
            content: Content to format (string, list, or dict)

        Returns:
            Formatted content with structure preserved
        """
        if isinstance(content, str):
            return {"text": content, "formatted": self.parse_inline_formatting(content)}
        elif isinstance(content, list):
            return {
                "list": content,
                "formatted_list": [
                    {
                        "text": item,
                        "formatted": (self.parse_inline_formatting(item) if isinstance(item, str) else item),
                    }
                    for item in content
                ],
            }
        elif isinstance(content, dict):
            # Recursively format dict content
            formatted_dict = {}
            for key, value in content.items():
                if isinstance(value, str):
                    formatted_dict[key] = {
                        "text": value,
                        "formatted": self.parse_inline_formatting(value),
                    }
                else:
                    formatted_dict[key] = value
            return formatted_dict
        else:
            # Return non-string/list/dict content as-is
            return {"raw": content}

    def format_table_data(self, table_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Format table data while preserving structure.

        Args:
            table_data: Table configuration with data

        Returns:
            Table data with formatting applied to cells
        """
        if not isinstance(table_data, dict) or "data" not in table_data:
            return table_data

        formatted_table = table_data.copy()

        if isinstance(table_data["data"], list):
            formatted_data = []
            for row in table_data["data"]:
                if isinstance(row, list):
                    formatted_row = []
                    for cell in row:
                        if isinstance(cell, str):
                            formatted_row.append({"text": cell, "formatted": self.parse_inline_formatting(cell)})
                        else:
                            # Keep non-string cells as-is
                            formatted_row.append(cell)
                    formatted_data.append(formatted_row)
                else:
                    # Keep non-list rows as-is
                    formatted_data.append(row)
            formatted_table["data"] = formatted_data

        return formatted_table

    def format_rich_content_blocks(self, rich_content: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """
        Format rich content blocks (headings, paragraphs, bullets).

        Args:
            rich_content: List of rich content blocks

        Returns:
            Rich content with formatting applied
        """
        if not rich_content:
            return []

        formatted_blocks = []

        for block in rich_content:
            if not isinstance(block, dict):
                formatted_blocks.append(block)
                continue

            formatted_block = block.copy()

            # Format headings
            if "heading" in block:
                formatted_block["heading_formatted"] = self.parse_inline_formatting(block["heading"])

            # Format paragraphs
            if "paragraph" in block:
                formatted_block["paragraph_formatted"] = self.parse_inline_formatting(block["paragraph"])

            # Format bullets
            if "bullets" in block and isinstance(block["bullets"], list):
                formatted_bullets = []
                for bullet in block["bullets"]:
                    if isinstance(bullet, str):
                        formatted_bullets.append({"text": bullet, "formatted": self.parse_inline_formatting(bullet)})
                    else:
                        formatted_bullets.append(bullet)
                formatted_block["bullets_formatted"] = formatted_bullets

            formatted_blocks.append(formatted_block)

        return formatted_blocks

    def format_slide_data(self, slide_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Format complete slide data with ALL formatting types preserved.

        This is the main entry point for slide formatting.

        Args:
            slide_data: Raw slide data dictionary

        Returns:
            Slide data with all formatting applied
        """
        if not isinstance(slide_data, dict):
            raise TypeError(f"slide_data must be a dictionary, got {type(slide_data).__name__}")

        # Create a copy to avoid modifying original
        formatted_data = slide_data.copy()

        # Format title and subtitle
        for field in ["title", "subtitle"]:
            if field in formatted_data and formatted_data[field]:
                formatted_data[f"{field}_formatted"] = self.parse_inline_formatting(formatted_data[field])

        # Format simple content lists to rich content
        if "content" in formatted_data and isinstance(formatted_data["content"], list):
            formatted_data["rich_content"] = self.format_simple_content_list(formatted_data["content"])
            # Keep original for compatibility
            # del formatted_data["content"]

        # Format rich content blocks
        if "rich_content" in formatted_data:
            formatted_data["rich_content_formatted"] = self.format_rich_content_blocks(formatted_data["rich_content"])

        # Format table data
        if "table" in formatted_data:
            formatted_data["table_formatted"] = self.format_table_data(formatted_data["table"])

        # Format all field content (for semantic field names like content_left, content_col1)
        content_fields = [key for key in formatted_data.keys() if key.startswith(("content_", "title_")) and not key.endswith("_formatted")]

        for field in content_fields:
            content = formatted_data[field]
            if content:  # Only format non-empty content
                formatted_data[f"{field}_formatted"] = self.format_field_content(content)

        return formatted_data

    def _detect_heading_level(self, line: str) -> int:
        """Detect heading level from line text (H1-H6)."""
        if line.startswith("###### "):
            return 6
        elif line.startswith("##### "):
            return 5
        elif line.startswith("#### "):
            return 4
        elif line.startswith("### "):
            return 3
        elif line.startswith("## "):
            return 2
        elif line.startswith("# "):
            return 1
        return 0

    def _apply_heading_formatting(self, paragraph, text_content, heading_level):
        """Apply heading formatting with font scaling and proper weights."""
        from pptx.util import Pt

        # Apply inline formatting first (bold, italic, underline)
        segments = self.parse_inline_formatting(text_content)
        for segment in segments:
            run = paragraph.add_run()
            run.text = segment["text"]
            if "bold" in segment["format"]:
                run.font.bold = True
            if "italic" in segment["format"]:
                run.font.italic = True
            if "underline" in segment["format"]:
                run.font.underline = True
            # Apply language and font formatting if specified
            self._apply_language_font_formatting(run)

        # Apply heading-specific font scaling and weight
        base_size = 18  # PowerPoint default content size

        # Calculate scaled font size
        scale_factors = {1: 1.33, 2: 1.11, 3: 1.0, 4: 0.89, 5: 0.78, 6: 0.67}
        heading_font_size = base_size * scale_factors.get(heading_level, 1.0)

        # Apply font size and weight to all runs
        for run in paragraph.runs:
            if heading_level in [1, 2]:
                # H1 and H2: Light weight
                run.font.bold = False
            else:
                # H3-H6: Bold weight
                run.font.bold = True
            run.font.size = Pt(heading_font_size)

    def _apply_multiline_formatting_with_headings(self, text, paragraph):
        """Handle multi-line text with heading detection."""
        lines = text.split("\n")

        # Get text frame for adding multiple paragraphs
        text_frame = paragraph._parent

        # Clear existing content
        text_frame.clear()

        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue

            # Create paragraph for this line
            if i == 0:
                current_paragraph = text_frame.paragraphs[0]
            else:
                current_paragraph = text_frame.add_paragraph()

            heading_level = self._detect_heading_level(line)

            if heading_level:
                # Extract heading text
                heading_text = line[heading_level + 1 :].strip()
                self._apply_heading_formatting(current_paragraph, heading_text, heading_level)
            else:
                # Regular text
                segments = self.parse_inline_formatting(line)
                for segment in segments:
                    run = current_paragraph.add_run()
                    run.text = segment["text"]
                    if "bold" in segment["format"]:
                        run.font.bold = True
                    if "italic" in segment["format"]:
                        run.font.italic = True
                    if "underline" in segment["format"]:
                        run.font.underline = True

    # REMOVED: parse_table_markdown_with_formatting() - complex markdown parsing
    # Replaced with plain text processing in TableHandler for 50%+ performance improvement
    # Table cells now contain plain text only, markdown treated as literal characters


# Global formatter instance factory
def get_content_formatter(language_code=None, font_name=None):
    """Get a ContentFormatter instance with optional formatting parameters."""
    return ContentFormatter(language_code=language_code, font_name=font_name)


# Default formatter instance for backward compatibility
content_formatter = ContentFormatter()


def format_slide_content(slide_data: Dict[str, Any]) -> Dict[str, Any]:
    """
    Convenience function for formatting slide content.

    Args:
        slide_data: Raw slide data

    Returns:
        Formatted slide data
    """
    return content_formatter.format_slide_data(slide_data)


def format_inline_text(text: str) -> List[Dict[str, Any]]:
    """
    Convenience function for formatting inline text.

    Args:
        text: Text with inline formatting

    Returns:
        List of formatted segments
    """
    return content_formatter.parse_inline_formatting(text)
