#!/usr/bin/env python3
"""
Deckbuilder Formatting Support

Comprehensive language and font formatting support for PowerPoint presentations.
Provides language ID constants, font validation, and unified presentation processing
that updates both master slides and content slides.
"""

import json
import os
import re
from difflib import get_close_matches
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.enum.lang import MSO_LANGUAGE_ID


class FormattingSupport:
    """Comprehensive formatting support for language and font settings"""

    # Language ID constants mapped to common locale codes
    # Using only verified constants from python-pptx MSO_LANGUAGE_ID
    LANGUAGE_IDS = {
        "en-US": MSO_LANGUAGE_ID.ENGLISH_US,
        "en-AU": MSO_LANGUAGE_ID.ENGLISH_AUS,
        "en-GB": MSO_LANGUAGE_ID.ENGLISH_UK,
        "en-CA": MSO_LANGUAGE_ID.ENGLISH_CANADIAN,
        "en-NZ": MSO_LANGUAGE_ID.ENGLISH_NEW_ZEALAND,
        "en-ZA": MSO_LANGUAGE_ID.ENGLISH_SOUTH_AFRICA,
        "es-ES": MSO_LANGUAGE_ID.SPANISH,
        "fr-FR": MSO_LANGUAGE_ID.FRENCH,
        "fr-CA": MSO_LANGUAGE_ID.FRENCH_CANADIAN,
        "de-DE": MSO_LANGUAGE_ID.GERMAN,
        "it-IT": MSO_LANGUAGE_ID.ITALIAN,
        "pt-PT": MSO_LANGUAGE_ID.PORTUGUESE,
        "pt-BR": MSO_LANGUAGE_ID.BRAZILIAN_PORTUGUESE,
        "nl-NL": MSO_LANGUAGE_ID.DUTCH,
        "sv-SE": MSO_LANGUAGE_ID.SWEDISH,
        "da-DK": MSO_LANGUAGE_ID.DANISH,
        "fi-FI": MSO_LANGUAGE_ID.FINNISH,
        "ru-RU": MSO_LANGUAGE_ID.RUSSIAN,
        "ja-JP": MSO_LANGUAGE_ID.JAPANESE,
        "ko-KR": MSO_LANGUAGE_ID.KOREAN,
    }

    # Common system fonts for validation
    COMMON_FONTS = [
        "Arial",
        "Calibri",
        "Segoe UI",
        "Times New Roman",
        "Verdana",
        "Tahoma",
        "Georgia",
        "Comic Sans MS",
        "Trebuchet MS",
        "Impact",
        "Palatino Linotype",
        "Book Antiqua",
        "Lucida Console",
        "Courier New",
        "Consolas",
        "Cambria",
        "Candara",
        "Corbel",
        "Franklin Gothic Medium",
        "Century Gothic",
        "Arial Black",
    ]

    def __init__(self):
        """Initialize formatting support"""
        self._language_mappings = {}

    # Language name to code mapping for backward compatibility
    LANGUAGE_NAME_TO_CODE = {
        "English (United States)": "en-US",
        "English (Australia)": "en-AU",
        "English (United Kingdom)": "en-GB",
        "English (Canada)": "en-CA",
        "English (New Zealand)": "en-NZ",
        "English (South Africa)": "en-ZA",
        "Spanish (Spain)": "es-ES",
        "French (France)": "fr-FR",
        "French (Canada)": "fr-CA",
        "German (Germany)": "de-DE",
        "Italian (Italy)": "it-IT",
        "Portuguese (Portugal)": "pt-PT",
        "Portuguese (Brazil)": "pt-BR",
        "Dutch (Netherlands)": "nl-NL",
        "Swedish (Sweden)": "sv-SE",
        "Danish (Denmark)": "da-DK",
        "Finnish (Finland)": "fi-FI",
        "Russian (Russia)": "ru-RU",
        "Japanese (Japan)": "ja-JP",
        "Korean (South Korea)": "ko-KR",
    }

    @classmethod
    def get_supported_languages(cls) -> Dict[str, str]:
        """
        Get dictionary of supported language codes with descriptions.

        Returns:
            Dict mapping language codes to human-readable descriptions
        """
        descriptions = {
            "en-US": "English (United States)",
            "en-AU": "English (Australia)",
            "en-GB": "English (United Kingdom)",
            "en-CA": "English (Canada)",
            "en-NZ": "English (New Zealand)",
            "en-ZA": "English (South Africa)",
            "es-ES": "Spanish (Spain)",
            "fr-FR": "French (France)",
            "fr-CA": "French (Canada)",
            "de-DE": "German (Germany)",
            "it-IT": "Italian (Italy)",
            "pt-PT": "Portuguese (Portugal)",
            "pt-BR": "Portuguese (Brazil)",
            "nl-NL": "Dutch (Netherlands)",
            "sv-SE": "Swedish (Sweden)",
            "da-DK": "Danish (Denmark)",
            "fi-FI": "Finnish (Finland)",
            "ru-RU": "Russian (Russia)",
            "ja-JP": "Japanese (Japan)",
            "ko-KR": "Korean (South Korea)",
        }
        return descriptions

    @classmethod
    def normalize_language_input(cls, language_input: str) -> Optional[str]:
        """
        Normalize language input to standard locale code format.

        Accepts both locale codes (en-AU) and full names (English (Australia)).

        Args:
            language_input: Language code or name

        Returns:
            Normalized locale code or None if not found
        """
        # Try direct code lookup first
        if language_input in cls.LANGUAGE_IDS:
            return language_input

        # Try full name lookup
        if language_input in cls.LANGUAGE_NAME_TO_CODE:
            return cls.LANGUAGE_NAME_TO_CODE[language_input]

        return None

    @classmethod
    def validate_language_code(cls, language_code: str) -> Tuple[bool, Optional[str], List[str]]:
        """
        Validate language code and provide suggestions for invalid codes.

        Accepts both locale codes (en-AU) and full names (English (Australia)).

        Args:
            language_code: Language code or name to validate

        Returns:
            Tuple of (is_valid, error_message, suggestions)
        """
        # Try to normalize the input
        normalized = cls.normalize_language_input(language_code)
        if normalized:
            return True, None, []

        # Generate suggestions for invalid codes from both formats
        all_valid_inputs = list(cls.LANGUAGE_IDS.keys()) + list(cls.LANGUAGE_NAME_TO_CODE.keys())
        suggestions = get_close_matches(language_code, all_valid_inputs, n=3, cutoff=0.6)

        error_msg = f"Unsupported language: '{language_code}'"
        return False, error_msg, suggestions

    @classmethod
    def validate_font_name(cls, font_name: str) -> Tuple[bool, Optional[str], List[str]]:
        """
        Validate font name and provide suggestions for common fonts.

        Args:
            font_name: Font name to validate

        Returns:
            Tuple of (is_valid, warning_message, suggestions)
        """
        # Always accept any font name (user might have custom fonts)
        # But provide suggestions if it's not a common font
        if font_name in cls.COMMON_FONTS:
            return True, None, []

        # Find similar common fonts
        suggestions = get_close_matches(font_name, cls.COMMON_FONTS, n=3, cutoff=0.6)

        warning_msg = f"Font '{font_name}' is not a common system font"
        return True, warning_msg, suggestions

    def apply_language_to_run(self, run, language_input: str) -> bool:
        """
        Apply language setting to a text run.

        Accepts both locale codes (en-AU) and full names (English (Australia)).

        Args:
            run: python-pptx Run object
            language_input: Language code or name (e.g., 'en-AU' or 'English (Australia)')

        Returns:
            True if language was applied successfully
        """
        try:
            # Normalize input to locale code
            language_code = self.normalize_language_input(language_input)
            if language_code and language_code in self.LANGUAGE_IDS:
                run.font.language_id = self.LANGUAGE_IDS[language_code]
                return True
            return False
        except Exception:
            return False

    def apply_font_to_run(self, run, font_name: str) -> bool:
        """
        Apply font family to a text run.

        Args:
            run: python-pptx Run object
            font_name: Font family name

        Returns:
            True if font was applied successfully
        """
        try:
            run.font.name = font_name
            return True
        except Exception:
            return False

    def load_language_mapping(self, language_code: str) -> Optional[Dict]:
        """
        Load language mapping configuration for text replacement.

        Args:
            language_code: Target language code (e.g., 'en-AU')

        Returns:
            Language mapping dictionary or None if not found
        """
        if language_code in self._language_mappings:
            return self._language_mappings[language_code]

        # Find language mapping file
        current_dir = Path(__file__).parent
        mapping_file = current_dir / "language_mappings" / f"{language_code}.json"

        if not mapping_file.exists():
            return None

        try:
            with open(mapping_file, "r", encoding="utf-8") as f:
                mapping_data = json.load(f)
            self._language_mappings[language_code] = mapping_data
            return mapping_data
        except Exception as e:
            print(f"⚠️  Warning: Could not load language mapping for {language_code}: {e}")
            return None

    def preserve_case(self, original_word: str, replacement_word: str) -> str:
        """
        Preserve the case pattern of the original word in the replacement.

        Args:
            original_word: Original word with case pattern
            replacement_word: Replacement word in lowercase

        Returns:
            Replacement word with preserved case pattern
        """
        if original_word.isupper():
            return replacement_word.upper()
        elif original_word.islower():
            return replacement_word.lower()
        elif original_word.istitle():
            # Handle title case for phrases (e.g., "Cell Phone" -> "Mobile Phone")
            return replacement_word.title()
        else:
            # Mixed case - attempt to preserve pattern character by character
            result = ""
            for i, char in enumerate(replacement_word):
                if i < len(original_word):
                    if original_word[i].isupper():
                        result += char.upper()
                    else:
                        result += char.lower()
                else:
                    result += char.lower()
            return result

    def is_context_exception(self, text: str, word_position: int, word: str, except_contexts: List[str]) -> bool:
        """
        Check if a word should be excluded from replacement based on context.

        Args:
            text: Full text containing the word
            word_position: Position of the word in the text
            word: The word to check
            except_contexts: List of context words that should prevent replacement

        Returns:
            True if word should be excluded from replacement
        """
        if not except_contexts:
            return False

        # Get surrounding context (30 characters before and after)
        start = max(0, word_position - 30)
        end = min(len(text), word_position + len(word) + 30)
        context = text[start:end].lower()

        # Check if any exception context is present
        for exception_context in except_contexts:
            if exception_context.lower() in context:
                return True

        return False

    def apply_text_replacements(self, text: str, language_code: str) -> str:
        """
        Apply text replacements based on language mapping configuration.

        Args:
            text: Text to process
            language_code: Target language code

        Returns:
            Text with language-specific replacements applied
        """
        if not text:
            return text

        mapping = self.load_language_mapping(language_code)
        if not mapping:
            return text

        result_text = text

        # Apply spelling pattern replacements
        if "spelling_patterns" in mapping:
            for source_word, target_word in mapping["spelling_patterns"].items():
                # Use word boundaries for exact matches
                pattern = r"\b" + re.escape(source_word) + r"\b"
                matches = list(re.finditer(pattern, result_text, re.IGNORECASE))

                # Process matches in reverse order to preserve positions
                for match in reversed(matches):
                    original_word = match.group()
                    replacement = self.preserve_case(original_word, target_word)
                    result_text = result_text[: match.start()] + replacement + result_text[match.end() :]

        # Apply conditional mappings (like program/programme)
        if "conditional_mappings" in mapping:
            for source_word, config in mapping["conditional_mappings"].items():
                target_word = config["to"]
                except_contexts = config.get("except_contexts", [])

                pattern = r"\b" + re.escape(source_word) + r"\b"
                matches = list(re.finditer(pattern, result_text, re.IGNORECASE))

                for match in reversed(matches):
                    # Check if this occurrence should be excluded
                    if not self.is_context_exception(result_text, match.start(), source_word, except_contexts):
                        original_word = match.group()
                        replacement = self.preserve_case(original_word, target_word)
                        result_text = result_text[: match.start()] + replacement + result_text[match.end() :]

        # Apply vocabulary replacements (phrases)
        if "vocabulary" in mapping:
            for source_phrase, target_phrase in mapping["vocabulary"].items():
                # Use word boundaries for phrase matching
                pattern = r"\b" + re.escape(source_phrase) + r"\b"
                matches = list(re.finditer(pattern, result_text, re.IGNORECASE))

                for match in reversed(matches):
                    original_phrase = match.group()
                    replacement = self.preserve_case(original_phrase, target_phrase)
                    result_text = result_text[: match.start()] + replacement + result_text[match.end() :]

        return result_text

    def process_text_frame(self, text_frame, language_code: Optional[str] = None, font_name: Optional[str] = None) -> Dict[str, int]:
        """
        Process all text runs in a text frame, applying text replacements, language and/or font settings.

        Args:
            text_frame: python-pptx TextFrame object
            language_code: Optional language code to apply
            font_name: Optional font name to apply

        Returns:
            Dictionary with processing statistics
        """
        stats = {"runs_processed": 0, "language_applied": 0, "font_applied": 0, "text_replaced": 0}

        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    stats["runs_processed"] += 1

                    # Apply text replacements first if language code is provided
                    if language_code and run.text:
                        original_text = run.text
                        replaced_text = self.apply_text_replacements(original_text, language_code)
                        if replaced_text != original_text:
                            run.text = replaced_text
                            stats["text_replaced"] += 1

                    # Apply language ID changes
                    if language_code:
                        if self.apply_language_to_run(run, language_code):
                            stats["language_applied"] += 1

                    # Apply font changes
                    if font_name:
                        if self.apply_font_to_run(run, font_name):
                            stats["font_applied"] += 1

        except Exception as e:
            # Log error but continue processing
            print(f"⚠️  Warning: Error processing text frame: {e}")

        return stats

    def process_shape(self, shape, language_code: Optional[str] = None, font_name: Optional[str] = None) -> Dict[str, int]:
        """
        Process a shape, applying formatting to its text content.

        Args:
            shape: python-pptx Shape object
            language_code: Optional language code to apply
            font_name: Optional font name to apply

        Returns:
            Dictionary with processing statistics
        """
        stats = {"runs_processed": 0, "language_applied": 0, "font_applied": 0, "text_replaced": 0}

        try:
            # Handle text frames
            if hasattr(shape, "text_frame") and shape.text_frame:
                frame_stats = self.process_text_frame(shape.text_frame, language_code, font_name)
                for key in stats:
                    stats[key] += frame_stats[key]

            # Handle tables
            elif hasattr(shape, "table"):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            cell_stats = self.process_text_frame(cell.text_frame, language_code, font_name)
                            for key in stats:
                                stats[key] += cell_stats[key]

            # Handle grouped shapes
            elif hasattr(shape, "shapes"):
                for sub_shape in shape.shapes:
                    sub_stats = self.process_shape(sub_shape, language_code, font_name)
                    for key in stats:
                        stats[key] += sub_stats[key]

        except Exception as e:
            print(f"⚠️  Warning: Error processing shape: {e}")

        return stats

    def update_theme_fonts(self, presentation, font_name: str) -> bool:
        """
        Update theme majorFont and minorFont to use the specified font.

        Args:
            presentation: python-pptx Presentation object
            font_name: Font name to apply to both major and minor theme fonts

        Returns:
            True if theme fonts were updated successfully, False otherwise
        """
        try:
            # Access theme part through presentation relationships
            theme_part = None
            for rel in presentation.part.rels.values():
                if "theme" in rel.target_ref:
                    theme_part = rel.target_part
                    break

            if not theme_part:
                return False

            # Parse theme XML from blob
            from lxml import etree  # nosec B410 - parsing trusted PowerPoint XML

            theme_xml_bytes = theme_part.blob
            theme_xml = etree.fromstring(theme_xml_bytes)  # nosec B320 - parsing trusted PowerPoint XML

            # Define DrawingML namespace
            DRAWINGML_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

            # Find and update majorFont latin typeface
            major_latin_elements = theme_xml.xpath(".//a:majorFont/a:latin", namespaces=DRAWINGML_NS)
            minor_latin_elements = theme_xml.xpath(".//a:minorFont/a:latin", namespaces=DRAWINGML_NS)

            if major_latin_elements and minor_latin_elements:
                # Update both major and minor fonts to use the same font
                major_latin_elements[0].set("typeface", font_name)
                minor_latin_elements[0].set("typeface", font_name)

                # Save the modified XML back to the theme part
                modified_xml_bytes = etree.tostring(theme_xml, encoding="utf-8", xml_declaration=True)
                theme_part._blob = modified_xml_bytes

                return True

            return False

        except Exception as e:
            print(f"⚠️  Warning: Could not update theme fonts: {e}")
            return False

    def update_presentation(
        self,
        presentation_path: str,
        language_code: Optional[str] = None,
        font_name: Optional[str] = None,
        output_path: Optional[str] = None,
        create_backup: bool = True,
    ) -> Dict[str, any]:
        """
        Update both master slides and content slides in a PowerPoint presentation.

        Args:
            presentation_path: Path to the PowerPoint file
            language_code: Optional language code to apply
            font_name: Optional font name to apply
            output_path: Optional output path (default: update in place)
            create_backup: Whether to create a backup file

        Returns:
            Dictionary with operation results and statistics
        """
        # Fallback to environment variables when parameters are None
        if language_code is None:
            language_code = get_default_language()
        if font_name is None:
            font_name = get_default_font()

        pptx_path = Path(presentation_path)

        if not pptx_path.exists():
            return {"success": False, "error": f"File not found: {presentation_path}", "stats": {}}

        # Create backup if requested
        backup_path = None
        if create_backup:
            backup_path = pptx_path.with_suffix(".bak.pptx")
            try:
                import shutil

                shutil.copy2(pptx_path, backup_path)
            except Exception as e:
                return {"success": False, "error": f"Failed to create backup: {e}", "stats": {}}

        try:
            # Load presentation
            prs = Presentation(str(pptx_path))

            total_stats = {
                "master_slides_processed": 0,
                "content_slides_processed": 0,
                "total_runs_processed": 0,
                "total_language_applied": 0,
                "total_font_applied": 0,
                "total_text_replaced": 0,
                "theme_fonts_updated": 0,
            }

            # Update theme fonts if font name is provided
            if font_name:
                if self.update_theme_fonts(prs, font_name):
                    total_stats["theme_fonts_updated"] = 1

            # Process master slides
            try:
                for slide_master in prs.slide_masters:
                    total_stats["master_slides_processed"] += 1

                    # Process master slide shapes
                    for shape in slide_master.shapes:
                        shape_stats = self.process_shape(shape, language_code, font_name)
                        total_stats["total_runs_processed"] += shape_stats["runs_processed"]
                        total_stats["total_language_applied"] += shape_stats["language_applied"]
                        total_stats["total_font_applied"] += shape_stats["font_applied"]
                        total_stats["total_text_replaced"] += shape_stats["text_replaced"]

                    # Process slide layouts
                    for slide_layout in slide_master.slide_layouts:
                        for shape in slide_layout.shapes:
                            shape_stats = self.process_shape(shape, language_code, font_name)
                            total_stats["total_runs_processed"] += shape_stats["runs_processed"]
                            total_stats["total_language_applied"] += shape_stats["language_applied"]
                            total_stats["total_font_applied"] += shape_stats["font_applied"]
                            total_stats["total_text_replaced"] += shape_stats["text_replaced"]

            except Exception as e:
                print(f"⚠️  Warning: Error processing master slides: {e}")

            # Process content slides
            for slide in prs.slides:
                total_stats["content_slides_processed"] += 1

                for shape in slide.shapes:
                    shape_stats = self.process_shape(shape, language_code, font_name)
                    total_stats["total_runs_processed"] += shape_stats["runs_processed"]
                    total_stats["total_language_applied"] += shape_stats["language_applied"]
                    total_stats["total_font_applied"] += shape_stats["font_applied"]
                    total_stats["total_text_replaced"] += shape_stats["text_replaced"]

            # Save presentation
            save_path = output_path if output_path else presentation_path
            prs.save(save_path)

            return {
                "success": True,
                "message": f"Presentation updated successfully: {save_path}",
                "backup_path": str(backup_path) if backup_path else None,
                "stats": total_stats,
            }

        except Exception as e:
            return {"success": False, "error": f"Failed to process presentation: {e}", "stats": {}}


def get_default_language() -> Optional[str]:
    """Get default language from environment variable"""
    return os.getenv("DECK_PROOFING_LANGUAGE")


def get_default_font() -> Optional[str]:
    """Get default font from environment variable"""
    return os.getenv("DECK_DEFAULT_FONT")


def print_supported_languages():
    """Print all supported languages in a user-friendly format"""
    languages = FormattingSupport.get_supported_languages()

    print("📋 Supported Proofing Languages:")
    for code, description in sorted(languages.items()):
        print(f"  {code:<6} {description}")
    print()
    print("💡 Usage examples:")
    print("  deckbuilder -l en-AU create presentation.md")
    print("  deckbuilder remap template.pptx --language en-US")
    print("  export DECK_PROOFING_LANGUAGE=en-GB")
