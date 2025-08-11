"""
PowerPoint Placeholder Type Constants and Mappings

This module defines semantic groupings of PowerPoint placeholder types for
generic content placement without hardcoding layout names. This allows the
deckbuilder to work with any PowerPoint template by detecting placeholder
types rather than relying on specific layout configurations.

Based on python-pptx PP_PLACEHOLDER_TYPE enumeration.
"""

from pptx.enum.shapes import PP_PLACEHOLDER_TYPE

# Title-related placeholders - for slide titles and headings
TITLE_PLACEHOLDERS = {
    PP_PLACEHOLDER_TYPE.TITLE,  # TITLE (1) - Standard slide title
    PP_PLACEHOLDER_TYPE.CENTER_TITLE,  # CENTER_TITLE (3) - Centered title (title slides)
    PP_PLACEHOLDER_TYPE.VERTICAL_TITLE,  # VERTICAL_TITLE (5) - Vertical orientation title
}

# Subtitle placeholders - for slide subtitles (typically on title slides)
SUBTITLE_PLACEHOLDERS = {PP_PLACEHOLDER_TYPE.SUBTITLE}  # SUBTITLE (4) - Subtitle text

# Main content placeholders - for primary slide content, bullets, paragraphs
CONTENT_PLACEHOLDERS = {
    PP_PLACEHOLDER_TYPE.BODY,  # BODY (2) - Main content area
    PP_PLACEHOLDER_TYPE.VERTICAL_BODY,  # VERTICAL_BODY (6) - Vertical text content
    PP_PLACEHOLDER_TYPE.OBJECT,  # OBJECT (7) - Often used for text content in modern templates
}

# Media and object placeholders - for rich content like images, charts, tables
MEDIA_PLACEHOLDERS = {
    PP_PLACEHOLDER_TYPE.PICTURE,  # PICTURE (18) - Image placeholders
    PP_PLACEHOLDER_TYPE.CHART,  # CHART (8) - Chart/graph placeholders
    PP_PLACEHOLDER_TYPE.TABLE,  # TABLE (12) - Data table placeholders
    PP_PLACEHOLDER_TYPE.MEDIA_CLIP,  # MEDIA_CLIP (10) - Video/audio content
    PP_PLACEHOLDER_TYPE.VERTICAL_OBJECT,  # VERTICAL_OBJECT (17) - Vertical objects
}

# Layout and metadata placeholders - for slide decorations and information
LAYOUT_PLACEHOLDERS = {
    PP_PLACEHOLDER_TYPE.HEADER,  # HEADER (14) - Page header text
    PP_PLACEHOLDER_TYPE.FOOTER,  # FOOTER (15) - Page footer text
    PP_PLACEHOLDER_TYPE.DATE,  # DATE (16) - Date/timestamp
    PP_PLACEHOLDER_TYPE.SLIDE_NUMBER,  # SLIDE_NUMBER (13) - Slide numbering
}

# Specialized placeholders - for specific content types
SPECIAL_PLACEHOLDERS = {
    PP_PLACEHOLDER_TYPE.ORG_CHART,  # ORG_CHART (11) - Organization charts
    PP_PLACEHOLDER_TYPE.BITMAP,  # BITMAP (9) - Bitmap images
    PP_PLACEHOLDER_TYPE.SLIDE_IMAGE,  # SLIDE_IMAGE (101) - Slide thumbnail images
    PP_PLACEHOLDER_TYPE.MIXED,  # MIXED (-2) - Mixed content types
}

# All placeholder types grouped by semantic function
ALL_PLACEHOLDER_GROUPS = {
    "title": TITLE_PLACEHOLDERS,
    "subtitle": SUBTITLE_PLACEHOLDERS,
    "content": CONTENT_PLACEHOLDERS,
    "media": MEDIA_PLACEHOLDERS,
    "layout": LAYOUT_PLACEHOLDERS,
    "special": SPECIAL_PLACEHOLDERS,
}


def get_placeholder_category(placeholder_type):
    """
    Determine the semantic category of a placeholder type.

    Args:
        placeholder_type: PP_PLACEHOLDER_TYPE enum value

    Returns:
        str: Category name ('title', 'subtitle', 'content', 'media', 'layout', 'special')
        None: If placeholder type is not recognized
    """
    if not isinstance(placeholder_type, PP_PLACEHOLDER_TYPE):
        return None
    for category, types in ALL_PLACEHOLDER_GROUPS.items():
        if _safe_enum_check(placeholder_type, types):
            return category
    return None


def _safe_enum_check(value, enum_set):
    """
    Safely check if a value is in an enum set, preventing TypeError on invalid types.

    This handles the Python 3.11 vs 3.12 difference where boolean values
    used with 'in' operator against Enum types can cause TypeError in 3.11.
    """
    if not isinstance(value, PP_PLACEHOLDER_TYPE):
        return False
    try:
        return value in enum_set
    except (TypeError, AttributeError):
        return False


def is_title_placeholder(placeholder_type):
    """Check if placeholder type is for titles."""
    return _safe_enum_check(placeholder_type, TITLE_PLACEHOLDERS)


def is_subtitle_placeholder(placeholder_type):
    """Check if placeholder type is for subtitles."""
    return _safe_enum_check(placeholder_type, SUBTITLE_PLACEHOLDERS)


def is_content_placeholder(placeholder_type):
    """Check if placeholder type is for main content."""
    return _safe_enum_check(placeholder_type, CONTENT_PLACEHOLDERS)


def is_media_placeholder(placeholder_type):
    """Check if placeholder type is for media/objects."""
    return _safe_enum_check(placeholder_type, MEDIA_PLACEHOLDERS)
