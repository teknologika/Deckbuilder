#!/usr/bin/env python3
"""
Deckbuilder Formatting Support

Comprehensive language and font formatting support for PowerPoint presentations.
Provides language ID constants, font validation, and unified presentation processing
that updates both master slides and content slides.
"""

import os
from difflib import get_close_matches
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.enum.lang import MSO_LANGUAGE_ID


class FormattingSupport:
    """Comprehensive formatting support for language and font settings"""

    # Language ID constants mapped to common locale codes
    # Using only verified constants from python-pptx MSO_LANGUAGE_ID
    LANGUAGE_IDS = {
        "en-US": MSO_LANGUAGE_ID.ENGLISH_US,
        "en-AU": MSO_LANGUAGE_ID.ENGLISH_AUS,
        "en-GB": MSO_LANGUAGE_ID.ENGLISH_UK,
        "en-CA": MSO_LANGUAGE_ID.ENGLISH_CANADIAN,
        "en-NZ": MSO_LANGUAGE_ID.ENGLISH_NEW_ZEALAND,
        "en-ZA": MSO_LANGUAGE_ID.ENGLISH_SOUTH_AFRICA,
        "es-ES": MSO_LANGUAGE_ID.SPANISH,
        "fr-FR": MSO_LANGUAGE_ID.FRENCH,
        "fr-CA": MSO_LANGUAGE_ID.FRENCH_CANADIAN,
        "de-DE": MSO_LANGUAGE_ID.GERMAN,
        "it-IT": MSO_LANGUAGE_ID.ITALIAN,
        "pt-PT": MSO_LANGUAGE_ID.PORTUGUESE,
        "pt-BR": MSO_LANGUAGE_ID.BRAZILIAN_PORTUGUESE,
        "nl-NL": MSO_LANGUAGE_ID.DUTCH,
        "sv-SE": MSO_LANGUAGE_ID.SWEDISH,
        "da-DK": MSO_LANGUAGE_ID.DANISH,
        "fi-FI": MSO_LANGUAGE_ID.FINNISH,
        "ru-RU": MSO_LANGUAGE_ID.RUSSIAN,
        "ja-JP": MSO_LANGUAGE_ID.JAPANESE,
        "ko-KR": MSO_LANGUAGE_ID.KOREAN,
    }

    # Common system fonts for validation
    COMMON_FONTS = [
        "Arial",
        "Calibri",
        "Segoe UI",
        "Times New Roman",
        "Verdana",
        "Tahoma",
        "Georgia",
        "Comic Sans MS",
        "Trebuchet MS",
        "Impact",
        "Palatino Linotype",
        "Book Antiqua",
        "Lucida Console",
        "Courier New",
        "Consolas",
        "Cambria",
        "Candara",
        "Corbel",
        "Franklin Gothic Medium",
        "Century Gothic",
        "Arial Black",
    ]

    def __init__(self):
        """Initialize formatting support"""
        pass

    # Language name to code mapping for backward compatibility
    LANGUAGE_NAME_TO_CODE = {
        "English (United States)": "en-US",
        "English (Australia)": "en-AU",
        "English (United Kingdom)": "en-GB",
        "English (Canada)": "en-CA",
        "English (New Zealand)": "en-NZ",
        "English (South Africa)": "en-ZA",
        "Spanish (Spain)": "es-ES",
        "French (France)": "fr-FR",
        "French (Canada)": "fr-CA",
        "German (Germany)": "de-DE",
        "Italian (Italy)": "it-IT",
        "Portuguese (Portugal)": "pt-PT",
        "Portuguese (Brazil)": "pt-BR",
        "Dutch (Netherlands)": "nl-NL",
        "Swedish (Sweden)": "sv-SE",
        "Danish (Denmark)": "da-DK",
        "Finnish (Finland)": "fi-FI",
        "Russian (Russia)": "ru-RU",
        "Japanese (Japan)": "ja-JP",
        "Korean (South Korea)": "ko-KR",
    }

    @classmethod
    def get_supported_languages(cls) -> Dict[str, str]:
        """
        Get dictionary of supported language codes with descriptions.

        Returns:
            Dict mapping language codes to human-readable descriptions
        """
        descriptions = {
            "en-US": "English (United States)",
            "en-AU": "English (Australia)",
            "en-GB": "English (United Kingdom)",
            "en-CA": "English (Canada)",
            "en-NZ": "English (New Zealand)",
            "en-ZA": "English (South Africa)",
            "es-ES": "Spanish (Spain)",
            "fr-FR": "French (France)",
            "fr-CA": "French (Canada)",
            "de-DE": "German (Germany)",
            "it-IT": "Italian (Italy)",
            "pt-PT": "Portuguese (Portugal)",
            "pt-BR": "Portuguese (Brazil)",
            "nl-NL": "Dutch (Netherlands)",
            "sv-SE": "Swedish (Sweden)",
            "da-DK": "Danish (Denmark)",
            "fi-FI": "Finnish (Finland)",
            "ru-RU": "Russian (Russia)",
            "ja-JP": "Japanese (Japan)",
            "ko-KR": "Korean (South Korea)",
        }
        return descriptions

    @classmethod
    def normalize_language_input(cls, language_input: str) -> Optional[str]:
        """
        Normalize language input to standard locale code format.

        Accepts both locale codes (en-AU) and full names (English (Australia)).

        Args:
            language_input: Language code or name

        Returns:
            Normalized locale code or None if not found
        """
        # Try direct code lookup first
        if language_input in cls.LANGUAGE_IDS:
            return language_input

        # Try full name lookup
        if language_input in cls.LANGUAGE_NAME_TO_CODE:
            return cls.LANGUAGE_NAME_TO_CODE[language_input]

        return None

    @classmethod
    def validate_language_code(cls, language_code: str) -> Tuple[bool, Optional[str], List[str]]:
        """
        Validate language code and provide suggestions for invalid codes.

        Accepts both locale codes (en-AU) and full names (English (Australia)).

        Args:
            language_code: Language code or name to validate

        Returns:
            Tuple of (is_valid, error_message, suggestions)
        """
        # Try to normalize the input
        normalized = cls.normalize_language_input(language_code)
        if normalized:
            return True, None, []

        # Generate suggestions for invalid codes from both formats
        all_valid_inputs = list(cls.LANGUAGE_IDS.keys()) + list(cls.LANGUAGE_NAME_TO_CODE.keys())
        suggestions = get_close_matches(language_code, all_valid_inputs, n=3, cutoff=0.6)

        error_msg = f"Unsupported language: '{language_code}'"
        return False, error_msg, suggestions

    @classmethod
    def validate_font_name(cls, font_name: str) -> Tuple[bool, Optional[str], List[str]]:
        """
        Validate font name and provide suggestions for common fonts.

        Args:
            font_name: Font name to validate

        Returns:
            Tuple of (is_valid, warning_message, suggestions)
        """
        # Always accept any font name (user might have custom fonts)
        # But provide suggestions if it's not a common font
        if font_name in cls.COMMON_FONTS:
            return True, None, []

        # Find similar common fonts
        suggestions = get_close_matches(font_name, cls.COMMON_FONTS, n=3, cutoff=0.6)

        warning_msg = f"Font '{font_name}' is not a common system font"
        return True, warning_msg, suggestions

    def apply_language_to_run(self, run, language_input: str) -> bool:
        """
        Apply language setting to a text run.

        Accepts both locale codes (en-AU) and full names (English (Australia)).

        Args:
            run: python-pptx Run object
            language_input: Language code or name (e.g., 'en-AU' or 'English (Australia)')

        Returns:
            True if language was applied successfully
        """
        try:
            # Normalize input to locale code
            language_code = self.normalize_language_input(language_input)
            if language_code and language_code in self.LANGUAGE_IDS:
                run.font.language_id = self.LANGUAGE_IDS[language_code]
                return True
            return False
        except Exception:
            return False

    def apply_font_to_run(self, run, font_name: str) -> bool:
        """
        Apply font family to a text run.

        Args:
            run: python-pptx Run object
            font_name: Font family name

        Returns:
            True if font was applied successfully
        """
        try:
            run.font.name = font_name
            return True
        except Exception:
            return False

    def process_text_frame(
        self, text_frame, language_code: Optional[str] = None, font_name: Optional[str] = None
    ) -> Dict[str, int]:
        """
        Process all text runs in a text frame, applying language and/or font settings.

        Args:
            text_frame: python-pptx TextFrame object
            language_code: Optional language code to apply
            font_name: Optional font name to apply

        Returns:
            Dictionary with processing statistics
        """
        stats = {"runs_processed": 0, "language_applied": 0, "font_applied": 0}

        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    stats["runs_processed"] += 1

                    if language_code:
                        if self.apply_language_to_run(run, language_code):
                            stats["language_applied"] += 1

                    if font_name:
                        if self.apply_font_to_run(run, font_name):
                            stats["font_applied"] += 1

        except Exception as e:
            # Log error but continue processing
            print(f"⚠️  Warning: Error processing text frame: {e}")

        return stats

    def process_shape(
        self, shape, language_code: Optional[str] = None, font_name: Optional[str] = None
    ) -> Dict[str, int]:
        """
        Process a shape, applying formatting to its text content.

        Args:
            shape: python-pptx Shape object
            language_code: Optional language code to apply
            font_name: Optional font name to apply

        Returns:
            Dictionary with processing statistics
        """
        stats = {"runs_processed": 0, "language_applied": 0, "font_applied": 0}

        try:
            # Handle text frames
            if hasattr(shape, "text_frame") and shape.text_frame:
                frame_stats = self.process_text_frame(shape.text_frame, language_code, font_name)
                for key in stats:
                    stats[key] += frame_stats[key]

            # Handle tables
            elif hasattr(shape, "table"):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            cell_stats = self.process_text_frame(
                                cell.text_frame, language_code, font_name
                            )
                            for key in stats:
                                stats[key] += cell_stats[key]

            # Handle grouped shapes
            elif hasattr(shape, "shapes"):
                for sub_shape in shape.shapes:
                    sub_stats = self.process_shape(sub_shape, language_code, font_name)
                    for key in stats:
                        stats[key] += sub_stats[key]

        except Exception as e:
            print(f"⚠️  Warning: Error processing shape: {e}")

        return stats

    def update_presentation(
        self,
        presentation_path: str,
        language_code: Optional[str] = None,
        font_name: Optional[str] = None,
        output_path: Optional[str] = None,
        create_backup: bool = True,
    ) -> Dict[str, any]:
        """
        Update both master slides and content slides in a PowerPoint presentation.

        Args:
            presentation_path: Path to the PowerPoint file
            language_code: Optional language code to apply
            font_name: Optional font name to apply
            output_path: Optional output path (default: update in place)
            create_backup: Whether to create a backup file

        Returns:
            Dictionary with operation results and statistics
        """
        pptx_path = Path(presentation_path)

        if not pptx_path.exists():
            return {"success": False, "error": f"File not found: {presentation_path}", "stats": {}}

        # Create backup if requested
        backup_path = None
        if create_backup:
            backup_path = pptx_path.with_suffix(".bak.pptx")
            try:
                import shutil

                shutil.copy2(pptx_path, backup_path)
            except Exception as e:
                return {"success": False, "error": f"Failed to create backup: {e}", "stats": {}}

        try:
            # Load presentation
            prs = Presentation(str(pptx_path))

            total_stats = {
                "master_slides_processed": 0,
                "content_slides_processed": 0,
                "total_runs_processed": 0,
                "total_language_applied": 0,
                "total_font_applied": 0,
            }

            # Process master slides
            try:
                for slide_master in prs.slide_masters:
                    total_stats["master_slides_processed"] += 1

                    # Process master slide shapes
                    for shape in slide_master.shapes:
                        shape_stats = self.process_shape(shape, language_code, font_name)
                        total_stats["total_runs_processed"] += shape_stats["runs_processed"]
                        total_stats["total_language_applied"] += shape_stats["language_applied"]
                        total_stats["total_font_applied"] += shape_stats["font_applied"]

                    # Process slide layouts
                    for slide_layout in slide_master.slide_layouts:
                        for shape in slide_layout.shapes:
                            shape_stats = self.process_shape(shape, language_code, font_name)
                            total_stats["total_runs_processed"] += shape_stats["runs_processed"]
                            total_stats["total_language_applied"] += shape_stats["language_applied"]
                            total_stats["total_font_applied"] += shape_stats["font_applied"]

            except Exception as e:
                print(f"⚠️  Warning: Error processing master slides: {e}")

            # Process content slides
            for slide in prs.slides:
                total_stats["content_slides_processed"] += 1

                for shape in slide.shapes:
                    shape_stats = self.process_shape(shape, language_code, font_name)
                    total_stats["total_runs_processed"] += shape_stats["runs_processed"]
                    total_stats["total_language_applied"] += shape_stats["language_applied"]
                    total_stats["total_font_applied"] += shape_stats["font_applied"]

            # Save presentation
            save_path = output_path if output_path else presentation_path
            prs.save(save_path)

            return {
                "success": True,
                "message": f"Presentation updated successfully: {save_path}",
                "backup_path": str(backup_path) if backup_path else None,
                "stats": total_stats,
            }

        except Exception as e:
            return {"success": False, "error": f"Failed to process presentation: {e}", "stats": {}}


def get_default_language() -> Optional[str]:
    """Get default language from environment variable"""
    return os.getenv("DECK_PROOFING_LANGUAGE")


def get_default_font() -> Optional[str]:
    """Get default font from environment variable"""
    return os.getenv("DECK_DEFAULT_FONT")


def print_supported_languages():
    """Print all supported languages in a user-friendly format"""
    languages = FormattingSupport.get_supported_languages()

    print("📋 Supported Proofing Languages:")
    for code, description in sorted(languages.items()):
        print(f"  {code:<6} {description}")
    print()
    print("💡 Usage examples:")
    print("  deckbuilder -l en-AU create presentation.md")
    print("  deckbuilder remap template.pptx --language en-US")
    print("  export DECK_PROOFING_LANGUAGE=en-GB")
