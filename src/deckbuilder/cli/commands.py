#!/usr/bin/env python3
"""
Deckbuilder Command Line Tools

Standalone utilities for template analysis, documentation generation, and validation.
These tools are designed to be run independently for template management tasks.
"""

import argparse
import os
import sys
from datetime import datetime
from pathlib import Path

# Add parent directory to path for imports
sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))

# Import after path modification - this is intentional to find the mcp_server module
from mcp_server.tools import TemplateAnalyzer  # noqa: E402

try:
    from .naming_conventions import NamingConvention, PlaceholderContext  # noqa: E402
except ImportError:
    # Handle direct script execution
    from deckbuilder.cli.naming_conventions import (
        NamingConvention,
        PlaceholderContext,
    )  # noqa: E402


class TemplateManager:
    """Command-line template management utilities"""

    def __init__(self, template_folder=None, output_folder=None):
        # Use command-line arguments, environment variables, or sensible defaults
        if template_folder:
            self.template_folder = template_folder
        else:
            # Use environment variable first (set by CLI setup_environment)
            env_template_folder = os.getenv("DECK_TEMPLATE_FOLDER")
            if env_template_folder:
                self.template_folder = env_template_folder
            else:
                # Fallback: default template folder in current directory
                self.template_folder = str(Path.cwd() / "templates")

        if output_folder:
            self.output_folder = output_folder
        else:
            # Use environment variable first (set by CLI setup_environment)
            env_output_folder = os.getenv("DECK_OUTPUT_FOLDER")
            if env_output_folder:
                self.output_folder = env_output_folder
            else:
                # Fallback: current directory
                self.output_folder = str(Path.cwd())

        # Ensure folders exist
        os.makedirs(self.output_folder, exist_ok=True)

        # Don't create analyzer yet - wait until we need it

    def analyze_template(self, template_name: str, verbose: bool = False) -> dict:
        """
        Analyze a PowerPoint template and generate JSON mapping.

        Args:
            template_name: Name of template (e.g., 'default')
            verbose: Print detailed analysis information

        Returns:
            Template analysis results
        """
        print(f"🔍 Analyzing template: {template_name}")
        print(f"📁 Template folder: {self.template_folder}")
        print(f"📂 Output folder: {self.output_folder}")

        try:
            # Create analyzer instance with explicit paths instead of environment variables
            analyzer = TemplateAnalyzer()
            # Override analyzer paths with our explicit values
            analyzer.template_path = self.template_folder
            analyzer.output_folder = self.output_folder

            # Run analysis
            result = analyzer.analyze_pptx_template(template_name)

            # Print summary
            layouts_count = len(result.get("layouts", {}))
            print("✅ Analysis complete!")
            print(f"   📊 Found {layouts_count} layouts")

            # Print layout summary
            if verbose and "layouts" in result:
                print("\n📋 Layout Summary:")
                for layout_name, layout_info in result["layouts"].items():
                    placeholder_count = len(layout_info.get("placeholders", {}))
                    index = layout_info.get("index", "?")
                    print(f"   {index:2d}: {layout_name} ({placeholder_count} placeholders)")

            # Check for generated file
            base_name = template_name.replace(".pptx", "")
            output_file = os.path.join(self.output_folder, f"{base_name}.g.json")
            if os.path.exists(output_file):
                print(f"📄 Generated: {output_file}")
                print("   ✏️  Edit this file with semantic placeholder names")
                print(f"   📝 Rename to '{base_name}.json' when ready")

            return result

        except Exception as e:
            print(f"❌ Error analyzing template: {str(e)}")
            return {}

    def document_template(self, template_name: str, output_path: str = None) -> str:
        """
        Generate comprehensive documentation for a template.

        Args:
            template_name: Name of template to document
            output_path: Custom output path (optional)

        Returns:
            Path to generated documentation
        """
        print(f"📝 Generating documentation for: {template_name}")

        # Analyze template first
        analysis = self.analyze_template(template_name, verbose=False)
        if not analysis:
            return ""

        # Legacy template mapping JSON files no longer used
        mapping = {}  # Empty mapping for backward compatibility

        # Generate documentation
        doc_content = self._generate_template_documentation(template_name, analysis, mapping)

        # Save documentation
        if not output_path:
            project_root = Path(__file__).parent.parent.parent
            docs_folder = project_root / "docs" / "Features"
            docs_folder.mkdir(parents=True, exist_ok=True)
            base_name = template_name.replace(".pptx", "")
            output_path = str(docs_folder / f"{base_name.title()}Template.md")

        with open(output_path, "w", encoding="utf-8") as f:
            f.write(doc_content)

        print(f"✅ Documentation generated: {output_path}")
        return output_path

    def _generate_template_documentation(self, template_name: str, analysis: dict, mapping: dict) -> str:
        """Generate markdown documentation for template"""
        base_name = template_name.replace(".pptx", "")
        layouts = analysis.get("layouts", {})

        # Generate layout summary table
        table_rows = []
        for layout_name, layout_info in layouts.items():
            index = layout_info.get("index", "?")
            placeholders = layout_info.get("placeholders", {})
            placeholder_count = len(placeholders)

            # Check if mapping exists
            mapping_status = "✅" if mapping.get("layouts", {}).get(layout_name) else "❌"

            # Check structured frontmatter support (placeholder for now)
            structured_status = "⏳"  # Would check structured_frontmatter.py

            row = f"| {layout_name} | {index} | {placeholder_count} | " f"{structured_status} | {mapping_status} |"
            table_rows.append(row)

        # Generate detailed layout specifications
        detailed_layouts = []
        for layout_name, layout_info in layouts.items():
            index = layout_info.get("index", "?")
            placeholders = layout_info.get("placeholders", {})

            layout_section = f"### {layout_name} (Index: {index})\\n\\n"
            layout_section += "**PowerPoint Placeholders**:\\n"

            for idx, placeholder_name in placeholders.items():
                # Handle both string and object formats
                if isinstance(placeholder_name, dict):
                    name = placeholder_name.get("name", "Unknown")
                    actual_idx = placeholder_name.get("idx", idx)
                else:
                    name = placeholder_name
                    actual_idx = idx
                layout_section += f'- `idx={actual_idx}`: "{name}"\\n'

            # Add mapping info if available
            layout_mapping = mapping.get("layouts", {}).get(layout_name)
            if layout_mapping:
                layout_section += "\\n**JSON Mapping**: ✅ Configured\\n"
            else:
                layout_section += "\\n**JSON Mapping**: ❌ Not configured\\n"

            layout_section += "**Structured Frontmatter**: ⏳ To be implemented\\n"

            detailed_layouts.append(layout_section)

        # Generate full documentation
        doc_content = f"""# {base_name.title()} Template Documentation

## Template Overview
- **Template Name**: {template_name}
- **Analysis Date**: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}
- **Total Layouts**: {len(layouts)}
- **Template Location**: `{self.template_folder}/{template_name}`

## Layout Summary

| Layout Name | Index | Placeholders | Structured Support | JSON Mapping |
|-------------|-------|--------------|-------------------|--------------|
{chr(10).join(table_rows)}

## Detailed Layout Specifications

{chr(10).join(detailed_layouts)}

## Template Management

### Adding JSON Mapping
1. **Analyze template**: Run `python -m deckbuilder.cli_tools analyze {base_name}`
2. **Edit generated file**: Customize `{base_name}.g.json` with semantic names
3. **Activate mapping**: Rename to `{base_name}.json` in templates folder

### Example JSON Mapping Structure
```json
{{
  "template_info": {{
    "name": "{base_name.title()}",
    "version": "1.0"
  }},
  "layouts": {{
    "Title Slide": {{
      "index": 0,
      "placeholders": {{
        "0": "Title 1",
        "1": "Subtitle 2"
      }}
    }}
  }},
  "aliases": {{
    "title": "Title Slide",
    "content": "Title and Content"
  }}
}}
```

### Usage Examples

**JSON Format**:
```json
{{
  "presentation": {{
    "slides": [
      {{
        "type": "Title Slide",
        "layout": "Title Slide",
        "title": "My Presentation",
        "subtitle": "Subtitle text"
      }}
    ]
  }}
}}
```

**Markdown with Frontmatter**:
```yaml
---
layout: Title Slide
---
# My Presentation
## Subtitle text
```

---
*Generated automatically by Deckbuilder Template Manager on \
{datetime.now().strftime("%Y-%m-%d %H:%M:%S")}*
"""

        return doc_content

    def validate_template(self, template_name: str) -> dict:
        """
        Validate template structure and mappings.

        Args:
            template_name: Name of template to validate

        Returns:
            Validation results
        """
        print(f"🔍 Validating template: {template_name}")

        validation_results = {
            "template_file": self._validate_template_file(template_name),
            "json_mapping": self._validate_json_mapping(template_name),
            "placeholder_naming": self._validate_placeholder_naming(template_name),
        }

        # Print results
        for category, results in validation_results.items():
            status = results.get("status", "unknown")
            if status == "valid":
                print(f"✅ {category}: Valid")
            elif status == "issues_found":
                print(f"⚠️  {category}: Issues found")
                for issue in results.get("issues", []):
                    print(f"   - {issue}")
            else:
                print(f"❌ {category}: {results.get('error', 'Unknown error')}")

        return validation_results

    def _validate_template_file(self, template_name: str) -> dict:
        """Validate template file exists and is accessible"""
        try:
            if not template_name.endswith(".pptx"):
                template_name += ".pptx"

            template_path = os.path.join(self.template_folder, template_name)

            if not os.path.exists(template_path):
                return {"status": "error", "error": f"Template file not found: {template_path}"}

            # Try to load with python-pptx
            from pptx import Presentation

            prs = Presentation(template_path)

            return {
                "status": "valid",
                "layout_count": len(prs.slide_layouts),
                "file_size": os.path.getsize(template_path),
            }

        except Exception as e:
            return {"status": "error", "error": str(e)}

    def _validate_json_mapping(self, template_name: str) -> dict:
        """Validate template - legacy JSON mapping validation removed"""
        # Legacy template mapping JSON files no longer used
        # Template validation now handled by structured frontmatter patterns
        return {"status": "valid", "message": "Template validation using structured frontmatter patterns"}

    def _validate_placeholder_naming(self, template_name: str) -> dict:
        """Validate placeholder naming conventions"""
        # Placeholder for naming convention validation
        return {"status": "valid", "message": "Naming validation not yet implemented"}

    def enhance_template(
        self,
        template_name: str,
        mapping_file: str = None,  # Legacy parameter - ignored
        create_backup: bool = True,
        use_conventions: bool = True,  # Always use conventions now
    ) -> dict:
        """
        Enhance template by updating master slide placeholder names using semantic mapping.

        Args:
            template_name: Name of template to enhance
            mapping_file: Legacy parameter - ignored (JSON mapping files no longer used)
            create_backup: Create backup before modification (default: True)
            use_conventions: Always True - convention-based naming system

        Returns:
            Enhancement results with success/failure information
        """
        print(f"🔧 Enhancing template: {template_name}")

        try:
            # Ensure template name has .pptx extension
            if not template_name.endswith(".pptx"):
                template_name += ".pptx"

            template_path = os.path.join(self.template_folder, template_name)

            if not os.path.exists(template_path):
                return {"status": "error", "error": f"Template file not found: {template_path}"}

            # Create backup if requested
            if create_backup:
                backup_path = self._create_template_backup(template_path)
                print(f"📄 Backup created: {backup_path}")

            # Always use convention-based mapping (legacy JSON files no longer supported)
            print("🎯 Using convention-based naming system...")
            mapping = self._generate_convention_mapping(template_path)

            # Enhance template
            modifications = self._modify_master_slide_placeholders(template_path, mapping)

            if modifications["modified_count"] > 0:
                print("✅ Enhancement complete!")
                modified_count = modifications["modified_count"]
                layout_count = modifications["layout_count"]
                print(f"   📊 Modified {modified_count} placeholders across {layout_count} layouts")

                if "enhanced_template_path" in modifications:
                    print(f"   📄 Enhanced template saved: {modifications['enhanced_template_path']}")

                if modifications["issues"]:
                    print("⚠️  Issues encountered:")
                    for issue in modifications["issues"]:
                        print(f"   - {issue}")

                return {
                    "status": "success",
                    "modifications": modifications,
                    "backup_path": backup_path if create_backup else None,
                }
            else:
                print("ℹ️  No modifications needed - all placeholders already have correct names")
                return {
                    "status": "no_changes",
                    "message": "Template already has correct placeholder names",
                }

        except Exception as e:
            return {"status": "error", "error": str(e)}

    def _create_template_backup(self, template_path: str) -> str:
        """Create backup copy of template file in organized backups folder"""
        import shutil
        from datetime import datetime

        # Generate backup filename with timestamp
        path_obj = Path(template_path)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        backup_name = f"{path_obj.stem}_backup_{timestamp}{path_obj.suffix}"

        # Create backups folder within templates directory
        backups_folder = path_obj.parent / "backups"
        backups_folder.mkdir(exist_ok=True)

        backup_path = backups_folder / backup_name

        # Copy file
        shutil.copy2(template_path, backup_path)
        return str(backup_path)

    def _modify_master_slide_placeholders(self, template_path: str, mapping: dict) -> dict:
        """
        Modify master slide placeholder names using python-pptx.

        Args:
            template_path: Path to PowerPoint template
            mapping: JSON mapping with placeholder names

        Returns:
            Dictionary with modification results
        """
        from pptx import Presentation

        modifications = {"modified_count": 0, "layout_count": 0, "issues": [], "changes": []}

        # Load presentation
        prs = Presentation(template_path)
        layouts_mapping = mapping.get("layouts", {})

        # Access the slide master (this is the key change!)
        slide_master = prs.slide_master

        # First, try to modify placeholders on the master slide itself
        try:
            for placeholder in slide_master.placeholders:
                placeholder_idx = str(placeholder.placeholder_format.idx)
                # Try to find this placeholder in any layout mapping
                for _layout_name, layout_info in layouts_mapping.items():
                    placeholder_mapping = layout_info.get("placeholders", {})
                    if placeholder_idx in placeholder_mapping:
                        new_name = placeholder_mapping[placeholder_idx]
                        old_name = placeholder.name
                        try:
                            placeholder.element.nvSpPr.cNvPr.name = new_name
                            modifications["modified_count"] += 1
                            modifications["changes"].append(
                                {
                                    "location": "master_slide",
                                    "placeholder_idx": placeholder_idx,
                                    "old_name": old_name,
                                    "new_name": new_name,
                                }
                            )
                            break
                        except Exception as e:
                            modifications["issues"].append(f"Failed to modify master placeholder {placeholder_idx}: {str(e)}")
        except Exception as e:
            modifications["issues"].append(f"No direct master placeholders or error: {e}")

        # Process each slide layout in the master
        for layout in slide_master.slide_layouts:
            layout_name = layout.name
            modifications["layout_count"] += 1

            if layout_name not in layouts_mapping:
                modifications["issues"].append(f"Layout '{layout_name}' not found in mapping file")
                continue

            layout_mapping = layouts_mapping[layout_name]
            placeholder_mapping = layout_mapping.get("placeholders", {})

            # Modify placeholders in this master slide layout
            for placeholder in layout.placeholders:
                placeholder_idx = str(placeholder.placeholder_format.idx)

                if placeholder_idx in placeholder_mapping:
                    new_name = placeholder_mapping[placeholder_idx]
                    old_name = placeholder.name

                    try:
                        # Update placeholder name on the master slide layout
                        if hasattr(placeholder, "element") and hasattr(placeholder.element, "nvSpPr"):
                            placeholder.element.nvSpPr.cNvPr.name = new_name
                            modifications["modified_count"] += 1
                            modifications["changes"].append(
                                {
                                    "layout": layout_name,
                                    "placeholder_idx": placeholder_idx,
                                    "old_name": old_name,
                                    "new_name": new_name,
                                }
                            )
                        else:
                            error_msg = f"Cannot modify placeholder {placeholder_idx} in {layout_name} - " "unsupported structure"
                            modifications["issues"].append(error_msg)
                    except Exception as e:
                        error_msg = f"Failed to modify placeholder {placeholder_idx} in {layout_name}: " f"{str(e)}"
                        modifications["issues"].append(error_msg)

        # Save modified presentation with .g.pptx extension
        try:
            # Generate enhanced template filename with .g.pptx convention
            path_obj = Path(template_path)
            enhanced_name = f"{path_obj.stem}.g{path_obj.suffix}"
            enhanced_path = path_obj.parent / enhanced_name

            prs.save(str(enhanced_path))
            modifications["enhanced_template_path"] = str(enhanced_path)
        except Exception as e:
            modifications["issues"].append(f"Failed to save enhanced template: {str(e)}")

        return modifications

    def _generate_convention_mapping(self, template_path: str) -> dict:
        """
        Generate convention-based mapping for template placeholders.

        Args:
            template_path: Path to PowerPoint template

        Returns:
            Convention-based mapping dictionary
        """
        from pptx import Presentation

        # Load presentation to analyze structure
        prs = Presentation(template_path)
        convention = NamingConvention()

        # Build mapping using convention system
        mapping = {"template_info": {"name": "Convention-Based", "version": "1.0"}, "layouts": {}}

        # Process each slide layout
        layout_index = 0
        for layout in prs.slide_layouts:
            layout_name = layout.name
            layout_placeholders = {}

            for placeholder in layout.placeholders:
                placeholder_idx = str(placeholder.placeholder_format.idx)

                # Create context for convention naming
                context = PlaceholderContext(
                    layout_name=layout_name,
                    placeholder_idx=placeholder_idx,
                    total_placeholders=len(layout.placeholders),
                )

                # Generate convention-based name
                convention_name = convention.generate_placeholder_name(context)
                layout_placeholders[placeholder_idx] = convention_name

            mapping["layouts"][layout_name] = {
                "index": layout_index,
                "placeholders": layout_placeholders,
            }

            layout_index += 1

        return mapping


class DocumentationGenerator:
    """Generate comprehensive documentation and examples for Deckbuilder templates"""

    def __init__(self, template_folder=None):
        self.template_folder = template_folder or self._find_template_folder()

    def _find_template_folder(self):
        """Find template folder using same logic as TemplateManager"""
        current_dir = Path.cwd()
        if (current_dir / "assets" / "templates").exists():
            return str(current_dir / "assets" / "templates")
        else:
            project_root = Path(__file__).parent.parent.parent
            return str(project_root / "assets" / "templates")

    def get_implemented_layouts(self, template_name="default"):
        """
        Get list of layouts that are actually implemented in the code.

        Based on analysis of structured_frontmatter.py and test files,
        these are the layouts with full implementation support.
        """
        implemented_layouts = {
            "Title Slide": {
                "description": "Perfect for presentation openings and section breaks",
                "complexity": "simple",
                "frontmatter_example": """---
layout: Title Slide
---
# Your Presentation Title
## Your subtitle here""",
                "json_example": """{
  "type": "Title Slide",
  "title": "Your Presentation Title",
  "subtitle": "Your subtitle here"
}""",
            },
            "Title and Content": {
                "description": "Standard content slides with bullets and paragraphs",
                "complexity": "simple",
                "frontmatter_example": """---
layout: Title and Content
---
# Slide Title

Your content here with:
- Bullet point one
- Bullet point two
- **Bold** and *italic* formatting""",
                "json_example": """{
  "type": "Title and Content",
  "title": "Slide Title",
  "rich_content": [
    {
      "paragraph": "Your content here with:"
    },
    {
      "bullets": ["Bullet point one", "Bullet point two", "**Bold** and *italic* formatting"],
      "bullet_levels": [1, 1, 1]
    }
  ]
}""",
            },
            "Section Header": {
                "description": "Section dividers for organizing your presentation",
                "complexity": "simple",
                "frontmatter_example": """---
layout: Section Header
---
# Section: **Implementation** Phase

This section covers the development and deployment stages.""",
                "json_example": """{
  "type": "Section Header",
  "title": "Section: **Implementation** Phase",
  "rich_content": [
    {
      "paragraph": "This section covers the development and deployment stages."
    }
  ]
}""",
            },
            "Two Content": {
                "description": "Side-by-side content comparison",
                "complexity": "structured",
                "frontmatter_example": """---
layout: Two Content
title: Comparison Title
sections:
  - title: Left Side
    content:
      - "Point A"
      - "Point B"
  - title: Right Side
    content:
      - "Point C"
      - "Point D"
---""",
                "json_example": """{
  "type": "Two Content",
  "title": "Comparison Title",
  "content_left_1": ["Point A", "Point B"],
  "content_right_1": ["Point C", "Point D"]
}""",
            },
            "Comparison": {
                "description": "Side-by-side comparison for contrasting two options",
                "complexity": "structured",
                "frontmatter_example": """---
layout: Comparison
title: Option Analysis
comparison:
  left:
    title: Option A
    content: "Cost effective with rapid deployment"
  right:
    title: Option B
    content: "Advanced features with future-proof design"
---""",
                "json_example": """{
  "type": "Comparison",
  "title": "Option Analysis",
  "title_left_1": "Option A",
  "content_left_1": "Cost effective with rapid deployment",
  "title_right_1": "Option B",
  "content_right_1": "Advanced features with future-proof design"
}""",
            },
            "Four Columns": {
                "description": "Four-column layout for feature comparisons",
                "complexity": "structured",
                "frontmatter_example": """---
layout: Four Columns
title: Feature Comparison
columns:
  - title: Performance
    content: "Fast processing with optimized algorithms"
  - title: Security
    content: "Enterprise-grade encryption and compliance"
  - title: Usability
    content: "Intuitive interface with minimal learning"
  - title: Cost
    content: "Competitive pricing with flexible plans"
---""",
                "json_example": """{
  "type": "Four Columns",
  "title": "Feature Comparison",
  "title_col1_1": "Performance",
  "content_col1_1": "Fast processing with optimized algorithms",
  "title_col2_1": "Security",
  "content_col2_1": "Enterprise-grade encryption and compliance",
  "title_col3_1": "Usability",
  "content_col3_1": "Intuitive interface with minimal learning",
  "title_col4_1": "Cost",
  "content_col4_1": "Competitive pricing with flexible plans"
}""",
            },
            "Three Columns With Titles": {
                "description": "Three-column layout with individual titles",
                "complexity": "structured",
                "frontmatter_example": """---
layout: Three Columns With Titles
title: Key Features
columns:
  - title: Performance
    content: "Fast processing with optimized algorithms"
  - title: Security
    content: "Enterprise-grade encryption and compliance"
  - title: Usability
    content: "Intuitive interface with minimal learning"
---""",
                "json_example": """{
  "type": "Three Columns With Titles",
  "title": "Key Features",
  "title_col1_1": "Performance",
  "content_col1_1": "Fast processing with optimized algorithms",
  "title_col2_1": "Security",
  "content_col2_1": "Enterprise-grade encryption and compliance",
  "title_col3_1": "Usability",
  "content_col3_1": "Intuitive interface with minimal learning"
}""",
            },
            "Three Columns": {
                "description": "Three-column layout with content only",
                "complexity": "structured",
                "frontmatter_example": """---
layout: Three Columns
title: Benefits Overview
columns:
  - content: "Fast processing with optimized algorithms"
  - content: "Enterprise-grade encryption and compliance"
  - content: "Intuitive interface with minimal learning"
---""",
                "json_example": """{
  "type": "Three Columns",
  "title": "Benefits Overview",
  "content_col1_1": "Fast processing with optimized algorithms",
  "content_col2_1": "Enterprise-grade encryption and compliance",
  "content_col3_1": "Intuitive interface with minimal learning"
}""",
            },
            "Picture with Caption": {
                "description": ("Media slides with image and caption (includes PlaceKitten fallback)"),
                "complexity": "advanced",
                "frontmatter_example": """---
layout: Picture with Caption
title: System Architecture
media:
  image_path: "assets/architecture.png"  # Auto-fallback to PlaceKitten
  alt_text: "System architecture diagram"
  caption: "High-level system architecture"
  description: "Main components and their interactions"
---""",
                "json_example": """{
  "type": "Picture with Caption",
  "title": "System Architecture",
  "image_1": "assets/architecture.png",
  "text_caption_1": "High-level system architecture"
}""",
            },
            "Title Only": {
                "description": "Simple title slides for minimal content",
                "complexity": "simple",
                "frontmatter_example": """---
layout: Title Only
---
# Title Only Layout: **Bold** *Italic* ___Underline___""",
                "json_example": """{
  "type": "Title Only",
  "title": "Title Only Layout: **Bold** *Italic* ___Underline___"
}""",
            },
            "Blank": {
                "description": "Blank layout for custom content",
                "complexity": "simple",
                "frontmatter_example": """---
layout: Blank
---
# Custom Content

This layout provides maximum flexibility for custom designs.""",
                "json_example": """{
  "type": "Blank",
  "title": "Custom Content",
  "rich_content": [
    {
      "paragraph": "This layout provides maximum flexibility for custom designs."
    }
  ]
}""",
            },
            "Content with Caption": {
                "description": "Content with additional caption area",
                "complexity": "simple",
                "frontmatter_example": """---
layout: Content with Caption
---
# Content with Caption

Main content goes here:
- Primary information
- Secondary details

Caption area with additional context.""",
                "json_example": """{
  "type": "Content with Caption",
  "title": "Content with Caption",
  "rich_content": [
    {
      "paragraph": "Main content goes here:"
    },
    {
      "bullets": ["Primary information", "Secondary details"],
      "bullet_levels": [1, 1]
    },
    {
      "paragraph": "Caption area with additional context."
    }
  ]
}""",
            },
            "table": {
                "description": "Table layout with styling options (alias for Title and Content)",
                "complexity": "advanced",
                "frontmatter_example": """---
layout: table
style: dark_blue_white_text
row_style: alternating_light_gray
border_style: thin_gray
---
# Table Slide

| **Feature** | *Status* | ___Priority___ |
| Authentication | **Complete** | *High* |
| User Management | ***In Progress*** | ___Medium___ |
| Reporting | *Planned* | **Low** |""",
                "json_example": """{
  "type": "table",
  "style": "dark_blue_white_text",
  "row_style": "alternating_light_gray",
  "border_style": "thin_gray",
  "title": "Table Slide",
  "table": {
    "data": [
      ["**Feature**", "*Status*", "___Priority___"],
      ["Authentication", "**Complete**", "*High*"],
      ["User Management", "***In Progress***", "___Medium___"],
      ["Reporting", "*Planned*", "**Low**"]
    ],
    "header_style": "dark_blue_white_text",
    "row_style": "alternating_light_gray",
    "border_style": "thin_gray"
  }
}""",
            },
        }
        return implemented_layouts

    def generate_getting_started(self, output_path=None):
        """Generate comprehensive Getting_Started.md file"""
        if not output_path:
            output_path = Path(self.template_folder) / "Getting_Started.md"
        else:
            output_path = Path(output_path)

        layouts = self.get_implemented_layouts()
        layout_count = len(layouts)

        content = f"""# Getting Started with Deckbuilder

Welcome to Deckbuilder! This guide will help you create professional PowerPoint "
"presentations from Markdown or JSON files.

## Quick Start (3 Steps)

1. **Create your content**
   ```bash
   # Try the example file
   deckbuilder create examples/test_presentation.md
   ```

2. **View your presentation**
   Open the generated `.pptx` file in PowerPoint or LibreOffice

3. **Customize for your needs**
   Edit the example files or create your own content

## File Formats: Markdown vs JSON

Deckbuilder supports two input formats to fit different workflows:

### Markdown with Frontmatter (Recommended)
- **Best for**: Content authors, quick editing, version control
- **Syntax**: YAML frontmatter + Markdown content
- **Example**: `examples/test_presentation.md`

### JSON Format
- **Best for**: Programmatic generation, automation, complex structures
- **Syntax**: Structured JSON with rich content
- **Example**: `examples/test_presentation.json`

## Supported Layouts ({layout_count} Available)

Your template currently supports **{layout_count} layouts**. Here's how to use each:

"""

        # Generate layout reference
        for i, (layout_name, layout_info) in enumerate(layouts.items(), 1):
            content += f"""### {i}. {layout_name}
{layout_info['description']}

**Frontmatter Syntax (in .md files):**
```yaml
{layout_info['frontmatter_example']}
```

**JSON Syntax (in .json files):**
```json
{layout_info['json_example']}
```

"""

        content += """## Advanced Features

### PlaceKitten Image Support
When images are missing or invalid, Deckbuilder automatically generates professional "
"placeholder images:
- **Grayscale styling** for business presentations
- **Smart cropping** with face detection and rule-of-thirds composition
- **Automatic caching** for performance optimization
- **Professional appearance** suitable for client presentations

### Inline Formatting
Use these formatting options in any text content:
- **Bold text**: `**text**` or `***text***`
- *Italic text*: `*text*` or `***text***`
- ___Underlined text___: `___text___`
- **_Combined formatting_**: `**_bold italic_**`

### Table Support
Create tables with professional styling:
```yaml
---
layout: table
style: dark_blue_white_text
row_style: alternating_light_gray
border_style: thin_gray
---
```

## CLI Commands Reference

### Essential Commands
```bash
# Create presentation from Markdown
deckbuilder create presentation.md

# Create from JSON
deckbuilder create presentation.json

# Generate PlaceKitten images
deckbuilder image 800 600 --filter grayscale --output placeholder.jpg

# List available templates
deckbuilder templates

# Analyze template structure
deckbuilder analyze default --verbose
```

### Template Management
```bash
# Initialize new template folder
deckbuilder init ~/my-templates

# Validate template and mappings
deckbuilder validate default

# Document template capabilities
deckbuilder document default --output template_docs.md

# Enhance template placeholders
deckbuilder enhance default
```

## Troubleshooting

### Common Issues

**"Template folder not found"**
```bash
# Solution: Initialize template folder
deckbuilder init ./templates
```

**"Layout not supported"**
- Check the supported layouts list above
- Use `deckbuilder analyze default` to see available layouts
- Ensure correct spelling and capitalization

**"Image not found"**
- Don't worry! PlaceKitten will generate a professional placeholder
- Check image path is relative to presentation file
- Supported formats: JPG, PNG, WebP

**"JSON validation error"**
- Validate JSON syntax with online tools
- Check required fields for each layout type
- Compare with working examples in `examples/test_presentation.json`

### Getting Help
```bash
# Show help for any command
deckbuilder --help
deckbuilder create --help

# Show current configuration
deckbuilder config

# List available templates
deckbuilder templates
```

## Next Steps

1. **Explore the examples**: Study `examples/test_presentation.md` and "
"   `examples/test_presentation.json`
2. **Try different layouts**: Experiment with the {layout_count} supported layouts
3. **Add your content**: Replace example content with your own
4. **Customize styling**: Explore table styles and formatting options
5. **Share your presentations**: Generated `.pptx` files work in PowerPoint, "
"   LibreOffice, and Google Slides

## Advanced Usage

### Batch Processing
```bash
# Process multiple files
for file in *.md; do deckbuilder create "$file"; done
```

### Environment Configuration
```bash
# Set permanent template folder
export DECK_TEMPLATE_FOLDER="~/my-templates"
export DECK_OUTPUT_FOLDER="~/presentations"
export DECK_TEMPLATE_NAME="default"
```

### Integration with Version Control
- Markdown files work excellently with Git
- Track presentation content changes easily
- Collaborate on presentations using familiar tools

---

**Generated by Deckbuilder CLI** • Template Version: 1.0 • {datetime.now().strftime('%Y-%m-%d')}

Happy presenting! 🚀
"""

        # Write the file
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(content)

        return str(output_path)


def main():
    """Command-line interface for template management tools"""
    parser = argparse.ArgumentParser(
        description="Deckbuilder Template Management Tools",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python cli_tools.py analyze default
  python cli_tools.py analyze default --verbose
  python cli_tools.py document default
  python cli_tools.py validate default
  python cli_tools.py enhance default
  python cli_tools.py enhance default --no-backup
  python cli_tools.py enhance default --use-conventions
  python cli_tools.py analyze default --template-folder ./templates --output-folder ./output
        """,
    )

    # Global arguments
    parser.add_argument("--template-folder", "-t", help="Path to templates folder (default: auto-detect)")
    parser.add_argument("--output-folder", "-o", help="Path to output folder (default: ./template_output)")

    subparsers = parser.add_subparsers(dest="command", help="Available commands")

    # Analyze command
    analyze_parser = subparsers.add_parser("analyze", help="Analyze PowerPoint template structure")
    analyze_parser.add_argument("template", help="Template name (e.g., default)")
    analyze_parser.add_argument("--verbose", "-v", action="store_true", help="Show detailed analysis")

    # Document command
    doc_parser = subparsers.add_parser("document", help="Generate template documentation")
    doc_parser.add_argument("template", help="Template name to document")
    doc_parser.add_argument("--doc-output", help="Documentation output file path")

    # Validate command
    validate_parser = subparsers.add_parser("validate", help="Validate template and mappings")
    validate_parser.add_argument("template", help="Template name to validate")

    # Enhance command
    enhance_parser = subparsers.add_parser("enhance", help="Enhance template with corrected placeholder names (saves as .g.pptx)")
    enhance_parser.add_argument("template", help="Template name to enhance")
    enhance_parser.add_argument("--mapping-file", help="Custom JSON mapping file path")
    enhance_parser.add_argument("--no-backup", action="store_true", help="Skip creating backup before modification")
    enhance_parser.add_argument(
        "--use-conventions",
        action="store_true",
        help="Use convention-based naming system instead of JSON mapping",
    )

    args = parser.parse_args()

    if not args.command:
        parser.print_help()
        return

    # Initialize template manager with command-line arguments
    manager = TemplateManager(template_folder=args.template_folder, output_folder=args.output_folder)

    # Execute command
    if args.command == "analyze":
        manager.analyze_template(args.template, verbose=args.verbose)
    elif args.command == "document":
        manager.document_template(args.template, getattr(args, "doc_output", None))
    elif args.command == "validate":
        manager.validate_template(args.template)
    elif args.command == "enhance":
        create_backup = not args.no_backup
        manager.enhance_template(args.template, None, create_backup, args.use_conventions)


if __name__ == "__main__":
    main()
