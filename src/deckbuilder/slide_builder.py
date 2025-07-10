from pptx.enum.shapes import PP_PLACEHOLDER_TYPE

from .placeholder_types import (
    is_content_placeholder,
    is_media_placeholder,
    is_subtitle_placeholder,
    is_title_placeholder,
)
from .logging_config import slide_builder_print, debug_print, error_print


class SlideBuilder:
    """Handles core slide creation, layout mapping, and placeholder management."""

    def __init__(self, layout_mapping=None):
        """
        Initialize the slide builder.

        Args:
            layout_mapping: Optional layout mapping dictionary
        """
        self.layout_mapping = layout_mapping
        self._current_slide_index = 0

    def clear_slides(self, prs):
        """Clear all slides from the presentation."""
        slide_count = len(prs.slides)
        for i in range(slide_count - 1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]

        # Reset slide index for consistent image selection
        self._current_slide_index = 0

    def add_slide(self, prs, slide_data: dict, content_formatter, image_placeholder_handler):
        """
        Add a single slide to the presentation based on slide data.

        Args:
            prs: PowerPoint presentation object
            slide_data: Dictionary containing slide information
            content_formatter: ContentFormatter instance for handling content
            image_placeholder_handler: ImagePlaceholderHandler instance for handling images
        """
        # Type validation: ensure slide_data is a dictionary
        if not isinstance(slide_data, dict):
            raise TypeError(f"slide_data must be a dictionary, got {type(slide_data).__name__}: {slide_data}")

        # Track slide index for consistent image selection
        self._current_slide_index = getattr(self, "_current_slide_index", 0) + 1

        # Auto-parse JSON formatting for inline formatting support
        slide_data = content_formatter.auto_parse_json_formatting(slide_data)

        # Get slide type and determine layout using JSON mapping
        # Prefer explicit "layout" field over "type" field
        layout_or_type = slide_data.get("layout", slide_data.get("type", "content"))

        # Use layout mapping if available
        if self.layout_mapping:
            aliases = self.layout_mapping.get("aliases", {})
            layouts = self.layout_mapping.get("layouts", {})

            # Get layout name from aliases (or use direct layout name if it exists in layouts)
            if layout_or_type in layouts:
                layout_name = layout_or_type
            else:
                layout_name = aliases.get(layout_or_type, layout_or_type)

            # Get layout index
            layout_info = layouts.get(layout_name, {})
            layout_index = layout_info.get("index", 1)
        else:
            # Fallback
            layout_name = layout_or_type  # Use the original layout name as fallback
            layout_index = 1

        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)

        # Copy descriptive placeholder names from template mapping
        self._copy_placeholder_names_from_mapping(slide, layout_name)

        # Add content to placeholders using template mapping + semantic detection
        self._apply_content_to_mapped_placeholders(slide, slide_data, layout_name, content_formatter, image_placeholder_handler)

        # All content should be processed through placeholders only - no legacy content blocks
        debug_print("  Slide completed using structured frontmatter placeholders only")

        return slide

    def add_slide_with_direct_mapping(self, prs, slide_data: dict, content_formatter, image_placeholder_handler):
        """
        Add slide using direct field mapping (no markdown conversion).

        Args:
            prs: PowerPoint presentation object
            slide_data: Dictionary containing slide information
            content_formatter: ContentFormatter instance
            image_placeholder_handler: ImagePlaceholderHandler instance
        """
        # Import the universal formatting module
        from .content_formatting import content_formatter as cf

        # Apply universal formatting to slide data
        formatted_slide = cf.format_slide_data(slide_data)

        # Add slide using the standard method
        return self.add_slide(prs, formatted_slide, content_formatter, image_placeholder_handler)

    def _copy_placeholder_names_from_mapping(self, slide, layout_name):
        """
        Copy descriptive placeholder names from template mapping to slide placeholders.

        This enhances the PowerPoint editing experience by providing meaningful placeholder
        names like "Col 1 Title Placeholder 2" instead of generic "Text Placeholder 2".

        Args:
            slide: PowerPoint slide object
            layout_name: Name of the PowerPoint layout
        """
        if not self.layout_mapping:
            return

        # Get layout info from template mapping
        layouts = self.layout_mapping.get("layouts", {})
        layout_info = layouts.get(layout_name, {})
        placeholder_mappings = layout_info.get("placeholders", {})

        # Update placeholder names to match template mapping
        for placeholder in slide.placeholders:
            placeholder_idx = str(placeholder.placeholder_format.idx)
            if placeholder_idx in placeholder_mappings:
                descriptive_name = placeholder_mappings[placeholder_idx]
                try:
                    # Update the placeholder name
                    try:
                        placeholder.element.nvSpPr.cNvPr.name = descriptive_name
                    except AttributeError:
                        pass  # Some placeholder types don't support name changes
                except Exception:
                    # Fallback: some placeholder types might not allow name changes
                    pass  # nosec - Continue processing other placeholders

    def _apply_content_to_mapped_placeholders(self, slide, slide_data, layout_name, content_formatter, image_placeholder_handler):
        """
        Apply content to placeholders using template JSON mappings + semantic detection.

        This unified method works with both JSON input and markdown frontmatter input:
        1. Looks up layout in template JSON mappings
        2. For each field in slide_data, finds corresponding placeholder index
        3. Gets actual placeholder and determines its semantic type
        4. Applies content using appropriate semantic handler

        Args:
            slide: PowerPoint slide object
            slide_data: Dictionary containing slide content (from JSON or markdown)
            layout_name: Name of the PowerPoint layout
            content_formatter: ContentFormatter instance
            image_placeholder_handler: ImagePlaceholderHandler instance
        """
        if not self.layout_mapping:
            # Fallback to basic semantic detection if no mapping available
            self._add_content_to_placeholders_fallback(slide, slide_data, content_formatter)
            return

        # Get layout info from template mapping
        layouts = self.layout_mapping.get("layouts", {})
        layout_info = layouts.get(layout_name, {})
        placeholder_mappings = layout_info.get("placeholders", {})

        # Create reverse mapping: field_name -> placeholder_index
        field_to_index = {}
        for placeholder_idx, field_name in placeholder_mappings.items():
            field_to_index[field_name] = int(placeholder_idx)

        # Enhanced debugging: Show all available placeholders and template mapping
        slide_builder_print(f"Layout '{layout_name}' - Template Mapping Analysis:")
        slide_builder_print(f"  Available placeholders in template: {list(field_to_index.keys())}")

        # Show actual PowerPoint placeholders
        actual_placeholders = []
        for ph in slide.placeholders:
            ph_type = ph.placeholder_format.type.name if hasattr(ph.placeholder_format.type, "name") else str(ph.placeholder_format.type)
            try:
                ph_name = getattr(ph.element.nvSpPr.cNvPr, "name", "unnamed")
            except AttributeError:
                ph_name = "unnamed"
            actual_placeholders.append(f"{ph.placeholder_format.idx}:{ph_type}({ph_name})")
        slide_builder_print(f"  PowerPoint placeholders found: {actual_placeholders}")

        # Process each field in slide_data using semantic detection
        # For canonical JSON format, process the placeholders object if it exists
        content_data = slide_data.get("placeholders", {}) if "placeholders" in slide_data else slide_data
        slide_builder_print(f"  Content fields to map: {list(content_data.keys())}")
        successful_mappings = []
        failed_mappings = []

        for field_name, field_value in content_data.items():
            # Skip non-content fields
            if field_name in ["type", "table", "layout"]:
                continue

            slide_builder_print(f"    Mapping field '{field_name}' (value: {str(field_value)[:50]}...)")

            # Find placeholder using semantic detection
            target_placeholder = None
            mapping_method = None

            # Handle title placeholders
            if field_name == "title":
                slide_builder_print("      Method: Title field resolution")

                # First try to find title field in template mapping
                title_field = field_to_index.get("title")
                if title_field is not None:
                    slide_builder_print(f"        Trying template mapping 'title' -> idx {title_field}")
                    for placeholder in slide.placeholders:
                        if placeholder.placeholder_format.idx == title_field:
                            target_placeholder = placeholder
                            mapping_method = f"Template mapping title -> idx {title_field}"
                            break
                else:
                    slide_builder_print("        No 'title' field in template mapping, trying title_top")
                    # Try title_top which is common in templates
                    title_top_idx = field_to_index.get("title_top")
                    if title_top_idx is not None:
                        slide_builder_print(f"        Trying template mapping 'title_top' -> idx {title_top_idx}")
                        for placeholder in slide.placeholders:
                            if placeholder.placeholder_format.idx == title_top_idx:
                                target_placeholder = placeholder
                                mapping_method = f"Template mapping title_top -> idx {title_top_idx}"
                                break

                # Final fallback to semantic title detection
                if not target_placeholder:
                    slide_builder_print("        Template mapping failed, trying semantic title detection")
                    for placeholder in slide.placeholders:
                        if is_title_placeholder(placeholder.placeholder_format.type):
                            target_placeholder = placeholder
                            mapping_method = f"Semantic title (idx {placeholder.placeholder_format.idx})"
                            break

            # Handle subtitle placeholders
            elif field_name == "subtitle":
                slide_builder_print("      Method: Semantic subtitle detection")
                for placeholder in slide.placeholders:
                    if is_subtitle_placeholder(placeholder.placeholder_format.type):
                        target_placeholder = placeholder
                        mapping_method = f"Semantic subtitle (idx {placeholder.placeholder_format.idx})"
                        break

            # Handle content placeholders
            elif field_name == "content":
                slide_builder_print("      Method: Content field resolution")

                # First try to find content field in template mapping
                content_field = field_to_index.get("content")
                if content_field is not None:
                    slide_builder_print(f"        Trying template mapping 'content' -> idx {content_field}")
                    for placeholder in slide.placeholders:
                        if placeholder.placeholder_format.idx == content_field:
                            target_placeholder = placeholder
                            mapping_method = f"Template mapping content -> idx {content_field}"
                            break
                else:
                    slide_builder_print("        No 'content' field in template mapping, trying content_1")
                    # Fallback: try to find content_1 specifically (for layouts like vertical)
                    content_1_idx = field_to_index.get("content_1")
                    if content_1_idx is not None:
                        slide_builder_print(f"        Trying template mapping 'content_1' -> idx {content_1_idx}")
                        for placeholder in slide.placeholders:
                            if placeholder.placeholder_format.idx == content_1_idx:
                                target_placeholder = placeholder
                                mapping_method = f"Template mapping content_1 -> idx {content_1_idx}"
                                break

                # Final fallback to semantic content placeholder detection
                if not target_placeholder:
                    slide_builder_print("        Template mapping failed, trying semantic content detection")
                    for placeholder in slide.placeholders:
                        if is_content_placeholder(placeholder.placeholder_format.type):
                            target_placeholder = placeholder
                            mapping_method = f"Semantic content (idx {placeholder.placeholder_format.idx})"
                            break

            # Handle image_path fields and image placeholder fields - find PICTURE placeholders
            elif field_name == "image_path" or field_name.endswith(".image_path") or "image" in field_name.lower():
                slide_builder_print("      Method: Image field detection")
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.PICTURE:
                        target_placeholder = placeholder
                        mapping_method = f"Image placeholder (idx {placeholder.placeholder_format.idx})"
                        slide_builder_print(f"        Found PICTURE placeholder at idx {placeholder.placeholder_format.idx}")
                        break
                if not target_placeholder:
                    slide_builder_print(f"        No PICTURE placeholder found for image field '{field_name}'")

            # Handle other fields by checking if they match placeholder names in JSON mapping
            else:
                slide_builder_print("      Method: Template mapping lookup")
                # Try to find by exact field name match in JSON mapping
                target_field = field_name

                # Handle common field name variations - be flexible for users
                resolved_field = self._resolve_field_name_variations(field_name, field_to_index)

                if resolved_field != field_name:
                    slide_builder_print(f"        Field name resolved: '{field_name}' -> '{resolved_field}'")
                    target_field = resolved_field
                else:
                    slide_builder_print(f"        Using field name as-is: '{field_name}'")

                if target_field in field_to_index:
                    placeholder_idx = field_to_index[target_field]
                    slide_builder_print(f"        Template mapping found: '{target_field}' -> idx {placeholder_idx}")
                    for placeholder in slide.placeholders:
                        if placeholder.placeholder_format.idx == placeholder_idx:
                            target_placeholder = placeholder
                            mapping_method = f"Template mapping '{target_field}' -> idx {placeholder_idx}"
                            slide_builder_print(f"        Matched placeholder at idx {placeholder_idx}")
                            break
                    if not target_placeholder:
                        slide_builder_print(f"        Template mapping failed: idx {placeholder_idx} not found in slide")
                else:
                    slide_builder_print(f"        No template mapping found for field '{target_field}'")

            if target_placeholder:
                slide_builder_print(f"    SUCCESS: '{field_name}' mapped using {mapping_method}")
                successful_mappings.append(f"{field_name} -> {mapping_method}")

                # Apply content based on placeholder's semantic type
                self._apply_content_by_semantic_type(
                    slide,
                    target_placeholder,
                    field_name,
                    field_value,
                    slide_data,
                    content_formatter,
                    image_placeholder_handler,
                )
            else:
                error_print(f"    FAILED: '{field_name}' could not be mapped to any placeholder")
                failed_mappings.append(field_name)

        # Summary of mapping results
        slide_builder_print("  Mapping Summary:")
        slide_builder_print(f"    Successful: {len(successful_mappings)} fields")
        for mapping in successful_mappings:
            slide_builder_print(f"      {mapping}")
        if failed_mappings:
            error_print(f"    Failed: {len(failed_mappings)} fields")
            for field in failed_mappings:
                error_print(f"      {field} (no suitable placeholder found)")

        # Process nested structures like media.image_path
        self._process_nested_image_fields(slide, slide_data, image_placeholder_handler)

    def _resolve_field_name_variations(self, field_name: str, field_to_index: dict) -> str:
        """
        Resolve field name variations to handle user flexibility.

        Users shouldn't have to know exact template field names - the engine should
        be smart enough to map common variations automatically.
        """
        # Return original if exact match exists
        if field_name in field_to_index:
            return field_name

        # Common variations mapping
        variations = {
            # Caption variations
            "text_caption": ["text_caption_1", "caption", "caption_1"],
            "caption": ["text_caption_1", "text_caption", "caption_1"],
            # Title variations - CRITICAL: map "title" to "title_top" for template compatibility
            "title": ["title_top", "title_top_1", "main_title"],
            "title_top": ["title", "title_top_1", "main_title"],
            "title_left": ["title_left_1", "left_title", "title_col1"],
            "title_right": ["title_right_1", "right_title", "title_col2"],
            # Content variations
            "content_left": ["content_left_1", "left_content", "content_col1"],
            "content_right": ["content_right_1", "right_content", "content_col2"],
            "content": ["content_1", "main_content", "body"],
            # Image variations
            "image": ["image_1", "image_path", "picture"],
            "image_1": ["image", "image_path", "picture"],
            "image_path": ["image", "image_1", "picture"],
            # Column variations
            "content_col1": ["content_left", "content_left_1", "col1_content"],
            "content_col2": ["content_right", "content_right_1", "col2_content"],
            "content_col3": ["content_col3_1", "col3_content"],
            "content_col4": ["content_col4_1", "col4_content"],
            "title_col1": ["title_left", "title_left_1", "col1_title"],
            "title_col2": ["title_right", "title_right_1", "col2_title"],
            "title_col3": ["title_col3_1", "col3_title"],
            "title_col4": ["title_col4_1", "col4_title"],
            # Item variations (for agenda, lists)
            "content_item1": ["content_item1_1", "item1_content", "item_1"],
            "content_item2": ["content_item2_1", "item2_content", "item_2"],
            "content_item3": ["content_item3_1", "item3_content", "item_3"],
            "content_item4": ["content_item4_1", "item4_content", "item_4"],
            "content_item5": ["content_item5_1", "item5_content", "item_5"],
            "content_item6": ["content_item6_1", "item6_content", "item_6"],
            # Number variations (for agenda)
            "number_item1": ["number_item1_1", "item1_number", "num_1"],
            "number_item2": ["number_item2_1", "item2_number", "num_2"],
            "number_item3": ["number_item3_1", "item3_number", "num_3"],
            "number_item4": ["number_item4_1", "item4_number", "num_4"],
            "number_item5": ["number_item5_1", "item5_number", "num_5"],
            "number_item6": ["number_item6_1", "item6_number", "num_6"],
            # SWOT Analysis variations
            "content_top_left": ["content_16", "strengths", "strength"],
            "content_top_right": ["content_17", "weaknesses", "weakness"],
            "content_bottom_left": ["content_18", "opportunities", "opportunity"],
            "content_bottom_right": ["content_19", "threats", "threat"],
            "content_16": ["content_top_left", "strengths", "strength"],
            "content_17": ["content_top_right", "weaknesses", "weakness"],
            "content_18": ["content_bottom_left", "opportunities", "opportunity"],
            "content_19": ["content_bottom_right", "threats", "threat"],
        }

        # Check if field_name has variations to try
        if field_name in variations:
            for variant in variations[field_name]:
                if variant in field_to_index:
                    return variant

        # Reverse lookup - check if template has a field that maps to this user field
        for template_field in field_to_index.keys():
            if template_field in variations:
                if field_name in variations[template_field]:
                    return template_field

        # Smart suffix handling - try adding/removing _1 suffix
        if field_name.endswith("_1"):
            base_name = field_name[:-2]
            if base_name in field_to_index:
                return base_name
        else:
            suffixed_name = field_name + "_1"
            if suffixed_name in field_to_index:
                return suffixed_name

        # Partial matching for complex fields
        for template_field in field_to_index.keys():
            # Check if field names are similar (contain each other)
            if field_name.lower() in template_field.lower() or template_field.lower() in field_name.lower():
                return template_field

        # Return original if no variations found
        return field_name

    def _add_content_to_placeholders_fallback(self, slide, slide_data, content_formatter):
        """
        Fallback method for basic semantic placeholder detection when no JSON mapping available.
        Uses inline formatting (**bold**, *italic*, ___underline___) processed at render time.
        """
        for shape in slide.placeholders:
            placeholder_type = shape.placeholder_format.type

            # Handle title placeholders
            if "title" in slide_data and is_title_placeholder(placeholder_type):
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text_frame = shape.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                    content_formatter.apply_inline_formatting(slide_data["title"], p)
                else:
                    shape.text = slide_data["title"]

            # Handle subtitle placeholders
            elif "subtitle" in slide_data and is_subtitle_placeholder(placeholder_type):
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text_frame = shape.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                    content_formatter.apply_inline_formatting(slide_data["subtitle"], p)
                else:
                    shape.text = slide_data["subtitle"]

            # Handle main content placeholders (for simple content)
            elif "content" in slide_data and is_content_placeholder(placeholder_type):
                content_formatter.add_simple_content_to_placeholder(shape, slide_data["content"])

    def _apply_content_by_semantic_type(
        self,
        slide,
        placeholder,
        field_name,
        field_value,
        slide_data,
        content_formatter,
        image_placeholder_handler,
    ):
        """
        Apply content to a placeholder based on its semantic type and the content type.
        Uses inline formatting (**bold**, *italic*, ___underline___) processed at render time.
        """
        # Set slide context for content formatter (needed for table processing)
        content_formatter._current_slide = slide

        placeholder_type = placeholder.placeholder_format.type

        # Apply content based on placeholder semantic type
        if is_title_placeholder(placeholder_type) or is_subtitle_placeholder(placeholder_type):
            # Title/subtitle placeholders - apply inline formatting directly
            if hasattr(placeholder, "text_frame") and placeholder.text_frame:
                # Use text frame for formatting support
                text_frame = placeholder.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                # Handle formatted content properly for titles
                if isinstance(field_value, list) and field_value and isinstance(field_value[0], dict) and "text" in field_value[0]:
                    content_formatter.apply_formatted_segments_to_paragraph(field_value, p)
                else:
                    content_formatter.apply_inline_formatting(str(field_value), p)
            else:
                # Fallback to simple text
                placeholder.text = str(field_value)

        elif is_content_placeholder(placeholder_type):
            # Content placeholders - handle text, lists, etc. with inline formatting
            content_formatter.add_simple_content_to_placeholder(placeholder, field_value)

        elif is_media_placeholder(placeholder_type):
            # Media placeholders - handle images, charts, etc.
            if placeholder_type == PP_PLACEHOLDER_TYPE.PICTURE:
                image_placeholder_handler.handle_image_placeholder(placeholder, field_name, field_value, slide_data)
            elif placeholder_type == PP_PLACEHOLDER_TYPE.OBJECT and hasattr(placeholder, "text_frame"):
                # OBJECT placeholders with text_frame should be treated as content placeholders
                content_formatter.add_simple_content_to_placeholder(placeholder, field_value)
            else:
                # Other media types - fallback to text for now
                if hasattr(placeholder, "text_frame") and placeholder.text_frame:
                    text_frame = placeholder.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                    if isinstance(field_value, list) and field_value and isinstance(field_value[0], dict) and "text" in field_value[0]:
                        content_formatter.apply_formatted_segments_to_paragraph(field_value, p)
                    else:
                        content_formatter.apply_inline_formatting(str(field_value), p)

        else:
            # Other placeholder types - apply inline formatting where possible
            if hasattr(placeholder, "text_frame") and placeholder.text_frame:
                text_frame = placeholder.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                if isinstance(field_value, list) and field_value and isinstance(field_value[0], dict) and "text" in field_value[0]:
                    content_formatter.apply_formatted_segments_to_paragraph(field_value, p)
                else:
                    content_formatter.apply_inline_formatting(str(field_value), p)
            else:
                placeholder.text = str(field_value)

    def _process_nested_image_fields(self, slide, slide_data, image_placeholder_handler):
        """
        Process nested image fields like media.image_path from structured frontmatter.

        Note: This method handles raw frontmatter with nested media structures.
        Structured frontmatter conversion flattens media.image_path to image_1,
        which is already handled by the main field processing loop.

        Args:
            slide: PowerPoint slide object
            slide_data: Dictionary containing slide content
            image_placeholder_handler: ImagePlaceholderHandler instance
        """
        # Skip if this appears to be structured frontmatter that was already converted
        # (indicated by presence of flattened image fields like image_1, image_path)
        has_converted_image_fields = any(field_name == "image_path" or field_name.endswith("_1") and "image" in field_name.lower() for field_name in slide_data.keys())

        if has_converted_image_fields:
            # Already processed by structured frontmatter conversion
            return

        # Check for media structure with image_path (raw frontmatter)
        if "media" in slide_data and isinstance(slide_data["media"], dict):
            media_data = slide_data["media"]
            image_path = media_data.get("image_path")

            if image_path:
                # Find the first PICTURE placeholder
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.PICTURE:
                        image_placeholder_handler.handle_image_placeholder(placeholder, "media.image_path", image_path, slide_data)
                        break
