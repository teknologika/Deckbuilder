#!/usr/bin/env python3
"""
Built-in End-to-End Validation System for Deckbuilder

Provides automatic validation that runs on every presentation generation to prevent
layout regressions and ensure JSON ↔ Template ↔ PPTX alignment.

GitHub Issue: https://github.com/teknologika/Deckbuilder/issues/36
"""

import json
import re
from pathlib import Path
from typing import Dict, List, Any, Optional
from pptx import Presentation
from .logging_config import validation_print, error_print, success_print


class ValidationError(Exception):
    """Validation error with detailed fix instructions."""

    def __init__(self, message: str, slide_num: Optional[int] = None, field_name: Optional[str] = None):
        self.slide_num = slide_num
        self.field_name = field_name
        super().__init__(message)


class PresentationValidator:
    """
    Built-in validation system that runs automatically on every presentation generation.

    Prevents layout regressions by validating:
    1. Pre-generation: JSON ↔ Template mapping alignment
    2. Post-generation: PPTX output ↔ JSON input verification
    """

    def __init__(self, presentation_data: Dict[str, Any], template_name: str, template_folder: str):
        self.presentation_data = presentation_data
        self.template_name = template_name
        self.template_folder = Path(template_folder)
        self.template_mapping = self._load_template_mapping()

    def _load_template_mapping(self) -> Dict[str, Any]:
        """Load template mapping JSON file."""
        mapping_file = self.template_folder / f"{self.template_name}.json"
        if not mapping_file.exists():
            raise ValidationError(f"Template mapping file not found: {mapping_file}\n" f"Fix: Create {self.template_name}.json in {self.template_folder}")

        with open(mapping_file, "r") as f:
            return json.load(f)

    def validate_markdown_to_json(self, markdown_content: str, converted_json: Dict[str, Any]):
        """
        Validate Markdown → JSON conversion before template validation.

        Ensures structured frontmatter is properly converted and no data is lost.
        """
        validation_print("🔍 Markdown → JSON validation: Frontmatter processing...")

        # Parse markdown sections manually to validate conversion
        markdown_sections = self._parse_markdown_sections(markdown_content)
        json_slides = converted_json.get("slides", [])

        # Validate slide count matches
        if len(markdown_sections) != len(json_slides):
            raise ValidationError(
                f"Markdown → JSON conversion error: {len(markdown_sections)} markdown sections "
                f"converted to {len(json_slides)} JSON slides\n"
                f"Fix: Check markdown section parsing and frontmatter conversion"
            )

        # Validate each section conversion
        for section_idx, (md_section, json_slide) in enumerate(zip(markdown_sections, json_slides)):
            self._validate_section_conversion(section_idx + 1, md_section, json_slide)

        success_print("✅ Markdown → JSON validation passed")

    def validate_pre_generation(self):
        """
        Validate JSON ↔ Template mapping alignment before generation.

        Raises ValidationError immediately on any mapping issues.
        """
        validation_print("🔍 JSON → Template validation: Mapping alignment...")

        # Debug: Show template mapping structure
        validation_print(f"[Validation] Template mapping loaded: {self.template_name}.json")
        layouts = self.template_mapping.get("layouts", {})
        validation_print(f"[Validation] Available layouts: {list(layouts.keys())}")
        validation_print(f"[Validation] Validating {len(self.presentation_data.get('slides', []))} slides")

        # Validate each slide's placeholders can be mapped
        for slide_idx, slide_data in enumerate(self.presentation_data.get("slides", [])):
            slide_num = slide_idx + 1
            layout_name = slide_data.get("layout")

            validation_print(f"[Validation] Slide {slide_num}: Checking layout '{layout_name}'")

            if not layout_name:
                raise ValidationError(f"Slide {slide_num}: Missing 'layout' field\n" f"Fix: Add 'layout' field with valid layout name")

            # Check layout exists in template mapping
            layouts = self.template_mapping.get("layouts", {})
            if layout_name not in layouts:
                available_layouts = list(layouts.keys())
                error_print(f"[Validation] ERROR: Layout '{layout_name}' not found in template mapping")
                raise ValidationError(
                    f"Slide {slide_num}: Unknown layout '{layout_name}'\n" f"Available layouts: {', '.join(available_layouts)}\n" f"Fix: Use one of the available layouts or update template mapping"
                )

            # Show slide content fields for debugging
            placeholders = slide_data.get("placeholders", {})
            validation_print(f"[Validation]   Placeholder fields: {list(placeholders.keys())}")

            # Legacy content blocks should not exist in structured frontmatter
            if "content" in slide_data:
                validation_print("[Validation]   WARNING: Legacy content blocks detected - should be converted to placeholders")

            # Validate placeholder mappings
            self._validate_slide_placeholders(slide_num, slide_data, layout_name)

        success_print("✅ Pre-generation validation passed")

    def _validate_slide_placeholders(self, slide_num: int, slide_data: Dict[str, Any], layout_name: str):
        """Validate that all placeholders in slide can be mapped to template."""
        layout_info = self.template_mapping["layouts"][layout_name]
        placeholder_mappings = layout_info.get("placeholders", {})

        validation_print(f"[Validation]   Template placeholders for '{layout_name}': {placeholder_mappings}")

        # Create reverse mapping: field_name -> placeholder_index
        field_to_index = {}
        for placeholder_idx, field_name in placeholder_mappings.items():
            field_to_index[field_name] = int(placeholder_idx)

        validation_print(f"[Validation]   Field-to-index mapping: {field_to_index}")

        # Check all placeholder fields in slide data
        placeholders = slide_data.get("placeholders", {})
        unmapped_fields = []
        mapped_fields = []

        for field_name in placeholders.keys():
            if field_name in ["style"]:  # Skip non-content fields
                continue

            # Check if field can be resolved
            if self._can_resolve_field_name(field_name, field_to_index):
                mapped_fields.append(field_name)
                validation_print(f"[Validation]     ✓ '{field_name}' can be mapped")
            else:
                unmapped_fields.append(field_name)
                validation_print(f"[Validation]     ✗ '{field_name}' cannot be mapped")

        if unmapped_fields:
            available_fields = list(field_to_index.keys())
            error_print(f"[Validation] ERROR: Unmapped fields found for slide {slide_num}")
            raise ValidationError(
                f"Slide {slide_num} ({layout_name}): Cannot map placeholder fields: {', '.join(unmapped_fields)}\n"
                f"Available template fields: {', '.join(available_fields)}\n"
                f"Fix: Update template mapping to include these fields or correct field names in JSON"
            )

    def _can_resolve_field_name(self, field_name: str, field_to_index: Dict[str, int]) -> bool:
        """
        Enhanced validation - check if field name can be resolved using same logic as slide_builder.

        This uses the same field name resolution logic as the content placement system
        to ensure validation and processing are consistent.
        """
        # Direct match
        if field_name in field_to_index:
            return True

        # Always allow semantic fields that have guaranteed fallbacks
        if field_name in ["title", "subtitle"]:
            return True  # Always handled by semantic detection

        # Enhanced field name resolution for common variations
        resolved_field = self._resolve_field_name_variations(field_name, field_to_index)
        if resolved_field != field_name and resolved_field in field_to_index:
            return True

        # Content field variations
        if field_name == "content":
            # Check if template has content, content_1, or other content variations
            for variant in ["content", "content_1", "main_content", "body"]:
                if variant in field_to_index:
                    return True

        return False

    def _resolve_field_name_variations(self, field_name: str, field_to_index: dict) -> str:
        """
        Resolve field name variations - same logic as slide_builder for consistency.
        """
        # Return original if exact match exists
        if field_name in field_to_index:
            return field_name

        # Common variations mapping - must match slide_builder.py exactly
        variations = {
            # Caption variations
            "text_caption": ["text_caption_1", "caption", "caption_1"],
            "caption": ["text_caption_1", "text_caption", "caption_1"],
            # Title variations - CRITICAL: map "title" to "title_top" for template compatibility
            "title": ["title_top", "title_top_1", "main_title"],
            "title_top": ["title", "title_top_1", "main_title"],
            "title_left": ["title_left_1", "left_title", "title_col1"],
            "title_right": ["title_right_1", "right_title", "title_col2"],
            # Content variations
            "content_left": ["content_left_1", "left_content", "content_col1"],
            "content_right": ["content_right_1", "right_content", "content_col2"],
            "content": ["content_1", "main_content", "body"],
            # Image variations
            "image": ["image_1", "image_path", "picture"],
            "image_1": ["image", "image_path", "picture"],
            "image_path": ["image", "image_1", "picture"],
            # SWOT Analysis variations
            "content_top_left": ["content_16", "strengths", "strength"],
            "content_top_right": ["content_17", "weaknesses", "weakness"],
            "content_bottom_left": ["content_18", "opportunities", "opportunity"],
            "content_bottom_right": ["content_19", "threats", "threat"],
            "content_16": ["content_top_left", "strengths", "strength"],
            "content_17": ["content_top_right", "weaknesses", "weakness"],
            "content_18": ["content_bottom_left", "opportunities", "opportunity"],
            "content_19": ["content_bottom_right", "threats", "threat"],
        }

        # Check if field_name has variations to try
        if field_name in variations:
            for variant in variations[field_name]:
                if variant in field_to_index:
                    return variant

        # Reverse lookup - check if template has a field that maps to this user field
        for template_field in field_to_index.keys():
            if template_field in variations:
                if field_name in variations[template_field]:
                    return template_field

        # Smart suffix handling - try adding/removing _1 suffix
        if field_name.endswith("_1"):
            base_name = field_name[:-2]
            if base_name in field_to_index:
                return base_name
        else:
            suffixed_name = field_name + "_1"
            if suffixed_name in field_to_index:
                return suffixed_name

        # Return original if no variations found
        return field_name

    def validate_post_generation(self, pptx_file_path: str):
        """
        Validate PPTX output ↔ JSON input after generation.

        Raises ValidationError if generated content doesn't match specification.
        """
        validation_print("🔍 Post-generation validation: PPTX ↔ JSON verification...")
        validation_print(f"[Validation] Loading generated PPTX: {pptx_file_path}")

        if not Path(pptx_file_path).exists():
            raise ValidationError(f"Generated PPTX file not found: {pptx_file_path}")

        # Load generated presentation
        prs = Presentation(pptx_file_path)

        # Validate slide count
        expected_slides = len(self.presentation_data.get("slides", []))
        actual_slides = len(prs.slides)

        validation_print(f"[Validation] Slide count check: expected={expected_slides}, actual={actual_slides}")

        if actual_slides != expected_slides:
            raise ValidationError(f"Slide count mismatch: expected {expected_slides}, got {actual_slides}\n" f"Fix: Check slide generation logic for dropped or duplicated slides")

        # Validate each slide content
        validation_errors = []
        validation_print(f"[Validation] Validating content for {actual_slides} slides...")

        for slide_idx, (slide, slide_spec) in enumerate(zip(prs.slides, self.presentation_data["slides"])):
            slide_num = slide_idx + 1
            layout_name = slide_spec.get("layout", "unknown")

            try:
                self._validate_slide_content(slide_num, slide, slide_spec)
                validation_print(f"[Validation] Slide {slide_num} ({layout_name}): Content validation passed")
            except ValidationError as e:
                error_print(f"[Validation] Slide {slide_num} ({layout_name}): Content validation failed")
                validation_errors.append(str(e))

        if validation_errors:
            error_summary = "\n".join(validation_errors)
            raise ValidationError(f"Post-generation validation failed:\n{error_summary}\n" f"Fix: Check placeholder mapping logic in slide_builder.py")

        success_print("✅ Post-generation validation passed")

    def _validate_slide_content(self, slide_num: int, slide, slide_spec: Dict[str, Any]):
        """Validate individual slide content against specification."""
        layout_name = slide_spec["layout"]
        actual_layout = slide.slide_layout.name

        # Validate layout
        if actual_layout != layout_name:
            raise ValidationError(f"Slide {slide_num}: Layout mismatch - expected '{layout_name}', got '{actual_layout}'")

        # Validate critical placeholders are not empty
        self._validate_critical_placeholders(slide_num, slide, slide_spec)

        # Content blocks validation removed - structured frontmatter uses placeholders only

    def _validate_critical_placeholders(self, slide_num: int, slide, slide_spec: Dict[str, Any]):
        """Validate that critical placeholders have content."""
        placeholders_spec = slide_spec.get("placeholders", {})
        layout_name = slide_spec["layout"]

        # Get actual placeholder content
        actual_content = {}
        all_placeholders = {}  # Track all placeholders for debugging

        for shape in slide.shapes:
            try:
                if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                    ph_type = shape.placeholder_format.type
                    ph_idx = shape.placeholder_format.idx
                    try:
                        ph_name = getattr(shape.element.nvSpPr.cNvPr, "name", "unnamed")
                    except AttributeError:
                        ph_name = "unnamed"

                    content = self._extract_shape_text(shape)
                    semantic_type = self._get_semantic_type(ph_type)

                    # Special handling for PICTURE placeholders
                    has_content = bool(content.strip())
                    if ph_type.name == "PICTURE":
                        # For picture placeholders, check if image is present
                        has_image = self._check_placeholder_has_image(shape)
                        has_content = has_image
                        if has_image:
                            content = "[IMAGE PRESENT]"

                    # Track all placeholders
                    all_placeholders[f"{semantic_type}_{ph_idx}"] = {
                        "content": content,
                        "name": ph_name,
                        "has_content": has_content,
                    }

                    if has_content:  # Count content or images
                        actual_content[f"{semantic_type}_{ph_idx}"] = content
            except ValueError:
                continue

        # Check critical fields that should have content
        critical_missing = []

        for field_name, expected_content in placeholders_spec.items():
            if field_name in ["style"]:  # Skip non-content fields
                continue

            if not expected_content or str(expected_content).strip() == "":
                continue  # Skip empty expected content

            # Check if this field has corresponding content in PPTX
            found_content = False
            expected_clean = self._strip_formatting_markers(str(expected_content))

            # Special handling for image fields
            if field_name in ["image", "image_1", "image_path"] or "image" in field_name.lower():
                # For image fields, just check if we have any content in actual_content
                # (which includes "[IMAGE PRESENT]" for successful image insertion)
                for _ph_key, actual_text in actual_content.items():
                    if "PICTURE" in _ph_key or "[IMAGE PRESENT]" in actual_text:
                        found_content = True
                        break
            # Special handling for table fields
            elif isinstance(expected_content, dict) and expected_content.get("type") == "table":
                # For table fields, check if we have table shapes in the slide or table content indicators
                # Check if we have table shapes in the slide
                table_shapes = [shape for shape in slide.shapes if hasattr(shape, "table")]
                if table_shapes:
                    found_content = True
                else:
                    # Fallback: check for table content indicators in text
                    for _ph_key, actual_text in actual_content.items():
                        if "table" in actual_text.lower() or "rows" in actual_text.lower():
                            found_content = True
                            break
            else:
                # Normal text content validation
                for _ph_key, actual_text in actual_content.items():
                    if expected_clean.lower() in actual_text.lower():
                        found_content = True
                        break

            if not found_content:
                critical_missing.append(f"{field_name} (expected: '{expected_clean[:30]}...')")

        # Special validation for vertical layouts
        if "vertical" in layout_name.lower():
            content_found = any("content" in key.lower() or "body" in key.lower() for key in actual_content.keys())
            if not content_found and placeholders_spec.get("content"):
                critical_missing.append("content (vertical layout missing main content)")

        if critical_missing:
            # Show detailed validation info only on failures
            error_print(f"[Validation] Slide {slide_num} ({layout_name}): VALIDATION FAILED")
            error_print(f"[Validation]   Expected placeholders: {list(placeholders_spec.keys())}")
            error_print("[Validation]   Found placeholders in slide:")
            for ph_key, ph_info in all_placeholders.items():
                status = "✓ HAS CONTENT" if ph_info["has_content"] else "✗ EMPTY"
                error_print(f"[Validation]     {ph_key} ('{ph_info['name']}'): {status}")
                if ph_info["has_content"]:
                    error_print(f"[Validation]       Content: '{ph_info['content'][:50]}{'...' if len(ph_info['content']) > 50 else ''}'")
            error_print(f"[Validation]   Non-empty placeholders: {list(actual_content.keys())}")
            error_print("[Validation]   Checking critical field mappings:")
            for field_name, expected_content in placeholders_spec.items():
                if field_name in ["style"]:
                    continue
                if not expected_content or str(expected_content).strip() == "":
                    error_print(f"[Validation]     '{field_name}': SKIPPED (empty expected content)")
                    continue
                expected_clean = self._strip_formatting_markers(str(expected_content))
                error_print(f"[Validation]     '{field_name}': Looking for '{expected_clean[:30]}{'...' if len(expected_clean) > 30 else ''}'")
                if field_name in critical_missing:
                    error_print("[Validation]       ✗ NOT FOUND in any placeholder")
                else:
                    error_print("[Validation]       ✓ FOUND")
            raise ValidationError(f"Slide {slide_num} ({layout_name}): Missing critical content: {', '.join(critical_missing)}")

    # _validate_content_blocks method removed - structured frontmatter uses placeholders only

    def _extract_shape_text(self, shape) -> str:
        """Extract all text from a shape."""
        text_parts = []

        if hasattr(shape, "text") and shape.text.strip():
            text_parts.append(shape.text.strip())
        elif hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text.strip())

        return "\n".join(text_parts)

    def _get_semantic_type(self, ph_type) -> str:
        """Get semantic type name for a placeholder."""
        return ph_type.name if hasattr(ph_type, "name") else str(ph_type)

    def _strip_formatting_markers(self, text: str) -> str:
        """Remove formatting markers from text for comparison."""
        # Remove formatting markers
        text = re.sub(r"\*\*\*(.*?)\*\*\*", r"\1", text)  # ***bold italic***
        text = re.sub(r"___(.*?)___", r"\1", text)  # ___underline___
        text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)  # **bold**
        text = re.sub(r"\*(.*?)\*", r"\1", text)  # *italic*
        return text.strip()

    def _check_placeholder_has_image(self, shape) -> bool:
        """Check if a PICTURE placeholder contains an image."""
        try:
            # Check if shape has image content
            if hasattr(shape, "image"):
                return shape.image is not None

            # Check if shape has fill with image
            if hasattr(shape, "fill") and shape.fill.type is not None:
                # PICTURE fill type indicates image is present
                return True

            # Check for image in shape element
            if hasattr(shape, "element") and shape.element is not None:
                # Look for image elements in the shape XML
                image_elements = shape.element.xpath(
                    ".//a:blip",
                    namespaces={"a": "http://schemas.openxmlformats.org/drawingml/2006/main"},
                )
                return len(image_elements) > 0

        except Exception:  # nosec B110
            # If we can't determine, assume no image for validation purposes
            pass

        return False

    def _parse_markdown_sections(self, markdown_content: str) -> List[Dict[str, Any]]:
        """Parse markdown into sections with frontmatter and content."""
        import yaml

        # Split by frontmatter delimiters
        sections = []
        blocks = re.split(r"^---\s*$", markdown_content, flags=re.MULTILINE)

        i = 0
        while i < len(blocks):
            # Skip empty blocks
            if not blocks[i].strip():
                i += 1
                continue

            # Look for frontmatter + content pairs
            if i + 1 < len(blocks):
                try:
                    frontmatter_raw = blocks[i].strip()
                    content_raw = blocks[i + 1].strip() if i + 1 < len(blocks) else ""

                    # Parse frontmatter
                    frontmatter = yaml.safe_load(frontmatter_raw) or {}

                    sections.append(
                        {
                            "frontmatter": frontmatter,
                            "content": content_raw,
                            "raw_frontmatter": frontmatter_raw,
                        }
                    )

                    i += 2
                except yaml.YAMLError as e:
                    raise ValidationError(f"YAML parsing error in markdown frontmatter: {e}\n" f"Fix: Check YAML syntax in frontmatter section")
            else:
                i += 1

        return sections

    def _validate_section_conversion(self, section_num: int, md_section: Dict[str, Any], json_slide: Dict[str, Any]):
        """Validate individual section conversion from markdown to JSON."""
        frontmatter = md_section["frontmatter"]

        # Validate layout field conversion
        expected_layout = frontmatter.get("layout")
        actual_layout = json_slide.get("layout")

        if expected_layout != actual_layout:
            raise ValidationError(
                f"Section {section_num}: Layout conversion failed\n" f"Markdown layout: '{expected_layout}'\n" f"JSON layout: '{actual_layout}'\n" f"Fix: Check frontmatter to JSON conversion logic"
            )

        # Validate critical frontmatter fields are preserved
        critical_fields = ["title", "layout"]
        for field in critical_fields:
            if field in frontmatter:
                # Check if field appears in JSON placeholders
                json_placeholders = json_slide.get("placeholders", {})
                if field not in json_placeholders and field != "layout":
                    raise ValidationError(
                        f"Section {section_num}: Frontmatter field '{field}' not found in JSON placeholders\n"
                        f"Markdown value: '{frontmatter[field]}'\n"
                        f"Fix: Check structured frontmatter conversion logic"
                    )
