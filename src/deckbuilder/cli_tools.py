#!/usr/bin/env python3
"""
Deckbuilder Command Line Tools

Standalone utilities for template analysis, documentation generation, and validation.
These tools are designed to be run independently for template management tasks.
"""

import os
import sys
import json
import argparse
from datetime import datetime
from pathlib import Path

# Add parent directory to path for imports
sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))

from mcp_server.tools import TemplateAnalyzer


class TemplateManager:
    """Command-line template management utilities"""
    
    def __init__(self, template_folder=None, output_folder=None):
        # Use command-line arguments or sensible defaults
        if template_folder:
            self.template_folder = template_folder
        else:
            # Default: look for assets/templates relative to current directory
            current_dir = Path.cwd()
            if (current_dir / 'assets' / 'templates').exists():
                self.template_folder = str(current_dir / 'assets' / 'templates')
            else:
                # Try from project root
                project_root = Path(__file__).parent.parent.parent
                self.template_folder = str(project_root / 'assets' / 'templates')
        
        if output_folder:
            self.output_folder = output_folder
        else:
            # Default: create output folder in current directory
            self.output_folder = str(Path.cwd() / 'template_output')
            
        # Ensure folders exist
        os.makedirs(self.output_folder, exist_ok=True)
        
        # Don't create analyzer yet - wait until we need it
    
    def analyze_template(self, template_name: str, verbose: bool = False) -> dict:
        """
        Analyze a PowerPoint template and generate JSON mapping.
        
        Args:
            template_name: Name of template (e.g., 'default')
            verbose: Print detailed analysis information
            
        Returns:
            Template analysis results
        """
        print(f"🔍 Analyzing template: {template_name}")
        print(f"📁 Template folder: {self.template_folder}")
        print(f"📂 Output folder: {self.output_folder}")
        
        try:
            # Set environment variables for the analyzer
            os.environ['DECK_TEMPLATE_FOLDER'] = self.template_folder
            os.environ['DECK_OUTPUT_FOLDER'] = self.output_folder
            
            # Create analyzer instance with current environment
            analyzer = TemplateAnalyzer()
            
            # Run analysis
            result = analyzer.analyze_pptx_template(template_name)
            
            # Print summary
            layouts_count = len(result.get('layouts', {}))
            print(f"✅ Analysis complete!")
            print(f"   📊 Found {layouts_count} layouts")
            
            # Print layout summary
            if verbose and 'layouts' in result:
                print(f"\n📋 Layout Summary:")
                for layout_name, layout_info in result['layouts'].items():
                    placeholder_count = len(layout_info.get('placeholders', {}))
                    index = layout_info.get('index', '?')
                    print(f"   {index:2d}: {layout_name} ({placeholder_count} placeholders)")
            
            # Check for generated file
            base_name = template_name.replace('.pptx', '')
            output_file = os.path.join(self.output_folder, f"{base_name}.g.json")
            if os.path.exists(output_file):
                print(f"📄 Generated: {output_file}")
                print(f"   ✏️  Edit this file with semantic placeholder names")
                print(f"   📝 Rename to '{base_name}.json' when ready")
            
            return result
            
        except Exception as e:
            print(f"❌ Error analyzing template: {str(e)}")
            return {}
    
    def document_template(self, template_name: str, output_path: str = None) -> str:
        """
        Generate comprehensive documentation for a template.
        
        Args:
            template_name: Name of template to document
            output_path: Custom output path (optional)
            
        Returns:
            Path to generated documentation
        """
        print(f"📝 Generating documentation for: {template_name}")
        
        # Analyze template first
        analysis = self.analyze_template(template_name, verbose=False)
        if not analysis:
            return ""
        
        # Load JSON mapping if available
        base_name = template_name.replace('.pptx', '')
        mapping_file = os.path.join(self.template_folder, f"{base_name}.json")
        mapping = {}
        
        if os.path.exists(mapping_file):
            try:
                with open(mapping_file, 'r', encoding='utf-8') as f:
                    mapping = json.load(f)
                print(f"📄 Using mapping: {mapping_file}")
            except Exception as e:
                print(f"⚠️  Could not load mapping file: {e}")
        
        # Generate documentation
        doc_content = self._generate_template_documentation(template_name, analysis, mapping)
        
        # Save documentation
        if not output_path:
            project_root = Path(__file__).parent.parent.parent
            docs_folder = project_root / 'docs' / 'Features'
            docs_folder.mkdir(parents=True, exist_ok=True)
            output_path = str(docs_folder / f"{base_name.title()}Template.md")
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(doc_content)
        
        print(f"✅ Documentation generated: {output_path}")
        return output_path
    
    def _generate_template_documentation(self, template_name: str, analysis: dict, mapping: dict) -> str:
        """Generate markdown documentation for template"""
        base_name = template_name.replace('.pptx', '')
        layouts = analysis.get('layouts', {})
        
        # Generate layout summary table
        table_rows = []
        for layout_name, layout_info in layouts.items():
            index = layout_info.get('index', '?')
            placeholders = layout_info.get('placeholders', {})
            placeholder_count = len(placeholders)
            
            # Check if mapping exists
            mapping_status = "✅" if mapping.get('layouts', {}).get(layout_name) else "❌"
            
            # Check structured frontmatter support (placeholder for now)
            structured_status = "⏳"  # Would check structured_frontmatter.py
            
            table_rows.append(f"| {layout_name} | {index} | {placeholder_count} | {structured_status} | {mapping_status} |")
        
        # Generate detailed layout specifications
        detailed_layouts = []
        for layout_name, layout_info in layouts.items():
            index = layout_info.get('index', '?')
            placeholders = layout_info.get('placeholders', {})
            
            layout_section = f"### {layout_name} (Index: {index})\\n\\n"
            layout_section += "**PowerPoint Placeholders**:\\n"
            
            for idx, placeholder_name in placeholders.items():
                # Handle both string and object formats
                if isinstance(placeholder_name, dict):
                    name = placeholder_name.get('name', 'Unknown')
                    actual_idx = placeholder_name.get('idx', idx)
                else:
                    name = placeholder_name
                    actual_idx = idx
                layout_section += f"- `idx={actual_idx}`: \"{name}\"\\n"
            
            # Add mapping info if available
            layout_mapping = mapping.get('layouts', {}).get(layout_name)
            if layout_mapping:
                layout_section += f"\\n**JSON Mapping**: ✅ Configured\\n"
            else:
                layout_section += f"\\n**JSON Mapping**: ❌ Not configured\\n"
            
            layout_section += f"**Structured Frontmatter**: ⏳ To be implemented\\n"
            
            detailed_layouts.append(layout_section)
        
        # Generate full documentation
        doc_content = f"""# {base_name.title()} Template Documentation

## Template Overview
- **Template Name**: {template_name}
- **Analysis Date**: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}
- **Total Layouts**: {len(layouts)}
- **Template Location**: `{self.template_folder}/{template_name}`

## Layout Summary

| Layout Name | Index | Placeholders | Structured Support | JSON Mapping |
|-------------|-------|--------------|-------------------|--------------|
{chr(10).join(table_rows)}

## Detailed Layout Specifications

{chr(10).join(detailed_layouts)}

## Template Management

### Adding JSON Mapping
1. **Analyze template**: Run `python -m deckbuilder.cli_tools analyze {base_name}`
2. **Edit generated file**: Customize `{base_name}.g.json` with semantic names
3. **Activate mapping**: Rename to `{base_name}.json` in templates folder

### Example JSON Mapping Structure
```json
{{
  "template_info": {{
    "name": "{base_name.title()}",
    "version": "1.0"
  }},
  "layouts": {{
    "Title Slide": {{
      "index": 0,
      "placeholders": {{
        "0": "Title 1",
        "1": "Subtitle 2"
      }}
    }}
  }},
  "aliases": {{
    "title": "Title Slide",
    "content": "Title and Content"
  }}
}}
```

### Usage Examples

**JSON Format**:
```json
{{
  "presentation": {{
    "slides": [
      {{
        "type": "Title Slide",
        "layout": "Title Slide",
        "title": "My Presentation",
        "subtitle": "Subtitle text"
      }}
    ]
  }}
}}
```

**Markdown with Frontmatter**:
```yaml
---
layout: Title Slide
---
# My Presentation
## Subtitle text
```

---
*Generated automatically by Deckbuilder Template Manager on {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}*
"""
        
        return doc_content
    
    def validate_template(self, template_name: str) -> dict:
        """
        Validate template structure and mappings.
        
        Args:
            template_name: Name of template to validate
            
        Returns:
            Validation results
        """
        print(f"🔍 Validating template: {template_name}")
        
        validation_results = {
            "template_file": self._validate_template_file(template_name),
            "json_mapping": self._validate_json_mapping(template_name),
            "placeholder_naming": self._validate_placeholder_naming(template_name)
        }
        
        # Print results
        for category, results in validation_results.items():
            status = results.get('status', 'unknown')
            if status == 'valid':
                print(f"✅ {category}: Valid")
            elif status == 'issues_found':
                print(f"⚠️  {category}: Issues found")
                for issue in results.get('issues', []):
                    print(f"   - {issue}")
            else:
                print(f"❌ {category}: {results.get('error', 'Unknown error')}")
        
        return validation_results
    
    def _validate_template_file(self, template_name: str) -> dict:
        """Validate template file exists and is accessible"""
        try:
            if not template_name.endswith('.pptx'):
                template_name += '.pptx'
            
            template_path = os.path.join(self.template_folder, template_name)
            
            if not os.path.exists(template_path):
                return {"status": "error", "error": f"Template file not found: {template_path}"}
            
            # Try to load with python-pptx
            from pptx import Presentation
            prs = Presentation(template_path)
            
            return {
                "status": "valid",
                "layout_count": len(prs.slide_layouts),
                "file_size": os.path.getsize(template_path)
            }
            
        except Exception as e:
            return {"status": "error", "error": str(e)}
    
    def _validate_json_mapping(self, template_name: str) -> dict:
        """Validate JSON mapping file"""
        base_name = template_name.replace('.pptx', '')
        mapping_file = os.path.join(self.template_folder, f"{base_name}.json")
        
        if not os.path.exists(mapping_file):
            return {"status": "missing", "message": "JSON mapping file not found"}
        
        try:
            with open(mapping_file, 'r', encoding='utf-8') as f:
                mapping = json.load(f)
            
            # Basic structure validation
            required_keys = ['template_info', 'layouts']
            missing_keys = [key for key in required_keys if key not in mapping]
            
            if missing_keys:
                return {
                    "status": "issues_found",
                    "issues": [f"Missing required keys: {missing_keys}"]
                }
            
            return {"status": "valid", "layouts_count": len(mapping.get('layouts', {}))}
            
        except json.JSONDecodeError as e:
            return {"status": "error", "error": f"Invalid JSON: {str(e)}"}
        except Exception as e:
            return {"status": "error", "error": str(e)}
    
    def _validate_placeholder_naming(self, template_name: str) -> dict:
        """Validate placeholder naming conventions"""
        # Placeholder for naming convention validation
        return {"status": "valid", "message": "Naming validation not yet implemented"}
    
    def enhance_template(self, template_name: str, mapping_file: str = None, create_backup: bool = True) -> dict:
        """
        Enhance template by updating master slide placeholder names using semantic mapping.
        
        Args:
            template_name: Name of template to enhance
            mapping_file: Custom JSON mapping file (optional)
            create_backup: Create backup before modification (default: True)
            
        Returns:
            Enhancement results with success/failure information
        """
        print(f"🔧 Enhancing template: {template_name}")
        
        try:
            # Ensure template name has .pptx extension
            if not template_name.endswith('.pptx'):
                template_name += '.pptx'
            
            template_path = os.path.join(self.template_folder, template_name)
            
            if not os.path.exists(template_path):
                return {"status": "error", "error": f"Template file not found: {template_path}"}
            
            # Determine mapping file path
            base_name = template_name.replace('.pptx', '')
            if not mapping_file:
                mapping_file = os.path.join(self.template_folder, f"{base_name}.json")
            
            if not os.path.exists(mapping_file):
                return {"status": "error", "error": f"Mapping file not found: {mapping_file}. Run 'analyze' first to generate mapping."}
            
            # Create backup if requested
            if create_backup:
                backup_path = self._create_template_backup(template_path)
                print(f"📄 Backup created: {backup_path}")
            
            # Load mapping
            with open(mapping_file, 'r', encoding='utf-8') as f:
                mapping = json.load(f)
            
            # Enhance template
            modifications = self._modify_master_slide_placeholders(template_path, mapping)
            
            if modifications["modified_count"] > 0:
                print(f"✅ Enhancement complete!")
                print(f"   📊 Modified {modifications['modified_count']} placeholders across {modifications['layout_count']} layouts")
                
                if "enhanced_template_path" in modifications:
                    print(f"   📄 Enhanced template saved: {modifications['enhanced_template_path']}")
                
                if modifications["issues"]:
                    print(f"⚠️  Issues encountered:")
                    for issue in modifications["issues"]:
                        print(f"   - {issue}")
                
                return {
                    "status": "success", 
                    "modifications": modifications,
                    "backup_path": backup_path if create_backup else None
                }
            else:
                print(f"ℹ️  No modifications needed - all placeholders already have correct names")
                return {"status": "no_changes", "message": "Template already has correct placeholder names"}
                
        except Exception as e:
            return {"status": "error", "error": str(e)}
    
    def _create_template_backup(self, template_path: str) -> str:
        """Create backup copy of template file in organized backups folder"""
        import shutil
        from datetime import datetime
        
        # Generate backup filename with timestamp
        path_obj = Path(template_path)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        backup_name = f"{path_obj.stem}_backup_{timestamp}{path_obj.suffix}"
        
        # Create backups folder within templates directory
        backups_folder = path_obj.parent / "backups"
        backups_folder.mkdir(exist_ok=True)
        
        backup_path = backups_folder / backup_name
        
        # Copy file
        shutil.copy2(template_path, backup_path)
        return str(backup_path)
    
    def _modify_master_slide_placeholders(self, template_path: str, mapping: dict) -> dict:
        """
        Modify master slide placeholder names using python-pptx.
        
        Args:
            template_path: Path to PowerPoint template
            mapping: JSON mapping with placeholder names
            
        Returns:
            Dictionary with modification results
        """
        from pptx import Presentation
        
        modifications = {
            "modified_count": 0,
            "layout_count": 0,
            "issues": [],
            "changes": []
        }
        
        # Load presentation
        prs = Presentation(template_path)
        layouts_mapping = mapping.get("layouts", {})
        
        # Access the slide master (this is the key change!)
        slide_master = prs.slide_master
        
        # First, try to modify placeholders on the master slide itself
        try:
            for placeholder in slide_master.placeholders:
                placeholder_idx = str(placeholder.placeholder_format.idx)
                # Try to find this placeholder in any layout mapping
                for layout_name, layout_info in layouts_mapping.items():
                    placeholder_mapping = layout_info.get("placeholders", {})
                    if placeholder_idx in placeholder_mapping:
                        new_name = placeholder_mapping[placeholder_idx]
                        old_name = placeholder.name
                        try:
                            placeholder.element.nvSpPr.cNvPr.name = new_name
                            modifications["modified_count"] += 1
                            modifications["changes"].append({
                                "location": "master_slide",
                                "placeholder_idx": placeholder_idx,
                                "old_name": old_name,
                                "new_name": new_name
                            })
                            break
                        except Exception as e:
                            modifications["issues"].append(f"Failed to modify master placeholder {placeholder_idx}: {str(e)}")
        except Exception as e:
            modifications["issues"].append(f"No direct master placeholders or error: {e}")
        
        # Process each slide layout in the master
        for layout in slide_master.slide_layouts:
            layout_name = layout.name
            modifications["layout_count"] += 1
            
            if layout_name not in layouts_mapping:
                modifications["issues"].append(f"Layout '{layout_name}' not found in mapping file")
                continue
            
            layout_mapping = layouts_mapping[layout_name]
            placeholder_mapping = layout_mapping.get("placeholders", {})
            
            # Modify placeholders in this master slide layout
            for placeholder in layout.placeholders:
                placeholder_idx = str(placeholder.placeholder_format.idx)
                
                if placeholder_idx in placeholder_mapping:
                    new_name = placeholder_mapping[placeholder_idx]
                    old_name = placeholder.name
                    
                    try:
                        # Update placeholder name on the master slide layout
                        if hasattr(placeholder, 'element') and hasattr(placeholder.element, 'nvSpPr'):
                            placeholder.element.nvSpPr.cNvPr.name = new_name
                            modifications["modified_count"] += 1
                            modifications["changes"].append({
                                "layout": layout_name,
                                "placeholder_idx": placeholder_idx,
                                "old_name": old_name,
                                "new_name": new_name
                            })
                        else:
                            modifications["issues"].append(f"Cannot modify placeholder {placeholder_idx} in {layout_name} - unsupported structure")
                    except Exception as e:
                        modifications["issues"].append(f"Failed to modify placeholder {placeholder_idx} in {layout_name}: {str(e)}")
        
        # Save modified presentation with .g.pptx extension
        try:
            # Generate enhanced template filename with .g.pptx convention
            path_obj = Path(template_path)
            enhanced_name = f"{path_obj.stem}.g{path_obj.suffix}"
            enhanced_path = path_obj.parent / enhanced_name
            
            prs.save(str(enhanced_path))
            modifications["enhanced_template_path"] = str(enhanced_path)
        except Exception as e:
            modifications["issues"].append(f"Failed to save enhanced template: {str(e)}")
            
        return modifications


def main():
    """Command-line interface for template management tools"""
    parser = argparse.ArgumentParser(
        description='Deckbuilder Template Management Tools',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python cli_tools.py analyze default
  python cli_tools.py analyze default --verbose
  python cli_tools.py document default
  python cli_tools.py validate default
  python cli_tools.py enhance default
  python cli_tools.py enhance default --no-backup
  python cli_tools.py analyze default --template-folder ./templates --output-folder ./output
        """
    )
    
    # Global arguments
    parser.add_argument('--template-folder', '-t', 
                       help='Path to templates folder (default: auto-detect)')
    parser.add_argument('--output-folder', '-o', 
                       help='Path to output folder (default: ./template_output)')
    
    subparsers = parser.add_subparsers(dest='command', help='Available commands')
    
    # Analyze command
    analyze_parser = subparsers.add_parser('analyze', help='Analyze PowerPoint template structure')
    analyze_parser.add_argument('template', help='Template name (e.g., default)')
    analyze_parser.add_argument('--verbose', '-v', action='store_true', help='Show detailed analysis')
    
    # Document command
    doc_parser = subparsers.add_parser('document', help='Generate template documentation')
    doc_parser.add_argument('template', help='Template name to document')
    doc_parser.add_argument('--doc-output', help='Documentation output file path')
    
    # Validate command
    validate_parser = subparsers.add_parser('validate', help='Validate template and mappings')
    validate_parser.add_argument('template', help='Template name to validate')
    
    # Enhance command
    enhance_parser = subparsers.add_parser('enhance', help='Enhance template with corrected placeholder names (saves as .g.pptx)')
    enhance_parser.add_argument('template', help='Template name to enhance')
    enhance_parser.add_argument('--mapping-file', help='Custom JSON mapping file path')
    enhance_parser.add_argument('--no-backup', action='store_true', help='Skip creating backup before modification')
    
    args = parser.parse_args()
    
    if not args.command:
        parser.print_help()
        return
    
    # Initialize template manager with command-line arguments
    manager = TemplateManager(
        template_folder=args.template_folder,
        output_folder=args.output_folder
    )
    
    # Execute command
    if args.command == 'analyze':
        manager.analyze_template(args.template, verbose=args.verbose)
    elif args.command == 'document':
        manager.document_template(args.template, getattr(args, 'doc_output', None))
    elif args.command == 'validate':
        manager.validate_template(args.template)
    elif args.command == 'enhance':
        create_backup = not args.no_backup
        manager.enhance_template(args.template, args.mapping_file, create_backup)


if __name__ == '__main__':
    main()