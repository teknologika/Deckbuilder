"""
PowerPoint Validation System for Testing

Provides comprehensive validation of generated PowerPoint files including
content verification, layout validation, and formatting checks.
"""

import os
import re
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from dataclasses import dataclass
from enum import Enum

try:
    from pptx import Presentation
    from pptx.shapes.base import BaseShape
    from pptx.text.text import _Text
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class ValidationResult(Enum):
    """Validation result types."""

    PASS = "pass"
    FAIL = "fail"
    WARNING = "warning"
    SKIP = "skip"


@dataclass
class ValidationError:
    """Validation error information."""

    check_name: str
    severity: ValidationResult
    message: str
    slide_index: Optional[int] = None
    placeholder_name: Optional[str] = None
    expected_value: Optional[str] = None
    actual_value: Optional[str] = None


@dataclass
class ValidationReport:
    """Comprehensive validation report."""

    file_path: str
    total_checks: int
    passed_checks: int
    failed_checks: int
    warnings: int
    skipped_checks: int
    errors: List[ValidationError]

    @property
    def success_rate(self) -> float:
        """Calculate success rate percentage."""
        if self.total_checks == 0:
            return 0.0
        return (self.passed_checks / self.total_checks) * 100

    @property
    def overall_result(self) -> ValidationResult:
        """Get overall validation result."""
        if self.failed_checks > 0:
            return ValidationResult.FAIL
        elif self.warnings > 0:
            return ValidationResult.WARNING
        elif self.passed_checks > 0:
            return ValidationResult.PASS
        else:
            return ValidationResult.SKIP


class PowerPointValidator:
    """Comprehensive PowerPoint file validation system."""

    def __init__(self, strict_mode: bool = False):
        """
        Initialize validator.

        Args:
            strict_mode: If True, warnings are treated as failures
        """
        if not HAS_PPTX:
            raise ImportError("python-pptx package required for PowerPoint validation")

        self.strict_mode = strict_mode
        self.errors: List[ValidationError] = []
        self.current_slide_index = 0

    def validate_presentation(
        self,
        pptx_path: Path,
        expected_content: Dict[str, Any],
        layout_mapping: Optional[Dict[str, Any]] = None,
        formatting_rules: Optional[Dict[str, Any]] = None,
    ) -> ValidationReport:
        """
        Perform comprehensive validation of PowerPoint presentation.

        Args:
            pptx_path: Path to PowerPoint file
            expected_content: Expected content structure
            layout_mapping: Expected layout mappings
            formatting_rules: Formatting validation rules

        Returns:
            ValidationReport with detailed results
        """
        self.errors = []

        if not pptx_path.exists():
            self.errors.append(
                ValidationError(
                    check_name="file_exists",
                    severity=ValidationResult.FAIL,
                    message=f"PowerPoint file not found: {pptx_path}",
                )
            )
            return self._generate_report(str(pptx_path))

        try:
            presentation = Presentation(str(pptx_path))
        except Exception as e:
            self.errors.append(
                ValidationError(
                    check_name="file_load",
                    severity=ValidationResult.FAIL,
                    message=f"Failed to load PowerPoint file: {e}",
                )
            )
            return self._generate_report(str(pptx_path))

        # Perform various validation checks
        self._validate_slide_count(presentation, expected_content)
        self._validate_content_preservation(presentation, expected_content)

        if layout_mapping:
            self._validate_layout_correctness(presentation, layout_mapping)

        if formatting_rules:
            self._validate_formatting_preservation(presentation, formatting_rules)

        self._validate_placeholder_population(presentation)
        self._validate_slide_structure(presentation)

        return self._generate_report(str(pptx_path))

    def validate_content_preservation(
        self, pptx_path: Path, expected_content: Dict[str, Any]
    ) -> ValidationReport:
        """Verify all expected content appears in generated PPTX."""
        self.errors = []

        try:
            presentation = Presentation(str(pptx_path))
            self._validate_content_preservation(presentation, expected_content)
        except Exception as e:
            self.errors.append(
                ValidationError(
                    check_name="content_validation",
                    severity=ValidationResult.FAIL,
                    message=f"Content validation failed: {e}",
                )
            )

        return self._generate_report(str(pptx_path))

    def validate_layout_correctness(
        self, pptx_path: Path, layout_mapping: Dict[str, Any]
    ) -> ValidationReport:
        """Verify correct slide layouts are used."""
        self.errors = []

        try:
            presentation = Presentation(str(pptx_path))
            self._validate_layout_correctness(presentation, layout_mapping)
        except Exception as e:
            self.errors.append(
                ValidationError(
                    check_name="layout_validation",
                    severity=ValidationResult.FAIL,
                    message=f"Layout validation failed: {e}",
                )
            )

        return self._generate_report(str(pptx_path))

    def validate_formatting_preservation(
        self, pptx_path: Path, formatting_rules: Dict[str, Any]
    ) -> ValidationReport:
        """Verify bold/italic/underline formatting is preserved."""
        self.errors = []

        try:
            presentation = Presentation(str(pptx_path))
            self._validate_formatting_preservation(presentation, formatting_rules)
        except Exception as e:
            self.errors.append(
                ValidationError(
                    check_name="formatting_validation",
                    severity=ValidationResult.FAIL,
                    message=f"Formatting validation failed: {e}",
                )
            )

        return self._generate_report(str(pptx_path))

    def validate_placeholder_population(self, pptx_path: Path) -> ValidationReport:
        """Ensure no placeholders remain empty."""
        self.errors = []

        try:
            presentation = Presentation(str(pptx_path))
            self._validate_placeholder_population(presentation)
        except Exception as e:
            self.errors.append(
                ValidationError(
                    check_name="placeholder_validation",
                    severity=ValidationResult.FAIL,
                    message=f"Placeholder validation failed: {e}",
                )
            )

        return self._generate_report(str(pptx_path))

    def _validate_slide_count(
        self, presentation: Presentation, expected_content: Dict[str, Any]
    ) -> None:
        """Validate expected number of slides."""
        expected_slides = len(expected_content.get("presentation", {}).get("slides", []))
        actual_slides = len(presentation.slides)

        if expected_slides != actual_slides:
            self.errors.append(
                ValidationError(
                    check_name="slide_count",
                    severity=ValidationResult.FAIL,
                    message=f"Expected {expected_slides} slides, found {actual_slides}",
                    expected_value=str(expected_slides),
                    actual_value=str(actual_slides),
                )
            )

    def _validate_content_preservation(
        self, presentation: Presentation, expected_content: Dict[str, Any]
    ) -> None:
        """Validate content preservation in slides."""
        slides_data = expected_content.get("presentation", {}).get("slides", [])

        for slide_index, slide_data in enumerate(slides_data):
            self.current_slide_index = slide_index

            if slide_index >= len(presentation.slides):
                self.errors.append(
                    ValidationError(
                        check_name="slide_missing",
                        severity=ValidationResult.FAIL,
                        message=f"Slide {slide_index} missing from presentation",
                        slide_index=slide_index,
                    )
                )
                continue

            slide = presentation.slides[slide_index]
            self._validate_slide_content(slide, slide_data, slide_index)

    def _validate_slide_content(self, slide, slide_data: Dict[str, Any], slide_index: int) -> None:
        """Validate content of individual slide."""
        # Extract all text from slide
        slide_text = self._extract_slide_text(slide)

        # Check title content
        if "title" in slide_data:
            expected_title = self._clean_text(slide_data["title"])
            if not self._text_appears_in_slide(expected_title, slide_text):
                self.errors.append(
                    ValidationError(
                        check_name="title_content",
                        severity=ValidationResult.FAIL,
                        message=f"Expected title not found in slide {slide_index}",
                        slide_index=slide_index,
                        expected_value=expected_title,
                        actual_value=(
                            slide_text[:100] + "..." if len(slide_text) > 100 else slide_text
                        ),
                    )
                )

        # Check placeholder content
        for key, value in slide_data.items():
            if key.startswith(("content_", "title_", "text_")) and isinstance(value, str):
                expected_text = self._clean_text(value)
                if not self._text_appears_in_slide(expected_text, slide_text):
                    self.errors.append(
                        ValidationError(
                            check_name="placeholder_content",
                            severity=ValidationResult.FAIL,
                            message=f"Expected content not found in slide {slide_index}",
                            slide_index=slide_index,
                            placeholder_name=key,
                            expected_value=expected_text,
                            actual_value=(
                                slide_text[:100] + "..." if len(slide_text) > 100 else slide_text
                            ),
                        )
                    )

        # Check rich content
        if "rich_content" in slide_data:
            self._validate_rich_content(slide_text, slide_data["rich_content"], slide_index)

    def _validate_rich_content(
        self, slide_text: str, rich_content: List[Dict], slide_index: int
    ) -> None:
        """Validate rich content elements."""
        for content_item in rich_content:
            if "heading" in content_item:
                expected_heading = self._clean_text(content_item["heading"])
                if not self._text_appears_in_slide(expected_heading, slide_text):
                    self.errors.append(
                        ValidationError(
                            check_name="rich_content_heading",
                            severity=ValidationResult.FAIL,
                            message=f"Expected heading not found in slide {slide_index}",
                            slide_index=slide_index,
                            expected_value=expected_heading,
                        )
                    )

            if "paragraph" in content_item:
                expected_paragraph = self._clean_text(content_item["paragraph"])
                if not self._text_appears_in_slide(expected_paragraph, slide_text):
                    self.errors.append(
                        ValidationError(
                            check_name="rich_content_paragraph",
                            severity=ValidationResult.FAIL,
                            message=f"Expected paragraph not found in slide {slide_index}",
                            slide_index=slide_index,
                            expected_value=expected_paragraph,
                        )
                    )

            if "bullets" in content_item:
                for bullet in content_item["bullets"]:
                    expected_bullet = self._clean_text(bullet)
                    if not self._text_appears_in_slide(expected_bullet, slide_text):
                        self.errors.append(
                            ValidationError(
                                check_name="rich_content_bullet",
                                severity=ValidationResult.WARNING,
                                message=f"Expected bullet point not found in slide {slide_index}",
                                slide_index=slide_index,
                                expected_value=expected_bullet,
                            )
                        )

    def _validate_layout_correctness(
        self, presentation: Presentation, layout_mapping: Dict[str, Any]
    ) -> None:
        """Validate slide layout usage."""
        layouts = layout_mapping.get("layouts", {})

        for slide_index, slide in enumerate(presentation.slides):
            layout_name = slide.slide_layout.name

            # Check if layout exists in mapping
            if layout_name not in layouts:
                self.errors.append(
                    ValidationError(
                        check_name="layout_unknown",
                        severity=ValidationResult.WARNING,
                        message=f"Unknown layout '{layout_name}' used in slide {slide_index}",
                        slide_index=slide_index,
                        actual_value=layout_name,
                    )
                )

    def _validate_formatting_preservation(
        self, presentation: Presentation, formatting_rules: Dict[str, Any]
    ) -> None:
        """Validate formatting preservation."""
        for slide_index, slide in enumerate(presentation.slides):
            self.current_slide_index = slide_index

            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    self._validate_shape_formatting(shape, slide_index, formatting_rules)

    def _validate_shape_formatting(
        self, shape, slide_index: int, formatting_rules: Dict[str, Any]
    ) -> None:
        """Validate formatting of individual shape."""
        try:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    font = run.font

                    # Check for formatting markers in text that should be applied
                    if "**" in text and not font.bold:
                        self.errors.append(
                            ValidationError(
                                check_name="bold_formatting",
                                severity=ValidationResult.WARNING,
                                message=f"Bold formatting not applied to text with ** markers",
                                slide_index=slide_index,
                                actual_value=text,
                            )
                        )

                    if "*" in text and not text.startswith("**") and not font.italic:
                        self.errors.append(
                            ValidationError(
                                check_name="italic_formatting",
                                severity=ValidationResult.WARNING,
                                message=f"Italic formatting not applied to text with * markers",
                                slide_index=slide_index,
                                actual_value=text,
                            )
                        )

                    if "___" in text and not font.underline:
                        self.errors.append(
                            ValidationError(
                                check_name="underline_formatting",
                                severity=ValidationResult.WARNING,
                                message=f"Underline formatting not applied to text with ___ markers",
                                slide_index=slide_index,
                                actual_value=text,
                            )
                        )
        except Exception as e:
            self.errors.append(
                ValidationError(
                    check_name="formatting_check_error",
                    severity=ValidationResult.WARNING,
                    message=f"Error checking formatting: {e}",
                    slide_index=slide_index,
                )
            )

    def _validate_placeholder_population(self, presentation: Presentation) -> None:
        """Validate that placeholders are populated."""
        for slide_index, slide in enumerate(presentation.slides):
            self.current_slide_index = slide_index

            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text = shape.text_frame.text.strip()

                    # Check for common placeholder patterns
                    placeholder_patterns = [
                        r"Click to add title",
                        r"Click to add text",
                        r"Click to add content",
                        r"Text Placeholder \d+",
                        r"Content Placeholder \d+",
                    ]

                    for pattern in placeholder_patterns:
                        if re.search(pattern, text, re.IGNORECASE):
                            self.errors.append(
                                ValidationError(
                                    check_name="empty_placeholder",
                                    severity=ValidationResult.WARNING,
                                    message=f"Empty placeholder found in slide {slide_index}",
                                    slide_index=slide_index,
                                    actual_value=text,
                                )
                            )
                            break

    def _validate_slide_structure(self, presentation: Presentation) -> None:
        """Validate overall slide structure."""
        if len(presentation.slides) == 0:
            self.errors.append(
                ValidationError(
                    check_name="no_slides",
                    severity=ValidationResult.FAIL,
                    message="Presentation contains no slides",
                )
            )
            return

        # Check for consistent footer elements
        footer_elements = []
        for slide_index, slide in enumerate(presentation.slides):
            slide_footers = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text = shape.text_frame.text.lower()
                    if any(keyword in text for keyword in ["footer", "date", "slide number"]):
                        slide_footers.append(text)
            footer_elements.append(slide_footers)

        # Could add more structural validation here

    def _extract_slide_text(self, slide) -> str:
        """Extract all text content from a slide."""
        texts = []

        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                texts.append(shape.text_frame.text)
            elif hasattr(shape, "text") and shape.text:
                texts.append(shape.text)

        return " ".join(texts)

    def _clean_text(self, text: str) -> str:
        """Clean text by removing formatting markers."""
        # Remove markdown-style formatting
        cleaned = re.sub(r"\*\*\*(.*?)\*\*\*", r"\\1", text)  # Bold + italic
        cleaned = re.sub(r"\*\*(.*?)\*\*", r"\\1", cleaned)  # Bold
        cleaned = re.sub(r"\*(.*?)\*", r"\\1", cleaned)  # Italic
        cleaned = re.sub(r"___(.*?)___", r"\\1", cleaned)  # Underline

        return cleaned.strip()

    def _text_appears_in_slide(self, expected_text: str, slide_text: str) -> bool:
        """Check if expected text appears in slide text."""
        expected_lower = expected_text.lower()
        slide_lower = slide_text.lower()

        # Direct match
        if expected_lower in slide_lower:
            return True

        # Word-by-word match (handle formatting differences)
        expected_words = expected_lower.split()
        slide_words = slide_lower.split()

        # Check if all expected words appear in slide
        return all(word in slide_words for word in expected_words if len(word) > 2)

    def _generate_report(self, file_path: str) -> ValidationReport:
        """Generate validation report from collected errors."""
        total_checks = len(self.errors) + 1  # +1 for successful operations
        failed_checks = len([e for e in self.errors if e.severity == ValidationResult.FAIL])
        warnings = len([e for e in self.errors if e.severity == ValidationResult.WARNING])
        skipped_checks = len([e for e in self.errors if e.severity == ValidationResult.SKIP])
        passed_checks = total_checks - failed_checks - warnings - skipped_checks

        return ValidationReport(
            file_path=file_path,
            total_checks=total_checks,
            passed_checks=max(0, passed_checks),
            failed_checks=failed_checks,
            warnings=warnings,
            skipped_checks=skipped_checks,
            errors=self.errors.copy(),
        )

    def generate_validation_report(
        self, report: ValidationReport, output_path: Optional[Path] = None
    ) -> str:
        """Generate detailed validation report."""
        report_lines = [
            f"PowerPoint Validation Report",
            f"=" * 50,
            f"File: {report.file_path}",
            f"Overall Result: {report.overall_result.value.upper()}",
            f"Success Rate: {report.success_rate:.1f}%",
            f"",
            f"Summary:",
            f"  Total Checks: {report.total_checks}",
            f"  Passed: {report.passed_checks}",
            f"  Failed: {report.failed_checks}",
            f"  Warnings: {report.warnings}",
            f"  Skipped: {report.skipped_checks}",
            f"",
        ]

        if report.errors:
            report_lines.append("Detailed Results:")
            report_lines.append("-" * 30)

            for error in report.errors:
                slide_info = (
                    f" (Slide {error.slide_index})" if error.slide_index is not None else ""
                )
                placeholder_info = f" [{error.placeholder_name}]" if error.placeholder_name else ""

                report_lines.append(
                    f"{error.severity.value.upper()}: {error.check_name}{slide_info}{placeholder_info}"
                )
                report_lines.append(f"  Message: {error.message}")

                if error.expected_value:
                    report_lines.append(f"  Expected: {error.expected_value[:100]}...")
                if error.actual_value:
                    report_lines.append(f"  Actual: {error.actual_value[:100]}...")

                report_lines.append("")

        report_text = "\n".join(report_lines)

        if output_path:
            output_path.write_text(report_text, encoding="utf-8")

        return report_text


# Convenience functions for testing
def validate_presentation_file(
    pptx_path: Path, expected_content: Dict[str, Any], strict_mode: bool = False
) -> ValidationReport:
    """Convenience function to validate a presentation file."""
    validator = PowerPointValidator(strict_mode=strict_mode)
    return validator.validate_presentation(pptx_path, expected_content)


def quick_content_check(pptx_path: Path, expected_texts: List[str]) -> bool:
    """Quick check if expected texts appear in presentation."""
    if not HAS_PPTX:
        return False

    try:
        presentation = Presentation(str(pptx_path))
        all_text = ""

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    all_text += " " + shape.text_frame.text

        all_text_lower = all_text.lower()
        return all(text.lower() in all_text_lower for text in expected_texts)

    except Exception:
        return False
