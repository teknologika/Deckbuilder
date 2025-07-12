"""
Test JSON Field Mapping with CORRECT Expected Behavior

These tests define the EXPECTED behavior for Issue #31:
- Semantic field names (content_left vs content_left_1)
- Direct JSON processing (no markdown conversion)
- ALL formatting types preserved
- Actual PowerPoint content validation

These tests SHOULD FAIL initially, then pass after implementation.
"""

import json
import os
import subprocess
import tempfile
from pathlib import Path
from pptx import Presentation
import sys

# Add src to path
sys.path.insert(0, str(Path(__file__).parent.parent.parent.parent / "src"))


class TestJSONFieldMappingCorrect:
    """Test CORRECT JSON field mapping behavior - these should work after fix"""

    def setup_method(self):
        """Setup for each test"""
        self.project_root = Path(__file__).parent.parent.parent.parent
        self.template_env_path = str(self.project_root / "src" / "deckbuilder" / "assets" / "templates")
        self.temp_dir = None

    def teardown_method(self):
        """Cleanup after each test"""
        if self.temp_dir and Path(self.temp_dir).exists():
            import shutil

            shutil.rmtree(self.temp_dir)

    def test_semantic_field_names_two_content(self):
        """Test Two Content layout with semantic field names (content_left, content_right)"""

        # CORRECT expected canonical JSON structure with structured frontmatter
        json_data = {
            "slides": [
                {
                    "layout": "Two Content",
                    "placeholders": {
                        "title": "Two Content Test",
                        "content_left": "Left item 1\nLeft item 2",
                        "content_right": "Right item 1\nRight item 2",
                    },
                    "content": [],
                }
            ]
        }

        with tempfile.TemporaryDirectory() as temp_dir:
            self.temp_dir = temp_dir

            # Write JSON input
            json_file = Path(temp_dir) / "two_content_test.json"
            with open(json_file, "w") as f:
                json.dump(json_data, f)

            # Generate PowerPoint
            env = os.environ.copy()
            env["DECK_TEMPLATE_FOLDER"] = self.template_env_path
            env["DECK_OUTPUT_FOLDER"] = temp_dir

            result = subprocess.run(
                [
                    "python",
                    str(self.project_root / "src" / "deckbuilder" / "cli.py"),
                    "create",
                    str(json_file),
                    "--output",
                    "two_content_test",
                ],
                capture_output=True,
                text=True,
                env=env,
                cwd=temp_dir,
            )

            print("\nTWO CONTENT SEMANTIC FIELD TEST:")
            print(f"CLI return code: {result.returncode}")
            print(f"CLI output: {result.stdout}")
            if result.stderr:
                print(f"CLI errors: {result.stderr}")

            # Should succeed without validation errors
            assert result.returncode == 0, f"CLI failed: {result.stderr}"

            # Validate PowerPoint content
            pptx_files = list(Path(temp_dir).glob("*.pptx"))
            assert len(pptx_files) > 0, "No PowerPoint file generated"

            prs = Presentation(str(pptx_files[0]))
            slide = prs.slides[0]

            # Extract all text from slide
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            slide_text_combined = " ".join(slide_text).lower()

            # Validate content placement
            title_found = "two content test" in slide_text_combined
            left_found = "left item 1" in slide_text_combined
            right_found = "right item 1" in slide_text_combined

            print(f"Slide text found: {slide_text}")
            print(f"Title found: {'✅' if title_found else '❌'}")
            print(f"Left content found: {'✅' if left_found else '❌'}")
            print(f"Right content found: {'✅' if right_found else '❌'}")

            # These SHOULD pass after implementation
            assert title_found, f"Title not found in slide. Found: {slide_text}"
            assert left_found, f"Left content not found in slide. Found: {slide_text}"
            assert right_found, f"Right content not found in slide. Found: {slide_text}"

    def test_semantic_field_names_four_columns(self):
        """Test Four Columns layout with semantic field names (content_col1, content_col2, etc.)"""

        # CORRECT expected canonical JSON structure with structured frontmatter
        json_data = {
            "slides": [
                {
                    "layout": "Four Columns",
                    "placeholders": {
                        "title": "Four Columns Test",
                        "content_col1": "Column 1 content",
                        "content_col2": "Column 2 content",
                        "content_col3": "Column 3 content",
                        "content_col4": "Column 4 content",
                    },
                    "content": [],
                }
            ]
        }

        with tempfile.TemporaryDirectory() as temp_dir:
            self.temp_dir = temp_dir

            json_file = Path(temp_dir) / "four_columns_test.json"
            with open(json_file, "w") as f:
                json.dump(json_data, f)

            env = os.environ.copy()
            env["DECK_TEMPLATE_FOLDER"] = self.template_env_path
            env["DECK_OUTPUT_FOLDER"] = temp_dir

            result = subprocess.run(
                [
                    "python",
                    str(self.project_root / "src" / "deckbuilder" / "cli.py"),
                    "create",
                    str(json_file),
                    "--output",
                    "four_columns_test",
                ],
                capture_output=True,
                text=True,
                env=env,
                cwd=temp_dir,
            )

            print("\nFOUR COLUMNS SEMANTIC FIELD TEST:")
            print(f"CLI return code: {result.returncode}")
            print(f"CLI output: {result.stdout}")

            assert result.returncode == 0, f"CLI failed: {result.stderr}"

            pptx_files = list(Path(temp_dir).glob("*.pptx"))
            assert len(pptx_files) > 0, "No PowerPoint file generated"

            prs = Presentation(str(pptx_files[0]))
            slide = prs.slides[0]

            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            slide_text_combined = " ".join(slide_text).lower()

            # Validate all column content
            title_found = "four columns test" in slide_text_combined
            col1_found = "column 1 content" in slide_text_combined
            col2_found = "column 2 content" in slide_text_combined
            col3_found = "column 3 content" in slide_text_combined
            col4_found = "column 4 content" in slide_text_combined

            print(f"Slide text found: {slide_text}")
            print(f"Title found: {'✅' if title_found else '❌'}")
            print(f"Column 1 found: {'✅' if col1_found else '❌'}")
            print(f"Column 2 found: {'✅' if col2_found else '❌'}")
            print(f"Column 3 found: {'✅' if col3_found else '❌'}")
            print(f"Column 4 found: {'✅' if col4_found else '❌'}")

            # These SHOULD pass after implementation
            assert title_found, f"Title not found. Found: {slide_text}"
            assert col1_found, f"Column 1 not found. Found: {slide_text}"
            assert col2_found, f"Column 2 not found. Found: {slide_text}"
            assert col3_found, f"Column 3 not found. Found: {slide_text}"
            assert col4_found, f"Column 4 not found. Found: {slide_text}"

    def test_all_formatting_types_preserved(self):
        """Test that ALL formatting types are preserved in JSON processing"""

        # CORRECT: All formatting types should work in structured frontmatter format
        json_data = {
            "slides": [
                {
                    "layout": "Title and Content",
                    "placeholders": {
                        "title": "**Formatting** Test with *All* Types",
                        "content": "## **Bold** Heading\nParagraph with *italic* and ___underline___ text.\n- ***Bold italic*** bullet\n- Regular bullet with **bold** text\n- Another with ___underlines___",
                    },
                    "content": [],
                },
                {
                    "layout": "Two Content",
                    "placeholders": {
                        "title": "Formatted **Two** Content",
                        "content_left": "**Bold** left item\n*Italic* left item",
                        "content_right": "___Underlined___ right item\n***Bold italic*** right item",
                    },
                    "content": [],
                },
            ]
        }

        with tempfile.TemporaryDirectory() as temp_dir:
            self.temp_dir = temp_dir

            json_file = Path(temp_dir) / "formatting_test.json"
            with open(json_file, "w") as f:
                json.dump(json_data, f)

            env = os.environ.copy()
            env["DECK_TEMPLATE_FOLDER"] = self.template_env_path
            env["DECK_OUTPUT_FOLDER"] = temp_dir

            result = subprocess.run(
                [
                    "python",
                    str(self.project_root / "src" / "deckbuilder" / "cli.py"),
                    "create",
                    str(json_file),
                    "--output",
                    "formatting_test",
                ],
                capture_output=True,
                text=True,
                env=env,
                cwd=temp_dir,
            )

            print("\nALL FORMATTING TYPES TEST:")
            print(f"CLI return code: {result.returncode}")
            print(f"CLI output: {result.stdout}")

            assert result.returncode == 0, f"CLI failed: {result.stderr}"

            pptx_files = list(Path(temp_dir).glob("*.pptx"))
            assert len(pptx_files) > 0, "No PowerPoint file generated"

            prs = Presentation(str(pptx_files[0]))

            # Check that content is present (formatting will be validated after implementation)
            all_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        all_text.append(shape.text.strip())

            combined_text = " ".join(all_text).lower()

            # Validate content is present (formatting preservation tested separately)
            formatting_content_found = "bold heading" in combined_text
            bullet_content_found = "bold italic bullet" in combined_text
            left_content_found = "bold left item" in combined_text
            right_content_found = "underlined right item" in combined_text

            print(f"All slide text: {all_text}")
            print(f"Formatting content found: {'✅' if formatting_content_found else '❌'}")
            print(f"Bullet content found: {'✅' if bullet_content_found else '❌'}")
            print(f"Left content found: {'✅' if left_content_found else '❌'}")
            print(f"Right content found: {'✅' if right_content_found else '❌'}")

            # Content should be present (formatting will be implemented)
            assert formatting_content_found, f"Formatting content not found. Found: {all_text}"
            assert bullet_content_found, f"Bullet content not found. Found: {all_text}"
            assert left_content_found, f"Left content not found. Found: {all_text}"
            assert right_content_found, f"Right content not found. Found: {all_text}"

    def test_no_markdown_conversion_errors(self):
        """Test that JSON processing doesn't produce markdown conversion errors"""

        # This should NOT produce "Missing required field" errors - canonical JSON format
        json_data = {
            "slides": [
                {
                    "layout": "Two Content",
                    "placeholders": {"title": "No Conversion Errors Test"},
                    "content": [
                        {
                            "type": "columns",
                            "columns": [
                                {"content": [{"type": "paragraph", "text": "Left content"}]},
                                {"content": [{"type": "paragraph", "text": "Right content"}]},
                            ],
                        }
                    ],
                }
            ]
        }

        with tempfile.TemporaryDirectory() as temp_dir:
            self.temp_dir = temp_dir

            json_file = Path(temp_dir) / "no_errors_test.json"
            with open(json_file, "w") as f:
                json.dump(json_data, f)

            env = os.environ.copy()
            env["DECK_TEMPLATE_FOLDER"] = self.template_env_path
            env["DECK_OUTPUT_FOLDER"] = temp_dir

            result = subprocess.run(
                [
                    "python",
                    str(self.project_root / "src" / "deckbuilder" / "cli.py"),
                    "create",
                    str(json_file),
                    "--output",
                    "no_errors_test",
                ],
                capture_output=True,
                text=True,
                env=env,
                cwd=temp_dir,
            )

            print("\nNO CONVERSION ERRORS TEST:")
            print(f"CLI return code: {result.returncode}")
            print(f"CLI output: {result.stdout}")
            print(f"CLI stderr: {result.stderr}")

            # Should NOT contain conversion error messages
            error_messages = [
                "Missing required field",
                "Error in structured frontmatter",
                "sections",
                "media",
            ]

            for error_msg in error_messages:
                assert error_msg not in result.stdout, f"Found conversion error: {error_msg}"
                assert error_msg not in result.stderr, f"Found conversion error: {error_msg}"

            # Should succeed
            assert result.returncode == 0, f"CLI failed: {result.stderr}"
