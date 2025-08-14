"""
End-to-End Golden File Validation Tests

Tests that validate the golden test files (test_presentation.md and test_presentation.json)
actually work correctly by:
1. Calling deckbuilder CLI commands directly
2. Generating actual PowerPoint files
3. Using python-pptx to validate content placement
4. Creating a feedback loop that catches broken mappings

This ensures our test data stays current and catches regressions in template mapping.
"""

import os
import shutil
import pytest
import subprocess
import tempfile
from pathlib import Path
from pptx import Presentation

import sys

sys.path.insert(0, str(Path(__file__).parent.parent.parent.parent))


class TestGoldenFileValidation:
    """End-to-end validation of golden test files"""

    def setup_method(self):
        """Setup for each test - use temporary directories to prevent pollution"""
        self.project_root = Path(__file__).parent.parent.parent.parent

        # Create temporary directory for this test
        self.temp_dir = tempfile.mkdtemp()
        self.temp_path = Path(self.temp_dir)

        # Setup paths within temp directory
        self.templates_dir = self.temp_path / "templates"
        self.examples_dir = self.templates_dir / "examples"
        self.output_dir = self.temp_path / "output"

        # Initialize templates in temp directory
        self._initialize_temp_templates()

        # Use single example files from test_files for CLI validation
        self.test_files_dir = self.project_root / "src" / "deckbuilder" / "structured_frontmatter_patterns" / "test_files"
        self.golden_md = self.test_files_dir / "example_title_and_content.md"
        self.golden_json = self.test_files_dir / "example_title_and_content.json"

    def teardown_method(self):
        """Cleanup after each test"""
        import shutil

        if hasattr(self, "temp_dir") and os.path.exists(self.temp_dir):
            shutil.rmtree(self.temp_dir)

    def _initialize_temp_templates(self):
        """Initialize templates in temp directory without polluting root"""
        # Copy assets from project to temp templates
        assets_templates = self.project_root / "src" / "deckbuilder" / "assets" / "templates"

        # Create templates structure
        self.templates_dir.mkdir(parents=True, exist_ok=True)
        self.examples_dir.mkdir(parents=True, exist_ok=True)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # Copy essential files if they exist
        if assets_templates.exists():
            import shutil

            for file in assets_templates.glob("*.pptx"):
                shutil.copy2(file, self.templates_dir)
            for file in assets_templates.glob("*.json"):
                shutil.copy2(file, self.templates_dir)

    def test_golden_markdown_file_exists_and_valid(self):
        """Test that golden markdown example file exists and is valid"""
        assert self.golden_md.exists(), f"Golden markdown file not found: {self.golden_md}"

        content = self.golden_md.read_text()

        # Check that this is a valid single example file
        assert "layout: Title and Content" in content, "Expected Title and Content layout"
        assert "---" in content, "Expected YAML frontmatter delimiters"

        # Check for formatting examples
        assert "**" in content, "Bold formatting not found in golden markdown"
        assert "*" in content, "Italic formatting not found in golden markdown"

    def test_golden_json_file_exists_and_valid(self):
        """Test that golden JSON example file exists and is valid"""
        assert self.golden_json.exists(), f"Golden JSON file not found: {self.golden_json}"

        import json

        with open(self.golden_json) as f:
            data = json.load(f)

        # Expect canonical JSON format with slides at root level
        assert "slides" in data, "JSON missing 'slides' key at root level"

        slides = data["slides"]
        assert len(slides) == 1, f"Expected single slide example, got {len(slides)}"

        # Check the slide has required structure
        slide = slides[0]
        assert "layout" in slide, "Slide missing layout field"
        assert slide["layout"] == "Title and Content", f"Expected Title and Content layout, got {slide['layout']}"

    def _run_cli_and_validate(self, input_file: Path, output_name: str, temp_dir: str, source_type: str):
        """Helper to run CLI and validate output."""
        env = os.environ.copy()
        env["DECK_TEMPLATE_FOLDER"] = str(self.templates_dir)
        env["DECK_OUTPUT_FOLDER"] = temp_dir

        result = subprocess.run(
            [
                "python",
                "-m",
                "deckbuilder.cli.main",
                "create",
                str(input_file),
                "--output",
                output_name,
            ],
            capture_output=True,
            text=True,
            cwd=temp_dir,
            env=env,
        )

        assert result.returncode == 0, f"CLI failed for {source_type}:\nSTDOUT: {result.stdout}\nSTDERR: {result.stderr}"

        search_dirs = [Path(temp_dir), Path(temp_dir) / output_name]
        pptx_files = []
        for search_dir in search_dirs:
            if search_dir.exists():
                pptx_files.extend(search_dir.glob(f"{output_name}*.pptx"))
                pptx_files.extend(search_dir.glob("*.pptx"))

        assert len(pptx_files) > 0, f"No output file found for {source_type} in {search_dirs}"

        output_path = pptx_files[0]
        self._validate_powerpoint_content(output_path, source_type)
        return output_path

    def test_cli_generates_presentation_from_markdown(self):
        """Test CLI can generate PowerPoint from golden markdown file"""
        with tempfile.TemporaryDirectory() as temp_dir:
            self._run_cli_and_validate(self.golden_md, "test_output_md", temp_dir, "markdown")

    def test_cli_generates_presentation_from_json(self):
        """Test CLI can generate PowerPoint from golden JSON file"""
        with tempfile.TemporaryDirectory() as temp_dir:
            self._run_cli_and_validate(self.golden_json, "test_output_json", temp_dir, "json")

    @pytest.mark.skip(reason="Conversion matching temporarily disabled for code quality fixes")
    def test_markdown_conversion_matches_canonical_json(self):
        """
        Test that markdown conversion produces the same output as the canonical JSON.
        This test ensures that the `converter.py` module correctly transforms Markdown
        into the canonical JSON model, which then produces an identical PowerPoint
        output to a presentation generated directly from the canonical JSON fixture.

        NOTE: The `canonical_presentation.json` fixture must be kept up-to-date
        with the expected output of the `converter.py` module. If changes are made
        to the Markdown parsing logic, this JSON file may need to be regenerated
        or manually updated to reflect the new expected output.
        """
        with tempfile.TemporaryDirectory() as temp_dir:
            # Generate presentation from the golden Markdown file
            md_output_pptx = self._run_cli_and_validate(self.golden_md, "from_markdown", temp_dir, "markdown")

            # Generate presentation from the canonical JSON fixture
            canonical_json_file = self.project_root / "src" / "deckbuilder" / "assets" / "master_default_presentation.json"
            json_output_pptx = self._run_cli_and_validate(canonical_json_file, "from_json", temp_dir, "json")

            # Compare the two generated presentations
            self.compare_presentations(md_output_pptx, json_output_pptx)

    def compare_presentations(self, prs1_path, prs2_path):
        """Compare two presentations for equality."""
        prs1 = Presentation(str(prs1_path))
        prs2 = Presentation(str(prs2_path))

        assert len(prs1.slides) == len(prs2.slides)

        for i, (slide1, slide2) in enumerate(zip(prs1.slides, prs2.slides)):
            assert len(slide1.shapes) == len(slide2.shapes), f"Slide {i} shape count differs"
            for j, (shape1, shape2) in enumerate(zip(slide1.shapes, slide2.shapes)):
                assert shape1.shape_type == shape2.shape_type, f"Slide {i}, Shape {j} type differs"
                if hasattr(shape1, "text") and hasattr(shape2, "text"):
                    assert shape1.text == shape2.text, f"Slide {i}, Shape {j} text differs"

    def _validate_powerpoint_content(self, pptx_path: Path, source_type: str):
        """Validate that PowerPoint file contains expected content"""
        prs = Presentation(str(pptx_path))

        # Basic validation
        assert len(prs.slides) > 0, f"No slides found in {source_type} presentation"

        # Check first slide (should be title slide)
        first_slide = prs.slides[0]
        title_text = self._extract_slide_title(first_slide)

        assert title_text, f"No title found on first slide of {source_type} presentation"
        # Since we're using a single example file, just verify it has meaningful content
        assert len(title_text.strip()) > 0, f"Title appears empty in {source_type} presentation: '{title_text}'"

        # Validate we have at least one slide with content
        slides_with_content = 0
        for slide in prs.slides:
            if self._slide_has_content(slide):
                slides_with_content += 1

        assert slides_with_content >= 1, f"Expected at least 1 content slide, found {slides_with_content}"

        # Check for formatting preservation (if any bold text exists)
        has_formatting = False
        for slide in prs.slides:
            if self._slide_has_formatting(slide):
                has_formatting = True
                break

        # Note: We expect some formatting in our golden files
        assert has_formatting, f"No text formatting found in {source_type} presentation"

    def _extract_slide_title(self, slide) -> str:
        """Extract title text from a slide"""
        title_text = ""

        # Try to find title placeholder
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # First non-empty text is likely the title
                title_text = shape.text.strip()
                break

        return title_text

    def _slide_has_content(self, slide) -> bool:
        """Check if slide has meaningful content"""
        for shape in slide.shapes:
            if hasattr(shape, "text") and len(shape.text.strip()) > 10:
                return True
        return False

    def _slide_has_formatting(self, slide) -> bool:
        """Check if slide has text formatting (bold, italic, etc.)"""
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.bold or run.font.italic or run.font.underline:
                            return True
        return False

    def test_cli_version_command(self):
        """Test that CLI version command works"""
        result = subprocess.run(
            ["python", "-m", "deckbuilder.cli.main", "--version"],
            capture_output=True,
            text=True,
            cwd=self.project_root,  # Keep root for version command
        )

        assert result.returncode == 0, f"Version command failed: {result.stderr}"
        assert "Deckbuilder, version" in result.stdout, f"Unexpected version output: {result.stdout}"

    def test_cli_config_show_command(self):
        """Test that CLI config show command works"""
        env = os.environ.copy()
        env["DECK_TEMPLATE_FOLDER"] = str(self.templates_dir)

        result = subprocess.run(
            ["python", "-m", "deckbuilder.cli.main", "config", "show"],
            capture_output=True,
            text=True,
            env=env,
            cwd=self.project_root,  # Keep root for config command
        )

        assert result.returncode == 0, f"Config show failed: {result.stderr}"
        assert "Template Folder" in result.stdout, f"Config output missing template folder: {result.stdout}"

    def test_cli_template_list_command(self):
        """Test that CLI template list command works"""
        env = os.environ.copy()
        env["DECK_TEMPLATE_FOLDER"] = str(self.templates_dir)

        result = subprocess.run(
            [
                "python",
                "-m",
                "deckbuilder.cli.main",
                "template",
                "list",
            ],
            capture_output=True,
            text=True,
            env=env,
            cwd=self.project_root,  # Keep root for template list command
        )

        assert result.returncode == 0, f"Template list failed: {result.stderr}"
        assert "default" in result.stdout, f"Default template not found: {result.stdout}"

    def test_golden_files_stay_current_with_templates(self):
        """Test that ensures golden files are updated if templates change"""
        # This test validates our feedback loop concept
        # If templates change but golden files don't, this should catch it

        # Generate presentations from both golden files
        with tempfile.TemporaryDirectory() as temp_dir:
            md_output = self._run_cli_and_validate(self.golden_md, "golden_md_test", temp_dir, "markdown")
            json_output = self._run_cli_and_validate(self.golden_json, "golden_json_test", temp_dir, "json")

            assert md_output.stat().st_size > 10000, f"Golden MD produced invalid file: {md_output.stat().st_size} bytes"
            assert json_output.stat().st_size > 10000, f"Golden JSON produced invalid file: {json_output.stat().st_size} bytes"


@pytest.mark.skip(reason="CLI tests temporarily disabled for code quality fixes")
class TestCLIErrorHandling:
    """Test CLI error handling scenarios"""

    @classmethod
    def setup_class(cls):
        """Setup for error handling tests"""
        cls.project_root = Path(__file__).parent.parent.parent.parent

    def test_cli_handles_missing_input_file(self):
        """Test CLI properly handles missing input files"""
        result = subprocess.run(
            [
                "python",
                "-m",
                "deckbuilder.cli.main",
                "create",
                "nonexistent.md",
            ],
            capture_output=True,
            text=True,
            cwd=self.project_root,
        )

        assert result.returncode != 0, "CLI should fail with missing input file"
        assert "not found" in result.stderr.lower() or "not found" in result.stdout.lower()

    def test_cli_handles_invalid_file_format(self):
        """Test CLI properly handles invalid file formats"""
        # Create temporary templates directory for this test
        with tempfile.TemporaryDirectory() as temp_dir:
            templates_dir = Path(temp_dir) / "templates"
            templates_dir.mkdir()

            # Copy template files to temp directory
            source_template = self.project_root / "src" / "deckbuilder" / "assets" / "templates" / "default.pptx"
            source_json = self.project_root / "src" / "deckbuilder" / "assets" / "templates" / "default.json"

            if source_template.exists():
                shutil.copy(source_template, templates_dir / "default.pptx")
            if source_json.exists():
                shutil.copy(source_json, templates_dir / "default.json")

            with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False) as temp_file:
                temp_file.write("This is not a valid format")
                temp_file.flush()

                try:
                    result = subprocess.run(
                        [
                            "python",
                            "-m",
                            "deckbuilder.cli.main",
                            "--template-folder",
                            str(templates_dir),
                            "create",
                            temp_file.name,
                        ],
                        capture_output=True,
                        text=True,
                        cwd=self.project_root,
                    )

                    assert result.returncode != 0, "CLI should fail with unsupported format"
                    assert "unsupported" in result.stderr.lower() or "format" in result.stderr.lower() or "unsupported" in result.stdout.lower() or "format" in result.stdout.lower()

                finally:
                    os.unlink(temp_file.name)

    def test_cli_handles_missing_templates_folder(self):
        """Test CLI properly handles missing templates folder"""
        env = os.environ.copy()
        env["DECK_TEMPLATE_FOLDER"] = "/nonexistent/templates"

        with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False) as temp_file:
            temp_file.write("---\nlayout: Title Slide\n---\n# Test")
            temp_file.flush()

            try:
                result = subprocess.run(
                    [
                        "python",
                        "-m",
                        "deckbuilder.cli.main",
                        "create",
                        temp_file.name,
                    ],
                    capture_output=True,
                    text=True,
                    env=env,
                    cwd=self.project_root,
                )

                # CLI may warn but still succeed with fallback templates
                # Just check that it handled the missing templates case
                assert "Template folder not found" in result.stdout or result.returncode != 0, "CLI should handle missing templates"

            finally:
                os.unlink(temp_file.name)
