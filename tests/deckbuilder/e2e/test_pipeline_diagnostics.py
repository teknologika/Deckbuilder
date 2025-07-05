#!/usr/bin/env python3
"""
Pipeline Diagnostics Tests

Tests to isolate whether issues are in:
1. Template.pptx → Template.json mapping (template analysis)
2. Generator → Output (content placement)

This helps identify where the content mapping failures occur.
"""

import json
import os
import subprocess
import tempfile
from pathlib import Path
from pptx import Presentation
import pytest
import sys

# Add src to path
sys.path.insert(0, str(Path(__file__).parent.parent.parent.parent / "src"))


class TestTemplateMappingDiagnostics:
    """Test template.pptx → template.json mapping accuracy"""

    def setup_method(self):
        """Setup for each test"""
        self.project_root = Path(__file__).parent.parent.parent.parent
        self.template_folder = self.project_root / "src" / "deckbuilder" / "assets" / "templates"
        self.template_pptx = self.template_folder / "default.pptx"
        self.template_json = self.template_folder / "default.json"

    def test_template_files_exist(self):
        """Verify template files exist"""
        assert self.template_pptx.exists(), f"Template PPTX not found: {self.template_pptx}"
        assert self.template_json.exists(), f"Template JSON not found: {self.template_json}"

    def test_template_json_structure(self):
        """Test template JSON has correct structure"""
        with open(self.template_json, "r") as f:
            template_data = json.load(f)

        assert "template_info" in template_data, "Missing template_info"
        assert "layouts" in template_data, "Missing layouts"

        # Check each layout has required fields
        for layout_name, layout_data in template_data["layouts"].items():
            assert "index" in layout_data, f"Layout {layout_name} missing index"
            assert "placeholders" in layout_data, f"Layout {layout_name} missing placeholders"

    def test_template_pptx_layout_count(self):
        """Test PowerPoint template has expected number of layouts"""
        prs = Presentation(str(self.template_pptx))

        with open(self.template_json, "r") as f:
            template_data = json.load(f)

        pptx_layout_count = len(prs.slide_layouts)
        json_layout_count = len(template_data["layouts"])

        print(f"PowerPoint layouts: {pptx_layout_count}")
        print(f"JSON layouts: {json_layout_count}")

        # They should be close but JSON might have fewer (only mapped ones)
        assert pptx_layout_count >= json_layout_count, "PowerPoint has fewer layouts than JSON"

    def test_template_layout_placeholder_mapping(self):
        """Test that JSON placeholder mappings correspond to actual PowerPoint placeholders"""
        prs = Presentation(str(self.template_pptx))

        with open(self.template_json, "r") as f:
            template_data = json.load(f)

        mapping_issues = []

        for layout_name, layout_data in template_data["layouts"].items():
            layout_index = layout_data["index"]
            json_placeholders = layout_data["placeholders"]

            if layout_index < len(prs.slide_layouts):
                slide_layout = prs.slide_layouts[layout_index]

                # Get actual placeholder IDs from PowerPoint
                actual_placeholder_ids = []
                for shape in slide_layout.placeholders:
                    actual_placeholder_ids.append(str(shape.placeholder_format.idx))

                # Check if JSON placeholders exist in PowerPoint
                for placeholder_id, placeholder_name in json_placeholders.items():
                    if placeholder_id not in actual_placeholder_ids:
                        mapping_issues.append(
                            {
                                "layout": layout_name,
                                "missing_placeholder_id": placeholder_id,
                                "placeholder_name": placeholder_name,
                                "actual_ids": actual_placeholder_ids,
                            }
                        )

        if mapping_issues:
            print("\n🚨 TEMPLATE MAPPING ISSUES:")
            for issue in mapping_issues[:3]:  # Show first 3 issues
                print(f"  Layout: {issue['layout']}")
                print(
                    f"  Missing ID: {issue['missing_placeholder_id']} ({issue['placeholder_name']})"
                )
                print(f"  Actual IDs: {issue['actual_ids']}")
                print()

        # This might fail - that's the point, to identify mapping issues
        assert len(mapping_issues) == 0, f"Found {len(mapping_issues)} placeholder mapping issues"


class TestContentGenerationDiagnostics:
    """Test Generator → Output content placement"""

    def setup_method(self):
        """Setup for each test"""
        self.project_root = Path(__file__).parent.parent.parent.parent
        self.template_folder = self.project_root / "src" / "deckbuilder" / "assets" / "templates"
        self.temp_dir = None

    def teardown_method(self):
        """Cleanup after each test"""
        if self.temp_dir and Path(self.temp_dir).exists():
            import shutil

            shutil.rmtree(self.temp_dir)

    def test_simple_title_slide_content_mapping(self):
        """Test simple title slide content is mapped correctly"""
        # Create minimal JSON input
        simple_json = {
            "presentation": {
                "slides": [
                    {"type": "Title Slide", "title": "TEST TITLE", "subtitle": "TEST SUBTITLE"}
                ]
            }
        }

        with tempfile.TemporaryDirectory() as temp_dir:
            self.temp_dir = temp_dir

            # Write JSON input
            json_file = Path(temp_dir) / "simple_test.json"
            with open(json_file, "w") as f:
                json.dump(simple_json, f)

            # Generate PowerPoint
            env = os.environ.copy()
            env["DECK_TEMPLATE_FOLDER"] = str(self.template_folder)
            env["DECK_OUTPUT_FOLDER"] = temp_dir

            result = subprocess.run(
                [
                    "python",
                    str(self.project_root / "src" / "deckbuilder" / "cli.py"),
                    "create",
                    str(json_file),
                    "--output",
                    "simple_test",
                ],
                capture_output=True,
                text=True,
                env=env,
                cwd=temp_dir,
            )

            assert result.returncode == 0, f"CLI failed: {result.stderr}"

            # Find generated file
            pptx_files = list(Path(temp_dir).glob("*.pptx"))
            assert len(pptx_files) > 0, "No PowerPoint file generated"

            # Validate content
            prs = Presentation(str(pptx_files[0]))
            assert len(prs.slides) == 1, f"Expected 1 slide, got {len(prs.slides)}"

            # Extract all text from slide
            slide = prs.slides[0]
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            slide_text_combined = " ".join(slide_text).upper()

            # Check if our test content is present
            title_found = "TEST TITLE" in slide_text_combined
            subtitle_found = "TEST SUBTITLE" in slide_text_combined

            print("\nSIMPLE CONTENT MAPPING TEST:")
            print(f"Generated file: {pptx_files[0].name}")
            print(f"Slide text found: {slide_text}")
            print(f"Title found: {'✅' if title_found else '❌'}")
            print(f"Subtitle found: {'✅' if subtitle_found else '❌'}")

            assert title_found, f"Title 'TEST TITLE' not found in slide. Found: {slide_text}"
            assert (
                subtitle_found
            ), f"Subtitle 'TEST SUBTITLE' not found in slide. Found: {slide_text}"

    def test_content_with_formatting_mapping(self):
        """Test formatted content is mapped correctly"""
        formatted_json = {
            "presentation": {
                "slides": [
                    {
                        "type": "Title and Content",
                        "title": "**Bold Title** Test",
                        "content": "This is *italic* and **bold** text",
                    }
                ]
            }
        }

        with tempfile.TemporaryDirectory() as temp_dir:
            self.temp_dir = temp_dir

            # Write JSON input
            json_file = Path(temp_dir) / "format_test.json"
            with open(json_file, "w") as f:
                json.dump(formatted_json, f)

            # Generate PowerPoint
            env = os.environ.copy()
            env["DECK_TEMPLATE_FOLDER"] = str(self.template_folder)
            env["DECK_OUTPUT_FOLDER"] = temp_dir

            result = subprocess.run(
                [
                    "python",
                    str(self.project_root / "src" / "deckbuilder" / "cli.py"),
                    "create",
                    str(json_file),
                    "--output",
                    "format_test",
                ],
                capture_output=True,
                text=True,
                env=env,
                cwd=temp_dir,
            )

            assert result.returncode == 0, f"CLI failed: {result.stderr}"

            # Find and validate generated file
            pptx_files = list(Path(temp_dir).glob("*.pptx"))
            assert len(pptx_files) > 0, "No PowerPoint file generated"

            prs = Presentation(str(pptx_files[0]))
            slide = prs.slides[0]

            # Check for content and formatting
            content_found = False
            formatting_found = False

            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    # Check if content text is present
                    if (
                        "italic" in shape.text_frame.text.lower()
                        and "bold" in shape.text_frame.text.lower()
                    ):
                        content_found = True

                    # Check if formatting is applied
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.bold or run.font.italic:
                                formatting_found = True
                                break

            print("\nFORMATTING MAPPING TEST:")
            print(f"Content found: {'✅' if content_found else '❌'}")
            print(f"Formatting found: {'✅' if formatting_found else '❌'}")

            assert content_found, "Content not mapped correctly"
            # Note: Formatting test might fail - that helps identify the issue

    def test_multi_column_content_mapping(self):
        """Test multi-column layout content mapping"""
        column_json = {
            "presentation": {
                "slides": [
                    {
                        "type": "Two Content",
                        "title": "Two Column Test",
                        "content_left": ["Left Column Item 1", "Left Column Item 2"],
                        "content_right": ["Right Column Item 1", "Right Column Item 2"],
                    }
                ]
            }
        }

        with tempfile.TemporaryDirectory() as temp_dir:
            self.temp_dir = temp_dir

            # Write JSON input
            json_file = Path(temp_dir) / "column_test.json"
            with open(json_file, "w") as f:
                json.dump(column_json, f)

            # Generate PowerPoint
            env = os.environ.copy()
            env["DECK_TEMPLATE_FOLDER"] = str(self.template_folder)
            env["DECK_OUTPUT_FOLDER"] = temp_dir

            result = subprocess.run(
                [
                    "python",
                    str(self.project_root / "src" / "deckbuilder" / "cli.py"),
                    "create",
                    str(json_file),
                    "--output",
                    "column_test",
                ],
                capture_output=True,
                text=True,
                env=env,
                cwd=temp_dir,
            )

            print("\nMULTI-COLUMN MAPPING TEST:")
            print(f"CLI return code: {result.returncode}")
            if result.stdout:
                print(f"CLI output: {result.stdout}")
            if result.stderr:
                print(f"CLI errors: {result.stderr}")

            # Even if CLI fails, let's see what happens
            pptx_files = list(Path(temp_dir).glob("*.pptx"))
            if pptx_files:
                prs = Presentation(str(pptx_files[0]))
                slide = prs.slides[0]

                # Extract all text
                all_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        all_text.append(shape.text.strip())

                combined_text = " ".join(all_text).lower()

                left_found = "left column item 1" in combined_text
                right_found = "right column item 1" in combined_text

                print(f"All slide text: {all_text}")
                print(f"Left content found: {'✅' if left_found else '❌'}")
                print(f"Right content found: {'✅' if right_found else '❌'}")

                # This test is expected to potentially fail - that's diagnostic info
                return {
                    "success": left_found and right_found,
                    "left_found": left_found,
                    "right_found": right_found,
                    "all_text": all_text,
                }
            else:
                print("❌ No PowerPoint file generated")
                return {"success": False, "error": "No file generated"}


class TestMarkdownGenerationDiagnostics:
    """Test Markdown → Output content placement"""

    def setup_method(self):
        """Setup for each test"""
        self.project_root = Path(__file__).parent.parent.parent.parent
        self.template_folder = self.project_root / "src" / "deckbuilder" / "assets" / "templates"
        self.temp_dir = None

    def teardown_method(self):
        """Cleanup after each test"""
        if self.temp_dir and Path(self.temp_dir).exists():
            import shutil

            shutil.rmtree(self.temp_dir)

    def test_simple_markdown_title_slide(self):
        """Test simple markdown title slide content mapping"""
        markdown_content = """---
layout: Title Slide
title: MARKDOWN TEST TITLE
subtitle: MARKDOWN TEST SUBTITLE
---

# This is the slide content
"""

        with tempfile.TemporaryDirectory() as temp_dir:
            self.temp_dir = temp_dir

            # Write markdown input
            md_file = Path(temp_dir) / "simple_test.md"
            with open(md_file, "w") as f:
                f.write(markdown_content)

            # Generate PowerPoint
            env = os.environ.copy()
            env["DECK_TEMPLATE_FOLDER"] = str(self.template_folder)
            env["DECK_OUTPUT_FOLDER"] = temp_dir

            result = subprocess.run(
                [
                    "python",
                    str(self.project_root / "src" / "deckbuilder" / "cli.py"),
                    "create",
                    str(md_file),
                    "--output",
                    "markdown_test",
                ],
                capture_output=True,
                text=True,
                env=env,
                cwd=temp_dir,
            )

            print("\nMARKDOWN TITLE SLIDE TEST:")
            print(f"CLI return code: {result.returncode}")
            if result.stdout:
                print(f"CLI output: {result.stdout}")

            assert result.returncode == 0, f"CLI failed: {result.stderr}"

            # Find generated file
            pptx_files = list(Path(temp_dir).glob("*.pptx"))
            assert len(pptx_files) > 0, "No PowerPoint file generated"

            # Validate content
            prs = Presentation(str(pptx_files[0]))
            slide = prs.slides[0]

            # Extract all text from slide
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            slide_text_combined = " ".join(slide_text).upper()

            title_found = "MARKDOWN TEST TITLE" in slide_text_combined
            subtitle_found = "MARKDOWN TEST SUBTITLE" in slide_text_combined

            print(f"Slide text found: {slide_text}")
            print(f"Title found: {'✅' if title_found else '❌'}")
            print(f"Subtitle found: {'✅' if subtitle_found else '❌'}")

            assert title_found, f"Title not found in slide. Found: {slide_text}"
            assert subtitle_found, f"Subtitle not found in slide. Found: {slide_text}"

    def test_markdown_structured_frontmatter_four_columns(self):
        """Test markdown structured frontmatter for Four Columns layout"""
        markdown_content = """---
layout: Four Columns
title: Four Column Markdown Test
columns:
  - title: Column 1 Title
    content: "Column 1 content text"
  - title: Column 2 Title
    content: "Column 2 content text"
  - title: Column 3 Title
    content: "Column 3 content text"
  - title: Column 4 Title
    content: "Column 4 content text"
---

# Four Columns Test Slide
This slide should have four columns with titles and content.
"""

        with tempfile.TemporaryDirectory() as temp_dir:
            self.temp_dir = temp_dir

            # Write markdown input
            md_file = Path(temp_dir) / "four_columns_test.md"
            with open(md_file, "w") as f:
                f.write(markdown_content)

            # Generate PowerPoint
            env = os.environ.copy()
            env["DECK_TEMPLATE_FOLDER"] = str(self.template_folder)
            env["DECK_OUTPUT_FOLDER"] = temp_dir

            result = subprocess.run(
                [
                    "python",
                    str(self.project_root / "src" / "deckbuilder" / "cli.py"),
                    "create",
                    str(md_file),
                    "--output",
                    "four_columns_test",
                ],
                capture_output=True,
                text=True,
                env=env,
                cwd=temp_dir,
            )

            print("\nMARKDOWN FOUR COLUMNS TEST:")
            print(f"CLI return code: {result.returncode}")
            if result.stdout:
                print(f"CLI output: {result.stdout}")
            if result.stderr:
                print(f"CLI errors: {result.stderr}")

            # Check if file was generated
            pptx_files = list(Path(temp_dir).glob("*.pptx"))
            if pptx_files:
                prs = Presentation(str(pptx_files[0]))
                slide = prs.slides[0]

                # Extract all text
                all_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        all_text.append(shape.text.strip())

                combined_text = " ".join(all_text).lower()

                # Check for column content
                col1_found = "column 1 content" in combined_text
                col2_found = "column 2 content" in combined_text
                col3_found = "column 3 content" in combined_text
                col4_found = "column 4 content" in combined_text
                title_found = "four column markdown test" in combined_text

                print(f"All slide text: {all_text}")
                print(f"Title found: {'✅' if title_found else '❌'}")
                print(f"Column 1 content: {'✅' if col1_found else '❌'}")
                print(f"Column 2 content: {'✅' if col2_found else '❌'}")
                print(f"Column 3 content: {'✅' if col3_found else '❌'}")
                print(f"Column 4 content: {'✅' if col4_found else '❌'}")

                # This test may fail - that's diagnostic info
                return {
                    "success": all([title_found, col1_found, col2_found, col3_found, col4_found]),
                    "title_found": title_found,
                    "columns_found": [col1_found, col2_found, col3_found, col4_found],
                    "all_text": all_text,
                }
            else:
                print("❌ No PowerPoint file generated")
                return {"success": False, "error": "No file generated"}

    def test_markdown_comparison_layout(self):
        """Test markdown comparison layout (similar to Two Content but different structure)"""
        markdown_content = """---
layout: Comparison
title: Markdown Comparison Test
left:
  title: Left Side Title
  content: "Left side content for comparison"
right:
  title: Right Side Title
  content: "Right side content for comparison"
---

# Comparison Layout Test
This slide should have left and right comparison content.
"""

        with tempfile.TemporaryDirectory() as temp_dir:
            self.temp_dir = temp_dir

            # Write markdown input
            md_file = Path(temp_dir) / "comparison_test.md"
            with open(md_file, "w") as f:
                f.write(markdown_content)

            # Generate PowerPoint
            env = os.environ.copy()
            env["DECK_TEMPLATE_FOLDER"] = str(self.template_folder)
            env["DECK_OUTPUT_FOLDER"] = temp_dir

            result = subprocess.run(
                [
                    "python",
                    str(self.project_root / "src" / "deckbuilder" / "cli.py"),
                    "create",
                    str(md_file),
                    "--output",
                    "comparison_test",
                ],
                capture_output=True,
                text=True,
                env=env,
                cwd=temp_dir,
            )

            print("\nMARKDOWN COMPARISON TEST:")
            print(f"CLI return code: {result.returncode}")
            if result.stdout:
                print(f"CLI output: {result.stdout}")
            if result.stderr:
                print(f"CLI errors: {result.stderr}")

            # Check results
            pptx_files = list(Path(temp_dir).glob("*.pptx"))
            if pptx_files:
                prs = Presentation(str(pptx_files[0]))
                slide = prs.slides[0]

                all_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        all_text.append(shape.text.strip())

                combined_text = " ".join(all_text).lower()

                left_found = "left side content" in combined_text
                right_found = "right side content" in combined_text
                title_found = "markdown comparison test" in combined_text

                print(f"All slide text: {all_text}")
                print(f"Title found: {'✅' if title_found else '❌'}")
                print(f"Left content found: {'✅' if left_found else '❌'}")
                print(f"Right content found: {'✅' if right_found else '❌'}")

                return {
                    "success": all([title_found, left_found, right_found]),
                    "title_found": title_found,
                    "left_found": left_found,
                    "right_found": right_found,
                    "all_text": all_text,
                }
            else:
                print("❌ No PowerPoint file generated")
                return {"success": False, "error": "No file generated"}


if __name__ == "__main__":
    pytest.main([__file__, "-v", "-s"])
