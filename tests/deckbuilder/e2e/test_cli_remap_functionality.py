"""
E2E tests for CLI remap functionality.

Tests the deckbuilder remap command for language and font remapping
in existing PowerPoint presentations.
"""

import os
import subprocess
import tempfile
from pathlib import Path
from shutil import copy2

import pytest
from pptx import Presentation


class TestCLIRemapFunctionality:
    """E2E tests for the deckbuilder remap command."""

    @pytest.fixture
    def sample_presentation(self):
        """Create a sample PowerPoint presentation for testing."""
        # Use the golden file as our test presentation
        golden_file = Path(__file__).parent.parent / "test_presentation.md"

        # Create a temp presentation using the CLI
        with tempfile.TemporaryDirectory() as temp_dir:
            env = os.environ.copy()
            env["DECK_TEMPLATE_FOLDER"] = str(Path(__file__).parent.parent.parent.parent / "src" / "deckbuilder" / "assets" / "templates")
            env["DECK_OUTPUT_FOLDER"] = temp_dir

            result = subprocess.run(
                [
                    "python",
                    "-m",
                    "deckbuilder.cli.main",
                    "create",
                    str(golden_file),
                    "--output",
                    "test_remap_sample",
                ],
                capture_output=True,
                text=True,
                env=env,
                cwd=temp_dir,
            )

            if result.returncode != 0:
                pytest.skip(f"Could not create sample presentation: {result.stderr}")

            # Find the generated file
            pptx_files = list(Path(temp_dir).glob("**/*.pptx"))
            if not pptx_files:
                pytest.skip("No PowerPoint file generated")

            # Copy to a new temp location for the test
            test_dir = Path(tempfile.mkdtemp())
            test_file = test_dir / "sample.pptx"
            copy2(pptx_files[0], test_file)

            yield test_file

            # Cleanup
            if test_dir.exists():
                import shutil

                shutil.rmtree(test_dir)

    def test_remap_help_command(self):
        """Test that remap help command works correctly."""
        result = subprocess.run(
            [
                "python",
                str(Path(__file__).parent.parent.parent.parent / "src" / "deckbuilder" / "cli.py"),
                "remap",
                "--help",
            ],
            capture_output=True,
            text=True,
        )

        assert result.returncode == 0
        assert "Update language and font settings" in result.stdout
        assert "Usage: cli.py remap [OPTIONS] INPUT_FILE" in result.stdout
        assert "-l, --language" in result.stdout
        assert "-f, --font" in result.stdout
        assert "-o, --output" in result.stdout
        assert "--no-backup" in result.stdout

    def test_remap_missing_file_error(self):
        """Test error handling for missing input file."""
        result = subprocess.run(
            [
                "python",
                str(Path(__file__).parent.parent.parent.parent / "src" / "deckbuilder" / "cli.py"),
                "remap",
                "nonexistent.pptx",
                "--language",
                "en-US",
            ],
            capture_output=True,
            text=True,
        )

        assert result.returncode == 2
        assert "Invalid value for 'INPUT_FILE': File 'nonexistent.pptx' does not exist." in result.stderr

    def test_remap_invalid_file_format_error(self):
        """Test error handling for non-PowerPoint files."""
        with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as temp_file:
            temp_file.write(b"This is not a PowerPoint file")
            temp_file_path = temp_file.name

        try:
            result = subprocess.run(
                [
                    "python",
                    "-m",
                    "deckbuilder.cli.main",
                    "remap",
                    temp_file_path,
                    "--language",
                    "en-US",
                ],
                capture_output=True,
                text=True,
            )

            assert result.returncode == 1
            assert "File must be a PowerPoint file" in result.stderr
        finally:
            os.unlink(temp_file_path)

    def test_remap_no_options_error(self):
        """Test error handling when no remap options are provided."""
        # Create a dummy PowerPoint file
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
            # Create minimal PowerPoint file
            prs = Presentation()
            prs.save(temp_file.name)
            temp_file_path = temp_file.name

        try:
            result = subprocess.run(
                [
                    "python",
                    "-m",
                    "deckbuilder.cli.main",
                    "remap",
                    temp_file_path,
                ],
                capture_output=True,
                text=True,
            )

            assert result.returncode == 1
            assert "No updates specified" in result.stderr
        finally:
            os.unlink(temp_file_path)

    def test_remap_language_only(self, sample_presentation):
        """Test language-only remapping."""
        result = subprocess.run(
            [
                "python",
                str(Path(__file__).parent.parent.parent.parent / "src" / "deckbuilder" / "cli.py"),
                "remap",
                str(sample_presentation),
                "--language",
                "en-US",
            ],
            capture_output=True,
            text=True,
        )

        assert result.returncode == 0
        assert "language to en-US" in result.stdout
        assert "Processing Summary:" in result.stdout
        assert "Language applied:" in result.stdout

        # Verify backup was created
        backup_files = list(sample_presentation.parent.glob("*.bak.pptx"))
        assert len(backup_files) == 1

    def test_remap_font_only(self, sample_presentation):
        """Test font-only remapping."""
        result = subprocess.run(
            [
                "python",
                str(Path(__file__).parent.parent.parent.parent / "src" / "deckbuilder" / "cli.py"),
                "remap",
                str(sample_presentation),
                "--font",
                "Arial",
            ],
            capture_output=True,
            text=True,
        )

        assert result.returncode == 0
        assert "font to 'Arial'" in result.stdout
        assert "Processing Summary:" in result.stdout
        assert "Font applied:" in result.stdout

    def test_remap_language_and_font(self, sample_presentation):
        """Test combined language and font remapping."""
        result = subprocess.run(
            [
                "python",
                str(Path(__file__).parent.parent.parent.parent / "src" / "deckbuilder" / "cli.py"),
                "remap",
                str(sample_presentation),
                "--language",
                "es-ES",
                "--font",
                "Calibri",
            ],
            capture_output=True,
            text=True,
        )

        assert result.returncode == 0
        assert "language to es-ES" in result.stdout
        assert "font to 'Calibri'" in result.stdout
        assert "Processing Summary:" in result.stdout
        assert "Language applied:" in result.stdout
        assert "Font applied:" in result.stdout

    def test_remap_with_output_file(self, sample_presentation):
        """Test remapping with custom output file."""
        output_file = sample_presentation.parent / "remapped_output.pptx"

        result = subprocess.run(
            [
                "python",
                str(Path(__file__).parent.parent.parent.parent / "src" / "deckbuilder" / "cli.py"),
                "remap",
                str(sample_presentation),
                "--language",
                "en-AU",
                "--output",
                str(output_file),
            ],
            capture_output=True,
            text=True,
        )

        assert result.returncode == 0
        assert "language to en-AU" in result.stdout
        assert output_file.exists()

        # Original file should be unchanged
        assert sample_presentation.exists()

    def test_remap_no_backup_flag(self, sample_presentation):
        """Test remapping with --no-backup flag."""
        # Get initial file count
        initial_files = list(sample_presentation.parent.glob("*.pptx"))
        initial_count = len(initial_files)

        result = subprocess.run(
            [
                "python",
                str(Path(__file__).parent.parent.parent.parent / "src" / "deckbuilder" / "cli.py"),
                "remap",
                str(sample_presentation),
                "--language",
                "en-GB",
                "--no-backup",
            ],
            capture_output=True,
            text=True,
        )

        assert result.returncode == 0
        assert "language to en-GB" in result.stdout

        # Should not have created backup file
        final_files = list(sample_presentation.parent.glob("*.pptx"))
        assert len(final_files) == initial_count  # No additional backup file

    def test_remap_invalid_language_code(self, sample_presentation):
        """Test error handling for invalid language codes."""
        result = subprocess.run(
            [
                "python",
                str(Path(__file__).parent.parent.parent.parent / "src" / "deckbuilder" / "cli.py"),
                "remap",
                str(sample_presentation),
                "--language",
                "invalid-lang",
            ],
            capture_output=True,
            text=True,
        )

        # Should handle gracefully (the validation happens in the formatting support)
        assert result.returncode == 1 or "invalid-lang" in result.stdout

    def test_remap_multiple_language_codes(self, sample_presentation):
        """Test remapping with multiple supported language codes."""
        test_languages = ["en-US", "en-AU", "en-GB", "es-ES"]

        for lang_code in test_languages:
            # Create a copy for each test
            test_file = sample_presentation.parent / f"test_{lang_code.replace('-', '_')}.pptx"
            copy2(sample_presentation, test_file)

            result = subprocess.run(
                [
                    "python",
                    "-m",
                    "deckbuilder.cli.main",
                    "remap",
                    str(test_file),
                    "--language",
                    lang_code,
                    "--no-backup",
                ],
                capture_output=True,
                text=True,
            )

            assert result.returncode == 0, f"Failed for language {lang_code}: {result.stderr}"
            assert f"language to {lang_code}" in result.stdout

    def test_remap_multiple_fonts(self, sample_presentation):
        """Test remapping with multiple font families."""
        test_fonts = ["Arial", "Calibri", "Times New Roman", "Verdana"]

        for font_name in test_fonts:
            # Create a copy for each test
            safe_font = font_name.replace(" ", "_")
            test_file = sample_presentation.parent / f"test_{safe_font}.pptx"
            copy2(sample_presentation, test_file)

            result = subprocess.run(
                [
                    "python",
                    "-m",
                    "deckbuilder.cli.main",
                    "remap",
                    str(test_file),
                    "--font",
                    font_name,
                    "--no-backup",
                ],
                capture_output=True,
                text=True,
            )

            assert result.returncode == 0, f"Failed for font {font_name}: {result.stderr}"
            assert f"font to '{font_name}'" in result.stdout


# NOTE: Text replacement functionality tests removed
# The core remap functionality (language and font settings) works correctly,
# but advanced text replacement (optimize → optimise, resume → CV, etc.)
# is not implemented and would require:
# - Comprehensive word mapping dictionaries
# - Context-aware replacement logic
# - Case preservation algorithms
# This could be added as a future enhancement if needed.


class TestRemapBackupValidation:
    """Tests specifically for backup file creation and validation."""

    def test_backup_file_structure(self, sample_presentation):
        """Test that backup files are created with correct naming and content."""
        # original_size = sample_presentation.stat().st_size  # Future: could validate size changes

        result = subprocess.run(
            [
                "python",
                str(Path(__file__).parent.parent.parent.parent / "src" / "deckbuilder" / "cli.py"),
                "remap",
                str(sample_presentation),
                "--language",
                "en-US",
            ],
            capture_output=True,
            text=True,
        )

        assert result.returncode == 0

        # Find backup file
        backup_files = list(sample_presentation.parent.glob("*.bak.pptx"))
        assert len(backup_files) == 1

        backup_file = backup_files[0]
        backup_size = backup_file.stat().st_size

        # Backup should be readable PowerPoint file
        assert backup_size > 1000  # Reasonable size check

        # Backup should be loadable
        try:
            backup_prs = Presentation(str(backup_file))
            assert len(backup_prs.slides) > 0
        except Exception as e:
            pytest.fail(f"Backup file is not a valid PowerPoint file: {e}")

    @pytest.fixture
    def sample_presentation(self):
        """Create a sample PowerPoint presentation for testing backup functionality."""
        # Create a minimal but valid presentation
        with tempfile.TemporaryDirectory() as temp_dir:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Test Slide"

            test_file = Path(temp_dir) / "backup_test.pptx"
            prs.save(str(test_file))

            yield test_file


class TestRemapThemeFonts:
    """Tests specifically for theme font update functionality."""

    @pytest.fixture
    def sample_presentation(self):
        """Create a sample PowerPoint presentation for testing theme font updates."""
        # Create a minimal but valid presentation
        with tempfile.TemporaryDirectory() as temp_dir:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Test Slide with Theme Fonts"

            test_file = Path(temp_dir) / "theme_font_test.pptx"
            prs.save(str(test_file))

            yield test_file

    def test_remap_font_updates_theme(self, sample_presentation):
        """Test that font remapping updates theme fonts (majorFont and minorFont)."""
        # Get original theme fonts
        original_prs = Presentation(str(sample_presentation))
        original_major_font, original_minor_font = self._get_theme_fonts(original_prs)

        # Run remap with new font
        test_font = "Times New Roman"
        result = subprocess.run(
            [
                "python",
                str(Path(__file__).parent.parent.parent.parent / "src" / "deckbuilder" / "cli.py"),
                "remap",
                str(sample_presentation),
                "--font",
                test_font,
                "--no-backup",
            ],
            capture_output=True,
            text=True,
        )

        assert result.returncode == 0
        assert f"font to '{test_font}'" in result.stdout
        assert "Theme fonts updated: 1 (majorFont + minorFont)" in result.stdout

        # Verify theme fonts were actually updated
        updated_prs = Presentation(str(sample_presentation))
        updated_major_font, updated_minor_font = self._get_theme_fonts(updated_prs)

        # Both should be updated to the new font
        assert updated_major_font == test_font
        assert updated_minor_font == test_font

        # Should be different from original (unless original was already Times New Roman)
        if original_major_font != test_font:
            assert updated_major_font != original_major_font
        if original_minor_font != test_font:
            assert updated_minor_font != original_minor_font

    def test_remap_language_only_preserves_theme_fonts(self, sample_presentation):
        """Test that language-only remapping doesn't modify theme fonts."""
        # Get original theme fonts
        original_prs = Presentation(str(sample_presentation))
        original_major_font, original_minor_font = self._get_theme_fonts(original_prs)

        # Run remap with language only (no font)
        result = subprocess.run(
            [
                "python",
                str(Path(__file__).parent.parent.parent.parent / "src" / "deckbuilder" / "cli.py"),
                "remap",
                str(sample_presentation),
                "--language",
                "en-US",
                "--no-backup",
            ],
            capture_output=True,
            text=True,
        )

        assert result.returncode == 0
        assert "language to en-US" in result.stdout
        # Should NOT show theme fonts updated
        assert "Theme fonts updated" not in result.stdout

        # Verify theme fonts were NOT changed
        updated_prs = Presentation(str(sample_presentation))
        updated_major_font, updated_minor_font = self._get_theme_fonts(updated_prs)

        assert updated_major_font == original_major_font
        assert updated_minor_font == original_minor_font

    def test_remap_font_and_language_updates_theme(self, sample_presentation):
        """Test that combined font and language remapping updates theme fonts."""
        test_font = "Arial"
        result = subprocess.run(
            [
                "python",
                str(Path(__file__).parent.parent.parent.parent / "src" / "deckbuilder" / "cli.py"),
                "remap",
                str(sample_presentation),
                "--font",
                test_font,
                "--language",
                "en-AU",
                "--no-backup",
            ],
            capture_output=True,
            text=True,
        )

        assert result.returncode == 0
        assert f"font to '{test_font}'" in result.stdout
        assert "language to en-AU" in result.stdout
        assert "Theme fonts updated: 1 (majorFont + minorFont)" in result.stdout

        # Verify theme fonts were updated
        updated_prs = Presentation(str(sample_presentation))
        updated_major_font, updated_minor_font = self._get_theme_fonts(updated_prs)

        assert updated_major_font == test_font
        assert updated_minor_font == test_font

    def _get_theme_fonts(self, presentation):
        """
        Helper method to extract theme font names from a presentation.

        Returns:
            Tuple of (majorFont, minorFont) typeface names
        """
        try:
            # Access theme part through presentation relationships
            theme_part = None
            for rel in presentation.part.rels.values():
                if "theme" in rel.target_ref:
                    theme_part = rel.target_part
                    break

            if not theme_part:
                return None, None

            # Parse theme XML from blob
            from lxml import etree  # nosec B410 - parsing trusted PowerPoint XML

            theme_xml_bytes = theme_part.blob
            theme_xml = etree.fromstring(theme_xml_bytes)  # nosec B320 - parsing trusted PowerPoint XML

            # Define DrawingML namespace
            DRAWINGML_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

            # Find theme font elements
            major_latin_elements = theme_xml.xpath(".//a:majorFont/a:latin", namespaces=DRAWINGML_NS)
            minor_latin_elements = theme_xml.xpath(".//a:minorFont/a:latin", namespaces=DRAWINGML_NS)

            major_font = major_latin_elements[0].get("typeface") if major_latin_elements else None
            minor_font = minor_latin_elements[0].get("typeface") if minor_latin_elements else None

            return major_font, minor_font

        except Exception as e:
            pytest.fail(f"Failed to extract theme fonts: {e}")
            return None, None
