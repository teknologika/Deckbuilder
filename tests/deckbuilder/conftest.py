"""
Deckbuilder-specific pytest configuration and fixtures.
"""

import os
import shutil
import tempfile
from pathlib import Path
from unittest.mock import Mock, patch

import pytest

# Import deckbuilder components
try:
    from deckbuilder.engine import Deckbuilder
    from deckbuilder.layout_intelligence import LayoutIntelligence
    from deckbuilder.naming_conventions import NamingConvention, PlaceholderContext
    from deckbuilder.structured_frontmatter import (
        StructuredFrontmatterConverter,
        StructuredFrontmatterRegistry,
        StructuredFrontmatterValidator,
    )
except ImportError:
    # Handle missing imports gracefully for early testing
    Deckbuilder = None
    NamingConvention = None


@pytest.fixture
def deckbuilder_temp_dir():
    """Temporary directory specific to deckbuilder tests."""
    temp_dir = tempfile.mkdtemp(prefix="deckbuilder_test_")
    yield Path(temp_dir)
    shutil.rmtree(temp_dir, ignore_errors=True)


@pytest.fixture
def mock_deckbuilder_env(deckbuilder_temp_dir):
    """Mock environment variables for deckbuilder."""
    templates_dir = deckbuilder_temp_dir / "templates"
    output_dir = deckbuilder_temp_dir / "output"
    templates_dir.mkdir()
    output_dir.mkdir()

    env_vars = {
        "DECK_TEMPLATE_FOLDER": str(templates_dir),
        "DECK_OUTPUT_FOLDER": str(output_dir),
        "DECK_TEMPLATE_NAME": "default",
    }

    with patch.dict(os.environ, env_vars):
        yield {"templates_dir": templates_dir, "output_dir": output_dir, "env_vars": env_vars}


@pytest.fixture
def default_template_json():
    """Default template JSON for testing."""
    return {
        "template_info": {"name": "Default", "version": "1.0"},
        "layouts": {
            "Title Slide": {
                "index": 0,
                "placeholders": {
                    "0": "title_top_1",
                    "1": "subtitle_1",
                    "10": "date_footer_1",
                    "11": "footer_footer_1",
                    "12": "slide_number_footer_1",
                },
            },
            "Title and Content": {
                "index": 1,
                "placeholders": {
                    "0": "title_top_1",
                    "1": "content_1",
                    "10": "date_footer_1",
                    "11": "footer_footer_1",
                    "12": "slide_number_footer_1",
                },
            },
            "Four Columns With Titles": {
                "index": 2,
                "placeholders": {
                    "0": "title_top_1",
                    "1": "title_col1_1",
                    "2": "content_col1_1",
                    "3": "title_col2_1",
                    "4": "content_col2_1",
                    "5": "title_col3_1",
                    "6": "content_col3_1",
                    "7": "title_col4_1",
                    "8": "content_col4_1",
                    "10": "date_footer_1",
                    "11": "footer_footer_1",
                    "12": "slide_number_footer_1",
                },
            },
            "Comparison": {
                "index": 3,
                "placeholders": {
                    "0": "title_top_1",
                    "1": "title_left_1",
                    "2": "content_left_1",
                    "3": "title_right_1",
                    "4": "content_right_1",
                    "10": "date_footer_1",
                    "11": "footer_footer_1",
                    "12": "slide_number_footer_1",
                },
            },
            "Two Content": {
                "index": 4,
                "placeholders": {
                    "0": "title_top_1",
                    "1": "content_left_1",
                    "2": "content_right_1",
                    "10": "date_footer_1",
                    "11": "footer_footer_1",
                    "12": "slide_number_footer_1",
                },
            },
        },
    }


@pytest.fixture
def naming_convention():
    """NamingConvention instance for testing."""
    if NamingConvention is None:
        pytest.skip("NamingConvention not available")
    return NamingConvention()


@pytest.fixture
def placeholder_context():
    """Sample PlaceholderContext for testing."""
    if PlaceholderContext is None:
        pytest.skip("PlaceholderContext not available")
    return PlaceholderContext(
        layout_name="Four Columns With Titles",
        placeholder_idx="1",
        placeholder_type="content",
    )


@pytest.fixture
def structured_frontmatter_registry(default_template_json):
    """StructuredFrontmatterRegistry for testing."""
    if StructuredFrontmatterRegistry is None:
        pytest.skip("StructuredFrontmatterRegistry not available")
    return StructuredFrontmatterRegistry(default_template_json)


@pytest.fixture
def structured_frontmatter_converter(default_template_json):
    """StructuredFrontmatterConverter for testing."""
    if StructuredFrontmatterConverter is None:
        pytest.skip("StructuredFrontmatterConverter not available")
    return StructuredFrontmatterConverter(default_template_json)


@pytest.fixture
def structured_frontmatter_validator():
    """StructuredFrontmatterValidator for testing."""
    if StructuredFrontmatterValidator is None:
        pytest.skip("StructuredFrontmatterValidator not available")
    return StructuredFrontmatterValidator()


@pytest.fixture
def layout_intelligence():
    """LayoutIntelligence instance for testing."""
    if LayoutIntelligence is None:
        pytest.skip("LayoutIntelligence not available")
    # Create a test layout intelligence file
    test_intelligence_data = {
        "content_patterns": {
            "intent_recognition": {
                "comparison": {"keywords": ["vs", "versus", "compare", "comparison"]},
                "overview": {"keywords": ["overview", "summary", "introduction"]},
            }
        },
        "layout_compatibility": {
            "Four Columns With Titles": {
                "optimal_for": ["multiple_columns", "comparison"],
                "confidence_factors": {"columns": 0.8, "titles": 0.6},
            }
        },
        "recommendation_engine": {
            "scoring_weights": {
                "content_structure": 0.4,
                "keyword_matching": 0.3,
                "intent_recognition": 0.2,
                "layout_compatibility": 0.1,
            },
            "minimum_confidence": 0.6,
        },
    }

    # Create temporary intelligence file
    with tempfile.NamedTemporaryFile(mode="w", suffix=".json", delete=False) as f:
        import json

        json.dump(test_intelligence_data, f)
        temp_file = f.name

    try:
        intelligence = LayoutIntelligence(temp_file)
        yield intelligence
    finally:
        os.unlink(temp_file)


@pytest.fixture
def sample_markdown_content():
    """Sample markdown content for testing."""
    return """# Test Presentation

## Overview
This is a test presentation with various formatting.

- **Bold** bullet point
- *Italic* text example
- ___Underlined___ content

Additional paragraph with mixed formatting."""


@pytest.fixture
def sample_structured_frontmatter():
    """Sample structured frontmatter YAML."""
    return {
        "layout": "Four Columns With Titles",
        "title": "Feature Comparison",
        "columns": [
            {"title": "Performance", "content": "Fast processing"},
            {"title": "Security", "content": "Enterprise encryption"},
            {"title": "Usability", "content": "Intuitive interface"},
            {"title": "Cost", "content": "Competitive pricing"},
        ],
    }


@pytest.fixture
def mock_pptx_presentation():
    """Mock PowerPoint presentation for testing."""
    mock_pres = Mock()
    mock_slide = Mock()
    mock_layout = Mock()
    mock_placeholder = Mock()

    # Set up basic mock structure
    mock_pres.slides = [mock_slide]
    mock_slide.slide_layout = mock_layout
    mock_layout.placeholders = [mock_placeholder]
    mock_placeholder.text = ""
    mock_placeholder.placeholder_format.idx = 1

    return mock_pres


@pytest.fixture
def sample_presentation_json():
    """Sample presentation JSON data."""
    return {
        "presentation": {
            "slides": [
                {
                    "type": "Title Slide",
                    "layout": "Title Slide",
                    "title": "Test Presentation",
                    "rich_content": [{"heading": "Testing Framework", "level": 2}],
                },
                {
                    "type": "Four Columns With Titles",
                    "title": "Feature Overview",
                    "title_col1_1": "Performance",
                    "content_col1_1": "Fast processing with optimized algorithms",
                    "title_col2_1": "Security",
                    "content_col2_1": "Enterprise-grade encryption",
                    "title_col3_1": "Usability",
                    "content_col3_1": "Intuitive interface",
                    "title_col4_1": "Cost",
                    "content_col4_1": "Competitive pricing",
                },
            ]
        }
    }


# Test helper functions specific to deckbuilder
def create_test_template_file(template_dir: Path, template_name: str, template_data: dict) -> Path:
    """Create a test template JSON file."""
    template_file = template_dir / f"{template_name}.json"
    import json

    with open(template_file, "w") as f:
        json.dump(template_data, f, indent=2)
    return template_file


def create_test_pptx_file(output_dir: Path, filename: str) -> Path:
    """Create a test PowerPoint file."""
    from pptx import Presentation

    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Test Slide"

    pptx_file = output_dir / filename
    pres.save(str(pptx_file))
    return pptx_file
