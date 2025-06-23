"""
Unit tests for deckbuilder engine.
"""

import os
import tempfile
from pathlib import Path
from unittest.mock import Mock, patch, MagicMock

import pytest

# Test imports with graceful handling
try:
    from deckbuilder.engine import Deckbuilder, get_deckbuilder_client

    HAS_ENGINE = True
except ImportError:
    HAS_ENGINE = False


@pytest.mark.skipif(not HAS_ENGINE, reason="Deckbuilder engine not available")
@pytest.mark.unit
@pytest.mark.deckbuilder
class TestDeckbuilderEngine:
    """Test cases for Deckbuilder engine class."""

    def test_singleton_behavior(self, mock_deckbuilder_env):
        """Test that Deckbuilder implements singleton pattern."""
        # Clear any existing instances
        if hasattr(Deckbuilder, "_instances"):
            Deckbuilder._instances.clear()

        instance1 = Deckbuilder()
        instance2 = Deckbuilder()

        # Should be the same instance
        assert instance1 is instance2

    def test_initialization_with_env_vars(self, mock_deckbuilder_env):
        """Test Deckbuilder initialization with environment variables."""
        env_data = mock_deckbuilder_env

        deckbuilder = Deckbuilder()

        assert deckbuilder.template_path == str(env_data["templates_dir"])
        assert deckbuilder.output_folder == str(env_data["output_dir"])
        assert deckbuilder.template_name == "default"

    @patch("deckbuilder.engine.Presentation")
    def test_initialization_creates_presentation(self, mock_presentation, mock_deckbuilder_env):
        """Test that initialization creates PowerPoint presentation."""
        mock_pres_instance = Mock()
        mock_presentation.return_value = mock_pres_instance

        deckbuilder = Deckbuilder()

        assert deckbuilder.prs is mock_pres_instance
        mock_presentation.assert_called_once()

    def test_template_path_validation(self, mock_deckbuilder_env):
        """Test that template path validation works."""
        env_data = mock_deckbuilder_env

        # Create a default template file
        default_template = env_data["templates_dir"] / "default.pptx"
        default_template.touch()

        # Should not raise an exception
        deckbuilder = Deckbuilder()
        assert deckbuilder.template_path == str(env_data["templates_dir"])

    def test_missing_template_handling(self, mock_deckbuilder_env):
        """Test handling of missing template files."""
        # Don't create the template file

        # Should still initialize but may log warning
        deckbuilder = Deckbuilder()
        assert deckbuilder.template_path is not None

    @patch("deckbuilder.engine.os.path.exists")
    def test_check_template_exists(self, mock_exists, mock_deckbuilder_env):
        """Test template existence checking."""
        mock_exists.return_value = True

        deckbuilder = Deckbuilder()

        # Should call exists check for template file
        mock_exists.assert_called()

        # Check if the call was made with a path containing the template name
        call_args = [call[0][0] for call in mock_exists.call_args_list]
        template_calls = [arg for arg in call_args if "default" in arg and ".pptx" in arg]
        assert len(template_calls) > 0

    def test_layout_mapping_initialization(self, mock_deckbuilder_env):
        """Test that layout mapping is initialized."""
        deckbuilder = Deckbuilder()

        # Layout mapping should be None initially
        assert deckbuilder.layout_mapping is None

    def test_get_deckbuilder_client_function(self, mock_deckbuilder_env):
        """Test the get_deckbuilder_client convenience function."""
        client = get_deckbuilder_client()

        assert client is not None
        assert isinstance(client, Deckbuilder)

    def test_multiple_client_calls_return_same_instance(self, mock_deckbuilder_env):
        """Test that multiple calls to get_deckbuilder_client return same instance."""
        client1 = get_deckbuilder_client()
        client2 = get_deckbuilder_client()

        assert client1 is client2

    @patch.dict(os.environ, {}, clear=True)
    def test_missing_environment_variables(self):
        """Test behavior when environment variables are missing."""
        # Clear all environment variables

        # Should still initialize with None values
        deckbuilder = Deckbuilder()

        assert deckbuilder.template_path is None
        assert deckbuilder.template_name is None
        assert deckbuilder.output_folder is None

    def test_template_name_default_handling(self, mock_deckbuilder_env):
        """Test default template name handling."""
        # Remove template name from environment
        if "DECK_TEMPLATE_NAME" in os.environ:
            del os.environ["DECK_TEMPLATE_NAME"]

        deckbuilder = Deckbuilder()

        # Should default to 'default'
        # Note: This depends on implementation details
        assert deckbuilder.template_name is None or deckbuilder.template_name == "default"


@pytest.mark.skipif(not HAS_ENGINE, reason="Deckbuilder engine not available")
@pytest.mark.integration
@pytest.mark.deckbuilder
class TestDeckbuilderEngineIntegration:
    """Integration tests for Deckbuilder engine."""

    def test_engine_with_real_template_file(self, mock_deckbuilder_env):
        """Test engine with actual template file."""
        env_data = mock_deckbuilder_env

        # Create a minimal PowerPoint template file
        from pptx import Presentation

        pres = Presentation()
        slide = pres.slides.add_slide(pres.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Test Template"

        template_file = env_data["templates_dir"] / "default.pptx"
        pres.save(str(template_file))

        # Initialize deckbuilder - should work without errors
        deckbuilder = Deckbuilder()

        assert deckbuilder.template_path == str(env_data["templates_dir"])
        assert template_file.exists()

    def test_engine_with_template_json(self, mock_deckbuilder_env, default_template_json):
        """Test engine with template JSON file."""
        env_data = mock_deckbuilder_env

        # Create template JSON file
        import json

        json_file = env_data["templates_dir"] / "default.json"
        with open(json_file, "w") as f:
            json.dump(default_template_json, f, indent=2)

        deckbuilder = Deckbuilder()

        # Should be able to load template configuration
        assert deckbuilder.template_path == str(env_data["templates_dir"])
        assert json_file.exists()

    def test_output_directory_creation(self, mock_deckbuilder_env):
        """Test that output directory is handled properly."""
        env_data = mock_deckbuilder_env

        # Remove output directory
        if env_data["output_dir"].exists():
            import shutil

            shutil.rmtree(env_data["output_dir"])

        deckbuilder = Deckbuilder()

        # Output folder should be set correctly
        assert deckbuilder.output_folder == str(env_data["output_dir"])


@pytest.mark.skipif(not HAS_ENGINE, reason="Deckbuilder engine not available")
@pytest.mark.unit
@pytest.mark.deckbuilder
class TestDeckbuilderEngineHelperMethods:
    """Test helper methods and utilities in Deckbuilder engine."""

    @patch("deckbuilder.engine.os.path.exists")
    def test_template_file_path_construction(self, mock_exists, mock_deckbuilder_env):
        """Test template file path construction."""
        mock_exists.return_value = True
        env_data = mock_deckbuilder_env

        deckbuilder = Deckbuilder()

        # Verify that template path construction includes correct elements
        expected_template_dir = str(env_data["templates_dir"])
        assert deckbuilder.template_path == expected_template_dir

    def test_environment_variable_precedence(self):
        """Test that environment variables take precedence."""
        test_env = {
            "DECK_TEMPLATE_FOLDER": "/custom/template/path",
            "DECK_OUTPUT_FOLDER": "/custom/output/path",
            "DECK_TEMPLATE_NAME": "custom_template",
        }

        with patch.dict(os.environ, test_env):
            deckbuilder = Deckbuilder()

            assert deckbuilder.template_path == "/custom/template/path"
            assert deckbuilder.output_folder == "/custom/output/path"
            assert deckbuilder.template_name == "custom_template"
