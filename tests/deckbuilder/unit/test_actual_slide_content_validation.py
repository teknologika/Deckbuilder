#!/usr/bin/env python3
"""
Actual Slide Content Validation Tests

These tests validate that content is ACTUALLY placed in PowerPoint slides
by examining the generated .pptx files with python-pptx.

This addresses the critical issue where tests pass but slides are blank.
"""

import json
import os
import secrets
import shutil
import sys
from pathlib import Path

import pytest
from pptx import Presentation

# Add src to path
project_root = Path(__file__).parent.parent.parent.parent
sys.path.insert(0, str(project_root / "src"))

from deckbuilder.engine import Deckbuilder  # noqa: E402


@pytest.fixture
def test_output_dir():
    """Create and return test output directory with unique hex string."""
    hex_id = secrets.token_hex(3)  # 3 bytes = 6 hex characters
    output_dir = Path(__file__).parent / "output" / f"test_{hex_id}"
    output_dir.mkdir(parents=True, exist_ok=True)

    yield output_dir

    # Clean up after test
    if output_dir.exists():
        shutil.rmtree(output_dir)


@pytest.fixture
def deckbuilder_with_env(test_output_dir):
    """Initialize Deckbuilder with proper environment variables."""
    # Set required environment variables
    original_env = {}
    env_vars = {
        "DECK_TEMPLATE_FOLDER": str(project_root / "src" / "deckbuilder" / "assets" / "templates"),
        "DECK_TEMPLATE_NAME": "default",
        "DECK_OUTPUT_FOLDER": str(test_output_dir),
    }

    # Store original values and set new ones
    for key, value in env_vars.items():
        original_env[key] = os.environ.get(key)
        os.environ[key] = value

    # Reset singleton to pick up new environment variables
    Deckbuilder.reset()

    deck = Deckbuilder()

    yield deck

    # Restore original environment
    for key, original_value in original_env.items():
        if original_value is None:
            os.environ.pop(key, None)
        else:
            os.environ[key] = original_value


class TestActualSlideContentValidation:
    """Test that content actually appears in PowerPoint slides"""

    def _generate_presentation_from_json(
        self, json_data: dict, filename: str, deckbuilder_with_env, test_output_dir
    ) -> Path:
        """Helper method to generate presentation and return path to generated file."""
        deck = deckbuilder_with_env

        # Convert legacy format to canonical format if needed
        if "presentation" in json_data and "slides" in json_data["presentation"]:
            # Legacy format - convert to canonical
            canonical_data = {"slides": []}
            for slide_data in json_data["presentation"]["slides"]:
                canonical_slide = {
                    "layout": slide_data.get("type", "Title and Content"),
                    "placeholders": {
                        key: value
                        for key, value in slide_data.items()
                        if key not in ["type", "layout"]
                    },
                    "content": [],
                }
                canonical_data["slides"].append(canonical_slide)
        else:
            # Already canonical format
            canonical_data = json_data

        result = deck.create_presentation(
            presentation_data=canonical_data, fileName=filename, templateName="default"
        )

        # Find generated file - look for .g.pptx first (what engine creates)
        output_files = list(test_output_dir.glob("*.g.pptx"))
        if not output_files:
            output_files = list(test_output_dir.glob("*.pptx"))
        assert (
            len(output_files) > 0
        ), f"No output file generated. Result: {result}. Files in dir: {list(test_output_dir.iterdir())}"

        return output_files[0]

    def _validate_slide_content(self, pptx_path: Path, expected_content: dict) -> dict:
        """
        Validate that specific content appears in slides.

        Args:
            pptx_path: Path to generated PowerPoint file
            expected_content: Dict mapping slide_index -> expected_text_list

        Returns:
            Dict with validation results
        """
        if not pptx_path.exists():
            return {"error": f"PPTX file not found: {pptx_path}"}

        prs = Presentation(str(pptx_path))
        results = {"total_slides": len(prs.slides), "slides": {}, "errors": []}

        for slide_idx, slide in enumerate(prs.slides):
            slide_text = []
            has_content = False

            # Extract all text from slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                    has_content = True
                elif hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_text.append(paragraph.text.strip())
                            has_content = True

            # Check expected content
            found_expected = []
            missing_expected = []

            if slide_idx in expected_content:
                for expected_text in expected_content[slide_idx]:
                    found = any(expected_text.lower() in text.lower() for text in slide_text)
                    if found:
                        found_expected.append(expected_text)
                    else:
                        missing_expected.append(expected_text)

            results["slides"][slide_idx] = {
                "has_content": has_content,
                "text": slide_text,
                "found_expected": found_expected,
                "missing_expected": missing_expected,
            }

            # Add errors for missing content
            if slide_idx in expected_content and missing_expected:
                results["errors"].append(f"Slide {slide_idx + 1}: Missing {missing_expected}")

        return results

    def test_simple_json_content_appears_in_slides(self, deckbuilder_with_env, test_output_dir):
        """Test that simple JSON content actually appears in PowerPoint slides"""
        json_data = {
            "presentation": {
                "slides": [
                    {
                        "type": "Title Slide",
                        "title": "TEST_TITLE_UNIQUE_123",
                        "subtitle": "TEST_SUBTITLE_UNIQUE_456",
                    },
                    {
                        "type": "Title and Content",
                        "title": "CONTENT_TITLE_UNIQUE_789",
                        "content": "TEST_CONTENT_UNIQUE_ABC",
                    },
                ]
            }
        }

        expected_content = {
            0: ["TEST_TITLE_UNIQUE_123", "TEST_SUBTITLE_UNIQUE_456"],  # Title slide
            1: ["CONTENT_TITLE_UNIQUE_789", "TEST_CONTENT_UNIQUE_ABC"],  # Content slide
        }

        pptx_path = self._generate_presentation_from_json(
            json_data, "simple_test", deckbuilder_with_env, test_output_dir
        )
        validation = self._validate_slide_content(pptx_path, expected_content)

        assert (
            validation["total_slides"] == 2
        ), f"Expected 2 slides, got {validation['total_slides']}"
        assert len(validation["errors"]) == 0, f"Content validation failed: {validation['errors']}"

        for slide_idx in [0, 1]:
            assert validation["slides"][slide_idx][
                "has_content"
            ], f"Slide {slide_idx + 1} has no content"
            assert (
                len(validation["slides"][slide_idx]["missing_expected"]) == 0
            ), f"Slide {slide_idx + 1} missing expected content"

    def test_semantic_field_names_content_appears(self, deckbuilder_with_env, test_output_dir):
        """Test that semantic field names actually place content in slides"""
        json_data = {
            "presentation": {
                "slides": [
                    {
                        "type": "Two Content",
                        "title": "TWO_CONTENT_TITLE_UNIQUE_999",
                        "content_left": "LEFT_CONTENT_UNIQUE_111",
                        "content_right": "RIGHT_CONTENT_UNIQUE_222",
                    }
                ]
            }
        }

        expected_content = {
            0: [
                "TWO_CONTENT_TITLE_UNIQUE_999",
                "LEFT_CONTENT_UNIQUE_111",
                "RIGHT_CONTENT_UNIQUE_222",
            ]
        }

        pptx_path = self._generate_presentation_from_json(
            json_data, "semantic_test", deckbuilder_with_env, test_output_dir
        )
        validation = self._validate_slide_content(pptx_path, expected_content)

        assert (
            validation["total_slides"] == 1
        ), f"Expected 1 slide, got {validation['total_slides']}"
        assert (
            len(validation["errors"]) == 0
        ), f"Semantic field validation failed: {validation['errors']}"
        assert validation["slides"][0]["has_content"], "Slide has no content"

    def test_rich_content_blocks_appear_in_slides(self, deckbuilder_with_env, test_output_dir):
        """Test that rich content blocks actually appear in PowerPoint slides"""
        json_data = {
            "presentation": {
                "slides": [
                    {
                        "type": "Title and Content",
                        "title": "RICH_CONTENT_TITLE_UNIQUE_888",
                        "rich_content": [
                            {"heading": "HEADING_UNIQUE_777", "level": 2},
                            {"paragraph": "PARAGRAPH_UNIQUE_666"},
                            {"bullets": ["BULLET_UNIQUE_555", "BULLET_UNIQUE_444"]},
                        ],
                    }
                ]
            }
        }

        expected_content = {
            0: [
                "RICH_CONTENT_TITLE_UNIQUE_888",
                "HEADING_UNIQUE_777",
                "PARAGRAPH_UNIQUE_666",
                "BULLET_UNIQUE_555",
                "BULLET_UNIQUE_444",
            ]
        }

        pptx_path = self._generate_presentation_from_json(
            json_data, "rich_content_test", deckbuilder_with_env, test_output_dir
        )
        validation = self._validate_slide_content(pptx_path, expected_content)

        assert (
            validation["total_slides"] == 1
        ), f"Expected 1 slide, got {validation['total_slides']}"
        assert (
            len(validation["errors"]) == 0
        ), f"Rich content validation failed: {validation['errors']}"
        assert validation["slides"][0]["has_content"], "Slide has no content"

    def test_comprehensive_layouts_file_content_appears(
        self, deckbuilder_with_env, test_output_dir
    ):
        """Test that the comprehensive layouts golden file actually produces content"""
        # Load the actual comprehensive layouts file
        golden_json_path = Path(__file__).parent.parent / "test_comprehensive_layouts.json"
        with open(golden_json_path, "r") as f:
            json_data = json.load(f)

        pptx_path = self._generate_presentation_from_json(
            json_data, "comprehensive_test", deckbuilder_with_env, test_output_dir
        )
        prs = Presentation(str(pptx_path))

        content_slides = 0
        blank_slides = 0

        # Check first 5 slides for content
        for i in range(min(5, len(prs.slides))):
            slide = prs.slides[i]
            has_content = any(
                shape.text.strip() for shape in slide.shapes if hasattr(shape, "text")
            ) or any(
                paragraph.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text_frame") and shape.text_frame
                for paragraph in shape.text_frame.paragraphs
            )

            if has_content:
                content_slides += 1
            else:
                blank_slides += 1

        assert (
            content_slides > blank_slides
        ), f"Too many blank slides: {blank_slides} blank vs {content_slides} with content"
        assert content_slides >= 3, f"Expected at least 3 slides with content, got {content_slides}"


if __name__ == "__main__":
    pytest.main([__file__, "-v", "-s"])
