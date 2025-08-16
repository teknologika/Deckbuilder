#!/usr/bin/env python
# -*- coding: utf-8 -*-

import unittest
from unittest.mock import MagicMock
from pptx import Presentation
from deckbuilder.core.slide_builder_legacy import SlideBuilder
from deckbuilder.content.formatter import ContentFormatter
from deckbuilder.image.placeholder import ImagePlaceholderHandler


class TestSpeakerNotes(unittest.TestCase):
    def setUp(self):
        self.prs = Presentation()
        self.content_formatter = ContentFormatter()
        self.placekitten_mock = MagicMock()
        self.image_placeholder_handler = ImagePlaceholderHandler(".", self.placekitten_mock)
        self.slide_builder = SlideBuilder()

    def test_add_speaker_notes(self):
        # Create a slide with speaker notes
        slide_data = {"layout": "Title and Content", "title": "Test Slide", "content": "This is a test slide.", "speaker_notes": "This is a speaker note."}
        slide = self.slide_builder.add_slide(self.prs, slide_data, self.content_formatter, self.image_placeholder_handler)

        # Check that the speaker notes were added correctly
        self.assertTrue(slide.has_notes_slide)
        notes_slide = slide.notes_slide
        self.assertEqual(notes_slide.notes_text_frame.text, "This is a speaker note.")


if __name__ == "__main__":
    unittest.main()
