"""
Unit tests for ContentProcessor

Tests content application with minimal font interference and proper newline support.
"""

import unittest
from unittest.mock import Mock, MagicMock
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE

from src.deckbuilder.core.content_processor import ContentProcessor


class TestContentProcessor(unittest.TestCase):
    """Test ContentProcessor functionality."""

    def setUp(self):
        """Set up test fixtures."""
        self.processor = ContentProcessor()
        
        # Mock dependencies
        self.mock_slide = Mock()
        self.mock_placeholder = Mock()
        self.mock_content_formatter = Mock()
        self.mock_image_handler = Mock()
        
        # Mock text frame for newline testing
        self.mock_text_frame = Mock()
        self.mock_paragraph = Mock()
        self.mock_text_frame.paragraphs = [self.mock_paragraph]
        self.mock_text_frame.add_paragraph.return_value = self.mock_paragraph

    def test_convert_newlines_to_paragraphs_single_line(self):
        """Test newline conversion with single line text."""
        text_content = "Single line without newlines"
        
        self.processor.convert_newlines_to_paragraphs(
            text_content, self.mock_text_frame, self.mock_content_formatter
        )
        
        # Should use existing paragraph with content formatter
        self.mock_content_formatter.apply_inline_formatting.assert_called_once_with(
            text_content, self.mock_paragraph
        )
        # Should not add new paragraphs
        self.mock_text_frame.add_paragraph.assert_not_called()

    def test_convert_newlines_to_paragraphs_multiple_lines(self):
        """Test newline conversion with multiple lines."""
        text_content = "First line\nSecond line\nThird line"
        expected_lines = ["First line", "Second line", "Third line"]
        
        self.processor.convert_newlines_to_paragraphs(
            text_content, self.mock_text_frame, self.mock_content_formatter
        )
        
        # Should clear text frame
        self.mock_text_frame.clear.assert_called_once()
        
        # Should add 2 new paragraphs (first one uses existing paragraph)
        self.assertEqual(self.mock_text_frame.add_paragraph.call_count, 2)
        
        # Should format each line
        self.assertEqual(self.mock_content_formatter.apply_inline_formatting.call_count, 3)
        
        # Check that each line was formatted
        call_args_list = self.mock_content_formatter.apply_inline_formatting.call_args_list
        for i, expected_line in enumerate(expected_lines):
            args, kwargs = call_args_list[i]
            self.assertEqual(args[0], expected_line)

    def test_convert_newlines_to_paragraphs_no_formatter(self):
        """Test newline conversion without content formatter."""
        text_content = "Line 1\nLine 2"
        
        self.processor.convert_newlines_to_paragraphs(
            text_content, self.mock_text_frame, None
        )
        
        # Should set text directly on paragraphs
        self.assertEqual(self.mock_paragraph.text, "Line 1")
        
        # Should add one new paragraph
        self.mock_text_frame.add_paragraph.assert_called_once()

    def test_convert_newlines_to_paragraphs_empty_text(self):
        """Test newline conversion with empty text."""
        self.processor.convert_newlines_to_paragraphs(
            "", self.mock_text_frame, self.mock_content_formatter
        )
        
        # Should set empty text
        self.mock_text_frame.text = ""

    def test_apply_content_with_newlines(self):
        """Test the wrapper method for applying content with newlines."""
        text_content = "Test\nwith\nnewlines"
        
        self.processor._apply_content_with_newlines(
            text_content, self.mock_text_frame, self.mock_content_formatter
        )
        
        # Should delegate to convert_newlines_to_paragraphs
        self.mock_text_frame.clear.assert_called_once()
        
    def test_apply_title_content_with_newlines(self):
        """Test title content application with newline support."""
        # Mock placeholder with text frame
        self.mock_placeholder.text_frame = self.mock_text_frame
        field_value = "Title Line 1\nTitle Line 2"
        
        self.processor._apply_title_content(
            self.mock_placeholder, field_value, self.mock_content_formatter
        )
        
        # Should clear and process with newlines
        self.mock_text_frame.clear.assert_called_once()
        
    def test_apply_title_content_formatted_segments(self):
        """Test title content with pre-formatted segments."""
        self.mock_placeholder.text_frame = self.mock_text_frame
        field_value = [{"text": "Formatted", "format": {"bold": True}}]
        
        self.processor._apply_title_content(
            self.mock_placeholder, field_value, self.mock_content_formatter
        )
        
        # Should use formatted segments method
        self.mock_content_formatter.apply_formatted_segments_to_paragraph.assert_called_once_with(
            field_value, self.mock_paragraph
        )

    def test_apply_title_content_no_text_frame(self):
        """Test title content fallback when no text frame available."""
        # Mock placeholder without text frame
        self.mock_placeholder.text_frame = None
        field_value = "Simple title"
        
        self.processor._apply_title_content(
            self.mock_placeholder, field_value, self.mock_content_formatter
        )
        
        # Should fallback to simple text
        self.assertEqual(self.mock_placeholder.text, "Simple title")

    def test_apply_content_placeholder_table_data(self):
        """Test content placeholder with table data."""
        field_value = {"type": "table", "data": [["Header"], ["Data"]]}
        
        with unittest.mock.patch('src.deckbuilder.core.content_processor.TableBuilder') as mock_table_builder:
            mock_builder_instance = Mock()
            mock_table_builder.return_value = mock_builder_instance
            
            self.processor._apply_content_placeholder_content(
                self.mock_slide, self.mock_placeholder, "content", field_value, self.mock_content_formatter
            )
            
            # Should create table via TableBuilder
            mock_table_builder.assert_called_once_with(self.mock_content_formatter)
            mock_builder_instance.add_table_to_slide.assert_called_once_with(self.mock_slide, field_value)

    def test_apply_content_placeholder_regular_content(self):
        """Test content placeholder with regular text content."""
        field_value = "Regular content"
        
        self.processor._apply_content_placeholder_content(
            self.mock_slide, self.mock_placeholder, "content", field_value, self.mock_content_formatter
        )
        
        # Should delegate to content formatter
        self.mock_content_formatter.add_content_to_placeholder.assert_called_once_with(
            self.mock_placeholder, field_value
        )


if __name__ == '__main__':
    unittest.main()