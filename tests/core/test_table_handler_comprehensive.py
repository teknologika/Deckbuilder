"""
Comprehensive tests for TableHandler functionality.

Tests table detection, creation, positioning, font handling, and duplication prevention.
"""

import pytest
from unittest.mock import Mock

# Removed unused typing imports

# Import the actual TableHandler
from deckbuilder.core.table_handler import TableHandler


class TestTableDetection:
    """Test table detection functionality."""

    def setup_method(self):
        """Set up test instance."""
        self.handler = TableHandler()

    def test_detect_valid_markdown_table(self):
        """Test detection of valid markdown table."""
        valid_table = """
| Header 1 | Header 2 | Header 3 |
|----------|----------|----------|
| Row 1 Col 1 | Row 1 Col 2 | Row 1 Col 3 |
| Row 2 Col 1 | Row 2 Col 2 | Row 2 Col 3 |
"""
        result = self.handler.detect_table_content(valid_table)
        assert result is True, "Should detect valid markdown table"

    def test_detect_table_without_separator(self):
        """Test detection of table without separator line."""
        table_no_separator = """
| Header 1 | Header 2 | Header 3 |
| Row 1 Col 1 | Row 1 Col 2 | Row 1 Col 3 |
| Row 2 Col 1 | Row 2 Col 2 | Row 2 Col 3 |
"""
        result = self.handler.detect_table_content(table_no_separator)
        assert result is True, "Should detect table even without separator"

    def test_reject_single_row_table(self):
        """Test rejection of single row (not a valid table)."""
        single_row = "| Header 1 | Header 2 | Header 3 |"
        result = self.handler.detect_table_content(single_row)
        assert result is False, "Should reject single row as not a valid table"

    def test_reject_non_table_content(self):
        """Test rejection of non-table content."""
        non_table = "This is just regular text content without any table structure."
        result = self.handler.detect_table_content(non_table)
        assert result is False, "Should reject non-table content"

    def test_reject_empty_content(self):
        """Test rejection of empty or None content."""
        assert self.handler.detect_table_content("") is False
        assert self.handler.detect_table_content(None) is False
        assert self.handler.detect_table_content("   ") is False


class TestTableParsing:
    """Test table parsing functionality."""

    def setup_method(self):
        """Set up test instance."""
        self.handler = TableHandler()

    def test_parse_simple_table(self):
        """Test parsing of simple markdown table."""
        table_markdown = """
| Name | Age | City |
|------|-----|------|
| Alice | 30 | NYC |
| Bob | 25 | LA |
"""
        result = self.handler.parse_table_structure(table_markdown)
        expected = [["Name", "Age", "City"], ["Alice", "30", "NYC"], ["Bob", "25", "LA"]]
        assert result == expected, f"Expected {expected}, got {result}"

    def test_parse_table_with_markdown_as_text(self):
        """Test that markdown formatting is treated as plain text."""
        table_with_markdown = """
| **Bold** | *Italic* | [Link](url) |
|----------|----------|-------------|
| **Bold Data** | *Italic Data* | [Another Link](url2) |
"""
        result = self.handler.parse_table_structure(table_with_markdown)
        expected = [["**Bold**", "*Italic*", "[Link](url)"], ["**Bold Data**", "*Italic Data*", "[Another Link](url2)"]]
        assert result == expected, "Markdown should be treated as plain text"

    def test_parse_irregular_table(self):
        """Test parsing table with irregular column counts."""
        irregular_table = """
| Col1 | Col2 | Col3 |
|------|------|------|
| A | B |
| X | Y | Z | Extra |
"""
        result = self.handler.parse_table_structure(irregular_table)
        expected = [["Col1", "Col2", "Col3"], ["A", "B"], ["X", "Y", "Z", "Extra"]]
        assert result == expected, "Should handle irregular column counts"

    def test_parse_empty_cells(self):
        """Test parsing table with empty cells."""
        table_with_empty = """
| Name | Value | Notes |
|------|-------|-------|
| Item1 |  | Some notes |
|  | 42 |  |
"""
        result = self.handler.parse_table_structure(table_with_empty)
        expected = [["Name", "Value", "Notes"], ["Item1", "", "Some notes"], ["", "42", ""]]
        assert result == expected, "Should handle empty cells correctly"


class TestTableCreation:
    """Test table creation functionality."""

    def setup_method(self):
        """Set up test instance."""
        self.handler = TableHandler()

        # Mock slide and shapes
        self.mock_slide = Mock()
        self.mock_shapes = Mock()
        self.mock_slide.shapes = self.mock_shapes

        # Mock table shape and table
        self.mock_table_shape = Mock()
        self.mock_table = Mock()
        self.mock_table_shape.table = self.mock_table
        self.mock_shapes.add_table.return_value = self.mock_table_shape

        # Mock table rows and cells
        self.mock_cells = []
        self.mock_rows = []

    def test_create_table_with_valid_data(self):
        """Test creating table with valid data."""
        table_data = [["Header1", "Header2"], ["Data1", "Data2"], ["Data3", "Data4"]]
        position = (100, 200)  # Mock position
        size = (500, 300)  # Mock size

        # Set up mock table structure
        rows = len(table_data)
        cols = len(table_data[0])

        # Create mock cells for each position
        mock_table_cells = []
        for _row_idx in range(rows):
            row_cells = []
            for _col_idx in range(cols):
                mock_cell = Mock()
                mock_cell.text = ""  # Will be set by our method
                row_cells.append(mock_cell)
            mock_table_cells.append(row_cells)

        # Create mock rows that return the cells
        mock_table_rows = []
        for row_idx in range(rows):
            mock_row = Mock()
            mock_row.cells = mock_table_cells[row_idx]
            mock_table_rows.append(mock_row)

        self.mock_table.rows = mock_table_rows

        # Call the method
        result = self.handler.create_table_from_data(self.mock_slide, table_data, position, size)

        # Verify table was created with correct parameters
        self.mock_shapes.add_table.assert_called_once_with(rows, cols, position[0], position[1], size[0], size[1])

        # Verify table content was set
        for row_idx, row_data in enumerate(table_data):
            for col_idx, cell_text in enumerate(row_data):
                expected_text = str(cell_text).strip()
                actual_text = mock_table_cells[row_idx][col_idx].text
                assert actual_text == expected_text, f"Cell [{row_idx}][{col_idx}] should be '{expected_text}', got '{actual_text}'"

        assert result == self.mock_table_shape, "Should return the created table shape"

    def test_create_table_with_empty_data(self):
        """Test creating table with empty data."""
        result = self.handler.create_table_from_data(self.mock_slide, [], (100, 200), (500, 300))
        assert result is None, "Should return None for empty table data"
        self.mock_shapes.add_table.assert_not_called()

    def test_create_table_with_invalid_data(self):
        """Test creating table with invalid data."""
        invalid_data = [[], []]  # Empty rows
        result = self.handler.create_table_from_data(self.mock_slide, invalid_data, (100, 200), (500, 300))
        assert result is None, "Should return None for invalid table data"


class TestTablePositioning:
    """Test table positioning functionality."""

    def setup_method(self):
        """Set up test instance."""
        self.handler = TableHandler()
        self.mock_slide = Mock()

    def test_default_positioning(self):
        """Test default table positioning."""
        position = self.handler.position_table_on_slide(self.mock_slide)

        # Should return tuple with left and top positions
        assert isinstance(position, tuple), "Should return position tuple"
        assert len(position) == 2, "Position should have left and top values"

        left, top = position
        # Default positioning should be reasonable
        assert left > 0, "Left position should be positive"
        assert top > 0, "Top position should be positive"

    def test_positioning_with_content_height(self):
        """Test table positioning considering existing content height."""
        from pptx.util import Cm

        content_height = Cm(5.0)  # 5cm of existing content
        position = self.handler.position_table_on_slide(self.mock_slide, content_height)

        left, top = position
        # Table should be positioned below the content
        assert top > content_height, "Table should be positioned below existing content"

    def test_positioning_with_large_content(self):
        """Test table positioning with large content (should cap at max position)."""
        from pptx.util import Cm

        large_content_height = Cm(20.0)  # Very large content
        position = self.handler.position_table_on_slide(self.mock_slide, large_content_height)

        left, top = position
        # Should cap at reasonable maximum to keep table on slide
        assert top <= Cm(17.0), "Table position should be capped to fit on slide"


class TestFontHandling:
    """Test font handling functionality."""

    def setup_method(self):
        """Set up test instance."""
        self.handler = TableHandler()

    def test_get_default_fonts_with_explicit_sizes(self):
        """Test getting fonts when explicitly specified."""
        slide_data = {"header_font_size": 14, "data_font_size": 10}

        result = self.handler.get_default_fonts(slide_data)
        expected = {"header_font_size": 14, "data_font_size": 10}
        assert result == expected, "Should return explicitly specified font sizes"

    def test_get_default_fonts_with_no_sizes(self):
        """Test getting fonts when none specified (should use template defaults)."""
        slide_data = {"title": "Test", "content": "Content"}

        result = self.handler.get_default_fonts(slide_data)
        assert result == {}, "Should return empty dict when no explicit fonts"

    def test_get_default_fonts_with_none_values(self):
        """Test handling None values in font specifications."""
        slide_data = {"header_font_size": None, "data_font_size": 12, "other_field": "value"}

        result = self.handler.get_default_fonts(slide_data)
        expected = {"data_font_size": 12}
        assert result == expected, "Should ignore None values and include only valid sizes"

    def test_validate_font_sizes_valid(self):
        """Test validation of valid font sizes."""
        font_config = {"header_font_size": 16, "data_font_size": 10}

        result = self.handler.validate_font_sizes(font_config)
        assert result == font_config, "Valid font sizes should pass through unchanged"

    def test_validate_font_sizes_boundary_values(self):
        """Test validation of boundary font sizes."""
        font_config = {"header_font_size": 30, "data_font_size": 5}  # Too large  # Too small

        result = self.handler.validate_font_sizes(font_config)
        expected = {"header_font_size": 24, "data_font_size": 8}  # Capped at maximum  # Raised to minimum
        assert result == expected, "Font sizes should be capped at boundaries"

    def test_validate_font_sizes_invalid_values(self):
        """Test validation removes invalid font sizes."""
        font_config = {"header_font_size": "invalid", "data_font_size": 12}

        result = self.handler.validate_font_sizes(font_config)
        expected = {"data_font_size": 12}
        assert result == expected, "Invalid font sizes should be removed"

    def test_get_font_size_for_row(self):
        """Test getting appropriate font size for different row types."""
        header_size = 14
        data_size = 10

        # Header row (index 0)
        header_font = self.handler.get_font_size_for_row(0, header_size, data_size)
        assert header_font == header_size, "Should return header font size for row 0"

        # Data rows (index > 0)
        data_font = self.handler.get_font_size_for_row(1, header_size, data_size)
        assert data_font == data_size, "Should return data font size for row > 0"

        data_font2 = self.handler.get_font_size_for_row(5, header_size, data_size)
        assert data_font2 == data_size, "Should return data font size for any row > 0"


class TestTableDuplicationPrevention:
    """Test table duplication prevention functionality."""

    def setup_method(self):
        """Set up test instance."""
        self.handler = TableHandler()
        self.mock_slide = Mock()

    def test_detect_existing_tables(self):
        """Test detection of existing tables on slide."""
        # Create mock shapes - some tables, some not
        mock_shape1 = Mock(spec=["shape_type"])  # Only has shape_type
        mock_shape1.shape_type = 19  # TABLE type

        mock_shape2 = Mock(spec=["shape_type"])  # Only has shape_type
        mock_shape2.shape_type = 1  # Not a table

        mock_shape3 = Mock(spec=["table", "shape_type"])  # Has both attributes
        mock_shape3.shape_type = 5  # Some other shape type, but has table attribute
        mock_shape3.table = Mock()  # Has table attribute

        mock_shape4 = Mock(spec=["shape_type"])  # Only has shape_type
        mock_shape4.shape_type = 3  # Not a table type

        self.mock_slide.shapes = [mock_shape1, mock_shape2, mock_shape3, mock_shape4]

        result = self.handler.detect_existing_tables(self.mock_slide)

        # Should find shapes 1 and 3 (both have table characteristics)
        assert len(result) == 2, f"Should find 2 tables, found {len(result)}"
        assert mock_shape1 in result, "Should find shape with table shape_type"
        assert mock_shape3 in result, "Should find shape with table attribute"

    def test_find_table_content_in_slide_data(self):
        """Test finding table content in various slide data formats."""
        # Test with markdown table in direct field
        slide_data_markdown = {
            "content": """
| Name | Value |
|------|-------|
| Test | 123 |
""",
            "title": "Test Slide",
        }

        result = self.handler.find_table_content_in_slide_data(slide_data_markdown)
        assert result is not None, "Should find table content"
        assert result["source_field"] == "content", "Should identify correct source field"
        assert result["markdown"] is not None, "Should return markdown content"

        # Test with structured table data
        slide_data_structured = {"placeholders": {"table_data": {"type": "table", "data": [["Header"], ["Data"]]}}}

        result = self.handler.find_table_content_in_slide_data(slide_data_structured)
        assert result is not None, "Should find structured table data"
        assert result["source_field"] == "table_data", "Should identify correct source field"
        assert result["table_data"] is not None, "Should return table data object"

        # Test with no table content
        slide_data_no_table = {"title": "Just Title", "content": "Regular text content"}

        result = self.handler.find_table_content_in_slide_data(slide_data_no_table)
        assert result is None, "Should return None when no table content found"


if __name__ == "__main__":
    # Run the tests
    pytest.main([__file__, "-v"])
