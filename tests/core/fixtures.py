"""
Shared test fixtures for core module testing

Provides reusable test data, mock objects, and helper functions for
comprehensive testing of the refactored slide building architecture.
"""

import pytest
from unittest.mock import Mock, MagicMock
from pptx.util import Cm
from typing import List, Dict, Any


# ==================== Slide and PowerPoint Fixtures ====================

@pytest.fixture
def mock_presentation():
    """Create a mock PowerPoint presentation object."""
    prs = Mock()
    prs.slides = Mock()
    prs.slide_layouts = []
    return prs

@pytest.fixture  
def mock_slide():
    """Create a mock PowerPoint slide object."""
    slide = Mock()
    slide.shapes = Mock()
    slide.placeholders = []
    return slide

@pytest.fixture
def mock_slide_layout():
    """Create a mock PowerPoint slide layout."""
    layout = Mock()
    layout.name = "Title and Content"
    layout.placeholders = []
    return layout

@pytest.fixture
def mock_placeholder():
    """Create a mock PowerPoint placeholder."""
    placeholder = Mock()
    placeholder.text_frame = Mock()
    placeholder.text_frame.text = ""
    placeholder.placeholder_format = Mock()
    placeholder.placeholder_format.idx = 1
    placeholder.placeholder_format.type = 1  # TITLE
    return placeholder

# ==================== Table Test Data Fixtures (DRY - Reading from Real Files) ====================

def _load_test_file_content(filename: str) -> str:
    """Helper to load content from structured frontmatter test files."""
    import os
    from pathlib import Path
    
    test_files_dir = Path(__file__).parent.parent.parent / "src" / "deckbuilder" / "structured_frontmatter_patterns" / "test_files"
    file_path = test_files_dir / filename
    
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

@pytest.fixture
def simple_table_markdown():
    """Simple table from real pattern test files."""
    content = _load_test_file_content("example_table_only.md")
    # Extract just the table_data section
    import re
    match = re.search(r'table_data: \|\s*\n((?:\s*\|.*\n)*)', content)
    if match:
        return match.group(1).strip()
    return ""

@pytest.fixture
def complex_table_markdown():
    """Complex table with markdown formatting from actual test files (treated as plain text)."""
    content = _load_test_file_content("example_table_with_content_above.md")
    # Extract just the table_data section
    import re
    match = re.search(r'table_data: \|\s*\n((?:\s*\|.*\n)*)', content)
    if match:
        return match.group(1).strip()
    return ""

@pytest.fixture 
def real_slide_data_table_only():
    """Real slide data from table_only test file."""
    import yaml
    content = _load_test_file_content("example_table_only.md")
    # Parse frontmatter
    if content.startswith('---'):
        _, frontmatter, _ = content.split('---', 2)
        return yaml.safe_load(frontmatter)
    return {}

@pytest.fixture
def real_slide_data_table_with_content():
    """Real slide data from table_with_content_above test file.""" 
    import yaml
    content = _load_test_file_content("example_table_with_content_above.md")
    # Parse frontmatter
    if content.startswith('---'):
        _, frontmatter, _ = content.split('---', 2)
        return yaml.safe_load(frontmatter)
    return {}

@pytest.fixture
def real_slide_data_comparison():
    """Real slide data from comparison test file."""
    import yaml
    content = _load_test_file_content("example_comparison.md")
    # Parse frontmatter  
    if content.startswith('---'):
        _, frontmatter, _ = content.split('---', 2)
        return yaml.safe_load(frontmatter)
    return {}

@pytest.fixture
def malformed_table_markdown():
    """Malformed table for error handling tests."""
    return """
    | Header 1 | Header 2
    Missing closing pipe
    |----------|----------|
    | Data 1   | Data 2   |
    """

@pytest.fixture
def large_table_data():
    """Large table dataset for performance testing."""
    return [
        [f"Header {i}" for i in range(10)]
    ] + [
        [f"Row {row}, Col {col}" for col in range(10)]
        for row in range(100)
    ]

@pytest.fixture
def table_with_mixed_content():
    """Table with mixed content types (all treated as plain text)."""
    return """
    | Type | Example | Notes |
    |------|---------|-------|
    | Bold | **This is bold** | Treated as literal |
    | Italic | *This is italic* | Treated as literal |
    | Link | [Link text](url) | Treated as literal |
    | Code | `code block` | Treated as literal |
    """

# ==================== Slide Data Fixtures (Using Real Test Files) ====================

@pytest.fixture
def basic_slide_data():
    """Basic slide data from real test files.""" 
    import yaml
    content = _load_test_file_content("example_title_and_content.md")
    if content.startswith('---'):
        _, frontmatter, _ = content.split('---', 2)
        return yaml.safe_load(frontmatter)
    return {}

@pytest.fixture
def slide_data_with_table():
    """Real slide data containing table markdown."""
    return real_slide_data_table_only()

@pytest.fixture
def slide_data_with_structured_table():
    """Real slide data with table from actual test files."""
    return real_slide_data_table_with_content()

@pytest.fixture
def slide_data_with_placeholders():
    """Real slide data with comparison layout from actual test files."""
    return real_slide_data_comparison()

# Legacy fixtures for backward compatibility 
@pytest.fixture
def basic_slide_data_simple():
    """Simple slide data for basic testing."""
    return {
        "layout": "Title and Content",
        "title": "Test Slide Title", 
        "content": "Test slide content"
    }

# ==================== Pattern and Template Fixtures ====================

@pytest.fixture
def mock_pattern_data():
    """Mock structured frontmatter pattern data."""
    return {
        "description": "Standard slide with title and content area",
        "yaml_pattern": {
            "layout": "Title and Content",
            "title": "str",
            "content": "str"
        },
        "validation": {
            "required_fields": ["title"],
            "optional_fields": ["content"],
            "field_types": {
                "title": "string",
                "content": "string"
            }
        }
    }

@pytest.fixture
def mock_comparison_pattern():
    """Mock comparison layout pattern."""
    return {
        "description": "Side-by-side comparison layout",
        "yaml_pattern": {
            "layout": "Comparison",
            "title": "str",
            "title_left": "str",
            "content_left": "str", 
            "title_right": "str",
            "content_right": "str"
        },
        "validation": {
            "required_fields": ["title", "title_left", "content_left", "title_right", "content_right"]
        }
    }

# ==================== Mock System Components ====================

@pytest.fixture
def mock_pattern_loader():
    """Mock PatternLoader for testing placeholder management."""
    loader = Mock()
    loader.get_pattern_for_layout.return_value = {
        "yaml_pattern": {
            "layout": "Title and Content",
            "title": "str",
            "content": "str"
        }
    }
    loader.load_patterns.return_value = {
        "Title and Content": {
            "yaml_pattern": {"layout": "Title and Content", "title": "str", "content": "str"}
        }
    }
    return loader

@pytest.fixture
def mock_content_formatter():
    """Mock ContentFormatter for testing content processing."""
    formatter = Mock()
    formatter.format_slide_data.return_value = {"layout": "Title and Content", "title": "Test", "content": "Content"}
    return formatter

@pytest.fixture
def mock_image_placeholder_handler():
    """Mock ImagePlaceholderHandler for testing."""
    handler = Mock()
    return handler

# ==================== Error Simulation Fixtures ====================

@pytest.fixture
def mock_slide_with_error_shapes():
    """Mock slide with shapes that raise errors when accessed.""" 
    slide = Mock()
    
    error_shape = Mock()
    error_shape.shape_type = property(lambda self: (_ for _ in ()).throw(Exception("Shape access error")))
    
    good_shape = Mock()
    good_shape.shape_type = 19  # TABLE type
    
    slide.shapes = [error_shape, good_shape]
    slide.placeholders = []
    return slide

# ==================== Performance Testing Fixtures ====================

@pytest.fixture
def performance_table_data():
    """Large dataset for performance testing."""
    return {
        "small_table": "| A | B |\n| 1 | 2 |",
        "medium_table": "\n".join([f"| Col {i} | Data {i} |" for i in range(50)]),
        "large_table": "\n".join([f"| Col {i} | Data {i} |" for i in range(500)]),
    }

# ==================== Helper Functions ====================

def create_mock_table_shape(rows: int, cols: int) -> Mock:
    """Helper to create mock table shape with specified dimensions."""
    table_shape = Mock()
    table = Mock()
    
    # Create mock rows and cells
    mock_rows = []
    for row_idx in range(rows):
        row = Mock()
        cells = []
        for col_idx in range(cols):
            cell = Mock()
            cell.text = f"Cell {row_idx},{col_idx}"
            cells.append(cell)
        row.cells = cells
        mock_rows.append(row)
    
    table.rows = mock_rows
    table_shape.table = table
    return table_shape

def create_mock_slide_with_placeholders(placeholder_types: List[int]) -> Mock:
    """Helper to create mock slide with specific placeholder types."""
    slide = Mock()
    placeholders = []
    
    for idx, ptype in enumerate(placeholder_types):
        placeholder = Mock()
        placeholder.placeholder_format = Mock()
        placeholder.placeholder_format.type = ptype
        placeholder.placeholder_format.idx = idx
        placeholder.text_frame = Mock()
        placeholder.text_frame.text = f"Placeholder {idx}"
        placeholders.append(placeholder)
    
    slide.placeholders = placeholders
    slide.shapes = Mock()
    return slide

# ==================== Validation Helpers ====================

def assert_table_structure_equal(actual: List[List[str]], expected: List[List[str]]) -> None:
    """Helper to assert table structures are equal with detailed error messages."""
    assert len(actual) == len(expected), f"Row count mismatch: got {len(actual)}, expected {len(expected)}"
    
    for row_idx, (actual_row, expected_row) in enumerate(zip(actual, expected)):
        assert len(actual_row) == len(expected_row), \
            f"Column count mismatch in row {row_idx}: got {len(actual_row)}, expected {len(expected_row)}"
        
        for col_idx, (actual_cell, expected_cell) in enumerate(zip(actual_row, expected_row)):
            assert actual_cell == expected_cell, \
                f"Cell mismatch at [{row_idx}][{col_idx}]: got '{actual_cell}', expected '{expected_cell}'"

def assert_placeholder_mapping_valid(mapping: Dict[str, Any], expected_fields: List[str]) -> None:
    """Helper to validate placeholder mapping results."""
    assert isinstance(mapping, dict), "Mapping should be a dictionary"
    
    for field in expected_fields:
        assert field in mapping, f"Field '{field}' should be in mapping"
        assert mapping[field] is not None, f"Mapping for '{field}' should not be None"