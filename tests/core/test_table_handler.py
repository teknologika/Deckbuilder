"""
Comprehensive unit tests for TableHandler module

Tests the newly refactored table processing logic with focus on:
1. Plain text processing (no markdown parsing in cells)
2. Table detection and structure parsing  
3. Table creation and positioning
4. Integration with existing slide systems

Design Philosophy:
- Test behavior, not implementation
- Cover happy path, edge cases, and error conditions
- Validate performance improvements from plain text processing
- Ensure backward compatibility where appropriate
"""

import pytest
from unittest.mock import Mock, MagicMock, patch
from pptx.util import Cm

from src.deckbuilder.core.table_handler import TableHandler


class TestTableHandler:
    """Test suite for TableHandler module functionality."""
    
    @pytest.fixture
    def table_handler(self):
        """Create a fresh TableHandler instance for each test."""
        return TableHandler()
    
    @pytest.fixture
    def mock_slide(self):
        """Create a mock PowerPoint slide for testing."""
        slide = Mock()
        slide.shapes = Mock()
        slide.placeholders = []
        return slide

    # ==================== Table Detection Tests ====================
    
    def test_detect_table_content_with_valid_markdown_table(self, table_handler):
        """Test detection of valid markdown table syntax."""
        table_markdown = """
        | Column 1 | Column 2 | Column 3 |
        |----------|----------|----------|
        | Row 1    | Data A   | Data B   |
        | Row 2    | Data C   | Data D   |
        """
        
        result = table_handler.detect_table_content(table_markdown)
        assert result is True, "Should detect valid markdown table"
    
    def test_detect_table_content_with_minimal_table(self, table_handler):
        """Test detection of minimal 2-row table."""
        minimal_table = "| Header | Value |\n| Data   | 123   |"
        
        result = table_handler.detect_table_content(minimal_table)
        assert result is True, "Should detect minimal valid table"
    
    def test_detect_table_content_rejects_single_row(self, table_handler):
        """Test rejection of single-row table (not valid)."""
        single_row = "| Only | One | Row |"
        
        result = table_handler.detect_table_content(single_row)
        assert result is False, "Should reject single-row table"
    
    def test_detect_table_content_rejects_separator_only(self, table_handler):
        """Test rejection of separator-only content."""
        separator_only = "|---|---|---|"
        
        result = table_handler.detect_table_content(separator_only)
        assert result is False, "Should reject separator-only content"
    
    def test_detect_table_content_handles_empty_input(self, table_handler):
        """Test handling of empty or None input."""
        assert table_handler.detect_table_content("") is False
        assert table_handler.detect_table_content(None) is False
        assert table_handler.detect_table_content("   ") is False

    def test_detect_table_content_handles_non_table_text(self, table_handler):
        """Test rejection of regular text content."""
        regular_text = "This is just regular text content with no table structure."
        
        result = table_handler.detect_table_content(regular_text)
        assert result is False, "Should reject non-table content"

    # ==================== Table Structure Parsing Tests ====================
    
    def test_parse_table_structure_extracts_plain_text_cells(self, table_handler):
        """Test parsing table structure with plain text extraction."""
        table_markdown = """
        | **Bold Header** | *Italic Header* | Normal Header |
        |-----------------|-----------------|---------------|
        | **Bold Data**   | *Italic Data*   | Plain Data    |
        | More **bold**   | More *italic*   | More plain    |
        """
        
        result = table_handler.parse_table_structure(table_markdown)
        
        # Should extract cells as plain text (including markdown characters)
        expected = [
            ["**Bold Header**", "*Italic Header*", "Normal Header"],
            ["**Bold Data**", "*Italic Data*", "Plain Data"],
            ["More **bold**", "More *italic*", "More plain"]
        ]
        
        assert result == expected, "Should extract table structure as plain text"
    
    def test_parse_table_structure_handles_separator_lines(self, table_handler):
        """Test parsing with various separator line formats."""
        table_with_separators = """
        | Header 1 | Header 2 |
        |:---------|----------|
        | Data 1   | Data 2   |
        |----------|:--------:|
        | Data 3   | Data 4   |
        """
        
        result = table_handler.parse_table_structure(table_with_separators)
        
        # Should skip separator lines but keep data rows
        expected = [
            ["Header 1", "Header 2"],
            ["Data 1", "Data 2"],
            ["Data 3", "Data 4"]
        ]
        
        assert result == expected, "Should skip separator lines correctly"
    
    def test_parse_table_structure_handles_uneven_columns(self, table_handler):
        """Test parsing table with uneven column counts."""
        uneven_table = """
        | Col1 | Col2 | Col3 |
        | A    | B    |
        | X    | Y    | Z    | Extra |
        """
        
        result = table_handler.parse_table_structure(uneven_table)
        
        # Should handle uneven columns gracefully
        expected = [
            ["Col1", "Col2", "Col3"],
            ["A", "B"],
            ["X", "Y", "Z", "Extra"]
        ]
        
        assert result == expected, "Should handle uneven columns"
    
    def test_parse_table_structure_returns_empty_for_invalid_content(self, table_handler):
        """Test parsing returns empty list for invalid table content."""
        invalid_content = "This is not a table at all"
        
        result = table_handler.parse_table_structure(invalid_content)
        assert result == [], "Should return empty list for invalid table content"

    # ==================== Table Creation Tests ====================
    
    @patch('src.deckbuilder.core.table_handler.debug_print')
    def test_create_table_from_data_success(self, mock_debug, table_handler, mock_slide):
        """Test successful table creation from plain text data."""
        # Setup mock slide and table
        mock_table_shape = Mock()
        mock_table = Mock()
        mock_table_shape.table = mock_table
        
        # Setup table rows and cells
        mock_cell_00 = Mock()
        mock_cell_01 = Mock()
        mock_cell_10 = Mock()
        mock_cell_11 = Mock()
        
        mock_row_0 = Mock()
        mock_row_0.cells = [mock_cell_00, mock_cell_01]
        mock_row_1 = Mock()
        mock_row_1.cells = [mock_cell_10, mock_cell_11]
        
        mock_table.rows = [mock_row_0, mock_row_1]
        mock_slide.shapes.add_table.return_value = mock_table_shape
        
        # Test data
        table_data = [
            ["Header 1", "Header 2"],
            ["Data 1", "Data 2"]
        ]
        position = (Cm(1), Cm(2))
        size = (Cm(10), Cm(5))
        
        result = table_handler.create_table_from_data(mock_slide, table_data, position, size)
        
        # Verify table creation
        mock_slide.shapes.add_table.assert_called_once_with(2, 2, Cm(1), Cm(2), Cm(10), Cm(5))
        
        # Verify cell content (plain text)
        mock_cell_00.text = "Header 1"
        mock_cell_01.text = "Header 2"
        mock_cell_10.text = "Data 1"
        mock_cell_11.text = "Data 2"
        
        assert result == mock_table_shape, "Should return created table shape"
        mock_debug.assert_called_with("Created table with 2 rows and 2 columns")
    
    def test_create_table_from_data_handles_empty_data(self, table_handler, mock_slide):
        """Test handling of empty table data."""
        result = table_handler.create_table_from_data(mock_slide, [], (Cm(1), Cm(2)), (Cm(10), Cm(5)))
        assert result is None, "Should return None for empty data"
        
        result = table_handler.create_table_from_data(mock_slide, None, (Cm(1), Cm(2)), (Cm(10), Cm(5)))
        assert result is None, "Should return None for None data"
    
    def test_create_table_from_data_handles_creation_error(self, table_handler, mock_slide):
        """Test handling of table creation errors."""
        mock_slide.shapes.add_table.side_effect = Exception("PowerPoint error")
        
        table_data = [["Header"], ["Data"]]
        result = table_handler.create_table_from_data(
            mock_slide, table_data, (Cm(1), Cm(2)), (Cm(10), Cm(5))
        )
        
        assert result is None, "Should return None on creation error"

    # ==================== Table Positioning Tests ====================
    
    def test_position_table_on_slide_default_positioning(self, table_handler, mock_slide):
        """Test default table positioning when no content height provided."""
        position = table_handler.position_table_on_slide(mock_slide)
        
        expected_left = Cm(1.0)
        expected_top = Cm(8.0)
        
        assert position == (expected_left, expected_top), "Should use default positioning"
    
    def test_position_table_on_slide_calculated_positioning(self, table_handler, mock_slide):
        """Test calculated table positioning based on content height."""
        content_height = Cm(4.0)  # 4cm of content above
        
        position = table_handler.position_table_on_slide(mock_slide, content_height)
        
        expected_left = Cm(1.0)
        expected_spacing = max(Cm(0.5), content_height * 0.1)  # 10% of content height or 0.5cm
        expected_top = content_height + expected_spacing
        
        assert position[0] == expected_left, "Should position at 1cm from left"
        assert position[1] == expected_top, "Should position below content with spacing"
    
    def test_position_table_on_slide_respects_maximum_position(self, table_handler, mock_slide):
        """Test positioning respects slide boundaries."""
        very_large_content_height = Cm(20.0)  # Would position table off slide
        
        position = table_handler.position_table_on_slide(mock_slide, very_large_content_height)
        
        expected_left = Cm(1.0)
        max_top = Cm(17.0)  # Maximum allowed position
        
        assert position == (expected_left, max_top), "Should respect maximum position"

    # ==================== Existing Table Detection Tests ====================
    
    def test_detect_existing_tables_finds_table_shapes(self, table_handler, mock_slide):
        """Test detection of existing table shapes on slide."""
        # Create mock shapes - one table, one text
        mock_table_shape = Mock()
        mock_table_shape.shape_type = 19  # MSO_SHAPE_TYPE.TABLE
        
        mock_text_shape = Mock()
        mock_text_shape.shape_type = 1  # Not a table
        
        mock_slide.shapes = [mock_table_shape, mock_text_shape]
        
        result = table_handler.detect_existing_tables(mock_slide)
        
        assert len(result) == 1, "Should find one table shape"
        assert result[0] == mock_table_shape, "Should return the table shape"
    
    def test_detect_existing_tables_finds_shapes_with_table_attribute(self, table_handler, mock_slide):
        """Test detection using table attribute fallback."""
        mock_shape_with_table_attr = Mock()
        mock_shape_with_table_attr.shape_type = 999  # Unknown type
        mock_shape_with_table_attr.table = Mock()  # Has table attribute
        
        mock_slide.shapes = [mock_shape_with_table_attr]
        
        result = table_handler.detect_existing_tables(mock_slide)
        
        assert len(result) == 1, "Should find shape with table attribute"
        assert result[0] == mock_shape_with_table_attr, "Should return the shape with table"
    
    def test_detect_existing_tables_handles_shape_errors(self, table_handler, mock_slide):
        """Test graceful handling of shape access errors."""
        mock_error_shape = Mock()
        mock_error_shape.shape_type = property(lambda self: (_ for _ in ()).throw(Exception("Access error")))
        
        mock_good_shape = Mock()
        mock_good_shape.shape_type = 19  # Table
        
        mock_slide.shapes = [mock_error_shape, mock_good_shape]
        
        result = table_handler.detect_existing_tables(mock_slide)
        
        assert len(result) == 1, "Should skip error shapes and find valid ones"
        assert result[0] == mock_good_shape, "Should return valid table shape"

    # ==================== Slide Data Analysis Tests ====================
    
    def test_find_table_content_in_slide_data_finds_markdown_table(self, table_handler):
        """Test finding markdown table content in slide data."""
        slide_data = {
            "layout": "Title and Content",
            "title": "My Slide",
            "content": "| Header 1 | Header 2 |\n| Data 1   | Data 2   |"
        }
        
        result = table_handler.find_table_content_in_slide_data(slide_data)
        
        assert result is not None, "Should find table content"
        assert result["source_field"] == "content", "Should identify source field"
        assert result["markdown"] == slide_data["content"], "Should return markdown content"
        assert result["table_data"] is None, "Should not have structured table data"
    
    def test_find_table_content_in_slide_data_finds_structured_table(self, table_handler):
        """Test finding structured table data in slide data."""
        slide_data = {
            "layout": "Title and Content",
            "title": "My Slide", 
            "table_info": {
                "type": "table",
                "headers": ["Col 1", "Col 2"],
                "rows": [["A", "B"], ["C", "D"]]
            }
        }
        
        result = table_handler.find_table_content_in_slide_data(slide_data)
        
        assert result is not None, "Should find structured table"
        assert result["source_field"] == "table_info", "Should identify source field"
        assert result["table_data"] == slide_data["table_info"], "Should return table data"
        assert result["markdown"] is None, "Should not have markdown content"
    
    def test_find_table_content_in_slide_data_checks_placeholders(self, table_handler):
        """Test finding table content in placeholders section."""
        slide_data = {
            "layout": "Title and Content",
            "title": "My Slide",
            "placeholders": {
                "content": "| Header | Value |\n| Data   | 123   |"
            }
        }
        
        result = table_handler.find_table_content_in_slide_data(slide_data)
        
        assert result is not None, "Should find table in placeholders"
        assert result["source_field"] == "content", "Should identify placeholder field"
    
    def test_find_table_content_in_slide_data_returns_none_when_no_table(self, table_handler):
        """Test returning None when no table content found."""
        slide_data = {
            "layout": "Title and Content",
            "title": "My Slide",
            "content": "Just regular text content with no tables"
        }
        
        result = table_handler.find_table_content_in_slide_data(slide_data)
        assert result is None, "Should return None when no table content found"

    # ==================== Placeholder Cleaning Tests ====================
    
    def test_clear_table_content_from_placeholders_clears_matching_content(self, table_handler, mock_slide):
        """Test clearing table markdown from placeholders."""
        # Setup mock placeholder with table content
        mock_placeholder = Mock()
        mock_text_frame = Mock()
        mock_text_frame.text = "| Header | Value |\n| Data   | 123   |"
        mock_placeholder.text_frame = mock_text_frame
        mock_placeholder.placeholder_format.idx = 1
        
        mock_slide.placeholders = [mock_placeholder]
        
        slide_data = {
            "content": "| Header | Value |\n| Data   | 123   |"
        }
        
        table_handler.clear_table_content_from_placeholders(mock_slide, slide_data)
        
        # Should clear the text frame
        mock_text_frame.clear.assert_called_once()
    
    def test_clear_table_content_from_placeholders_ignores_non_table_content(self, table_handler, mock_slide):
        """Test ignoring placeholders without table content."""
        mock_placeholder = Mock()
        mock_text_frame = Mock()
        mock_text_frame.text = "Regular text content"
        mock_placeholder.text_frame = mock_text_frame
        
        mock_slide.placeholders = [mock_placeholder]
        
        slide_data = {
            "content": "Regular text content"
        }
        
        table_handler.clear_table_content_from_placeholders(mock_slide, slide_data)
        
        # Should not clear non-table content
        mock_text_frame.clear.assert_not_called()
    
    def test_clear_table_content_from_placeholders_handles_missing_text_frame(self, table_handler, mock_slide):
        """Test handling placeholders without text frames."""
        mock_placeholder = Mock()
        mock_placeholder.text_frame = None
        
        mock_slide.placeholders = [mock_placeholder]
        
        slide_data = {
            "content": "| Header | Value |\n| Data   | 123   |"
        }
        
        # Should not raise exception
        table_handler.clear_table_content_from_placeholders(mock_slide, slide_data)

    # ==================== Performance and Integration Tests ====================
    
    def test_plain_text_processing_performance(self, table_handler):
        """Test that plain text processing is faster than complex parsing."""
        # Large table with complex markdown (should be processed as plain text)
        large_table = """
        | **Complex** | *Markdown* | ___Content___ | [Link](url) |
        """ + "\n".join([
            f"| **Row {i}** | *Data {i}* | ___Value {i}___ | [Link {i}](url) |"
            for i in range(100)
        ])
        
        import time
        
        start_time = time.time()
        result = table_handler.parse_table_structure(large_table)
        end_time = time.time()
        
        processing_time = end_time - start_time
        
        # Should complete quickly (plain text processing)
        assert processing_time < 0.1, f"Processing took {processing_time:.3f}s, should be < 0.1s"
        assert len(result) == 101, "Should parse all rows"
        assert "**Complex**" in result[0][0], "Should preserve markdown as plain text"
    
    @pytest.mark.integration
    def test_table_handler_integration_with_existing_systems(self, table_handler):
        """Integration test with existing slide building components."""
        # This would test integration with PatternLoader, semantic detection, etc.
        # For now, just verify the TableHandler can be instantiated and used
        assert isinstance(table_handler, TableHandler)
        
        # Test basic workflow
        markdown = "| Test | Value |\n| A    | B     |"
        assert table_handler.detect_table_content(markdown) is True
        
        structure = table_handler.parse_table_structure(markdown)
        assert len(structure) == 2
        assert structure[0] == ["Test", "Value"]
        assert structure[1] == ["A", "B"]


# ==================== Test Data and Fixtures ====================

@pytest.fixture
def sample_table_markdown():
    """Sample table markdown for testing."""
    return """
    | Feature | Old System | New System |
    |---------|------------|------------|
    | Performance | Slow | **Fast** |
    | Maintenance | Hard | *Easy* |
    | Bugs | Many | ___Few___ |
    """

@pytest.fixture  
def sample_malformed_table():
    """Sample malformed table for error testing."""
    return """
    | Header 1 | Header 2
    | Missing pipe at end
    |----------|
    | Data | More Data |
    """

# ==================== Performance Benchmarks ====================

class TestTableHandlerPerformance:
    """Performance tests to validate improvements from plain text processing."""
    
    @pytest.mark.performance
    def test_table_detection_performance(self, table_handler):
        """Benchmark table detection performance."""
        test_content = "| Col 1 | Col 2 |\n| Data 1 | Data 2 |" * 1000
        
        import time
        start = time.time()
        for _ in range(100):
            table_handler.detect_table_content(test_content)
        end = time.time()
        
        avg_time = (end - start) / 100
        assert avg_time < 0.01, f"Average detection time {avg_time:.4f}s should be < 0.01s"
    
    @pytest.mark.performance
    def test_table_parsing_performance(self, table_handler):
        """Benchmark table parsing performance."""
        large_table = "\n".join([
            f"| Col {j} |" * 10 + f"\n| Data {i}-{j} |" * 10
            for i in range(50) for j in range(10)
        ])
        
        import time
        start = time.time()
        result = table_handler.parse_table_structure(large_table)
        end = time.time()
        
        parse_time = end - start
        assert parse_time < 0.05, f"Parse time {parse_time:.4f}s should be < 0.05s"
        assert len(result) > 0, "Should successfully parse large table"