"""
Test harness for structured frontmatter test files using black-box testing approach.

This test framework follows the specification from docs/Features/Structured_Frontmatter_Test_Files.md:
1. Execute deckbuilder CLI commands in black-box style
2. Validate output PPTX 1:1 against input JSON using python-pptx
3. DRY test harness - write once, use many principle
"""

import json
import subprocess
from pathlib import Path
from pptx import Presentation
from typing import Dict, Any, List, Tuple


class StructuredFrontmatterTester:
    """Write-once, use-many test harness for structured frontmatter examples."""

    def __init__(self):
        self.project_root = Path(__file__).parent.parent
        self.test_files_dir = self.project_root / "src/deckbuilder/structured_frontmatter_patterns/test_files"
        self.output_dir = self.project_root / "tests/output"
        self.venv_activate = self.project_root / ".venv/bin/activate"

        # Ensure output directory exists
        self.output_dir.mkdir(exist_ok=True)

        # Track created files for cleanup
        self.created_files = []

    def run_deckbuilder_command(self, input_file: Path, output_name: str) -> Tuple[bool, str, Path]:
        """
        Execute deckbuilder CLI command in black-box style.

        Args:
            input_file: Path to input JSON file
            output_name: Base name for output file

        Returns:
            Tuple of (success, output_text, output_pptx_path)
        """
        # Try to use venv if available, otherwise use python -m
        if self.venv_activate.exists():
            cmd = ["bash", "-c", f"source {self.venv_activate} && deckbuilder create {input_file} --output {self.output_dir}/{output_name}"]
        else:
            # CI environment - use python -m approach
            cmd = ["python", "-m", "deckbuilder.cli", "create", str(input_file), "--output", f"{self.output_dir}/{output_name}"]

        try:
            result = subprocess.run(cmd, capture_output=True, text=True, cwd=self.project_root, timeout=60)

            # Find the generated PPTX file
            output_pptx = None
            if result.returncode == 0:
                # Parse output to find generated filename
                output_lines = result.stdout.split("\n")
                for line in output_lines:
                    if "✅ Presentation complete:" in line:
                        # Extract filename from: "✅ Presentation complete: filename.pptx (X slides)"
                        filename = line.split("✅ Presentation complete: ")[1].split(" (")[0]
                        output_pptx = self.output_dir / filename
                        # Track for cleanup
                        if output_pptx and output_pptx.exists():
                            self.created_files.append(output_pptx)
                        break

            return result.returncode == 0, result.stdout + result.stderr, output_pptx

        except subprocess.TimeoutExpired:
            return False, "Command timed out", None
        except Exception as e:
            return False, f"Command failed: {str(e)}", None

    def load_input_json(self, json_file: Path) -> Dict[str, Any]:
        """Load and parse input JSON file."""
        with open(json_file, "r", encoding="utf-8") as f:
            return json.load(f)

    def validate_pptx_against_json(self, pptx_path: Path, input_data: Dict[str, Any]) -> List[str]:
        """
        Validate PPTX output 1:1 against input JSON using python-pptx.

        Args:
            pptx_path: Path to generated PPTX file
            input_data: Original JSON input data

        Returns:
            List of validation errors (empty if all valid)
        """
        errors = []

        if not pptx_path or not pptx_path.exists():
            errors.append(f"PPTX file not found: {pptx_path}")
            return errors

        try:
            prs = Presentation(str(pptx_path))
            slides_data = input_data.get("slides", [])

            # Validate slide count
            if len(prs.slides) != len(slides_data):
                errors.append(f"Slide count mismatch: expected {len(slides_data)}, got {len(prs.slides)}")
                return errors  # Can't continue validation without matching slide count

            # Validate each slide
            for slide_idx, (pptx_slide, json_slide) in enumerate(zip(prs.slides, slides_data)):
                slide_errors = self._validate_slide_content(slide_idx, pptx_slide, json_slide)
                errors.extend(slide_errors)

        except Exception as e:
            errors.append(f"Error reading PPTX file: {str(e)}")

        return errors

    def _validate_slide_content(self, slide_idx: int, pptx_slide, json_slide: Dict[str, Any]) -> List[str]:
        """Validate individual slide content against JSON data."""
        errors = []
        slide_prefix = f"Slide {slide_idx + 1}"

        # Validate layout
        expected_layout = json_slide.get("layout")
        actual_layout = pptx_slide.slide_layout.name
        if actual_layout != expected_layout:
            errors.append(f"{slide_prefix}: Layout mismatch - expected '{expected_layout}', got '{actual_layout}'")

        # Validate placeholders content
        placeholders = json_slide.get("placeholders", {})

        # Skip validation for styling/configuration fields that don't have visible content
        skip_fields = {"style", "row_style", "border_style", "row_height", "table_width", "column_widths", "header_font_size", "data_font_size", "custom_colors"}

        for field_name, expected_content in placeholders.items():
            # Skip styling/configuration fields
            if field_name in skip_fields:
                continue

            # Skip image and table_data fields for now (would need special validation)
            if field_name in ["image", "table_data"]:
                continue

            # Find corresponding placeholder in PPTX
            found_placeholder = False

            for shape in pptx_slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    # Simple text content validation - could be enhanced for rich formatting
                    if self._content_matches(shape.text, expected_content):
                        found_placeholder = True
                        break
                elif hasattr(shape, "placeholder_format"):
                    # Check placeholder by name/type
                    try:
                        placeholder_name = shape.element.nvSpPr.cNvPr.name
                        if placeholder_name and field_name.lower() in placeholder_name.lower():
                            if hasattr(shape, "text_frame") and shape.text_frame:
                                if self._content_matches(shape.text_frame.text, expected_content):
                                    found_placeholder = True
                                    break
                    except AttributeError:
                        continue

            if not found_placeholder and expected_content:  # Only error if content is not empty
                errors.append(f"{slide_prefix}: Content not found for field '{field_name}': '{expected_content}'")

        return errors

    def _content_matches(self, actual_text: str, expected_content) -> bool:
        """
        Compare actual PPTX text content with expected JSON content.

        This is a simplified comparison - could be enhanced to handle:
        - Rich formatting preservation
        - Bullet point structures
        - Complex formatting patterns
        """
        if not actual_text and not expected_content:
            return True

        if not actual_text or not expected_content:
            return False

        # Convert expected_content to string if it's not already
        expected_str = str(expected_content)

        # Remove markdown formatting for basic comparison
        expected_clean = expected_str
        for marker in ["**", "*", "___", "***"]:
            expected_clean = expected_clean.replace(marker, "")

        # Normalize whitespace
        actual_normalized = " ".join(actual_text.split())
        expected_normalized = " ".join(expected_clean.split())

        # Check if expected content is contained in actual (allows for some formatting differences)
        return expected_normalized.lower() in actual_normalized.lower()

    def test_example_file(self, json_file: Path, test_name: str) -> Tuple[bool, List[str]]:
        """
        Test a single example file end-to-end.

        Args:
            json_file: Path to JSON test file
            test_name: Name for test output files

        Returns:
            Tuple of (success, errors)
        """
        # Step 1: Load input JSON
        try:
            input_data = self.load_input_json(json_file)
        except Exception as e:
            return False, [f"Failed to load JSON file: {str(e)}"]

        # Step 2: Run deckbuilder command
        success, output, pptx_path = self.run_deckbuilder_command(json_file, test_name)
        if not success:
            return False, [f"Deckbuilder command failed: {output}"]

        # Step 3: Validate PPTX against JSON
        validation_errors = self.validate_pptx_against_json(pptx_path, input_data)

        return len(validation_errors) == 0, validation_errors

    def cleanup_test_files(self) -> int:
        """
        Clean up all test-generated PPTX files.

        Returns:
            Number of files cleaned up
        """
        cleaned_count = 0
        for file_path in self.created_files:
            try:
                if file_path.exists():
                    file_path.unlink()
                    cleaned_count += 1
            except Exception as e:
                print(f"Warning: Failed to delete {file_path}: {e}")

        self.created_files.clear()
        return cleaned_count


# Test instances following pytest conventions
class TestBasicCategoryExamples:
    """Test all basic category structured frontmatter examples."""

    def setup_method(self):
        """Set up test harness for each test method."""
        self.tester = StructuredFrontmatterTester()
        self.test_files_dir = self.tester.test_files_dir

    def test_title_slide_example(self):
        """Test title slide example."""
        json_file = self.test_files_dir / "example_title_slide.json"
        success, errors = self.tester.test_example_file(json_file, "test_title_slide")

        assert success, f"Title slide test failed: {errors}"

    def test_title_and_content_example(self):
        """Test title and content example."""
        json_file = self.test_files_dir / "example_title_and_content.json"
        success, errors = self.tester.test_example_file(json_file, "test_title_content")

        assert success, f"Title and content test failed: {errors}"

    def test_title_only_example(self):
        """Test title only example."""
        json_file = self.test_files_dir / "example_title_only.json"
        success, errors = self.tester.test_example_file(json_file, "test_title_only")

        assert success, f"Title only test failed: {errors}"

    def test_blank_example(self):
        """Test blank example."""
        json_file = self.test_files_dir / "example_blank.json"
        success, errors = self.tester.test_example_file(json_file, "test_blank")

        assert success, f"Blank test failed: {errors}"


class TestExpandedLayoutExamples:
    """Test expanded layout structured frontmatter examples."""

    def setup_method(self):
        """Set up test harness for each test method."""
        self.tester = StructuredFrontmatterTester()
        self.test_files_dir = self.tester.test_files_dir

    def test_two_content_example(self):
        """Test two content layout example."""
        json_file = self.test_files_dir / "example_two_content.json"
        success, errors = self.tester.test_example_file(json_file, "test_two_content")

        assert success, f"Two content test failed: {errors}"

    def test_comparison_example(self):
        """Test comparison layout example."""
        json_file = self.test_files_dir / "example_comparison.json"
        success, errors = self.tester.test_example_file(json_file, "test_comparison")

        assert success, f"Comparison test failed: {errors}"

    def test_three_columns_example(self):
        """Test three columns layout example."""
        json_file = self.test_files_dir / "example_three_columns.json"
        success, errors = self.tester.test_example_file(json_file, "test_three_columns")

        assert success, f"Three columns test failed: {errors}"

    def test_picture_with_caption_example(self):
        """Test picture with caption layout example."""
        json_file = self.test_files_dir / "example_picture_with_caption.json"
        success, errors = self.tester.test_example_file(json_file, "test_picture_caption")

        assert success, f"Picture with caption test failed: {errors}"

    def test_table_only_example(self):
        """Test table only layout example."""
        json_file = self.test_files_dir / "example_table_only.json"
        success, errors = self.tester.test_example_file(json_file, "test_table_only")

        assert success, f"Table only test failed: {errors}"

    def test_swot_analysis_example(self):
        """Test SWOT analysis layout example."""
        json_file = self.test_files_dir / "example_swot_analysis.json"
        success, errors = self.tester.test_example_file(json_file, "test_swot_analysis")

        assert success, f"SWOT analysis test failed: {errors}"

    def test_before_and_after_example(self):
        """Test before and after layout example."""
        json_file = self.test_files_dir / "example_before_and_after.json"
        success, errors = self.tester.test_example_file(json_file, "test_before_after")

        assert success, f"Before and after test failed: {errors}"

    def test_big_number_example(self):
        """Test big number layout example."""
        json_file = self.test_files_dir / "example_big_number.json"
        success, errors = self.tester.test_example_file(json_file, "test_big_number")

        assert success, f"Big number test failed: {errors}"

    def test_four_columns_example(self):
        """Test four columns layout example."""
        json_file = self.test_files_dir / "example_four_columns.json"
        success, errors = self.tester.test_example_file(json_file, "test_four_columns")

        assert success, f"Four columns test failed: {errors}"

    def test_key_metrics_example(self):
        """Test key metrics layout example."""
        json_file = self.test_files_dir / "example_key_metrics.json"
        success, errors = self.tester.test_example_file(json_file, "test_key_metrics")

        assert success, f"Key metrics test failed: {errors}"

    def test_problem_solution_example(self):
        """Test problem solution layout example."""
        json_file = self.test_files_dir / "example_problem_solution.json"
        success, errors = self.tester.test_example_file(json_file, "test_problem_solution")

        assert success, f"Problem solution test failed: {errors}"

    def test_process_steps_example(self):
        """Test process steps layout example."""
        json_file = self.test_files_dir / "example_process_steps.json"
        success, errors = self.tester.test_example_file(json_file, "test_process_steps")

        assert success, f"Process steps test failed: {errors}"

    def test_pros_and_cons_example(self):
        """Test pros and cons layout example."""
        json_file = self.test_files_dir / "example_pros_and_cons.json"
        success, errors = self.tester.test_example_file(json_file, "test_pros_cons")

        assert success, f"Pros and cons test failed: {errors}"

    def test_section_header_example(self):
        """Test section header layout example."""
        json_file = self.test_files_dir / "example_section_header.json"
        success, errors = self.tester.test_example_file(json_file, "test_section_header")

        assert success, f"Section header test failed: {errors}"

    def test_timeline_example(self):
        """Test timeline layout example."""
        json_file = self.test_files_dir / "example_timeline.json"
        success, errors = self.tester.test_example_file(json_file, "test_timeline")

        assert success, f"Timeline test failed: {errors}"

    def test_team_members_example(self):
        """Test team members layout example."""
        json_file = self.test_files_dir / "example_team_members.json"
        success, errors = self.tester.test_example_file(json_file, "test_team_members")

        assert success, f"Team members test failed: {errors}"

    def test_agenda_6_textboxes_example(self):
        """Test agenda 6 textboxes layout example."""
        json_file = self.test_files_dir / "example_agenda_6_textboxes.json"
        success, errors = self.tester.test_example_file(json_file, "test_agenda_6")

        assert success, f"Agenda 6 textboxes test failed: {errors}"

    def test_table_with_content_above_example(self):
        """Test table with content above layout example."""
        json_file = self.test_files_dir / "example_table_with_content_above.json"
        success, errors = self.tester.test_example_file(json_file, "test_table_content_above")

        assert success, f"Table with content above test failed: {errors}"


if __name__ == "__main__":
    # Allow running tests directly
    tester = StructuredFrontmatterTester()
    test_files_dir = tester.test_files_dir

    # Test all examples - organized by category
    basic_tests = [
        ("example_title_slide.json", "manual_title_slide"),
        ("example_title_and_content.json", "manual_title_content"),
        ("example_title_only.json", "manual_title_only"),
        ("example_blank.json", "manual_blank"),
    ]

    expanded_tests = [
        ("example_two_content.json", "manual_two_content"),
        ("example_comparison.json", "manual_comparison"),
        ("example_three_columns.json", "manual_three_columns"),
        ("example_picture_with_caption.json", "manual_picture_caption"),
        ("example_table_only.json", "manual_table_only"),
        ("example_swot_analysis.json", "manual_swot_analysis"),
        ("example_before_and_after.json", "manual_before_after"),
        ("example_big_number.json", "manual_big_number"),
        ("example_four_columns.json", "manual_four_columns"),
        ("example_key_metrics.json", "manual_key_metrics"),
        ("example_problem_solution.json", "manual_problem_solution"),
        ("example_process_steps.json", "manual_process_steps"),
        ("example_pros_and_cons.json", "manual_pros_cons"),
        ("example_section_header.json", "manual_section_header"),
        ("example_timeline.json", "manual_timeline"),
        ("example_team_members.json", "manual_team_members"),
        ("example_agenda_6_textboxes.json", "manual_agenda_6"),
        ("example_table_with_content_above.json", "manual_table_content_above"),
    ]

    print("🧪 Running Manual Test Suite for All Structured Frontmatter Examples")
    print("=" * 70)

    total_passed = 0
    total_failed = 0

    # Run basic category tests
    print(f"\n📁 Basic Category Examples ({len(basic_tests)} tests)")
    print("-" * 50)

    for test_file, test_name in basic_tests:
        json_file = test_files_dir / test_file
        print(f"📋 Testing: {test_file}")

        success, errors = tester.test_example_file(json_file, test_name)

        if success:
            print("  ✅ PASSED")
            total_passed += 1
        else:
            print("  ❌ FAILED")
            for error in errors:
                print(f"    • {error}")
            total_failed += 1

    # Run expanded category tests
    print(f"\n📁 Expanded Layout Examples ({len(expanded_tests)} tests)")
    print("-" * 50)

    for test_file, test_name in expanded_tests:
        json_file = test_files_dir / test_file
        print(f"📋 Testing: {test_file}")

        success, errors = tester.test_example_file(json_file, test_name)

        if success:
            print("  ✅ PASSED")
            total_passed += 1
        else:
            print("  ❌ FAILED")
            for error in errors:
                print(f"    • {error}")
            total_failed += 1

    # Final results
    total_tests = len(basic_tests) + len(expanded_tests)
    success_rate = (total_passed / total_tests) * 100 if total_tests > 0 else 0

    print("\n📊 Final Test Results")
    print("=" * 70)
    print(f"  Total Tests: {total_tests}")
    print(f"  ✅ Passed: {total_passed}")
    print(f"  ❌ Failed: {total_failed}")
    print(f"  📈 Success Rate: {success_rate:.1f}%")
    print("=" * 70)

    # Cleanup test files
    cleaned_count = tester.cleanup_test_files()
    print(f"\n🧹 Cleanup Complete: {cleaned_count} test files removed")
    print("=" * 70)
