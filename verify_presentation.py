#!/usr/bin/env python3
"""
Verification script to analyze generated PowerPoint presentations using python-pptx.
This script validates the dynamic multi-shape creation system.
"""

import sys
from pptx import Presentation

def analyze_presentation(pptx_path: str):
    """Analyze a PowerPoint presentation and report on its structure."""
    try:
        presentation = Presentation(pptx_path)
        print(f"📊 Analyzing presentation: {pptx_path}")
        print(f"Total slides: {len(presentation.slides)}")
        print("=" * 60)
        
        for slide_idx, slide in enumerate(presentation.slides, 1):
            print(f"\n🎯 Slide {slide_idx}: {slide.slide_layout.name}")
            print(f"   Shapes: {len(slide.shapes)}")
            
            # Analyze shapes
            shape_types = []
            table_count = 0
            text_shape_count = 0
            placeholder_count = 0
            
            for shape_idx, shape in enumerate(slide.shapes):
                shape_info = f"   [{shape_idx+1}]"
                
                if shape.shape_type.name == "PLACEHOLDER":
                    shape_info += f" PLACEHOLDER ({shape.placeholder_format.type.name})"
                    placeholder_count += 1
                    if hasattr(shape, 'text') and shape.text:
                        preview = shape.text[:50] + "..." if len(shape.text) > 50 else shape.text
                        shape_info += f" - '{preview}'"
                elif shape.shape_type.name == "TABLE":
                    shape_info += f" TABLE ({shape.table.rows} rows × {shape.table.columns} cols)"
                    table_count += 1
                elif shape.shape_type.name == "TEXT_BOX":
                    shape_info += f" TEXT_BOX"
                    text_shape_count += 1
                    if hasattr(shape, 'text') and shape.text:
                        preview = shape.text[:50] + "..." if len(shape.text) > 50 else shape.text
                        shape_info += f" - '{preview}'"
                else:
                    shape_info += f" {shape.shape_type.name}"
                
                shape_types.append(shape.shape_type.name)
                print(shape_info)
            
            # Summary for this slide
            print(f"   📈 Summary: {placeholder_count} placeholders, {table_count} tables, {text_shape_count} text boxes")
            
            # Special analysis for dynamic multi-shape slides
            if table_count > 0 and text_shape_count > 0:
                print(f"   ✅ DYNAMIC MULTI-SHAPE DETECTED: Mixed content with separate shapes")
            elif table_count > 0:
                print(f"   📋 Table-only slide")
            else:
                print(f"   📝 Text-only slide")
        
        print("\n" + "=" * 60)
        print("🎯 DYNAMIC MULTI-SHAPE ANALYSIS COMPLETE")
        
        # Look for evidence of the dynamic system working
        total_tables = sum(1 for slide in presentation.slides for shape in slide.shapes if shape.shape_type.name == "TABLE")
        total_text_boxes = sum(1 for slide in presentation.slides for shape in slide.shapes if shape.shape_type.name == "TEXT_BOX")
        
        print(f"📊 PRESENTATION SUMMARY:")
        print(f"   Total slides: {len(presentation.slides)}")
        print(f"   Total tables: {total_tables}")
        print(f"   Total dynamic text boxes: {total_text_boxes}")
        
        if total_tables > 0 and total_text_boxes > 0:
            print(f"   ✅ SUCCESS: Dynamic multi-shape creation is working")
        elif total_tables > 0:
            print(f"   ⚠️  Tables found but no additional text boxes created")
        else:
            print(f"   ❌ No tables or dynamic shapes detected")
        
        return True
        
    except Exception as e:
        print(f"❌ Error analyzing presentation: {e}")
        return False

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python verify_presentation.py <path_to_pptx>")
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    success = analyze_presentation(pptx_path)
    sys.exit(0 if success else 1)