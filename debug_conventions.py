#!/usr/bin/env python3
"""
Debug script for convention-based naming
"""

import os
import sys
sys.path.insert(0, 'src')

from pptx import Presentation
from src.deckbuilder.naming_conventions import NamingConvention, PlaceholderContext

def test_convention_mapping():
    """Test the convention mapping generation"""
    template_path = "assets/templates/default.pptx"
    
    print(f"Loading template: {template_path}")
    
    if not os.path.exists(template_path):
        print(f"ERROR: Template not found at {template_path}")
        return
    
    try:
        # Load presentation
        prs = Presentation(template_path)
        convention = NamingConvention()
        
        print(f"✅ Template loaded successfully")
        print(f"📊 Found {len(prs.slide_layouts)} slide layouts")
        
        # Test convention mapping generation
        mapping = {
            "template_info": {
                "name": "Convention-Based",
                "version": "1.0"
            },
            "layouts": {}
        }
        
        layout_count = 0
        for layout in prs.slide_layouts:
            layout_name = layout.name
            layout_placeholders = {}
            
            print(f"\n🔍 Processing layout: {layout_name}")
            print(f"   📊 Placeholders: {len(layout.placeholders)}")
            
            for placeholder in layout.placeholders:
                try:
                    placeholder_idx = str(placeholder.placeholder_format.idx)
                    
                    # Create context for convention naming
                    context = PlaceholderContext(
                        layout_name=layout_name,
                        placeholder_idx=placeholder_idx,
                        total_placeholders=len(layout.placeholders)
                    )
                    
                    # Generate convention-based name
                    convention_name = convention.generate_placeholder_name(context)
                    layout_placeholders[placeholder_idx] = convention_name
                    
                    print(f"   - idx={placeholder_idx}: {placeholder.name} → {convention_name}")
                    
                except Exception as e:
                    print(f"   ❌ Error processing placeholder: {e}")
            
            # Debug layout attributes
            print(f"   🔍 Layout attributes: {[attr for attr in dir(layout) if not attr.startswith('_')]}")
            
            # Try to get the layout index - use a simple counter for now
            layout_index = layout_count
            
            mapping["layouts"][layout_name] = {
                "index": layout_index,
                "placeholders": layout_placeholders
            }
            
            layout_count += 1
        
        print(f"\n✅ Convention mapping generated successfully!")
        print(f"📊 Processed {layout_count} layouts")
        
        # Show sample mapping
        print(f"\n📋 Sample mapping (first layout):")
        first_layout = list(mapping["layouts"].keys())[0]
        first_layout_data = mapping["layouts"][first_layout]
        print(f"Layout: {first_layout}")
        for idx, name in first_layout_data["placeholders"].items():
            print(f"  {idx}: {name}")
        
        return mapping
        
    except Exception as e:
        print(f"❌ ERROR: {e}")
        import traceback
        traceback.print_exc()
        return None

if __name__ == "__main__":
    test_convention_mapping()